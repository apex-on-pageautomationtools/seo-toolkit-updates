"""
Health Audit Module for SEO Toolkit Pro
Adapted from james-seo-tools generate_health_report.py
- No Patchright (uses Selenium CDP for screenshots)
- GSC screenshots (manual_action, security_issues) captured if GSC account is connected
- 4 format themes: James (docx), Neon (docx), Sigma (pptx), Omega (pptx)
"""

import os
import re
import tempfile
import time
import urllib.request as _ur
import urllib.error as _ue
from pathlib import Path
from datetime import datetime
from html.parser import HTMLParser
import engine

# ---------- HTTP helpers ----------
_UA = "Mozilla/5.0 (compatible; SEOToolkitPro-HealthBot/1.0)"


def _fetch_html(url, timeout=15):
    try:
        req = _ur.Request(url, headers={"User-Agent": _UA})
        with _ur.urlopen(req, timeout=timeout) as r:
            return r.read().decode("utf-8", "ignore")
    except Exception:
        return None


# Per-run cache of browser-rendered pages so the same URL is only rendered once
# across the canonical / double-meta / meta-robots / meta-source checks.
_RENDER_CACHE = {}


def _render_html(url, driver=None, wait=1.2):
    """Return the page as a real browser renders it (so JS-set title / meta /
    canonical are seen - Google renders JS, so this is the ACTUAL state that
    matters). Falls back to raw HTTP when no browser/driver is available."""
    if url in _RENDER_CACHE:
        return _RENDER_CACHE[url]
    html = None
    if driver is not None:
        import time
        try:
            driver.get(url)
            for _ in range(16):                       # wait up to ~4s for load
                try:
                    if driver.execute_script("return document.readyState") == "complete":
                        break
                except Exception:
                    break
                time.sleep(0.25)
            time.sleep(wait)                          # short settle for hydration
            h = driver.page_source or ""
            if h and len(h) > 200:
                html = h
        except Exception:
            html = None
    if not html:
        html = _fetch_html(url)
    _RENDER_CACHE[url] = html
    return html


def _google_via_selenium(driver, query, num=30):
    """Run a Google search through the Selenium browser to avoid bot blocks.
    Returns page HTML or None. Uses gentle delays to stay under the radar."""
    import time, random
    from urllib.parse import quote_plus
    try:
        url = f"https://www.google.com/search?q={quote_plus(query)}&num={num}"
        driver.get(url)
        time.sleep(random.uniform(2.5, 4.5))
        return driver.page_source
    except Exception:
        return None


def _get_status(url, timeout=10):
    for method in ("HEAD", "GET"):
        try:
            req = _ur.Request(url, headers={"User-Agent": _UA}, method=method)
            with _ur.urlopen(req, timeout=timeout) as r:
                return r.status
        except _ue.HTTPError as e:
            if method == "HEAD" and e.code in (403, 405):
                continue
            return e.code
        except Exception:
            return 0
    return 0


def _parse_meta(html):
    result = {"canonical": [], "title": [], "description": [], "robots": []}
    for m in re.finditer(r'<link[^>]+rel=["\']canonical["\'][^>]*href=["\']([^"\']+)["\']', html, re.I):
        result["canonical"].append(m.group(1).strip())
    for m in re.finditer(r'<link[^>]+href=["\']([^"\']+)["\'][^>]*rel=["\']canonical["\']', html, re.I):
        result["canonical"].append(m.group(1).strip())
    for m in re.finditer(r'<title[^>]*>(.*?)</title>', html, re.I | re.S):
        t = m.group(1).strip()
        if t:
            result["title"].append(t)
    for m in re.finditer(r'<meta[^>]+name=["\']description["\'][^>]*content=["\']([^"\']*)["\']', html, re.I):
        result["description"].append(m.group(1).strip())
    for m in re.finditer(r'<meta[^>]+content=["\']([^"\']*)["\'][^>]*name=["\']description["\']', html, re.I):
        result["description"].append(m.group(1).strip())
    for m in re.finditer(r'<meta[^>]+name=["\']robots["\'][^>]*content=["\']([^"\']*)["\']', html, re.I):
        result["robots"].append(m.group(1).strip())
    for m in re.finditer(r'<meta[^>]+content=["\']([^"\']*)["\'][^>]*name=["\']robots["\']', html, re.I):
        result["robots"].append(m.group(1).strip())
    return result


# ---------- Checks ----------
def check_status200(target_pages, domain):
    if not target_pages:
        target_pages = [f"https://{domain}/"]
    results = []
    for url in target_pages:
        url = url.strip()
        if not url:
            continue
        if not url.startswith("http"):
            url = "https://" + url
        code = _get_status(url)
        results.append({"url": url, "status": str(code), "ok": code == 200})
    return results


def _detect_live_homepage(domain):
    for prefix in (f"https://{domain}/", f"https://www.{domain}/",
                   f"http://{domain}/", f"http://www.{domain}/"):
        try:
            req = _ur.Request(prefix, headers={"User-Agent": _UA})
            with _ur.urlopen(req, timeout=10) as r:
                return r.url
        except Exception:
            continue
    return f"https://{domain}/"


def check_canonical(target_pages, domain, driver=None):
    pages = list(target_pages) if target_pages else []
    homepage = _detect_live_homepage(domain)
    if not any(p.rstrip("/") == homepage.rstrip("/") for p in pages):
        pages = [homepage] + pages
    results = []
    for url in pages:
        url = url.strip()
        if not url:
            continue
        if not url.startswith("http"):
            url = "https://" + url
        html = _render_html(url, driver)
        if html is None:
            results.append({"url": url, "canonical": "N/A", "ok": False, "note": "Could not fetch"})
            continue
        meta = _parse_meta(html)
        canonicals = list(dict.fromkeys(meta["canonical"]))
        if not canonicals:
            results.append({"url": url, "canonical": "Missing", "ok": False, "note": "No canonical tag"})
        elif len(canonicals) > 1:
            results.append({"url": url, "canonical": canonicals[0], "ok": False, "note": f"{len(canonicals)} canonical tags"})
        else:
            canon = canonicals[0].rstrip("/")
            page_url = url.rstrip("/")
            ok = canon == page_url
            note = "Self-referencing"
            if not ok:
                # A canonical that only differs by scheme (http/https) or a "www."
                # prefix from the page it's on - same host otherwise, same path -
                # is normal, correct canonicalization to one preferred version
                # (e.g. the non-www page pointing to the www version), not an issue.
                canon_host_path = _strip_scheme_www(canon)
                page_host_path = _strip_scheme_www(page_url)
                if canon_host_path == page_host_path:
                    ok = True
                    note = "Self-referencing (points to preferred www/https version)"
                else:
                    note = "Points elsewhere"
            results.append({"url": url, "canonical": canonicals[0], "ok": ok, "note": note})
    return results


def _strip_scheme_www(url):
    """Normalize a URL to host+path for comparing canonical variants that only
    differ by scheme (http/https) or a 'www.' prefix - not real canonical issues."""
    from urllib.parse import urlsplit
    parts = urlsplit(url)
    host = parts.netloc.lower()
    if host.startswith("www."):
        host = host[4:]
    return host + parts.path.rstrip("/")


def check_double_meta(target_pages, domain, driver=None):
    pages = list(target_pages) if target_pages else []
    if not pages:
        pages = [f"https://{domain}/"]
    results = []
    for url in pages:
        url = url.strip()
        if not url:
            continue
        if not url.startswith("http"):
            url = "https://" + url
        html = _render_html(url, driver)
        if html is None:
            results.append({"url": url, "title_count": "N/A", "desc_count": "N/A", "ok": False, "note": "Could not fetch"})
            continue
        meta = _parse_meta(html)
        tc, dc = len(meta["title"]), len(meta["description"])
        ok = (tc == 1 and dc == 1)
        notes = []
        if tc == 0:
            notes.append("Missing title")
        elif tc > 1:
            notes.append(f"Title ×{tc}")
        if dc == 0:
            notes.append("Missing meta desc")
        elif dc > 1:
            notes.append(f"Meta desc ×{dc}")
        results.append({"url": url, "title_count": str(tc), "desc_count": str(dc),
                         "ok": ok, "note": ", ".join(notes) if notes else "OK"})
    return results


def check_meta_suggestions(target_pages, domain, driver=None):
    """Extract the actual title + meta description per page, instead of
    screenshotting the raw page source."""
    pages = list(target_pages) if target_pages else []
    if not pages:
        pages = [f"https://{domain}/"]
    results = []
    for url in pages:
        url = url.strip()
        if not url:
            continue
        if not url.startswith("http"):
            url = "https://" + url
        html = _render_html(url, driver)
        if html is None:
            results.append({"url": url, "title": "N/A", "description": "N/A",
                            "ok": False, "note": "Could not fetch"})
            continue
        meta = _parse_meta(html)
        title = meta["title"][0] if meta["title"] else "Missing"
        desc = meta["description"][0] if meta["description"] else "Missing"
        ok = bool(meta["title"]) and bool(meta["description"])
        results.append({"url": url, "title": title, "description": desc, "ok": ok})
    return results


def check_meta_robots(target_pages, domain, driver=None):
    pages = list(target_pages) if target_pages else []
    if not pages:
        pages = [f"https://{domain}/"]
    results = []
    for url in pages:
        url = url.strip()
        if not url:
            continue
        if not url.startswith("http"):
            url = "https://" + url
        html = _render_html(url, driver)
        if html is None:
            results.append({"url": url, "robots_value": "N/A", "ok": False, "note": "Could not fetch"})
            continue
        meta = _parse_meta(html)
        robots_values = meta["robots"]
        if not robots_values:
            results.append({"url": url, "robots_value": "Not set (default: index,follow)", "ok": True, "note": "OK"})
            continue
        val = robots_values[0]
        val_lower = val.lower()
        issues = []
        if "noindex" in val_lower:
            issues.append("noindex")
        if "nofollow" in val_lower:
            issues.append("nofollow")
        ok = not issues
        results.append({"url": url, "robots_value": val, "ok": ok,
                         "note": ", ".join(issues) if issues else "OK"})
    return results


def check_versions_full(domain):
    """Determine whether the site runs as a single canonical version or as several
    independent ones. Each variant is followed through its FULL redirect chain to the
    final destination - a site that 301-redirects http/non-www/etc. to one canonical URL
    (e.g. http://x -> https://x -> https://www.x/) is a single version, NOT multiple.
    Only variants that independently land on DIFFERENT HTTP 200 URLs count as multiple."""
    variants = [
        f"http://{domain}/", f"http://www.{domain}/",
        f"https://{domain}/", f"https://www.{domain}/",
    ]
    import time as _time
    results = []  # (variant, final_url, final_code, redirected)
    for url in variants:
        final_url, final_code = None, 0
        # A site can rate-limit (429) or hiccup (503) when we fire 4 requests back
        # to back - retry with a short backoff before treating the variant as down,
        # otherwise a transient rate limit gets misreported as "site unreachable".
        for attempt in range(3):
            try:
                req = _ur.Request(url, headers={"User-Agent": _UA})
                # urlopen follows redirects by default, so geturl() is the FINAL
                # landing URL after the whole chain (http -> https -> www ...).
                with _ur.urlopen(req, timeout=12) as r:
                    final_url = (r.geturl() or url).rstrip("/")
                    final_code = r.status
                break
            except _ue.HTTPError as he:
                final_url, final_code = url.rstrip("/"), he.code
                if he.code in (429, 503) and attempt < 2:
                    _time.sleep(2 * (attempt + 1))
                    continue
                break
            except Exception:
                final_url, final_code = None, 0
                if attempt < 2:
                    _time.sleep(2 * (attempt + 1))
                    continue
                break
        redirected = bool(final_url and final_url != url.rstrip("/"))
        results.append((url, final_url, final_code, redirected))

    # Distinct URLs that actually serve content (HTTP 200) after following redirects.
    final_200 = sorted({fu for _, fu, code, _ in results if code == 200 and fu})

    if len(final_200) == 1:
        is_ok = True
        canonical = final_200[0]
        if any(red for _, _, code, red in results if code == 200):
            summary = (f"No issue found. All variants correctly redirect to a single "
                       f"canonical version: {canonical}")
        else:
            summary = f"No issue found. The site runs on a single version: {canonical}"
    elif len(final_200) >= 2:
        is_ok = False
        summary = (f"Issue found - the site is reachable at {len(final_200)} independent "
                   f"versions ({', '.join(final_200)}) that each return HTTP 200 without "
                   f"redirecting to one canonical URL. Fix: 301-redirect all variants to "
                   f"a single canonical URL.")
    else:
        is_ok = False
        summary = ("Could not verify the site version(s) - no variant returned HTTP 200. "
                   "Please check the site's availability manually.")

    rows = []
    for url, fu, code, redirected in results:
        if code == 200 and not redirected:
            rows.append((url, "200", "OK (canonical)", is_ok))
        elif code == 200 and redirected:
            rows.append((url, "200", f"redirects -> {fu}", True))
        elif code:
            rows.append((url, str(code), f"-> {fu or '?'}", True))
        else:
            rows.append((url, "ERR", "unreachable", False))
    return summary, rows


def check_dummy_content(domain, driver=None):
    query = f"site:{domain} lorem"
    if driver:
        html = _google_via_selenium(driver, query, num=10)
    else:
        html = _fetch_html(f"https://www.google.com/search?q=site:{domain}+lorem")
    if html is None:
        return "Could not check Google for dummy content. Please check manually."
    no_results_markers = ["did not match any documents", "No results found", "no results", "0 results"]
    for marker in no_results_markers:
        if marker.lower() in html.lower():
            return f"No dummy content found. Google search for 'site:{domain} lorem' returned no results."
    count_match = re.search(r'About ([\d,]+) results', html)
    count = count_match.group(1) if count_match else "some"
    titles = re.findall(r'<h3[^>]*>(.*?)</h3>', html, re.S)
    titles = [re.sub(r'<[^>]+>', '', t).strip() for t in titles[:5] if t.strip()]
    if titles:
        result = f"Possible dummy content found - Google shows {count} result(s) for 'site:{domain} lorem'.\nPages found:\n"
        result += "\n".join(f"  - {t}" for t in titles)
    else:
        result = f"Google returned {count} result(s) for 'site:{domain} lorem'. Please verify manually."
    return result


def _suggest_redirect_for_broken(url, headers):
    """For a dead link, try a couple of cheap same-URL variants a real 404 often
    has a working replacement for (trailing slash toggled, query string dropped,
    a trailing /index.html or /index.php stripped) and return the first one that
    actually resolves (200). Mirrors what a human checking brokenlinkcheck.com
    results would try by hand - not a full crawl, just the obvious candidates."""
    from urllib.parse import urlsplit, urlunsplit
    parts = urlsplit(url)
    candidates = []
    path = parts.path
    if path.endswith("/index.html") or path.endswith("/index.php"):
        candidates.append(path.rsplit("/", 1)[0] + "/")
    if path.endswith("/") and path != "/":
        candidates.append(path.rstrip("/"))
    elif not path.endswith("/"):
        candidates.append(path + "/")
    if parts.query:
        candidates.append(path)  # same path, query string dropped
    seen = set()
    for cand_path in candidates:
        cand = urlunsplit((parts.scheme, parts.netloc, cand_path, "", ""))
        if cand == url or cand in seen:
            continue
        seen.add(cand)
        try:
            req = _ur.Request(cand, headers=headers, method="HEAD")
            if _ur.urlopen(req, timeout=8).status == 200:
                return cand
        except Exception:
            continue
    return None


def check_broken_links(domain, target_pages=None):
    """Check all links on every target page (or homepage if none given).
    No artificial limit - every link on each page is checked. Each broken link
    records which page(s) it was found on (like brokenlinkcheck.com shows) and,
    where a working replacement URL can be found, a suggested redirect target."""
    from urllib.parse import urljoin
    headers = {"User-Agent": _UA}

    pages = [p.strip() for p in (target_pages or []) if p.strip()]
    if not pages:
        pages = [f"https://{domain}/"]
    # Normalise
    pages = [p if p.startswith("http") else f"https://{p}" for p in pages]

    all_links = []
    seen = set()
    found_on = {}
    for page_url in pages:
        try:
            html = _ur.urlopen(_ur.Request(page_url, headers=headers), timeout=20).read().decode("utf-8", "ignore")
        except Exception:
            continue
        for m in re.finditer(r'href=["\']([^"\']+)["\']', html):
            href = m.group(1).strip()
            if href.startswith(("mailto:", "tel:", "javascript:", "#", "data:")):
                continue
            u = urljoin(page_url, href)
            if u.startswith("http"):
                found_on.setdefault(u, []).append(page_url)
                if u not in seen:
                    seen.add(u)
                    all_links.append(u)

    def _check_one(u):
        code = None
        for method in ("HEAD", "GET"):
            try:
                req = _ur.Request(u, headers=headers, method=method)
                code = _ur.urlopen(req, timeout=12).status
                break
            except _ue.HTTPError as he:
                code = he.code
                if method == "HEAD" and code in (403, 405):
                    continue
                break
            except Exception:
                code = 0
        return u, code

    # Links are checked independently of each other, so a small thread pool checks
    # them concurrently instead of one HEAD/GET at a time - same per-link logic and
    # same 404/410-only broken criteria below, just not serialized.
    import concurrent.futures
    broken = []
    if all_links:
        with concurrent.futures.ThreadPoolExecutor(max_workers=10) as pool:
            dead = []
            for u, code in pool.map(_check_one, all_links):
                # Only genuinely dead links count as broken: 404 (Not Found) / 410 (Gone).
                # 429 (rate-limited), 403/401/405 (bot-blocked), 5xx (temporary server
                # errors) and timeouts are NOT broken links - they're the server
                # refusing the automated request, so flagging them gives false
                # "broken link" counts.
                if code in (404, 410):
                    dead.append((u, code))
            redirects = list(pool.map(lambda uc: _suggest_redirect_for_broken(uc[0], headers), dead))
            for (u, code), redirect in zip(dead, redirects):
                broken.append({"url": u, "code": code, "found_on": found_on.get(u, [])[:3],
                                "suggested_redirect": redirect})
    return len(all_links), broken


def check_sucuri(domain):
    """Fetch Sucuri SiteCheck results page and parse the actual verdict."""
    import re
    url = f"https://sitecheck.sucuri.net/results/https/{domain}"
    html = _fetch_html(url, timeout=30)
    if html is None:
        return {"ok": None, "summary": "Could not reach Sucuri SiteCheck. Check manually.", "details": [], "issues": []}
    details = []
    issues = []
    ok = True
    html_lower = html.lower()

    # Malware - check the actual result text, not template text
    has_no_malware = "didn't detect any malware" in html_lower or "no malware found" in html_lower
    has_malware = bool(re.search(r'Warning:\s*Malware\s+Detected', html))
    # The raw HTML contains BOTH templates; check which is actually shown
    # "No malware detected by scan (Low Risk)" = clean; presence of this confirms clean
    clean_scan = "no malware detected by scan" in html_lower
    if has_no_malware or clean_scan:
        details.append("No malware found")
        malware_ok = True
    elif has_malware:
        ok = False
        details.append("Malware detected!")
        issues.append("Malware detected on the website")
        malware_ok = False
    else:
        details.append("Malware status unclear")
        malware_ok = True

    # Blacklist
    if "is not blacklisted" in html_lower or "site is not blacklisted" in html_lower:
        details.append("Not blacklisted")
        bl_ok = True
    elif re.search(r'site\s+is\s+blacklisted', html, re.IGNORECASE):
        ok = False
        details.append("Domain is BLACKLISTED")
        issues.append("Domain is blacklisted")
        bl_ok = False
    else:
        bl_ok = True

    # Blacklist details - count "Domain clean by" entries
    bl_clean = len(re.findall(r'Domain clean by', html))
    if bl_clean:
        details.append(f"{bl_clean} blacklist checks clean")

    # Risk level - infer from malware/blacklist status
    if not malware_ok or not bl_ok:
        risk_level = "High"
    elif malware_ok and bl_ok:
        risk_level = "Low"
    else:
        risk_level = "Medium"
    details.append(f"Risk level: {risk_level}")

    # Security headers (from Hardening Improvements section)
    if "missing security header" in html_lower:
        header_issues = re.findall(r'Missing security header[^<]*<a[^>]*>([^<]+)</a>', html, re.IGNORECASE)
        header_issues += re.findall(r'Missing\s+<a[^>]*>([^<]+?)</a>\s*security header', html, re.IGNORECASE)
        for h in header_issues:
            issues.append(f"Missing: {h.strip()}")
        if not header_issues:
            issues.append("Missing security headers")

    # WAF
    if "no website application firewall" in html_lower or "please install a cloud-based waf" in html_lower:
        issues.append("No WAF/Firewall detected")
    elif "firewall detected" in html_lower:
        details.append("WAF/Firewall detected")

    # PHP version leaked
    if "leaked php version" in html_lower or "expose_php" in html_lower:
        issues.append("PHP version exposed in headers")

    # Missing Strict-Transport-Security
    if "strict-transport-security" in html_lower and "missing" in html_lower:
        issues.append("Missing: Strict-Transport-Security header")

    # Missing CSP
    if "content-security-policy directive" in html_lower and "missing" in html_lower:
        issues.append("Missing: Content-Security-Policy directive")

    summary_parts = []
    if ok:
        summary_parts.append("No critical issues found.")
    else:
        summary_parts.append("Critical issues detected!")
    summary_parts.append(f"Risk: {risk_level}.")
    if issues:
        summary_parts.append("Hardening: " + "; ".join(issues[:6]) + ".")

    return {"ok": ok, "summary": " ".join(summary_parts), "details": details, "issues": issues}


def check_robots_txt(domain):
    """Actually fetch and validate robots.txt."""
    url = f"https://{domain}/robots.txt"
    try:
        req = _ur.Request(url, headers={"User-Agent": _UA})
        with _ur.urlopen(req, timeout=10) as r:
            code = r.status
            content = r.read().decode("utf-8", "ignore")
    except _ue.HTTPError as e:
        return {"ok": False, "status": e.code,
                "summary": f"robots.txt returned HTTP {e.code} - file is missing or inaccessible."}
    except Exception as e:
        return {"ok": False, "status": 0,
                "summary": f"Could not fetch robots.txt: {e}"}
    if not content.strip():
        return {"ok": False, "status": code,
                "summary": "robots.txt exists but is empty."}
    has_user_agent = bool(re.search(r'(?i)^user-agent:', content, re.MULTILINE))
    has_sitemap = bool(re.search(r'(?i)^sitemap:', content, re.MULTILINE))
    has_disallow = bool(re.search(r'(?i)^disallow:', content, re.MULTILINE))
    # Is the whole site blocked from crawling? (Disallow: / under User-agent: * or Googlebot)
    root_blocked = False
    cur_ua = None
    for ln in content.splitlines():
        t = ln.split("#", 1)[0].strip()
        low = t.lower()
        if low.startswith("user-agent:"):
            cur_ua = low.split(":", 1)[1].strip()
        elif low.startswith("disallow:") and cur_ua in ("*", "googlebot"):
            if t.split(":", 1)[1].strip() == "/":
                root_blocked = True
    issues = []
    if root_blocked:
        issues.append("Main site is BLOCKED from crawling (Disallow: /)")
    if not has_user_agent:
        issues.append("Missing User-agent directive")
    if not has_sitemap:
        issues.append("No Sitemap reference in robots.txt")
    line_count = len([l for l in content.splitlines() if l.strip()])
    summary = f"robots.txt found ({line_count} lines). "
    if root_blocked:
        summary += "WARNING: the site is disallowed for search crawlers. "
    if issues:
        summary += "Issues: " + "; ".join(issues) + "."
    else:
        summary += "Main site is crawlable" + (", sitemap referenced." if has_sitemap else ".")
    return {"ok": not issues, "status": code, "summary": summary, "found": True,
            "root_blocked": root_blocked, "site_crawlable": not root_blocked,
            "has_sitemap_ref": has_sitemap, "has_disallow": has_disallow, "lines": line_count}


def check_sitemap(domain, robots_data=None):
    """Fetch and validate the sitemap. Tries /sitemap.xml then /sitemap_index.xml
    (many sites - e.g. Yoast - expose only a sitemap index). `robots_data` is
    accepted for compatibility with callers that pass robots.txt results; when it
    carries explicit Sitemap: URLs those are tried first."""
    candidates = []
    # Prefer any Sitemap: URLs surfaced by the robots.txt check.
    if isinstance(robots_data, dict):
        for k in ("sitemaps", "sitemap_urls", "sitemaps_found"):
            v = robots_data.get(k)
            if isinstance(v, (list, tuple)):
                candidates.extend([u for u in v if isinstance(u, str) and u.startswith("http")])
            elif isinstance(v, str) and v.startswith("http"):
                candidates.append(v)
    elif isinstance(robots_data, str):
        candidates.extend(re.findall(r'(?im)^\s*sitemap:\s*(\S+)', robots_data))
    # Common fallbacks.
    for u in (f"https://{domain}/sitemap.xml", f"https://{domain}/sitemap_index.xml"):
        if u not in candidates:
            candidates.append(u)

    last = {"ok": False, "status": 0, "summary": "No sitemap found."}
    for url in candidates:
        name = url.rsplit("/", 1)[-1] or "sitemap.xml"
        # A raw connection/timeout failure (NOT a real HTTP 404/etc.) gets one retry
        # after a short pause before giving up on this candidate - a transient
        # network blip shouldn't permanently report "sitemap not found" for a site
        # whose sitemap is actually fine, and this result isn't checked again later
        # in the same run.
        content = None
        for attempt in (1, 2):
            try:
                req = _ur.Request(url, headers={"User-Agent": _UA})
                with _ur.urlopen(req, timeout=15) as r:
                    code = r.status
                    content = r.read().decode("utf-8", "ignore")
                break
            except _ue.HTTPError as e:
                last = {"ok": False, "status": e.code,
                        "summary": f"{name} returned HTTP {e.code} - file is missing or inaccessible."}
                break
            except Exception as e:
                last = {"ok": False, "status": 0, "summary": f"Could not fetch {name}: {e}"}
                if attempt == 1:
                    time.sleep(3)
        if content is None:
            continue
        if not content.strip():
            last = {"ok": False, "status": code, "summary": f"{name} exists but is empty."}
            continue
        is_xml = content.strip().startswith("<?xml") or "<urlset" in content or "<sitemapindex" in content
        if not is_xml:
            last = {"ok": False, "status": code,
                    "summary": f"{name} exists but does not appear to be valid XML."}
            continue
        url_count = len(re.findall(r'<loc>', content))
        is_index = "<sitemapindex" in content
        if is_index:
            sub_count = len(re.findall(r'<sitemap>', content))
            summary = f"Sitemap index found with {sub_count} sub-sitemap(s)."
        else:
            summary = f"Sitemap found with {url_count} URL(s)."
        return {"ok": True, "status": code, "summary": summary,
                "url_count": url_count, "is_index": is_index, "url": url}
    return last


def check_serp(domain, driver=None):
    """Check Google SERP for site:{domain} - count indexed pages AND detect
    hacked/spam content (pharma, casino, Japanese SEO spam, cloaked pages)."""
    query = f"site:{domain}"
    if driver:
        html = _google_via_selenium(driver, query, num=30)
    else:
        html = _fetch_html(f"https://www.google.com/search?q=site:{domain}&num=30")
    if html is None:
        return {"ok": None, "summary": "Could not check Google SERP. Please check manually.",
                "result_count": "unknown", "spam_found": []}

    count_match = re.search(r'About ([\d,]+) results', html)
    count = count_match.group(1) if count_match else None

    no_results = any(m.lower() in html.lower()
                     for m in ["did not match any documents", "0 results", "No results found"])
    if no_results:
        return {"ok": False,
                "summary": f"No pages found indexed for site:{domain}. The website may not be indexed.",
                "result_count": "0", "spam_found": []}

    # Extract titles and snippets for spam/hack detection
    titles = re.findall(r'<h3[^>]*>(.*?)</h3>', html, re.S)
    titles = [re.sub(r'<[^>]+>', '', t).strip() for t in titles if t.strip()]
    snippets = re.findall(r'<span[^>]*class="[^"]*st[^"]*"[^>]*>(.*?)</span>', html, re.S)
    snippets = [re.sub(r'<[^>]+>', '', s).strip() for s in snippets if s.strip()]
    all_text = " ".join(titles + snippets).lower()

    # Hack/spam indicators
    SPAM_PATTERNS = [
        (r'\b(viagra|cialis|levitra|pharmacy|pharma|pills|tablets|medication)\b', "Pharma spam"),
        (r'\b(casino|poker|slots|gambling|bet online|betting)\b', "Casino/gambling spam"),
        (r'\b(buy cheap|cheap online|discount|wholesale|replica|counterfeit)\b', "Commercial spam"),
        (r'\b(hacked by|pwned|defaced|owned by)\b', "Hacked/defaced indicator"),
        (r'[぀-ヿ一-鿿]{5,}', "Japanese/Chinese SEO spam"),
        (r'\b(payday loan|loan online|cash advance|credit card hack)\b', "Financial spam"),
        (r'\b(free download|crack|keygen|serial key|torrent)\b', "Piracy/malware spam"),
        (r'\b(escort|adult|xxx|porn)\b', "Adult content injection"),
    ]
    spam_found = []
    for pattern, label in SPAM_PATTERNS:
        matches = re.findall(pattern, all_text, re.I)
        if matches:
            # Find which titles contain this spam
            flagged_titles = [t for t in titles if re.search(pattern, t.lower(), re.I)]
            spam_found.append({
                "type": label,
                "count": len(matches),
                "example_titles": flagged_titles[:3],
            })

    if spam_found:
        spam_types = ", ".join(s["type"] for s in spam_found)
        summary = (f"WARNING: Potential hacked/spam content detected in SERP results! "
                   f"Types found: {spam_types}. "
                   f"Google shows {count or 'multiple'} result(s) for site:{domain}. "
                   f"Immediate investigation recommended.")
        return {"ok": False, "summary": summary, "result_count": count or "unknown",
                "spam_found": spam_found}

    summary = f"No spam or hacked content detected. Google shows {count or 'multiple'} result(s) for site:{domain}."
    return {"ok": True, "summary": summary, "result_count": count or "unknown", "spam_found": []}


def check_blank_pages(target_pages, domain, driver=None):
    """Fetch target pages and check if they have minimal content (blank).
    Uses the RENDERED page so JS/SPA sites aren't falsely flagged as blank (their
    raw HTML is an empty shell). Only runs when target pages are provided."""
    pages = [p.strip() for p in (target_pages or []) if p.strip()]
    if not pages:
        return [], 0
    results = []
    for url in pages:
        if not url.startswith("http"):
            url = "https://" + url
        html = _render_html(url, driver)
        if html is None:
            results.append({"url": url, "ok": False, "note": "Could not fetch"})
            continue
        text = re.sub(r'<script[^>]*>.*?</script>', '', html, flags=re.S | re.I)
        text = re.sub(r'<style[^>]*>.*?</style>', '', text, flags=re.S | re.I)
        text = re.sub(r'<[^>]+>', '', text)
        text = re.sub(r'\s+', ' ', text).strip()
        is_blank = len(text) < 100
        results.append({"url": url, "ok": not is_blank, "text_len": len(text),
                         "note": f"Only {len(text)} chars of text - appears blank" if is_blank else "OK"})
    blank_count = sum(1 for r in results if not r["ok"])
    return results, blank_count


def check_pagespeed(domain, api_key=None):
    """Fetch real PageSpeed Insights data via the PSI API. Requires a real key
    (Keys sheet psi_api_key, synced into CONFIG) - no hardcoded fallback key
    (a real Google API key was previously hardcoded here directly in source,
    committed to a PUBLIC repo - treat it as compromised/rotated, never
    embed a real key in source again). Returns empty/None scores if no key
    is configured, same as any other key-gated feature in this app."""
    import json
    key = (api_key or "").strip()
    if not key:
        return {"mobile": None, "desktop": None}
    base_text = ("The time it takes to fully display the content on a specific page; it reports on the "
                 "performance of a page on both mobile and desktop devices. ")
    scores = {}
    for strategy in ("mobile", "desktop"):
        api_url = (f"https://www.googleapis.com/pagespeedonline/v5/runPagespeed"
                   f"?url=https://{domain}/&strategy={strategy}"
                   f"&category=PERFORMANCE&category=SEO&category=BEST_PRACTICES&category=ACCESSIBILITY"
                   f"&key={key}")
        try:
            req = _ur.Request(api_url, headers={"User-Agent": _UA})
            with _ur.urlopen(req, timeout=90) as r:
                data = json.loads(r.read().decode("utf-8"))
            cats = data.get("lighthouseResult", {}).get("categories", {})
            scores[strategy] = {
                "performance": int(cats.get("performance", {}).get("score", 0) * 100),
                "seo": int(cats.get("seo", {}).get("score", 0) * 100),
                "best_practices": int(cats.get("best-practices", {}).get("score", 0) * 100),
                "accessibility": int(cats.get("accessibility", {}).get("score", 0) * 100),
            }
        except Exception:
            scores[strategy] = None

    parts = [base_text]
    for strategy in ("mobile", "desktop"):
        s = scores.get(strategy)
        if s:
            parts.append(f"{strategy.title()}: Performance {s['performance']}/100, "
                         f"SEO {s['seo']}/100, Best Practices {s['best_practices']}/100, "
                         f"Accessibility {s['accessibility']}/100.")
        else:
            parts.append(f"{strategy.title()}: Could not retrieve scores.")

    return {"summary": " ".join(parts), "scores": scores}


# ---------- Table image renderer (PIL) ----------
def _render_table_image(title, headers, rows, col_widths, path):
    from PIL import Image, ImageDraw, ImageFont
    try:
        bold_font = ImageFont.truetype("C:/Windows/Fonts/arialbd.ttf", 20)
        header_font = ImageFont.truetype("C:/Windows/Fonts/arialbd.ttf", 16)
        cell_font = ImageFont.truetype("C:/Windows/Fonts/arial.ttf", 15)
        title_font = ImageFont.truetype("C:/Windows/Fonts/arialbd.ttf", 24)
    except Exception:
        bold_font = header_font = cell_font = title_font = ImageFont.load_default()

    W = sum(col_widths) + 2
    ROW_H = 28
    TITLE_H = 50
    HEADER_H = 32
    H = TITLE_H + HEADER_H + ROW_H * max(len(rows), 1) + 20
    img = Image.new("RGB", (W, H), "#FFFFFF")
    d = ImageDraw.Draw(img)
    d.text((16, 14), title, fill=(31, 41, 55), font=title_font)
    y = TITLE_H
    x = 0
    d.rectangle([0, y, W, y + HEADER_H], fill="#F3F4F6")
    for i, h in enumerate(headers):
        d.text((x + 8, y + 8), h, fill=(75, 85, 99), font=header_font)
        x += col_widths[i]
        if i < len(headers) - 1:
            d.line([(x, y), (x, y + HEADER_H)], fill="#D1D5DB", width=1)
    d.line([(0, y + HEADER_H), (W, y + HEADER_H)], fill="#D1D5DB", width=1)
    for ri, row in enumerate(rows):
        y = TITLE_H + HEADER_H + ri * ROW_H
        bg = "#F9FAFB" if ri % 2 == 0 else "#FFFFFF"
        d.rectangle([0, y, W, y + ROW_H], fill=bg)
        x = 0
        ok = row[-1] if isinstance(row[-1], bool) else True
        for ci, cell in enumerate(row[:-1]):
            color = (31, 41, 55)
            if ci == len(row) - 2:
                color = (22, 128, 67) if ok else (200, 40, 60)
            text = str(cell)
            max_chars = col_widths[ci] // 8
            if len(text) > max_chars:
                text = text[:max_chars - 1] + "..."
            d.text((x + 8, y + 6), text, fill=color, font=cell_font)
            x += col_widths[ci]
            if ci < len(row) - 2:
                d.line([(x, y), (x, y + ROW_H)], fill="#E5E7EB", width=1)
        d.line([(0, y + ROW_H), (W, y + ROW_H)], fill="#E5E7EB", width=1)
    if not rows:
        d.text((16, TITLE_H + HEADER_H + 8), "No data.", fill=(150, 150, 150), font=cell_font)
    img.save(str(path))


def _render_broken_links_image(domain, checked, broken, path):
    """`broken` is a list of dicts: {"url","code","found_on","suggested_redirect"}
    (older tuple form (url, code) is still accepted for compatibility)."""
    from PIL import Image, ImageDraw, ImageFont
    rows = [b if isinstance(b, dict) else {"url": b[0], "code": b[1], "found_on": [], "suggested_redirect": None}
            for b in broken[:25]]
    ROW_H = 46
    W = 1180
    H = 130 + (len(rows) + 1) * ROW_H + 30
    img = Image.new("RGB", (W, max(H, 190)), "#FFFFFF")
    d = ImageDraw.Draw(img)
    try:
        bold = ImageFont.truetype("C:/Windows/Fonts/arialbd.ttf", 26)
        f = ImageFont.truetype("C:/Windows/Fonts/arial.ttf", 18)
        fs = ImageFont.truetype("C:/Windows/Fonts/arial.ttf", 15)
    except Exception:
        bold = f = fs = ImageFont.load_default()
    d.text((26, 24), f"Broken Links Check - {domain}", fill=(31, 41, 55), font=bold)
    if not broken:
        d.text((26, 74), f"Checked {checked} links - No broken links found.", fill=(22, 128, 67), font=f)
    else:
        d.text((26, 74), f"Checked {checked} links - {len(broken)} broken:", fill=(200, 40, 60), font=f)
        y = 116
        d.text((26, y), "Broken URL / Status", fill=(120, 120, 120), font=fs)
        d.text((26, y + 22), "Found On / Suggested Redirect", fill=(120, 120, 120), font=fs)
        y += 28 + 22
        for row in rows:
            u, code = row["url"], row["code"]
            found_on = ", ".join(row.get("found_on") or []) or "-"
            redirect = row.get("suggested_redirect")
            d.text((26, y), f"{u[:110]}", fill=(40, 40, 40), font=fs)
            d.text((W - 90, y), str(code), fill=(200, 40, 60), font=fs)
            y += 22
            note = f"Found on: {found_on[:90]}"
            if redirect:
                note += f"   ->  Suggested redirect: {redirect[:90]}"
            d.text((26, y), note, fill=(90, 90, 90), font=fs)
            y += ROW_H - 22
    img.save(str(path))


# ---------- Checkpoint definitions ----------
INTRO = [
    ("Website health checkup and analysis",
     ": As a process of routine analysis, we have checked the website "
     "according to different points in order to determine if any issue "
     "persists in the website or not."),
    ("",
     "And after doing the analysis, we found that there is no health related "
     "issue in the website."),
]

CHECKPOINTS = [
    {"key": "sucuri", "label": "Sucuri Check",
     "body": ": We first of all check the website in the tool Sucuri, to check "
             "for any malware or website blacklisting issue.",
     "extra": ("URL:  ", "https://sitecheck.sucuri.net/results/https/{domain}"),
     "status": ("Status", ": We have not find issue in the website."),
     "capture": "url"},

    {"key": "manual_action", "label": "Manual Action:",
     "body": " The Manual Actions report lists manually detected issues with a "
             "page or site that are mostly attempts to manipulate our search "
             "index, but are not necessarily dangerous for users.",
     "status": ("Status -", " Not found"),
     "capture": "gsc"},

    {"key": "security_issues", "label": "Security Issues:",
     "body": " The Security Issues report lists indications that your site was "
             "hacked, or behavior on your site that could potentially harm a "
             "visitor or their computer: for example, phishing attacks or "
             "installing malware or unwanted software on the user's computer.",
     "status": ("Status:  ", "Not Found"),
     "capture": "gsc"},

    {"key": "robots", "label": "Robots.txt:",
     "body": " Robots.txt file has pages that needs to be blocked and are not "
             "useful for the website. Thus, we have checked the robots file is "
             "running fine for now.",
     "capture": "url"},

    {"key": "sitemap", "label": "Sitemap.xml: ",
     "body": "Sitemap file must have all the web-pages that are being present "
             "in the website so that all the pages can be identified by the "
             "Search Engine. We have checked the sitemap file as it is having "
             "all the relevant pages or not in the website.",
     "capture": "url"},

    {"key": "status200", "label": "All Target page's status is 200:- ",
     "body": "{status200_result}", "capture": "computed"},

    {"key": "versions", "label": "The website is running with multiple versions: - ",
     "body": "{versions_result}", "capture": "computed"},

    {"key": "meta_source", "label": "Meta Suggestions:- ",
     "body": "{meta_source_result}", "capture": "url"},

    {"key": "canonical", "label": "Right Canonical tag: - ",
     "body": "{canonical_result}", "capture": "computed"},

    {"key": "double_meta", "label": "Check for Double Meta Tags: - ",
     "body": "{double_meta_result}", "capture": "computed"},

    {"key": "meta_robots", "label": "Check Meta Robots Tag on Web-Pages: - ",
     "body": "{meta_robots_result}", "capture": "computed"},

    {"key": "layout", "label": "Check Website Layout: -  ",
     "body": "The screenshot below is the latest archived snapshot from web.archive.org. "
             "Compare it with the current live site to confirm no unexpected layout changes "
             "have occurred recently.",
     "capture": "url"},

    {"key": "dummy", "label": "Check Dummy Content: ",
     "body": "{dummy_content_result}", "capture": "url"},

    {"key": "broken_links", "label": "Check for Broken Links:- ",
     "body": "{broken_links_result}", "capture": "computed"},

    {"key": "serp", "label": "Any un-necessary SERP issues: - ",
     "body": "There is no unnecessary pages found in the SERP",
     "capture": "url"},

    {"key": "blank_page", "label": "Check Blank Page: ",
     "body": "We check our all-targeted page and didn't find any blank targeted page.",
     "capture": "url"},

    {"key": "pagespeed", "label": "Website Page speed:- ",
     "body": "{pagespeed_result}", "capture": "computed"},
]

CHECKPOINT_BY_KEY = {cp["key"]: cp for cp in CHECKPOINTS}

JAMES_KEYS = [cp["key"] for cp in CHECKPOINTS]

OMEGA_KEYS = [
    "sucuri", "manual_action", "security_issues", "robots", "sitemap", "serp", "double_meta",
    "blank_page", "dummy", "meta_source", "layout", "versions", "status200",
]

NEON_KEYS = [
    # Neon sequence: SEO-first, security + performance last (distinct from James)
    "robots", "sitemap", "meta_source", "canonical", "meta_robots",
    "serp", "status200", "versions", "layout", "blank_page",
    "broken_links", "sucuri", "manual_action", "security_issues", "pagespeed",
]

SIGMA_KEYS = [
    "sucuri", "manual_action", "security_issues", "robots", "sitemap", "status200",
    "versions", "serp", "meta_source", "double_meta", "layout",
    "blank_page", "dummy", "pagespeed",
]

XENONK_KEYS = [
    # Xenon K sequence: exactly JAMES_KEYS minus "broken_links" and "pagespeed"
    # (the reference report has no checkpoint for either), order preserved.
    k for k in JAMES_KEYS if k not in ("broken_links", "pagespeed")
]

FORMAT_INFO = {
    "james":  {"keys": JAMES_KEYS,   "ext": "docx", "label": "James (Full DOCX)"},
    "sigma":  {"keys": SIGMA_KEYS,   "ext": "pptx", "label": "Sigma (PPTX)"},
    "omega":  {"keys": OMEGA_KEYS,   "ext": "pptx", "label": "Omega (PPTX)"},
    "neon":   {"keys": NEON_KEYS,    "ext": "docx", "label": "Neon (DOCX)"},
    "xenonk": {"keys": XENONK_KEYS,  "ext": "docx", "label": "Xenon K"},
}

# Screenshot URLs for each checkpoint
SCREENSHOT_URLS = {
    "sucuri": "https://sitecheck.sucuri.net/results/https/{domain}",
    "robots": "https://{domain}/robots.txt",
    "sitemap": "https://{domain}/sitemap.xml",
    "meta_source": "https://{domain}/",
    "layout": "https://web.archive.org/web/2999/https://{domain}/",
    "dummy": "https://www.google.com/search?q=site:{domain}+lorem",
    "serp": "https://www.google.com/search?q=site:{domain}",
    "blank_page": "https://{domain}/",
}


def _highlight_source_snippet(driver, patterns):
    """Fetch the page's REAL raw HTML (same-origin fetch of its own URL - the
    actual bytes the server sent) and render a genuine view-source-style
    widget: a fake address bar reading "view-source:<url>", real line
    numbers, real surrounding source lines, and every line matching one of
    `patterns` highlighted - matching what a real, zoomed view-source:
    screenshot looks like, instead of a plain colored banner of extracted
    values. Contiguous/nearby matches are shown as one block; matches far
    apart in the source get their own block with a "..." gap between them.
    Returns the widget's on-page bounding box (for a tight screenshot crop),
    or None if the fetch/render failed."""
    import time
    script = """
        var callback = arguments[arguments.length - 1];
        var patterns = arguments[0];
        var old = document.getElementById('__source_highlight_widget');
        if (old) old.remove();
        fetch(location.href).then(function(r){ return r.text(); }).then(function(html){
            var lines = html.split(/\\r\\n|\\r|\\n/);
            var matchIdx = [];
            for (var p = 0; p < patterns.length; p++) {
                var re = new RegExp(patterns[p], 'i');
                for (var i = 0; i < lines.length; i++) {
                    if (re.test(lines[i])) { matchIdx.push(i); break; }
                }
            }
            matchIdx.sort(function(a, b){ return a - b; });
            var esc = function(s){ return String(s).replace(/&/g,'&amp;').replace(/</g,'&lt;').replace(/>/g,'&gt;'); };

            var box = document.createElement('div');
            box.id = '__source_highlight_widget';
            box.style.cssText = 'position:fixed;top:20px;left:20px;z-index:2147483647;'
                + 'background:#fff;border:1px solid #ccc;border-radius:6px;'
                + 'box-shadow:0 4px 16px rgba(0,0,0,.25);font:15px/1.6 Consolas,Menlo,monospace;'
                + 'overflow:hidden;max-width:900px;';

            var addr = document.createElement('div');
            addr.style.cssText = 'background:#f1f3f4;padding:8px 14px;border-bottom:1px solid #ddd;'
                + 'font:13px Arial,sans-serif;color:#333;';
            addr.textContent = 'view-source:' + location.href;
            box.appendChild(addr);

            var code = document.createElement('div');
            code.style.cssText = 'padding:10px 14px;white-space:pre;';
            var out = '';
            if (!matchIdx.length) {
                out = '<div>(tag not found in page source)</div>';
            } else {
                // Merge match lines that are close together (within 4 lines)
                // into one contiguous block, each padded by 1 line of context.
                var blocks = [];
                var curStart = matchIdx[0], curEnd = matchIdx[0];
                for (var i = 1; i < matchIdx.length; i++) {
                    if (matchIdx[i] - curEnd <= 4) { curEnd = matchIdx[i]; }
                    else { blocks.push([curStart, curEnd]); curStart = curEnd = matchIdx[i]; }
                }
                blocks.push([curStart, curEnd]);
                var matchSet = {};
                for (var i = 0; i < matchIdx.length; i++) { matchSet[matchIdx[i]] = true; }
                for (var b = 0; b < blocks.length; b++) {
                    if (b > 0) { out += '<div style="color:#999;padding:2px 6px 2px 34px;">...</div>'; }
                    var s = Math.max(0, blocks[b][0] - 1);
                    var e = Math.min(lines.length - 1, blocks[b][1] + 1);
                    for (var i = s; i <= e; i++) {
                        var bg = matchSet[i] ? 'background:#f9a825;' : '';
                        out += '<div style="' + bg + 'padding:2px 6px;">'
                            + '<span style="color:#999;display:inline-block;width:28px;">' + (i + 1) + '</span>'
                            + '<span>' + esc(lines[i]) + '</span></div>';
                    }
                }
            }
            code.innerHTML = out;
            box.appendChild(code);

            document.body.insertBefore(box, document.body.firstChild);
            window.scrollTo(0, 0);

            setTimeout(function() {
                var r = box.getBoundingClientRect();
                callback({x: r.x, y: r.y, width: r.width, height: r.height});
            }, 80);
        }).catch(function(){ callback(null); });
    """
    try:
        driver.set_script_timeout(15)
        rect = driver.execute_async_script(script, patterns)
        time.sleep(0.2)
        return rect
    except Exception:
        return None


def _latest_wayback_snapshot_url(domain, log_fn=None):
    """Resolve the actual latest-archived snapshot URL via the Wayback
    Availability API, instead of guessing a timestamp. Returns None if
    the domain has never been archived (caller should skip the
    screenshot rather than capture archive.org's "not found" page)."""
    import json as _json
    log_fn = log_fn or print
    try:
        req = _ur.Request(
            f"https://archive.org/wayback/available?url=https://{domain}/",
            headers={"User-Agent": _UA},
        )
        with _ur.urlopen(req, timeout=15) as resp:
            data = _json.loads(resp.read().decode("utf-8", "replace"))
        snap = (data.get("archived_snapshots") or {}).get("closest") or {}
        if snap.get("available") and snap.get("url"):
            return snap["url"]
    except Exception as e:
        log_fn(f"  [warn] layout: Wayback availability lookup failed ({e})")
    return None


# ---------- Selenium screenshot capture ----------
def capture_screenshots_selenium(driver, domain, out_dir, keys, log_fn=None):
    """Capture screenshots using an existing Selenium WebDriver via CDP.
    Returns dict of {key: image_path}."""
    captured = {}
    if log_fn is None:
        log_fn = print

    os.makedirs(out_dir, exist_ok=True)

    for key in keys:
        url_tmpl = SCREENSHOT_URLS.get(key)
        if not url_tmpl:
            continue

        url = url_tmpl.format(domain=domain)
        if key == "layout":
            # The old guessed URL (timestamp "2") resolved to the
            # OLDEST archived snapshot, not the latest as the report
            # caption claims - and if that earliest capture was broken
            # or never completed, the screenshot came out blank.
            # Resolve the real latest snapshot via the Availability API.
            real_url = _latest_wayback_snapshot_url(domain, log_fn)
            if not real_url:
                log_fn("  [warn] layout: no archived snapshot found for this domain, skipping screenshot")
                continue
            url = real_url
        path = os.path.join(out_dir, f"{key}.png")

        try:
            log_fn(f"  Capturing [{key}]...")
            import time
            highlight_rect = None
            try:
                driver.get(url)
            except Exception as e:
                # A slow/redirecting page (archive.org's Wayback lookup is the
                # worst offender) can throw a page-load timeout. Retry once
                # instead of abandoning this checkpoint's screenshot entirely.
                log_fn(f"  [warn] {key} navigation failed ({e}), retrying once...")
                time.sleep(2)
                driver.get(url)
            if key == "sucuri":
                time.sleep(15)
            elif key == "layout":
                time.sleep(8)
            else:
                time.sleep(4)

            if key in ("serp", "dummy"):
                # These checkpoints hit a raw Google search URL, which can
                # trip Google's bot-detection (reCAPTCHA / "unusual traffic")
                # and leave a CAPTCHA baked into the saved screenshot instead
                # of real SERP results. Detect it (reusing engine.py's
                # existing classify_page, already proven for the ranking
                # checker) and retry once with backoff before giving up.
                state = engine.classify_page(driver.page_source)
                if state in ("captcha", "soft_block"):
                    log_fn(f"  [warn] {key}: Google bot-check detected ({state}), retrying after backoff...")
                    time.sleep(20)
                    try:
                        driver.get(url)
                        time.sleep(5)
                    except Exception:
                        pass
                    state = engine.classify_page(driver.page_source)
                    if state in ("captcha", "soft_block"):
                        log_fn(f"  [warn] {key}: still bot-checked after retry, skipping screenshot to avoid capturing a CAPTCHA")
                        continue

            if key == "sucuri":
                # Sucuri shows a cookie-consent bar over the verdict. Click "Accept"
                # so it isn't in the screenshot; then remove any leftover banner.
                try:
                    driver.execute_script("""
                        var els = document.querySelectorAll('button, a, input[type=button], input[type=submit], span, div');
                        for (var i = 0; i < els.length; i++) {
                            var t = (els[i].textContent || els[i].value || '').trim().toLowerCase();
                            if (t === 'accept' || t === 'accept all' || t === 'i accept' ||
                                t === 'allow all' || t === 'allow' || t === 'got it' || t === 'agree') {
                                try { els[i].click(); } catch (e) {}
                                return;
                            }
                        }
                    """)
                    time.sleep(1.2)
                except Exception:
                    pass
                driver.execute_script(
                    "document.querySelectorAll('.cookie-banner,.consent-banner,[class*=cookie],[class*=consent],[id*=cookie],[id*=consent],#cookie-law-info-bar').forEach(function(e){e.remove();});"
                    "window.scrollTo(0, 0);"
                )
                time.sleep(0.6)
            elif key == "meta_source":
                # A view-source: URL used to be navigated here and the
                # screenshot just scrolled near the meta tag - two real
                # problems: view-source: is a privileged scheme some
                # automated/webdriver-controlled browsers restrict or refuse
                # to navigate to at all (confirmed live: "meta suggestions
                # not being found" across every format tested), and a plain
                # DOM-value banner (title/description text in a colored box)
                # doesn't actually show the SOURCE the way a real view-source:
                # screenshot does. This fetches the page's REAL raw HTML (a
                # same-origin fetch of its own URL - the actual bytes the
                # server sent) and renders a genuine view-source-style widget:
                # a fake address bar, real line numbers, real surrounding
                # source lines, with the <title> and meta description lines
                # highlighted - then the screenshot crops tight to just that
                # widget instead of a full-viewport shot.
                highlight_rect = _highlight_source_snippet(
                    driver, [r'<title\b[^>]*>.*?</title>', r'<meta\b[^>]*name=["\']description["\'][^>]*>'])
            else:
                driver.execute_script("window.scrollTo(0, 0);")
            time.sleep(0.5)

            # Use CDP for the screenshot - a plain viewport-only capture (no
            # captureBeyondViewport/clip), so every checkpoint gets a tight
            # crop of what's actually on screen after scrolling, not an
            # oversized capture of the full document. meta_source instead
            # crops tight to the injected source-viewer widget (see above).
            if key == "meta_source" and highlight_rect:
                margin = 12
                cdp_params = {"format": "png", "captureBeyondViewport": False,
                              "clip": {"x": max(0, highlight_rect["x"] - margin),
                                       "y": max(0, highlight_rect["y"] - margin),
                                       "width": highlight_rect["width"] + margin * 2,
                                       "height": highlight_rect["height"] + margin * 2,
                                       "scale": 1}}
            else:
                cdp_params = {"format": "png", "captureBeyondViewport": False}
            try:
                result = driver.execute_cdp_cmd("Page.captureScreenshot", cdp_params)
                import base64
                with open(path, "wb") as f:
                    f.write(base64.b64decode(result["data"]))
                captured[key] = path
                log_fn(f"  -> [{key}] captured")
            except Exception as e:
                log_fn(f"  [warn] CDP screenshot failed for {key}: {e}")
                try:
                    driver.save_screenshot(path)
                    captured[key] = path
                    log_fn(f"  -> [{key}] captured (fallback)")
                except Exception:
                    pass
        except Exception as e:
            log_fn(f"  [warn] capture '{key}' failed: {e}")
            # A hung/crashed renderer (confirmed real case: a Wayback Machine
            # "layout" snapshot timed out mid-render - archived pages can be
            # huge/broken and hang far longer than a normal page) can leave
            # the shared driver session wedged for every screenshot captured
            # AFTER this one too, since they all reuse the same driver -
            # each subsequent driver.get() then also fails, which is exactly
            # why one bad checkpoint was reported as "most of the screenshots
            # were missing" even though nothing else was actually wrong.
            # Force the browser back to a known-good blank page before
            # continuing so a wedged renderer can't poison the rest of the run.
            try:
                driver.get("about:blank")
            except Exception:
                pass

    return captured


# ---------- Prepare all health data ----------
def prepare_health_data(domain, captured=None, target_pages=None, log_fn=None, psi_api_key=None, driver=None, prefetched_futures=None):
    if log_fn is None:
        log_fn = print
    domain = re.sub(r'^\s*https?://', '', str(domain or '')).strip().strip('/').split('/')[0] or str(domain)
    captured = captured or {}
    target_pages = [p.strip() for p in (target_pages or []) if p.strip()]
    _RENDER_CACHE.clear()      # fresh per-run cache of browser-rendered pages

    _tmp = Path(tempfile.gettempdir())

    # --- Start slow checks in background threads (or reuse prefetched futures) ---
    import concurrent.futures
    executor = None

    if prefetched_futures and "sucuri" in prefetched_futures:
        fut_sucuri = prefetched_futures["sucuri"]
        fut_psi = prefetched_futures.get("psi")
        log_fn("Sucuri" + (" & PageSpeed" if fut_psi else "") + " already running from early start...")
    else:
        executor = concurrent.futures.ThreadPoolExecutor(max_workers=2)
        fut_sucuri = executor.submit(check_sucuri, domain)
        if psi_api_key:
            log_fn("Starting Sucuri & PageSpeed checks in background...")
            fut_psi = executor.submit(check_pagespeed, domain, psi_api_key)
        else:
            log_fn("Starting Sucuri in background (PageSpeed skipped)...")
            fut_psi = None

    # --- Fast checks run sequentially while slow ones are in progress ---
    _total = 11  # robots.txt/sitemap.xml/status200 now share one progress step (run in parallel)
    _step = [0]
    def _prog(msg):
        _step[0] += 1
        log_fn(f"[{_step[0]}/{_total}] {msg}")

    # robots.txt, sitemap.xml and the status200 check are each independent HTTP round
    # trips (none of them touch the shared Selenium driver, unlike canonical/double-meta
    # /meta-robots below), so run them concurrently instead of one after another - this
    # is the main reason a Health Audit was slow just to get through its first checks.
    _prog("Checking robots.txt, sitemap.xml & status200 (in parallel)...")
    with concurrent.futures.ThreadPoolExecutor(max_workers=3) as fast_pool:
        fut_robots = fast_pool.submit(check_robots_txt, domain)
        fut_sitemap = fast_pool.submit(check_sitemap, domain)
        fut_s200 = fast_pool.submit(check_status200, target_pages, domain)
        robots = fut_robots.result()
        sitemap = fut_sitemap.result()
        s200 = fut_s200.result()

    _prog("Running canonical check...")
    canon = check_canonical(target_pages, domain, driver=driver)

    _prog("Running double meta check...")
    dmeta = check_double_meta(target_pages, domain, driver=driver)

    _prog("Running meta robots check...")
    mrobots = check_meta_robots(target_pages, domain, driver=driver)

    _prog("Running URL versions check...")
    versions_summary, versions_rows = check_versions_full(domain)

    _prog("Running broken links check...")
    bl_checked, bl_broken = check_broken_links(domain, target_pages)

    _prog("Running dummy content check...")
    dummy_result = check_dummy_content(domain, driver=driver)

    _prog("Checking SERP for spam/hacked content...")
    serp = check_serp(domain, driver=driver)

    _prog("Checking blank pages...")
    blank_results, blank_count = check_blank_pages(target_pages, domain, driver=driver)

    _prog("Checking meta suggestions...")
    meta_sugg = check_meta_suggestions(target_pages, domain, driver=driver)

    # --- Collect slow checks ---
    if fut_psi:
        _prog("Waiting for Sucuri & PageSpeed results...")
    else:
        _prog("Waiting for Sucuri results...")
    sucuri = fut_sucuri.result(timeout=120)
    psi = fut_psi.result(timeout=120) if fut_psi else {"summary": "PageSpeed Insights - skipped", "scores": {}}
    if executor:
        executor.shutdown(wait=False)
    log_fn("Sucuri" + (" & PageSpeed" if fut_psi else "") + " done.")

    # Render table images
    if s200:
        p = _tmp / f"ha_status200_{domain}.png"
        rows = [(r["url"].replace("https://", "").replace("http://", ""),
                 r["status"], "OK" if r["ok"] else "Issue", r["ok"]) for r in s200]
        _render_table_image(f"HTTP Status Check - {domain}", ["Page URL", "Status", "Result"],
                            rows, [780, 120, 140], p)
        captured["status200"] = str(p)

    if canon:
        p = _tmp / f"ha_canonical_{domain}.png"
        rows = [(r["url"],
                 r["canonical"][:60] + "..." if len(r["canonical"]) > 60 else r["canonical"],
                 r["note"], r["ok"]) for r in canon]
        _render_table_image(f"Canonical Tag Check - {domain}", ["Page URL", "Canonical", "Result"],
                            rows, [540, 360, 180], p)
        captured["canonical"] = str(p)

    if psi.get("scores"):
        p = _tmp / f"ha_pagespeed_{domain}.png"
        rows = []
        for strategy in ("mobile", "desktop"):
            sc = psi["scores"].get(strategy)
            if sc:
                rows.append((strategy.title(), f"{sc['performance']}/100", f"{sc['seo']}/100",
                            f"{sc['best_practices']}/100", f"{sc['accessibility']}/100",
                            sc['performance'] >= 50))
        if rows:
            _render_table_image(f"PageSpeed Insights - {domain}",
                                ["Device", "Performance", "SEO", "Best Practices", "Accessibility"],
                                rows, [180, 220, 180, 260, 220], p)
            captured["pagespeed"] = str(p)

    if dmeta:
        p = _tmp / f"ha_doublemeta_{domain}.png"
        rows = [(r["url"].replace("https://", "").replace("http://", ""),
                 r["title_count"], r["desc_count"], r["note"], r["ok"]) for r in dmeta]
        _render_table_image(f"Meta Tags Check - {domain}", ["Page URL", "Title", "Meta Desc", "Result"],
                            rows, [600, 90, 110, 240], p)
        captured["double_meta"] = str(p)

    if mrobots:
        p = _tmp / f"ha_metarobots_{domain}.png"
        rows = [(r["url"].replace("https://", "").replace("http://", ""),
                 r["robots_value"], r["note"], r["ok"]) for r in mrobots]
        _render_table_image(f"Meta Robots Check - {domain}", ["Page URL", "Robots Value", "Result"],
                            rows, [540, 360, 180], p)
        captured["meta_robots"] = str(p)

    if versions_rows:
        p = _tmp / f"ha_versions_{domain}.png"
        _render_table_image(f"URL Versions Check - {domain}", ["URL Variant", "HTTP", "Result"],
                            versions_rows, [620, 90, 360], p)
        captured["versions"] = str(p)

    if bl_checked > 0:
        p = _tmp / f"ha_brokenlinks_{domain}.png"
        _render_broken_links_image(domain, bl_checked, bl_broken, p)
        captured["broken_links"] = str(p)

    # Render blank page table if we have results
    if blank_results:
        p = _tmp / f"ha_blankpage_{domain}.png"
        rows = [(r["url"].replace("https://", "").replace("http://", ""),
                 str(r.get("text_len", "N/A")), r["note"], r["ok"]) for r in blank_results]
        _render_table_image(f"Blank Page Check - {domain}", ["Page URL", "Text Length", "Result"],
                            rows, [700, 120, 220], p)
        captured["blank_page"] = str(p)

    # Render SERP spam table if spam found
    if serp.get("spam_found"):
        p = _tmp / f"ha_serp_spam_{domain}.png"
        rows = [(s["type"], str(s["count"]),
                 s["example_titles"][0][:60] + "..." if s["example_titles"] else "N/A", False)
                for s in serp["spam_found"]]
        _render_table_image(f"SERP Spam/Hack Detection - {domain}", ["Spam Type", "Hits", "Example Title"],
                            rows, [300, 80, 660], p)
        captured["serp"] = str(p)

    # Build summaries
    def _bl_summary():
        if bl_checked == 0:
            return "Could not check for broken links."
        if not bl_broken:
            return f"Checked {bl_checked} links - no broken links found."
        return f"Checked {bl_checked} links - {len(bl_broken)} broken link(s) found."

    def _s200_summary():
        if not s200:
            return "No target pages provided."
        bad = [r for r in s200 if not r["ok"]]
        if not bad:
            return f"All {len(s200)} page(s) returned HTTP 200."
        return f"{len(bad)} of {len(s200)} page(s) have non-200 status."

    def _canon_summary():
        if not canon:
            return "No pages checked."
        bad = [r for r in canon if not r["ok"]]
        if not bad:
            return f"All {len(canon)} page(s) have correct self-referencing canonical tags."
        return f"{len(bad)} of {len(canon)} page(s) have canonical issues."

    def _dmeta_summary():
        if not dmeta:
            return "No pages checked."
        bad = [r for r in dmeta if not r["ok"]]
        if not bad:
            return f"All {len(dmeta)} page(s) have title and meta description present once each."
        return f"{len(bad)} of {len(dmeta)} page(s) have meta tag issues."

    def _mrobots_summary():
        if not mrobots:
            return "No pages checked."
        bad = [r for r in mrobots if not r["ok"]]
        if not bad:
            return f"All {len(mrobots)} page(s) have correct meta robots (index, follow)."
        return f"{len(bad)} of {len(mrobots)} page(s) have robots tag issues."

    def _blank_summary():
        if not blank_results:
            return "No target pages provided for blank page check."
        if blank_count == 0:
            return f"All {len(blank_results)} target page(s) have content - no blank pages found."
        return f"{blank_count} of {len(blank_results)} target page(s) appear blank (less than 100 chars of text)."

    def _meta_sugg_summary():
        if not meta_sugg:
            return "No pages checked."
        bad = [r for r in meta_sugg if not r["ok"]]
        if not bad:
            return f"All {len(meta_sugg)} page(s) have a title and meta description. See screenshot below."
        return f"{len(bad)} of {len(meta_sugg)} page(s) are missing a title or meta description."

    computed = {
        "sucuri_result": sucuri["summary"] + (" Details: " + "; ".join(sucuri["details"]) if sucuri.get("details") else ""),
        "robots_result": robots["summary"],
        "sitemap_result": sitemap["summary"],
        "versions_result": versions_summary,
        "status200_result": _s200_summary(),
        "canonical_result": _canon_summary(),
        "double_meta_result": _dmeta_summary(),
        "meta_robots_result": _mrobots_summary(),
        "meta_source_result": _meta_sugg_summary(),
        "broken_links_result": _bl_summary(),
        "dummy_content_result": dummy_result,
        "serp_result": serp["summary"],
        "blank_page_result": _blank_summary(),
        "pagespeed_result": psi["summary"],
    }
    return captured, computed


# ---------- Safe save path ----------
def _safe_save_path(base_path):
    p = Path(base_path)
    if not p.exists():
        return p
    stem, suffix = p.stem, p.suffix
    parent = p.parent
    n = 1
    while True:
        candidate = parent / f"{stem} ({n}){suffix}"
        if not candidate.exists():
            return candidate
        n += 1


# ---------- DOCX builder ----------
def build_health_docx(domain, captured, computed, out_dir, include_keys=None, voice="we", page_border=False, header_fill=None):
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml import parse_xml, OxmlElement
    from docx.oxml.ns import nsdecls, qn
    from docx.shared import Inches, Pt, RGBColor

    FONT_NAME = "Candara"
    IMAGE_WIDTH_IN = 6.3
    IMAGE_MAX_H_IN = 4.6

    def _set_cell_border(cell, sz=24, color="111111"):
        """Solid border on all 4 sides of a table cell. `sz` is in EIGHTHS of a point
        (3pt = 24). A bordered INLINE picture gets its top edge clipped by Word's line
        box until the image is resized by hand - a table cell grows to fit its image,
        so the border is never clipped."""
        tcPr = cell._tc.get_or_add_tcPr()
        for old in tcPr.findall(qn('w:tcBorders')):
            tcPr.remove(old)
        borders = OxmlElement('w:tcBorders')
        for edge in ('top', 'left', 'bottom', 'right'):
            el = OxmlElement('w:' + edge)
            el.set(qn('w:val'), 'single')
            el.set(qn('w:sz'), str(sz))
            el.set(qn('w:space'), '0')
            el.set(qn('w:color'), color)
            borders.append(el)
        tcPr.append(borders)

    def _zero_cell_margins(cell):
        tcPr = cell._tc.get_or_add_tcPr()
        mar = OxmlElement('w:tcMar')
        for edge in ('top', 'start', 'bottom', 'end'):
            el = OxmlElement('w:' + edge)
            el.set(qn('w:w'), '0')
            el.set(qn('w:type'), 'dxa')
            mar.append(el)
        tcPr.append(mar)

    def _add_bordered_picture(img_path):
        """Screenshot framed with a 3pt border via a single-cell table (matches the
        reference report exactly) - a bordered inline picture gets clipped by Word
        until resized by hand, so we frame with a table cell border instead."""
        from docx.enum.table import WD_TABLE_ALIGNMENT
        from PIL import Image as PILImage
        with PILImage.open(img_path) as im:
            iw, ih = im.size
        if (IMAGE_WIDTH_IN * ih / iw) > IMAGE_MAX_H_IN:
            disp_w = IMAGE_MAX_H_IN * iw / ih
            pic_kwargs = {"height": Inches(IMAGE_MAX_H_IN)}
        else:
            disp_w = IMAGE_WIDTH_IN
            pic_kwargs = {"width": Inches(IMAGE_WIDTH_IN)}

        table = doc.add_table(rows=1, cols=1)
        table.alignment = WD_TABLE_ALIGNMENT.LEFT
        table.autofit = False
        table.allow_autofit = False
        col_w = Inches(disp_w)
        table.columns[0].width = col_w
        cell = table.cell(0, 0)
        cell.width = col_w
        _zero_cell_margins(cell)
        _set_cell_border(cell, sz=24, color="111111")

        cp = cell.paragraphs[0]
        cp.paragraph_format.space_before = Pt(0)
        cp.paragraph_format.space_after = Pt(0)
        cp.paragraph_format.line_spacing = 1.0
        run = cp.add_run()
        run.add_picture(img_path, **pic_kwargs)

        # Small gap after the framed image - also gives Word the paragraph it
        # needs to render cleanly after a table.
        spacer = doc.add_paragraph()
        spacer.add_run("").font.size = Pt(6)
        spacer.paragraph_format.space_after = Pt(6)
        spacer.paragraph_format.space_before = Pt(0)

    use_keys = include_keys or JAMES_KEYS
    cps = [CHECKPOINT_BY_KEY[k] for k in use_keys if k in CHECKPOINT_BY_KEY]
    doc = Document()

    for sname in ("Normal", "List Paragraph"):
        try:
            st = doc.styles[sname]
            st.font.name = FONT_NAME
            st.font.size = Pt(11)
        except KeyError:
            pass

    sec = doc.sections[0]
    sec.page_width, sec.page_height = Inches(8.5), Inches(11.0)
    sec.left_margin = sec.right_margin = Inches(1.0)
    sec.top_margin = sec.bottom_margin = Inches(1.0)

    if page_border:
        sectPr = sec._sectPr
        pgBorders = parse_xml(
            '<w:pgBorders xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"'
            ' w:offsetFrom="page">'
            '<w:top w:val="single" w:sz="12" w:space="24" w:color="000000"/>'
            '<w:left w:val="single" w:sz="12" w:space="24" w:color="000000"/>'
            '<w:bottom w:val="single" w:sz="12" w:space="24" w:color="000000"/>'
            '<w:right w:val="single" w:sz="12" w:space="24" w:color="000000"/>'
            '</w:pgBorders>'
        )
        pgSz = sectPr.find('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}pgSz')
        if pgSz is not None:
            pgSz.addnext(pgBorders)
        else:
            sectPr.append(pgBorders)

    # Header
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(domain)
    run.bold = True
    run.font.name = FONT_NAME
    run.font.size = Pt(24)

    def _voice(text):
        if voice == "i":
            return (text.replace("We first of all check", "I first of all check")
                    .replace("we have checked", "I have checked")
                    .replace("We have checked", "I have checked")
                    .replace("We have not", "I have not")
                    .replace("We check ", "I check ")
                    .replace("we found", "I found"))
        return text

    def _add_label_body(label, body):
        p = doc.add_paragraph(style="List Paragraph")
        p.paragraph_format.left_indent = Inches(0)
        if label:
            r = p.add_run(label)
            r.font.name = FONT_NAME
            r.font.size = Pt(11)
            r.bold = True
        if body:
            r = p.add_run(body)
            r.font.name = FONT_NAME
            r.font.size = Pt(11)
        return p

    def _shaded_header(text, fill):
        """A full-width shaded section-header bar (white bold) - used by Neon
        so its sections read as coloured bars instead of plain bold labels."""
        from docx.oxml import OxmlElement
        from docx.oxml.ns import qn
        p = doc.add_paragraph()
        p.paragraph_format.space_before = Pt(10)
        p.paragraph_format.space_after = Pt(4)
        r = p.add_run(text)
        r.font.name = FONT_NAME
        r.font.size = Pt(13)
        r.bold = True
        r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        shd = OxmlElement('w:shd')
        shd.set(qn('w:fill'), fill)
        shd.set(qn('w:val'), 'clear')
        p._p.get_or_add_pPr().append(shd)
        return p

    # Intro
    for label, body in INTRO:
        _add_label_body(label, _voice(body))
    doc.add_paragraph()

    placed = missing = 0
    for n, cp in enumerate(cps, start=1):
        body = cp["body"]
        for k, v in computed.items():
            body = body.replace("{" + k + "}", v)
        try:
            body = body.format(domain=domain)
        except Exception:
            pass
        if header_fill:
            # Neon style: coloured header bar with the section title, body below.
            _shaded_header(f"{n}. {cp['label'].strip().rstrip(':').strip()}", header_fill)
            btext = _voice(body)
            if btext.strip():
                bp = _add_label_body("", btext)
                if cp.get("red"):
                    for r in bp.runs:
                        r.font.color.rgb = RGBColor(0xCC, 0x00, 0x00)
                        r.bold = True
        else:
            p = _add_label_body(f"{n}. {cp['label']}", _voice(body))
            if cp.get("red") and p.runs:
                for r in p.runs:
                    if not r.bold:
                        r.font.color.rgb = RGBColor(0xCC, 0x00, 0x00)
                        r.bold = True

        extra = cp.get("extra")
        if extra:
            p = doc.add_paragraph(style="List Paragraph")
            p.paragraph_format.left_indent = Inches(0)
            r = p.add_run(extra[0])
            r.font.name = FONT_NAME
            r.font.size = Pt(11)
            r.bold = True
            r = p.add_run(extra[1].format(domain=domain))
            r.font.name = FONT_NAME
            r.font.size = Pt(11)

        # Prefer the REAL status capture_gsc_screenshots() actually detected
        # (No Issues Detected / Issue Found / could-not-check-no-access) over
        # the static CHECKPOINTS placeholder text - previously every report
        # showed the exact same hardcoded "Status - Not found" regardless of
        # what was actually found, since this line never looked at the real
        # detection result at all (confirmed real case: a domain GSC wasn't
        # even connected for still showed "Status - Not found", identical to
        # what a genuinely clean site would show).
        real_status = (captured.get("_statuses") or {}).get(cp["key"])
        status = cp.get("status")
        if real_status or status:
            p = doc.add_paragraph(style="List Paragraph")
            p.paragraph_format.left_indent = Inches(0)
            r = p.add_run(status[0] if status else "Status: ")
            r.font.name = FONT_NAME
            r.font.size = Pt(11)
            r.bold = True
            r = p.add_run(real_status if real_status else status[1])
            r.font.name = FONT_NAME
            r.font.size = Pt(11)

        # Image
        img = captured.get(cp["key"])
        if img and os.path.exists(img):
            try:
                _add_bordered_picture(img)
                placed += 1
            except Exception:
                missing += 1
        else:
            missing += 1

        if cp["key"] == "layout" and img and os.path.exists(img):
            p = doc.add_paragraph()
            r = p.add_run("Compare this archived snapshot with the current live site.")
            r.font.name = FONT_NAME
            r.font.size = Pt(11)
            r.font.color.rgb = RGBColor(0xCC, 0x00, 0x00)
            r.bold = True

        doc.add_paragraph()

    gen_date = datetime.now().strftime("%d-%B-%Y")
    out_path = Path(out_dir) / f"{domain} Website health checkup analysis report {gen_date}.docx"
    out_path = _safe_save_path(out_path)
    doc.save(str(out_path))
    return str(out_path), placed, missing


# ---------- DOCX builder (Xenon K) ----------
def build_health_docx_xenonk(domain, captured, computed, out_dir, include_keys=None):
    """Dedicated builder matching the "Xenon K" client reference exactly:
    a navy (1F3864) title banner with light-blue (BFD3EC) centred domain text and a
    blue (2E75B6) rule underneath, a very light blue (D9E2F3) page background tint,
    a thin auto-colour page border, Calibri body text (10.5pt, #2A2A2A - the
    reference's Word doc-default, left untouched rather than forced), navy
    (1F3864) bold checkpoint labels, and blue (2E75B6) bold intro label - all
    confirmed byte-for-byte against
    "Xenon format -Health Audit Report.docx" (word/document.xml + styles.xml).

    Screenshots are framed with the same 6.3in-wide, 3pt/#111111-bordered single-cell
    table used by build_health_docx's _add_bordered_picture (that table's column
    width - 9072 dxa - is an exact match for the reference's own screenshot tables).
    Two reference images (the Sucuri screenshot and the Wayback "Check Website
    Layout" screenshot) are instead a raw inline picture with a picture-level
    3pt border + drop shadow baked into the DrawingML; visually near-identical to
    the bordered-table treatment used everywhere else, so - to keep one proven
    code path - this builder frames those two the same bordered-table way as the
    rest rather than hand-building that DrawingML border+shadow XML. The only
    visible difference from the reference is the missing drop-shadow on those 2
    images.
    """
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml import parse_xml, OxmlElement
    from docx.oxml.ns import nsdecls, qn
    from docx.shared import Inches, Pt, RGBColor

    FONT_NAME = "Calibri"
    BODY_COLOR = RGBColor(0x2A, 0x2A, 0x2A)
    LABEL_COLOR = RGBColor(0x1F, 0x38, 0x64)
    INTRO_LABEL_COLOR = RGBColor(0x2E, 0x75, 0xB6)
    TITLE_TEXT_COLOR = RGBColor(0xBF, 0xD3, 0xEC)
    TITLE_BG = "1F3864"
    TITLE_RULE_COLOR = "2E75B6"
    PAGE_BG = "D9E2F3"
    IMAGE_WIDTH_IN = 6.3
    IMAGE_MAX_H_IN = 4.6

    def _set_cell_border(cell, sz=24, color="111111"):
        tcPr = cell._tc.get_or_add_tcPr()
        for old in tcPr.findall(qn('w:tcBorders')):
            tcPr.remove(old)
        borders = OxmlElement('w:tcBorders')
        for edge in ('top', 'left', 'bottom', 'right'):
            el = OxmlElement('w:' + edge)
            el.set(qn('w:val'), 'single')
            el.set(qn('w:sz'), str(sz))
            el.set(qn('w:space'), '0')
            el.set(qn('w:color'), color)
            borders.append(el)
        tcPr.append(borders)

    def _zero_cell_margins(cell):
        tcPr = cell._tc.get_or_add_tcPr()
        mar = OxmlElement('w:tcMar')
        for edge in ('top', 'start', 'bottom', 'end'):
            el = OxmlElement('w:' + edge)
            el.set(qn('w:w'), '0')
            el.set(qn('w:type'), 'dxa')
            mar.append(el)
        tcPr.append(mar)

    def _add_bordered_picture(img_path):
        """Screenshot framed with a 3pt/#111111 border via a single-cell table -
        matches the reference report's screenshot tables exactly (9072 dxa = 6.3in)."""
        from docx.enum.table import WD_TABLE_ALIGNMENT
        from PIL import Image as PILImage
        with PILImage.open(img_path) as im:
            iw, ih = im.size
        if (IMAGE_WIDTH_IN * ih / iw) > IMAGE_MAX_H_IN:
            disp_w = IMAGE_MAX_H_IN * iw / ih
            pic_kwargs = {"height": Inches(IMAGE_MAX_H_IN)}
        else:
            disp_w = IMAGE_WIDTH_IN
            pic_kwargs = {"width": Inches(IMAGE_WIDTH_IN)}

        table = doc.add_table(rows=1, cols=1)
        table.alignment = WD_TABLE_ALIGNMENT.LEFT
        table.autofit = False
        table.allow_autofit = False
        col_w = Inches(disp_w)
        table.columns[0].width = col_w
        cell = table.cell(0, 0)
        cell.width = col_w
        _zero_cell_margins(cell)
        _set_cell_border(cell, sz=24, color="111111")

        cp = cell.paragraphs[0]
        cp.paragraph_format.space_before = Pt(0)
        cp.paragraph_format.space_after = Pt(0)
        cp.paragraph_format.line_spacing = 1.0
        run = cp.add_run()
        run.add_picture(img_path, **pic_kwargs)

        spacer = doc.add_paragraph()
        spacer.add_run("").font.size = Pt(6)
        spacer.paragraph_format.space_after = Pt(6)
        spacer.paragraph_format.space_before = Pt(0)

    use_keys = include_keys or XENONK_KEYS
    cps = [CHECKPOINT_BY_KEY[k] for k in use_keys if k in CHECKPOINT_BY_KEY]
    doc = Document()

    # Reference doc-default run properties (word/styles.xml docDefaults):
    # Calibri, 10.5pt (sz 21 half-points), color #2A2A2A. Left un-set at the
    # style level too (matches reference, whose "Normal"/"List Paragraph"
    # style entries carry no explicit rPr/pPr of their own).
    for sname in ("Normal", "List Paragraph"):
        try:
            st = doc.styles[sname]
            st.font.name = FONT_NAME
            st.font.size = Pt(10.5)
            st.font.color.rgb = BODY_COLOR
        except KeyError:
            pass

    sec = doc.sections[0]
    sec.page_width, sec.page_height = Inches(8.5), Inches(11.0)
    sec.left_margin = sec.right_margin = Inches(1.0)
    sec.top_margin = sec.bottom_margin = Inches(0.75)

    # Thin auto-colour page border, offset 24pt from the page edge (matches
    # the reference's <w:pgBorders w:offsetFrom="page"> sz=4/color=auto exactly).
    sectPr = sec._sectPr
    pgBorders = parse_xml(
        '<w:pgBorders xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"'
        ' w:offsetFrom="page">'
        '<w:top w:val="single" w:sz="4" w:space="24" w:color="auto"/>'
        '<w:left w:val="single" w:sz="4" w:space="24" w:color="auto"/>'
        '<w:bottom w:val="single" w:sz="4" w:space="24" w:color="auto"/>'
        '<w:right w:val="single" w:sz="4" w:space="24" w:color="auto"/>'
        '</w:pgBorders>'
    )
    pgSz = sectPr.find('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}pgSz')
    if pgSz is not None:
        pgSz.addnext(pgBorders)
    else:
        sectPr.append(pgBorders)

    # Very light blue page background tint (matches reference's <w:background>).
    body_bg = parse_xml(
        f'<w:background {nsdecls("w")} w:color="{PAGE_BG}"/>'
    )
    doc.element.insert(0, body_bg)

    # Title banner: navy fill, centred light-blue 24pt text, blue rule underneath.
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    pPr = p._p.get_or_add_pPr()
    shd = OxmlElement('w:shd')
    shd.set(qn('w:val'), 'clear')
    shd.set(qn('w:color'), 'auto')
    shd.set(qn('w:fill'), TITLE_BG)
    pPr.append(shd)
    pBdr = OxmlElement('w:pBdr')
    bottom = OxmlElement('w:bottom')
    bottom.set(qn('w:val'), 'single')
    bottom.set(qn('w:sz'), '24')
    bottom.set(qn('w:space'), '2')
    bottom.set(qn('w:color'), TITLE_RULE_COLOR)
    pBdr.append(bottom)
    pPr.append(pBdr)
    run = p.add_run(domain)
    run.font.size = Pt(24)
    run.font.color.rgb = TITLE_TEXT_COLOR
    p.paragraph_format.space_after = Pt(10)

    def _add_label_body(label, body, label_color=LABEL_COLOR):
        p = doc.add_paragraph(style="List Paragraph")
        p.paragraph_format.left_indent = Inches(0)
        if label:
            r = p.add_run(label)
            r.bold = True
            r.font.color.rgb = label_color
        if body:
            r = p.add_run(body)
        return p

    # Intro (identical wording/formatting to the reference: bold blue label,
    # default-coloured body, single hard line-break between the two sentences).
    intro_p = doc.add_paragraph()
    ip_fmt = intro_p.paragraph_format
    ip_fmt.left_indent = Inches(0.083)
    ip_fmt.right_indent = Inches(0.083)
    ip_fmt.space_before = Pt(3)
    ip_fmt.space_after = Pt(6)
    ip_fmt.line_spacing = 1.15
    for i, (label, body) in enumerate(INTRO):
        if label:
            r = intro_p.add_run(label)
            r.bold = True
            r.font.color.rgb = INTRO_LABEL_COLOR
        if body:
            if i > 0:
                intro_p.add_run().add_break()
            intro_p.add_run(body)
    doc.add_paragraph()

    placed = missing = 0
    for n, cp in enumerate(cps, start=1):
        body = cp["body"]
        for k, v in computed.items():
            body = body.replace("{" + k + "}", v)
        try:
            body = body.format(domain=domain)
        except Exception:
            pass

        _add_label_body(f"{n}. {cp['label']}", body)

        extra = cp.get("extra")
        if extra:
            p = doc.add_paragraph(style="List Paragraph")
            p.paragraph_format.left_indent = Inches(0)
            r = p.add_run(extra[0])
            r.bold = True
            r.font.color.rgb = LABEL_COLOR
            p.add_run(extra[1].format(domain=domain))

        real_status = (captured.get("_statuses") or {}).get(cp["key"])
        status = cp.get("status")
        if real_status or status:
            p = doc.add_paragraph(style="List Paragraph")
            p.paragraph_format.left_indent = Inches(0)
            r = p.add_run(status[0] if status else "Status: ")
            r.bold = True
            r.font.color.rgb = LABEL_COLOR
            p.add_run(real_status if real_status else status[1])

        img = captured.get(cp["key"])
        if img and os.path.exists(img):
            try:
                _add_bordered_picture(img)
                placed += 1
            except Exception:
                missing += 1
        else:
            missing += 1

        doc.add_paragraph()

    gen_date = datetime.now().strftime("%d-%B-%Y")
    out_path = Path(out_dir) / f"{domain} Website health checkup analysis report {gen_date}.docx"
    out_path = _safe_save_path(out_path)
    doc.save(str(out_path))
    return str(out_path), placed, missing


# ---------- PPTX builder (Sigma) ----------
def build_health_pptx_sigma(domain, captured, computed, out_dir, include_keys=None):
    from pptx import Presentation
    from pptx.util import Inches as PInches, Pt as PPt
    from pptx.dml.color import RGBColor as PRGB
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    CREAM   = PRGB(0xEF, 0xED, 0xE3)
    GOLD    = PRGB(0xE6, 0xC0, 0x69)
    GOLD_DK = PRGB(0xB8, 0x92, 0x3B)
    INK     = PRGB(0x19, 0x1B, 0x0E)
    INK_SOFT = PRGB(0x6B, 0x6A, 0x5C)
    GREEN   = PRGB(0x00, 0xB0, 0x50)
    RED     = PRGB(0xC0, 0x00, 0x00)
    FONT    = "Calibri"

    use_keys = include_keys or SIGMA_KEYS
    cps = [CHECKPOINT_BY_KEY[k] for k in use_keys if k in CHECKPOINT_BY_KEY]

    prs = Presentation()
    prs.slide_width = PInches(13.333)
    prs.slide_height = PInches(7.5)
    blank = prs.slide_layouts[6]

    def rect(slide, x, y, w, h, rgb, line_rgb=None):
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PInches(x), PInches(y), PInches(w), PInches(h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = rgb
        if line_rgb is None:
            shp.line.fill.background()
        else:
            shp.line.color.rgb = line_rgb
            shp.line.width = PPt(1)
        shp.shadow.inherit = False
        return shp

    def bg(slide):
        shp = rect(slide, 0, 0, 13.333, 7.5, CREAM)
        spTree = slide.shapes._spTree
        spTree.remove(shp._element)
        spTree.insert(2, shp._element)

    def text(slide, runs, x, y, w, h, size, rgb=INK, bold=False, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP):
        tb = slide.shapes.add_textbox(PInches(x), PInches(y), PInches(w), PInches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        p = tf.paragraphs[0]
        p.alignment = align
        if not isinstance(runs, list):
            runs = [{"text": runs, "color": rgb, "bold": bold}]
        for seg in runs:
            r = p.add_run()
            r.text = seg["text"]
            r.font.size = PPt(size)
            r.font.name = FONT
            r.font.bold = seg.get("bold", bold)
            r.font.color.rgb = seg.get("color", rgb)
        return tb

    def add_image_fit(slide, img_path, bx, by, bw, bh):
        try:
            from PIL import Image
            if not img_path or not os.path.exists(img_path) or os.path.getsize(img_path) == 0:
                return False
            with Image.open(img_path) as im:
                im.verify()
            with Image.open(img_path) as im:
                iw, ih = im.size
        except Exception:
            return False
        ar = (iw / ih) if ih else 1.5
        box_ar = bw / bh
        if ar > box_ar:
            w = bw
            h = bw / ar
        else:
            h = bh
            w = bh * ar
        x = bx + (bw - w) / 2.0
        y = by + (bh - h) / 2.0
        rect(slide, x - 0.03, y - 0.03, w + 0.06, h + 0.06, GOLD)
        try:
            slide.shapes.add_picture(str(img_path), PInches(x), PInches(y),
                                     width=PInches(w), height=PInches(h))
            return True
        except Exception:
            return False

    # Title slide
    s = prs.slides.add_slide(blank)
    bg(s)
    rect(s, 0, 0, 13.333, 0.35, GOLD)
    rect(s, 0, 7.15, 13.333, 0.35, GOLD)
    text(s, domain, 1.0, 2.4, 11.33, 1.0, 40, INK, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, "Website Health Checkup & Analysis Report", 1.0, 3.5, 11.33, 0.6, 22,
         GOLD_DK, bold=True, align=PP_ALIGN.CENTER)
    gen_date = datetime.now().strftime("%d-%B-%Y")
    text(s, f"Generated: {gen_date}", 1.0, 4.2, 11.33, 0.45, 13, INK_SOFT,
         align=PP_ALIGN.CENTER)

    # Intro slide
    s = prs.slides.add_slide(blank)
    bg(s)
    rect(s, 11.5, 0, 1.83, 0.28, GOLD)
    text(s, "INTRODUCTION", 0.689, 0.3, 12.0, 0.7, 32, INK, bold=True)
    rect(s, 0.689, 1.0, 5.5, 0.05, GOLD)
    intro_body = " ".join(b.lstrip(": ").strip() if b.startswith(":") else b.strip() for _, b in INTRO)
    text(s, intro_body, 1.1, 1.8, 11.1, 3.4, 18, INK)

    # Checkpoint slides
    placed = missing = 0
    for n, cp in enumerate(cps, start=1):
        s = prs.slides.add_slide(blank)
        bg(s)
        rect(s, 0, 0, 13.333, 0.18, GOLD)
        rect(s, 11.5, 0, 1.83, 0.28, GOLD)
        text(s, f"{n}. {cp['label'].strip()}", 0.689, 0.32, 12.0, 0.7, 26, INK, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)
        rect(s, 0.689, 1.02, 5.5, 0.05, GOLD)

        body = cp["body"]
        for k, v in computed.items():
            body = body.replace("{" + k + "}", v)
        try:
            body = body.format(domain=domain)
        except Exception:
            pass
        body_color = PRGB(0xCC, 0x00, 0x00) if cp.get("red") else INK_SOFT
        text(s, body.strip(), 0.689, 1.12, 11.95, 0.95, 15, body_color, bold=bool(cp.get("red")))

        line_y = 2.05
        extra = cp.get("extra")
        if extra:
            text(s, [
                {"text": extra[0], "color": INK, "bold": True},
                {"text": extra[1].format(domain=domain), "color": GREEN, "bold": False},
            ], 0.689, line_y, 11.95, 0.35, 14)
            line_y += 0.4

        # Real detected status (see build_health_docx's comment on this) takes
        # priority over the static CHECKPOINTS placeholder.
        real_status = (captured.get("_statuses") or {}).get(cp["key"])
        status = cp.get("status")
        if real_status or status:
            text(s, [
                {"text": (status[0].strip() if status else "Status:") + " ", "color": INK, "bold": True},
                {"text": (real_status or status[1]).strip(), "color": GREEN, "bold": True},
            ], 0.689, line_y, 11.95, 0.35, 14)
            line_y += 0.4

        img_box_y = max(2.55, line_y + 0.05)
        img = captured.get(cp["key"])
        if img and add_image_fit(s, img, 0.7, img_box_y, 11.93, 7.2 - img_box_y):
            placed += 1
        else:
            missing += 1

    out_path = Path(out_dir) / f"{domain} Website health checkup analysis report {gen_date}.pptx"
    out_path = _safe_save_path(out_path)
    prs.save(str(out_path))
    return str(out_path), placed, missing


# ---------- PPTX builder (Omega) ----------
def build_health_pptx_omega(domain, captured, computed, out_dir):
    from pptx import Presentation as PRS
    from pptx.util import Inches as PInches, Pt as PPt
    from pptx.dml.color import RGBColor as PRGB
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    BLUE    = PRGB(0x1F, 0x49, 0x7D)
    INK     = PRGB(0x1F, 0x29, 0x37)
    INK_SOFT = PRGB(0x5A, 0x5A, 0x5A)
    GREEN   = PRGB(0x00, 0xB0, 0x50)
    FONT    = "Calibri"

    cps = [CHECKPOINT_BY_KEY[k] for k in OMEGA_KEYS if k in CHECKPOINT_BY_KEY]

    prs = PRS()
    prs.slide_width = PInches(13.333)
    prs.slide_height = PInches(7.5)
    blank = prs.slide_layouts[6]

    def rect(slide, x, y, w, h, rgb, line_rgb=None):
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PInches(x), PInches(y), PInches(w), PInches(h))
        shp.fill.solid()
        shp.fill.fore_color.rgb = rgb
        if line_rgb is None:
            shp.line.fill.background()
        else:
            shp.line.color.rgb = line_rgb
            shp.line.width = PPt(1)
        shp.shadow.inherit = False
        return shp

    def text(slide, runs, x, y, w, h, size, rgb=INK, bold=False, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP):
        tb = slide.shapes.add_textbox(PInches(x), PInches(y), PInches(w), PInches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        p = tf.paragraphs[0]
        p.alignment = align
        if not isinstance(runs, list):
            runs = [{"text": runs, "color": rgb, "bold": bold}]
        for seg in runs:
            r = p.add_run()
            r.text = seg["text"]
            r.font.size = PPt(size)
            r.font.name = FONT
            r.font.bold = seg.get("bold", bold)
            r.font.color.rgb = seg.get("color", rgb)
        return tb

    def add_image_fit(slide, img_path, bx, by, bw, bh):
        try:
            from PIL import Image
            if not img_path or not os.path.exists(img_path) or os.path.getsize(img_path) == 0:
                return False
            with Image.open(img_path) as im:
                im.verify()
            with Image.open(img_path) as im:
                iw, ih = im.size
        except Exception:
            return False
        ar = (iw / ih) if ih else 1.5
        box_ar = bw / bh
        if ar > box_ar:
            w = bw
            h = bw / ar
        else:
            h = bh
            w = bh * ar
        x = bx + (bw - w) / 2.0
        y = by + (bh - h) / 2.0
        rect(slide, x - 0.03, y - 0.03, w + 0.06, h + 0.06, BLUE)
        try:
            slide.shapes.add_picture(str(img_path), PInches(x), PInches(y),
                                     width=PInches(w), height=PInches(h))
            return True
        except Exception:
            return False

    BG = PRGB(0xF8, 0xF7, 0xF4)
    ACCENT = PRGB(0xB0, 0x8D, 0x57)

    def slide_bg(slide, header=False):
        """Background fill + top/bottom accent bars so every slide carries the
        report's design instead of the default blank-white pptx layout."""
        rect(slide, 0, 0, 13.333, 7.5, BG)
        if header:
            rect(slide, 0, 0, 13.333, 0.12, ACCENT)
        rect(slide, 0, 7.4, 13.333, 0.1, BLUE)

    # Title
    s = prs.slides.add_slide(blank)
    slide_bg(s, header=True)
    text(s, domain, 1.0, 2.2, 11.33, 1.0, 36, BLUE, bold=True,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    intro_body = " ".join(b.lstrip(": ").strip() if b.startswith(":") else b.strip() for _, b in INTRO)
    text(s, intro_body, 1.5, 3.5, 10.33, 2.5, 16, INK_SOFT, align=PP_ALIGN.CENTER)

    placed = missing = 0
    for n, cp in enumerate(cps, start=1):
        s = prs.slides.add_slide(blank)
        slide_bg(s)
        body = cp["body"]
        for k, v in computed.items():
            body = body.replace("{" + k + "}", v)
        try:
            body = body.format(domain=domain)
        except Exception:
            pass

        text(s, f"{n}. {cp['label'].strip()}", 0.5, 0.3, 12.3, 0.7, 24, BLUE, bold=True)
        body_color = PRGB(0xCC, 0x00, 0x00) if cp.get("red") else INK_SOFT
        text(s, body.strip(), 0.5, 1.05, 12.3, 0.95, 14, body_color, bold=bool(cp.get("red")))

        line_y = 2.0
        extra = cp.get("extra")
        if extra:
            text(s, [
                {"text": extra[0], "color": INK, "bold": True},
                {"text": extra[1].format(domain=domain), "color": GREEN, "bold": False},
            ], 0.5, line_y, 12.3, 0.35, 13)
            line_y += 0.35
        # Real detected status (see build_health_docx's comment on this) takes
        # priority over the static CHECKPOINTS placeholder.
        real_status = (captured.get("_statuses") or {}).get(cp["key"])
        status = cp.get("status")
        if real_status or status:
            text(s, [
                {"text": (status[0].strip() if status else "Status:") + " ", "color": INK, "bold": True},
                {"text": (real_status or status[1]).strip(), "color": GREEN, "bold": True},
            ], 0.5, line_y, 12.3, 0.35, 13)
            line_y += 0.35

        img_box_y = max(2.4, line_y + 0.1)
        img = captured.get(cp["key"])
        if img and add_image_fit(s, img, 0.5, img_box_y, 12.3, 7.2 - img_box_y):
            placed += 1
        else:
            missing += 1

    gen_date = datetime.now().strftime("%d-%B-%Y")
    out_path = Path(out_dir) / f"{domain} Health Analysis Report {gen_date}.pptx"
    out_path = _safe_save_path(out_path)
    prs.save(str(out_path))
    return str(out_path), placed, missing


# ---------- Main entry point ----------
def run_health_audit(domain, fmt="james", target_pages=None, out_dir=None,
                     driver=None, no_capture=False, log_fn=None, psi_api_key=None,
                     prefetched_futures=None):
    """Run a full health audit. Returns the output file path."""
    if log_fn is None:
        log_fn = print
    # Normalize to a bare host - a full URL (https://x.com/) breaks temp image
    # paths (".png" basename -> PIL "unknown file extension") and the filename.
    domain = re.sub(r'^\s*https?://', '', str(domain or '')).strip().strip('/').split('/')[0] or str(domain)

    # Build EXACTLY the selected format - never silently fall back to another one.
    fmt = str(fmt or "").strip().lower()
    fi = FORMAT_INFO.get(fmt)
    if not fi:
        raise ValueError(f"Unknown health audit format '{fmt}'. "
                         f"Available: {', '.join(sorted(FORMAT_INFO))}")
    use_keys = list(fi["keys"])
    if not psi_api_key and "pagespeed" in use_keys:
        use_keys.remove("pagespeed")
        log_fn("PageSpeed Insights skipped (not enabled).")

    total_steps = 3 if (not no_capture and driver is not None) else 2
    log_fn(f"[1/{total_steps}] Running health checks for {domain}...")
    captured, computed = prepare_health_data(domain, target_pages=target_pages, log_fn=log_fn, psi_api_key=psi_api_key, driver=driver, prefetched_futures=prefetched_futures)

    if not no_capture and driver is not None:
        log_fn(f"[2/{total_steps}] Capturing screenshots...")
        ss_dir = os.path.join(out_dir or tempfile.gettempdir(), "health_screenshots")
        # "blank_page" already has a real rendered results table from
        # prepare_health_data() above when target pages were checked - a raw
        # homepage screenshot here would silently overwrite that table (or,
        # if the raw capture fails, leave the checkpoint with nothing at
        # all). Only fall back to a raw screenshot when there's no table.
        capture_keys = [k for k in use_keys if not (k == "blank_page" and "blank_page" in captured)]
        ss = capture_screenshots_selenium(driver, domain, ss_dir, capture_keys, log_fn)
        captured.update(ss)

        gsc_keys = {"manual_action", "security_issues"} & set(use_keys)
        uncaptured_gsc = gsc_keys - set(captured.keys())
        if uncaptured_gsc:
            try:
                import gsc_audit
                accounts = gsc_audit.list_accounts()
                connected = [a for a in accounts if a.get("has_refresh")]
                # Previously just used connected[0] - the FIRST connected
                # account, regardless of whether it actually has access to
                # THIS domain. Confirmed real case: graceperfumes.ae's
                # manual_action/security_issues screenshots both failed with
                # "the signed-in account doesn't have access to this
                # property", even though the pre-run "GSC Connected" check
                # (a separate, correct code path - api_gsc_check_for_domain
                # in web_app_batch.py) had already confirmed
                # seo.digitall12@gmail.com genuinely has Full access - this
                # function was silently picking whichever connected account
                # happened to be first instead of the one already confirmed
                # to own this domain. Now checks each connected account's
                # real properties the same way that pre-run check does, and
                # only picks one actually covering this domain.
                gsc_email = None
                for _acct in connected:
                    try:
                        _tok0 = gsc_audit.get_access_token(_acct["email"])
                        _props = gsc_audit.list_properties(_tok0)
                        if any(domain in (p.get("siteUrl") or "").lower() for p in _props):
                            gsc_email = _acct["email"]
                            break
                    except Exception:
                        continue
                if not gsc_email and connected:
                    # No connected account was confirmed to actually own this
                    # domain - fall back to the first one rather than skipping
                    # entirely; resolve_property() below still surfaces the
                    # real "no access" reason if this guess is wrong too.
                    gsc_email = connected[0]["email"]
                # Resolve the REAL property (URL-prefix vs sc-domain) via the API so
                # the manual-actions / security pages load the correct property —
                # hardcoding sc-domain broke it for URL-prefix properties.
                prop_url = f"sc-domain:{domain}"
                if gsc_email:
                    try:
                        _tok = gsc_audit.get_access_token(gsc_email)
                        prop_url = gsc_audit.resolve_property(_tok, domain) or prop_url
                        log_fn(f"  GSC property: {prop_url}")
                    except Exception as _e:
                        log_fn(f"  Could not resolve GSC property ({_e}); using {prop_url}")
                gsc_pages = []
                if "manual_action" in uncaptured_gsc:
                    gsc_pages.append({"key": "manual_action", "page": "manual-actions", "wait": 12})
                if "security_issues" in uncaptured_gsc:
                    gsc_pages.append({"key": "security_issues", "page": "security-issues", "wait": 8})
                if gsc_pages:
                    gsc_captured = False
                    sessions = gsc_audit.list_sessions()
                    # A session's "accounts" list is only populated by a cookie scan, which
                    # may not have run - so a perfectly-good logged-in session can look
                    # "untagged". Try email-matched sessions FIRST, then every other session.
                    # capture_gsc_with_session verifies it's actually logged in (and retries
                    # in a visible window if Google bounces the headless relaunch to login).
                    def _email_match(s):
                        return bool(gsc_email) and gsc_email.lower() in [a.lower() for a in s.get("accounts", [])]
                    ordered = [s for s in sessions if _email_match(s)] + [s for s in sessions if not _email_match(s)]
                    for sess in ordered:
                        tag = "matched" if _email_match(sess) else "untagged"
                        log_fn(f"  Capturing GSC screenshots via session {sess.get('label', sess['id'])} ({tag})...")
                        gsc_ss = gsc_audit.capture_gsc_with_session(
                            sess["id"], prop_url, gsc_email or "", ss_dir,
                            pages=gsc_pages, log_fn=log_fn)
                        if isinstance(gsc_ss, dict) and "error" not in gsc_ss:
                            captured.update(gsc_ss)
                            gsc_captured = True
                            break
                        elif isinstance(gsc_ss, dict) and gsc_ss.get("error") == "session_expired":
                            log_fn(f"  Session {sess.get('label', sess['id'])} needs re-login in GSC Sessions.")
                    if not gsc_captured:
                        # A connected GSC *account* (OAuth token, used for API calls
                        # elsewhere in this report) is a separate thing from a GSC
                        # *browser session* (a logged-in Selenium profile, required
                        # here because Manual Actions/Security Issues have no public
                        # API and must be scraped from the UI). Reporting these as a
                        # plain "Not found" when the real reason is "couldn't check"
                        # is misleading, so put the real reason in the report body.
                        if not sessions:
                            reason = ("Could not verify - no GSC Browser Session is set up. "
                                      "Create one in GSC Audit > Browser Sessions (this is separate "
                                      "from connecting the GSC account).")
                            log_fn("  No GSC browser session found - create one in GSC Audit > Browser Sessions")
                        elif not gsc_email:
                            reason = "Could not verify - no GSC account is connected."
                            log_fn("  GSC not connected - skipping manual_action/security_issues screenshots")
                        else:
                            reason = ("Could not verify - the GSC Browser Session is not logged in "
                                      "for this account. Re-create it in GSC Audit > Browser Sessions.")
                            log_fn("  No valid GSC browser session found - create one in GSC Audit > Browser Sessions")
                        statuses = captured.setdefault("_statuses", {})
                        for gp in gsc_pages:
                            statuses[gp["key"]] = reason
                        try:
                            import json as _json, urllib.request as _ur2, urllib.parse as _up2
                            _bd = os.path.dirname(os.path.abspath(__file__))
                            _cfg_path = os.path.join(_bd, "config.json")
                            if os.path.exists(_cfg_path):
                                with open(_cfg_path, "r") as _cf:
                                    _cfg = _json.load(_cf)
                                _wurl = (_cfg.get("auth_api_url") or "").strip()
                                if _wurl:
                                    _qs = _up2.urlencode({"action": "get_config"})
                                    _req = _ur2.Request(f"{_wurl}?{_qs}")
                                    with _ur2.urlopen(_req, timeout=10) as _r:
                                        _data = _json.loads(_r.read())
                                    _mapping = _data.get("mapping", {}) if _data.get("success") else {}
                                    _clean = domain.lower().replace("www.", "")
                                    _info = _mapping.get(_clean)
                                    if _info:
                                        log_fn(f"  >> Domain found in GSC account: {_info.get('email', '?')} "
                                               f"(access: {_info.get('accessLevel', '?')}). "
                                               f"Create a browser session and log in with this account.")
                        except Exception:
                            pass
            except Exception as e:
                log_fn(f"  GSC screenshot capture skipped: {e}")
    else:
        log_fn(f"[2/3] Skipping screenshots (no browser or --no-capture)")

    log_fn(f"[{total_steps}/{total_steps}] Building {fi['ext'].upper()} report ({fmt})...")
    if not out_dir:
        out_dir = tempfile.gettempdir()
    os.makedirs(out_dir, exist_ok=True)

    if fmt == "sigma":
        path, placed, miss = build_health_pptx_sigma(domain, captured, computed, out_dir, use_keys)
    elif fmt == "omega":
        path, placed, miss = build_health_pptx_omega(domain, captured, computed, out_dir)
    elif fmt == "neon":
        path, placed, miss = build_health_docx(domain, captured, computed, out_dir, use_keys, voice="i", header_fill="215868")
    elif fmt == "xenonk":
        path, placed, miss = build_health_docx_xenonk(domain, captured, computed, out_dir, use_keys)
    else:
        path, placed, miss = build_health_docx(domain, captured, computed, out_dir, use_keys, page_border=True)

    log_fn(f"[DONE] {placed} images placed, {miss} text-only. Output: {path}")
    return path
