"""
SEO Toolkit Pro - GSC Audit Module
Handles Google OAuth, GSC API data, screenshots, and 5 PPTX report formats.
"""

import os
import re
import sys
import json
import time
import hashlib
import tempfile
import threading
import urllib.request
import urllib.error

# One lock per browser-session profile directory. Chrome locks a --user-data-dir, so two
# GSC screenshot captures on the SAME account/session must serialize - but captures on
# DIFFERENT accounts run fully in parallel (different profile dirs -> different locks).
_PROFILE_LOCKS = {}
_PROFILE_LOCKS_GUARD = threading.Lock()

def _profile_lock(profile_dir):
    key = os.path.abspath(profile_dir)
    with _PROFILE_LOCKS_GUARD:
        lk = _PROFILE_LOCKS.get(key)
        if lk is None:
            lk = threading.Lock()
            _PROFILE_LOCKS[key] = lk
        return lk
import urllib.parse
from datetime import datetime, timedelta
from io import BytesIO

BUNDLE_DIR = os.path.dirname(os.path.abspath(__file__))

def _gsc_auth_file():
    script_dir = BUNDLE_DIR
    pf = os.environ.get("ProgramFiles", "C:\\Program Files").lower()
    pfx86 = os.environ.get("ProgramFiles(x86)", "C:\\Program Files (x86)").lower()
    if script_dir.lower().startswith(pf) or script_dir.lower().startswith(pfx86):
        appdata = os.path.join(os.environ.get("LOCALAPPDATA", os.path.expanduser("~")), "SEO Toolkit Pro")
        os.makedirs(appdata, exist_ok=True)
        return os.path.join(appdata, ".gsc_accounts")
    return os.path.join(script_dir, ".gsc_accounts")

GSC_AUTH_FILE = _gsc_auth_file()

OAUTH_AUTH_URL = "https://accounts.google.com/o/oauth2/v2/auth"
OAUTH_TOKEN_URL = "https://oauth2.googleapis.com/token"
USERINFO_URL = "https://www.googleapis.com/oauth2/v3/userinfo"
SCOPES = "https://www.googleapis.com/auth/webmasters https://www.googleapis.com/auth/userinfo.email"

GSC_API_BASE = "https://www.googleapis.com/webmasters/v3"
SEARCH_ANALYTICS_URL = "https://searchconsole.googleapis.com/webmasters/v3"
URL_INSPECTION_URL = "https://searchconsole.googleapis.com/v1/urlInspection/index:inspect"

# ---------------------------------------------------------------------------
# Account storage
# ---------------------------------------------------------------------------

def _load_accounts():
    try:
        if os.path.exists(GSC_AUTH_FILE):
            with open(GSC_AUTH_FILE, "r", encoding="utf-8") as f:
                return json.load(f)
    except Exception:
        pass
    return {}


def _save_accounts(accounts):
    try:
        with open(GSC_AUTH_FILE, "w", encoding="utf-8") as f:
            json.dump(accounts, f, indent=2)
    except Exception:
        pass


def list_accounts():
    accs = _load_accounts()
    return [{"email": k, "has_refresh": bool(v.get("refresh_token"))}
            for k, v in accs.items()]


def remove_account(email):
    accs = _load_accounts()
    accs.pop(email.lower(), None)
    _save_accounts(accs)


# ---------------------------------------------------------------------------
# OAuth via Selenium popup
# ---------------------------------------------------------------------------

def get_oauth_url(client_id, redirect_uri="urn:ietf:wg:oauth:2.0:oob"):
    params = urllib.parse.urlencode({
        "client_id": client_id,
        "redirect_uri": redirect_uri,
        "response_type": "code",
        "scope": SCOPES,
        "access_type": "offline",
        "prompt": "consent",
    })
    return f"{OAUTH_AUTH_URL}?{params}"


def oauth_login_selenium(driver, client_id, client_secret, log_fn=None):
    """Open Google OAuth in the Selenium browser and wait for the user to authorize.
    Returns the email of the connected account."""
    if log_fn is None:
        log_fn = print

    redirect_uri = "http://localhost:19876/oauth_callback"
    params = urllib.parse.urlencode({
        "client_id": client_id,
        "redirect_uri": redirect_uri,
        "response_type": "code",
        "scope": SCOPES,
        "access_type": "offline",
        "prompt": "consent",
    })
    auth_url = f"{OAUTH_AUTH_URL}?{params}"

    original_window = driver.current_window_handle
    driver.execute_script("window.open(arguments[0], '_blank', 'width=500,height=700');", auth_url)

    # Wait for the OAuth redirect
    import time as _t
    new_handles = [h for h in driver.window_handles if h != original_window]
    if new_handles:
        driver.switch_to.window(new_handles[0])

    log_fn("  Waiting for Google authorization (log in and click Allow)...")

    code = None
    for _ in range(300):  # 5 minutes max
        _t.sleep(1)
        try:
            current_url = driver.current_url
            if "oauth_callback" in current_url and "code=" in current_url:
                parsed = urllib.parse.urlparse(current_url)
                qs = urllib.parse.parse_qs(parsed.query)
                code = qs.get("code", [None])[0]
                break
            if "approval_code" in current_url or "/oauthchooseaccount" not in current_url:
                # Check if the page shows an authorization code
                try:
                    page_text = driver.find_element("tag name", "body").text
                    if "4/" in page_text:
                        for line in page_text.split("\n"):
                            line = line.strip()
                            if line.startswith("4/"):
                                code = line
                                break
                except Exception:
                    pass
        except Exception:
            # Window might be closed
            break

    # Close the OAuth popup
    try:
        if driver.current_window_handle != original_window:
            driver.close()
    except Exception:
        pass
    try:
        driver.switch_to.window(original_window)
    except Exception:
        pass

    if not code:
        raise Exception("OAuth authorization failed - no code received. Please try again.")

    # Exchange code for tokens
    tokens = _exchange_code(code, client_id, client_secret, redirect_uri)
    email = _get_user_email(tokens["access_token"])

    accs = _load_accounts()
    accs[email] = {
        "email": email,
        "refresh_token": tokens.get("refresh_token", ""),
        "access_token": tokens["access_token"],
        "expires_at": time.time() + tokens.get("expires_in", 3600) - 60,
    }
    _save_accounts(accs)
    log_fn(f"  Connected: {email}")
    return email


def _exchange_code(code, client_id, client_secret, redirect_uri):
    data = urllib.parse.urlencode({
        "code": code,
        "client_id": client_id,
        "client_secret": client_secret,
        "redirect_uri": redirect_uri,
        "grant_type": "authorization_code",
    }).encode()
    req = urllib.request.Request(OAUTH_TOKEN_URL, data=data,
                                headers={"Content-Type": "application/x-www-form-urlencoded"})
    with urllib.request.urlopen(req, timeout=15) as r:
        tokens = json.loads(r.read().decode())
    if "error" in tokens:
        raise Exception(f"Token exchange failed: {tokens.get('error_description', tokens['error'])}")
    return tokens


def _get_user_email(access_token):
    req = urllib.request.Request(USERINFO_URL,
                                headers={"Authorization": f"Bearer {access_token}"})
    with urllib.request.urlopen(req, timeout=10) as r:
        info = json.loads(r.read().decode())
    return info.get("email", "").lower()


def _refresh_token(email):
    """Refresh an expired access token. Returns the new access_token."""
    accs = _load_accounts()
    acc = accs.get(email)
    if not acc or not acc.get("refresh_token"):
        raise Exception(f"No refresh token for {email}. Please reconnect the account.")

    config = _load_gsc_config()
    data = urllib.parse.urlencode({
        "client_id": config.get("gsc_client_id", config.get("client_id", "")),
        "client_secret": config.get("gsc_client_secret", config.get("client_secret", "")),
        "refresh_token": acc["refresh_token"],
        "grant_type": "refresh_token",
    }).encode()
    req = urllib.request.Request(OAUTH_TOKEN_URL, data=data,
                                headers={"Content-Type": "application/x-www-form-urlencoded"})
    with urllib.request.urlopen(req, timeout=15) as r:
        tokens = json.loads(r.read().decode())
    if "error" in tokens:
        remove_account(email)
        raise Exception(f"Session expired for {email}. Please reconnect.")

    acc["access_token"] = tokens["access_token"]
    acc["expires_at"] = time.time() + tokens.get("expires_in", 3600) - 60
    accs[email] = acc
    _save_accounts(accs)
    return tokens["access_token"]


def get_access_token(email):
    """Get a valid access token, refreshing if needed."""
    accs = _load_accounts()
    acc = accs.get(email.lower())
    if not acc:
        raise Exception(f"Account {email} not connected.")
    if acc.get("access_token") and acc.get("expires_at", 0) > time.time():
        return acc["access_token"]
    return _refresh_token(email.lower())


def _load_gsc_config():
    script_dir = os.path.dirname(os.path.abspath(__file__))
    pf = os.environ.get("ProgramFiles", "C:\\Program Files").lower()
    pfx86 = os.environ.get("ProgramFiles(x86)", "C:\\Program Files (x86)").lower()
    if script_dir.lower().startswith(pf) or script_dir.lower().startswith(pfx86):
        cfg_path = os.path.join(os.environ.get("LOCALAPPDATA", os.path.expanduser("~")),
                                "SEO Toolkit Pro", "config.json")
    else:
        cfg_path = os.path.join(script_dir, "config.json")
    try:
        if os.path.exists(cfg_path):
            with open(cfg_path, "r", encoding="utf-8") as f:
                return json.load(f)
    except Exception:
        pass
    return {}


# ---------------------------------------------------------------------------
# GSC Browser Sessions - persistent Chrome user-data-dirs per Google account
# ---------------------------------------------------------------------------

def _sessions_dir():
    appdata = os.path.join(
        os.environ.get("LOCALAPPDATA", os.path.expanduser("~")),
        "SEO Toolkit Pro", "gsc_sessions")
    os.makedirs(appdata, exist_ok=True)
    return appdata


def _session_meta_path(session_id):
    return os.path.join(_sessions_dir(), session_id, "meta.json")


def _load_session_meta(session_id):
    p = _session_meta_path(session_id)
    try:
        if os.path.exists(p):
            with open(p, "r", encoding="utf-8") as f:
                return json.load(f)
    except Exception:
        pass
    return {}


def _save_session_meta(session_id, meta):
    p = _session_meta_path(session_id)
    os.makedirs(os.path.dirname(p), exist_ok=True)
    with open(p, "w", encoding="utf-8") as f:
        json.dump(meta, f, indent=2)


def list_sessions():
    sdir = _sessions_dir()
    sessions = []
    if not os.path.isdir(sdir):
        return sessions
    for name in sorted(os.listdir(sdir)):
        meta_path = os.path.join(sdir, name, "meta.json")
        if os.path.isfile(meta_path):
            try:
                with open(meta_path, "r", encoding="utf-8") as f:
                    meta = json.load(f)
                meta["id"] = name
                meta["profile_dir"] = os.path.join(sdir, name, "chrome_profile")
                sessions.append(meta)
            except Exception:
                pass
    return sessions


def create_session(label=None):
    sid = f"s{int(time.time()*1000)}"
    profile_dir = os.path.join(_sessions_dir(), sid, "chrome_profile")
    os.makedirs(profile_dir, exist_ok=True)
    meta = {
        "label": label or f"Session {sid[-4:]}",
        "accounts": [],
        "created": datetime.now().isoformat(),
        "last_login": None,
    }
    _save_session_meta(sid, meta)
    return {"id": sid, **meta, "profile_dir": profile_dir}


def remove_session(session_id):
    import shutil
    sdir = os.path.join(_sessions_dir(), session_id)
    if os.path.isdir(sdir):
        shutil.rmtree(sdir, ignore_errors=True)
        return True
    return False


def set_session_account(session_id, email):
    """Tag a browser session with the Google account it's logged into (so the
    Manual Action / Security capture can find it by email)."""
    meta = _load_session_meta(session_id)
    meta["accounts"] = [email]
    meta["label"] = email
    meta["last_login"] = datetime.now().isoformat()
    _save_session_meta(session_id, meta)


def launch_session_browser(session_id, browser_pref="edge", log_fn=None):
    """Launch Chrome/Edge with the session's user-data-dir so the user can
    log into Google accounts.  Returns the driver (caller must quit it)."""
    if log_fn is None:
        log_fn = print
    import engine
    profile_dir = os.path.join(_sessions_dir(), session_id, "chrome_profile")
    os.makedirs(profile_dir, exist_ok=True)
    driver = engine.build_driver(
        profile_dir, proxy=None, headless=False,
        country="us", extra_extensions=[],
        logger=log_fn, browser_pref=browser_pref,
    )
    driver.get("https://accounts.google.com/")
    return driver


def scan_session_cookies(session_id, driver=None, log_fn=None):
    """Scan cookies in the session browser to detect logged-in Google accounts.
    If driver is provided, reads cookies from it; otherwise reads the cookie
    file from the profile dir."""
    if log_fn is None:
        log_fn = print
    accounts = []
    if driver:
        try:
            driver.get("https://myaccount.google.com/")
            time.sleep(3)
            cookies = driver.get_cookies()
            google_cookies = [c for c in cookies if ".google.com" in c.get("domain", "")]
            emails = set()
            for c in google_cookies:
                if c.get("name") == "SAPISID" or c.get("name") == "SID":
                    pass
            try:
                page_text = driver.find_element("tag name", "body").text
                import re
                found = re.findall(r'[\w.+-]+@[\w-]+\.[\w.]+', page_text)
                for e in found:
                    if e.endswith(('.com', '.org', '.net', '.in', '.io', '.co')):
                        emails.add(e.lower())
            except Exception:
                pass
            accounts = list(emails)
        except Exception as e:
            log_fn(f"  Cookie scan error: {e}")
    meta = _load_session_meta(session_id)
    if accounts:
        meta["accounts"] = accounts
        meta["last_login"] = datetime.now().isoformat()
    _save_session_meta(session_id, meta)
    return {"accounts": accounts, "session_id": session_id}


def check_session_valid(session_id, driver=None):
    """Quick check if a session's Google cookies are still valid.
    Returns True if we can load a Google page without hitting the login screen."""
    if driver is None:
        return False
    try:
        driver.get("https://search.google.com/search-console")
        time.sleep(4)
        url = driver.current_url
        if "accounts.google.com" in url or "signin" in url.lower():
            return False
        return True
    except Exception:
        return False


def find_session_for_email(email):
    """Find a session that has this Google account logged in."""
    email = email.lower().strip()
    for s in list_sessions():
        if email in [a.lower() for a in s.get("accounts", [])]:
            return s
    return None


def _trim_session_cache(profile_dir):
    """Delete a session profile's browser CACHE (not cookies/login), so each
    per-account session stays small on disk - the login (cookies) is only a few MB;
    the bulk is cache we don't need to keep."""
    import shutil
    cache_dirs = ["Cache", "Code Cache", "GPUCache", "DawnCache", "DawnGraphiteCache",
                  "GrShaderCache", "ShaderCache", "Service Worker/CacheStorage",
                  "Service Worker/ScriptCache", "Service Worker/Database"]
    for base in (profile_dir, os.path.join(profile_dir, "Default")):
        for c in cache_dirs:
            p = os.path.join(base, *c.split("/"))
            try:
                if os.path.isdir(p):
                    shutil.rmtree(p, ignore_errors=True)
            except Exception:
                pass


def capture_gsc_with_session(session_id, property_url, email, out_dir,
                              pages=None, browser_pref="edge", log_fn=None):
    """Launch the GSC-login session browser and capture GSC screenshots from it.
    Uses the exact profile where the user logged into Google Search Console. Google
    often bounces a HEADLESS relaunch of a logged-in profile to the sign-in page, so
    if the headless attempt is bounced we retry once in a VISIBLE window (matching how
    the session was created) - that keeps the login alive and the screenshots real."""
    if log_fn is None:
        log_fn = print
    import engine
    profile_dir = os.path.join(_sessions_dir(), session_id, "chrome_profile")

    def _attempt(headless):
        driver = None
        try:
            driver = engine.build_driver(
                profile_dir, proxy=None, headless=headless,
                country="us", extra_extensions=[],
                logger=log_fn, browser_pref=browser_pref,
            )
            driver.get("https://search.google.com/search-console")
            time.sleep(3)
            cur = (driver.current_url or "").lower()
            if "accounts.google.com" in cur or "/signin" in cur or "servicelogin" in cur:
                return {"error": "session_expired", "session_id": session_id}
            return capture_gsc_screenshots(driver, property_url, email, out_dir,
                                            pages=pages, log_fn=log_fn)
        except Exception as e:
            log_fn(f"  Session capture error ({'headless' if headless else 'visible'}): {e}")
            return {"error": str(e)}
        finally:
            if driver:
                try:
                    driver.quit()
                except Exception:
                    pass

    # Serialize captures that share this session's profile (Chrome locks the profile
    # dir); captures for other accounts hold different locks and run in parallel.
    lock = _profile_lock(profile_dir)
    acquired = lock.acquire(timeout=180)
    if not acquired:
        log_fn("  Another capture is using this Google session - timed out waiting.")
        return {"error": "session_busy", "session_id": session_id}
    try:
        result = _attempt(headless=True)
        if isinstance(result, dict) and result.get("error") == "session_expired":
            log_fn("  Headless GSC session bounced to login - retrying in a visible window...")
            result = _attempt(headless=False)
        _trim_session_cache(profile_dir)  # keep the session small: drop cache, keep the login
        return result
    finally:
        lock.release()


# ---------------------------------------------------------------------------
# GSC API helpers
# ---------------------------------------------------------------------------

def _raise_with_api_detail(e, url):
    """urllib's default HTTPError message is just 'HTTP Error 403: Forbidden' -
    it never reads the response body, which is where Google's API actually
    explains WHY (wrong scope, no permission on this property, API not
    enabled, etc.) - that detail is essential for diagnosing a 403/404 that
    only affects one domain/account, not a generic unhelpful message."""
    try:
        body = e.read().decode("utf-8", "ignore")
        detail = json.loads(body).get("error", {})
        reason = detail.get("message") or body[:300]
    except Exception:
        reason = str(e)
    raise Exception(f"GSC API {e.code} on {url.split('?')[0]}: {reason}") from e


def _api_get(url, token, timeout=15):
    req = urllib.request.Request(url, headers={
        "Authorization": f"Bearer {token}",
        "User-Agent": "SEOToolkitPro-GSC/1.0",
    })
    try:
        with urllib.request.urlopen(req, timeout=timeout) as r:
            return json.loads(r.read().decode())
    except urllib.error.HTTPError as e:
        _raise_with_api_detail(e, url)


def _fmt_date(val):
    """Format an ISO datetime to just the date part. '2026-06-25T12:31:13.883Z' → '2026-06-25'."""
    if not val or val == "N/A":
        return val or "N/A"
    s = str(val)
    if "T" in s:
        return s.split("T")[0]
    return s


def _detect_sitemap_type(sm):
    """Detect sitemap type - 'Index' for sitemap indexes, otherwise the API type."""
    api_type = sm.get("type", "")
    path = sm.get("path", "").lower()
    if api_type and api_type.lower() not in ("", "sitemap"):
        return api_type
    if "index" in path or path.endswith("sitemap_index.xml"):
        return "Index Sitemap"
    if api_type:
        return api_type
    return "Sitemap"


def _api_post(url, token, body, timeout=30):
    data = json.dumps(body).encode()
    req = urllib.request.Request(url, data=data, headers={
        "Authorization": f"Bearer {token}",
        "Content-Type": "application/json",
        "User-Agent": "SEOToolkitPro-GSC/1.0",
    })
    try:
        with urllib.request.urlopen(req, timeout=timeout) as r:
            return json.loads(r.read().decode())
    except urllib.error.HTTPError as e:
        _raise_with_api_detail(e, url)


def list_properties(token):
    """List all GSC properties accessible by this token."""
    resp = _api_get(f"{GSC_API_BASE}/sites", token)
    return resp.get("siteEntry", [])


_property_cache = {}  # domain -> (siteUrl, expires_at) - a token's property list is
                      # essentially static minute to minute, but resolve_property() was
                      # called with a live sites.list API request on every single call
                      # (e.g. every Health Audit run re-resolving the same domain).
_PROPERTY_CACHE_TTL = 600  # seconds


def resolve_property(token, domain):
    """Find the best GSC property URL for a domain. Cached per domain for 10 minutes
    so repeat lookups (e.g. Health Audit re-resolving a domain GSC Audit already
    resolved this session) skip the live sites.list API call entirely."""
    domain_clean = domain.lower().replace("www.", "").strip("/")
    cached = _property_cache.get(domain_clean)
    if cached and cached[1] > time.time():
        return cached[0]

    props = list_properties(token)
    # Prefer domain property, then URL prefix
    result = None
    for p in props:
        url = p.get("siteUrl", "")
        if url == f"sc-domain:{domain_clean}":
            result = url
            break
    if not result:
        for p in props:
            url = p.get("siteUrl", "")
            if domain_clean in url.lower():
                result = url
                break
    if not result and props:
        result = props[0].get("siteUrl", "")
    if not result:
        raise Exception(f"No GSC property found for '{domain}'. Make sure the domain is added to Google Search Console.")
    _property_cache[domain_clean] = (result, time.time() + _PROPERTY_CACHE_TTL)
    return result


def fetch_sitemaps(token, property_url):
    encoded = urllib.parse.quote(property_url, safe="")
    resp = _api_get(f"{GSC_API_BASE}/sites/{encoded}/sitemaps", token, timeout=20)
    return resp.get("sitemap", [])


def fetch_search_analytics(token, property_url, start_date, end_date, dimensions=None, row_limit=25):
    encoded = urllib.parse.quote(property_url, safe="")
    body = {
        "startDate": start_date,
        "endDate": end_date,
        "dimensions": dimensions or ["query"],
        "rowLimit": row_limit,
    }
    resp = _api_post(
        f"{SEARCH_ANALYTICS_URL}/sites/{encoded}/searchAnalytics/query",
        token, body, timeout=30
    )
    return resp.get("rows", [])


def fetch_performance_daily(token, property_url, start_date, end_date):
    return fetch_search_analytics(token, property_url, start_date, end_date,
                                  dimensions=["date"], row_limit=100)


def fetch_top_queries(token, property_url, start_date, end_date, limit=25):
    return fetch_search_analytics(token, property_url, start_date, end_date,
                                  dimensions=["query"], row_limit=limit)


def fetch_top_pages(token, property_url, start_date, end_date, limit=25):
    return fetch_search_analytics(token, property_url, start_date, end_date,
                                  dimensions=["page"], row_limit=limit)


def fetch_image_perf(token, property_url, start_date, end_date, limit=10):
    encoded = urllib.parse.quote(property_url, safe="")
    body = {
        "startDate": start_date,
        "endDate": end_date,
        "dimensions": ["query"],
        "rowLimit": limit,
        "searchType": "image",
    }
    resp = _api_post(
        f"{SEARCH_ANALYTICS_URL}/sites/{encoded}/searchAnalytics/query",
        token, body, timeout=30
    )
    return resp.get("rows", [])


def inspect_url(token, property_url, url_to_inspect):
    body = {
        "inspectionUrl": url_to_inspect,
        "siteUrl": property_url,
    }
    resp = _api_post(URL_INSPECTION_URL, token, body, timeout=30)
    result = resp.get("inspectionResult", {})
    return {
        "indexStatusResult": result.get("indexStatusResult", {}),
        "richResultsResult": result.get("richResultsResult", {}),
    }


def fetch_all_api_data(token, property_url, start_date, end_date, log_fn=None):
    if log_fn is None:
        log_fn = print
    data = {}
    log_fn("  Fetching sitemaps...")
    data["sitemaps"] = fetch_sitemaps(token, property_url)
    log_fn("  Fetching performance data...")
    data["perfDaily"] = fetch_performance_daily(token, property_url, start_date, end_date)
    log_fn("  Fetching top queries...")
    data["topQueries"] = fetch_top_queries(token, property_url, start_date, end_date)
    log_fn("  Fetching top pages...")
    data["topPages"] = fetch_top_pages(token, property_url, start_date, end_date)
    log_fn("  Fetching image performance...")
    data["imagePerf"] = fetch_image_perf(token, property_url, start_date, end_date)
    return data


def run_inspections(token, property_url, top_pages, log_fn=None):
    if log_fn is None:
        log_fn = print
    urls = []
    for p in top_pages[:10]:
        keys = p.get("keys", [])
        if keys:
            urls.append(keys[0])
    if not urls and property_url.startswith("http"):
        urls = [property_url.rstrip("/") + "/"]
    elif not urls and property_url.startswith("sc-domain:"):
        urls = [f"https://{property_url.replace('sc-domain:', '')}/"]

    inspections = []
    for i, url in enumerate(urls):
        log_fn(f"  Inspecting URL {i+1}/{len(urls)}: {url[:60]}...")
        try:
            result = inspect_url(token, property_url, url)
            inspections.append({"url": url, **result})
        except Exception as e:
            inspections.append({"url": url, "error": str(e)})
    return inspections


# ---------------------------------------------------------------------------
# GSC page screenshots via Selenium
# ---------------------------------------------------------------------------

def build_gsc_url(page, property_url, email=None):
    """Build a GSC console URL for a specific page."""
    resource_id = property_url
    base = "https://search.google.com/search-console"
    qs = f"resource_id={urllib.parse.quote(resource_id, safe='')}"
    return f"{base}/{page}?{qs}"


def _detect_page_status(driver, key):
    """Read page text to detect status for manual actions, security issues, removals."""
    try:
        body = driver.find_element("tag name", "body").text.lower()
        if key == "manual":
            if "no issues detected" in body or "no manual actions" in body:
                return "No Issues Detected"
            if "manual action" in body and ("issue" in body or "action" in body):
                return "Issue Found"
        elif key == "security":
            if "no issues detected" in body or "no security issues" in body:
                return "No Issues Detected"
            if "security issue" in body or "malware" in body or "hacked" in body:
                return "Issue Found"
        elif key == "removals":
            # Only real table-row status values mean there ARE removals. The
            # page's chrome/help text always contains the words "removal" and
            # "request", and every empty tab renders "no requests submitted", so
            # check the positive row markers FIRST, then the empty state.
            if ("temporarily removed" in body or "request canceled" in body
                    or "request denied" in body or "request approved" in body):
                return "Removals Found"
            if ("no requests submitted" in body or "no items" in body
                    or "no removals" in body or "nothing here" in body):
                return "No Removals Found"
            return "No Removals Found"
    except Exception:
        pass
    return ""


def _extract_removals_detail(driver):
    """Rich, per-URL breakdown of active GSC removal requests (URL / Type /
    Requested date / Status), scraped from the Removals page's rendered text -
    ported from the gsc-audit-studio extension's innerText-parsing strategy
    (D:\\Working Extensions\\gsc-audit-studio\\dashboard.js), which produced the
    detailed alert emails this is meant to restore. Returns None if nothing
    found (caller falls back to the short _detect_page_status string)."""
    try:
        text = driver.find_element("tag name", "body").text
    except Exception:
        return None
    if not text:
        return None
    lower = text.lower()

    url_matches = re.findall(r"https?://\S+", text)
    removal_urls = [u for u in url_matches
                    if "google.com" not in u and "support." not in u and "goo.gl" not in u]
    has_active = bool(re.search(r"processing request|temporarily remove|clear cached|outdated|filtered", lower))
    if not removal_urls or not has_active:
        return None

    lines = []
    for url in removal_urls[:5]:
        idx = text.find(url)
        context = text[idx:idx + 200]
        status_m = re.search(r"(Processing request|Active|Approved|Expired|"
                              r"Temporarily remove URL|Clear cached URL)", context, re.I)
        date_m = re.search(r"(Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec) \d+, \d{4}", context, re.I)
        type_m = re.search(r"(Temporarily remove URL|Clear cached URL)", context, re.I)
        line = f"  • {url}"
        if type_m:
            line += f"\n    Type: {type_m.group(0)}"
        if date_m:
            line += f"  | Requested: {date_m.group(0)}"
        if status_m:
            line += f"  | Status: {status_m.group(0)}"
        lines.append(line)

    return ("Active removal request(s) found in GSC Removals tool:\n\n"
            + "\n".join(lines) + "\n\nReview and confirm these removals are intentional.")


# Robots.txt analysis - ported from gsc-audit-studio's analyseRobotsTxt(), a
# pure HTTP fetch + parse (no DOM/Selenium needed), same rule set: flags a
# hard block of the whole site, an inaccessible/erroring robots.txt (Google
# treats that as a full crawl block too), and short/suspicious Disallow
# patterns that block important paths (search, tag, author pages etc.) while
# staying quiet about common, safe blocks (wp-admin, login, checkout...).
_ROBOTS_SAFE_DISALLOWS = ("/wp-admin", "/wp-includes", "/cgi-bin", "/admin", "/login",
                          "/cart", "/checkout", "/account", "/my-account", "/wp-json", "/xmlrpc.php")


def _analyse_robots_txt(text):
    issues = []
    if not text or not text.strip():
        return issues
    lines = [l.replace("\r", "").strip() for l in text.split("\n")]
    blocks = []
    current = None
    for line in lines:
        if not line or line.startswith("#"):
            continue
        lc = line.lower()
        if lc.startswith("user-agent:"):
            agent = line.split(":", 1)[1].strip().lower()
            if current is None:
                current = {"agents": [agent], "disallows": []}
            elif not current["disallows"]:
                current["agents"].append(agent)
            else:
                blocks.append(current)
                current = {"agents": [agent], "disallows": []}
        elif lc.startswith("disallow:") and current is not None:
            current["disallows"].append(line.split(":", 1)[1].strip())
    if current is not None:
        blocks.append(current)

    for block in blocks:
        applies_all = "*" in block["agents"]
        applies_google = any(a in ("googlebot", "googlebot-news", "googlebot-image") for a in block["agents"])
        if not applies_all and not applies_google:
            continue
        agent_label = "all crawlers (*)" if applies_all else ", ".join(block["agents"])
        for disallow in block["disallows"]:
            if disallow == "/":
                issues.append({"type": "Robots.txt Blocks All Crawling",
                               "detail": f'robots.txt has "Disallow: /" for {agent_label}. This blocks Google '
                                         f'from crawling ALL pages on the site. This is a critical SEO issue.'})
            elif disallow and not any(disallow.startswith(s) for s in _ROBOTS_SAFE_DISALLOWS) and len(disallow) <= 3:
                issues.append({"type": "Robots.txt Suspicious Block",
                               "detail": f'robots.txt has "Disallow: {disallow}" for {agent_label}. This may be '
                                         f'blocking important pages from being crawled.'})
    return issues


def check_robots_txt(property_url, log_fn=None):
    """Fetch and analyse this property's robots.txt. Never raises - a network
    failure just yields no issues rather than breaking the whole audit."""
    if log_fn is None:
        log_fn = print
    try:
        base = ("https://" + property_url.replace("sc-domain:", "")
                 if property_url.startswith("sc-domain:")
                 else (property_url if property_url.endswith("/") else property_url + "/"))
        robots_url = urllib.parse.urljoin(base, "/robots.txt")
        req = urllib.request.Request(robots_url, headers={"User-Agent": "SEOToolkitPro-GSC/1.0"})
        with urllib.request.urlopen(req, timeout=15) as r:
            code = r.getcode()
            body = r.read().decode("utf-8", "ignore")
    except urllib.error.HTTPError as e:
        code = e.code
        body = ""
    except Exception as e:
        log_fn(f"  [warn] robots.txt check failed: {e}")
        return []

    if code == 404:
        return []   # missing robots.txt = fully open, that's fine
    if code in (401, 403):
        return [{"type": "Robots.txt Access Denied",
                 "detail": f"robots.txt returned HTTP {code}. Google treats this as a complete block "
                           f"and will not crawl ANY page on the site until robots.txt is accessible."}]
    if code >= 500:
        return [{"type": "Robots.txt Server Error",
                 "detail": f"robots.txt returned HTTP {code} (server error). Google may temporarily "
                           f"stop crawling until the server error is resolved."}]
    if code >= 400:
        return [{"type": "Robots.txt Error",
                 "detail": f"robots.txt returned HTTP {code}. Google may have trouble reading crawl "
                           f"rules for this site."}]
    return _analyse_robots_txt(body)


def _looks_like_signin(driver):
    """True if the current page is a Google sign-in / account-chooser screen
    rather than real GSC content. The browser can pass the ONE upfront login
    check (in capture_gsc_with_session) and still bounce to sign-in on a
    LATER page - e.g. when the active Google account in a multi-login
    profile isn't the one with access to this specific property. Without
    this check that sign-in screenshot gets saved and reported as a normal
    capture, which is exactly the "screenshots not captured correctly" bug."""
    try:
        cur = (driver.current_url or "").lower()
        if "accounts.google.com" in cur or "/signin" in cur or "servicelogin" in cur:
            return True
        body = driver.find_element("tag name", "body").text.lower()
        if "sign in" in body and ("to continue to google search console" in body
                                   or "choose an account" in body):
            return True
    except Exception:
        pass
    return False


def capture_gsc_screenshots(driver, property_url, email, out_dir, pages=None, log_fn=None):
    """Capture GSC page screenshots using the existing Selenium browser.
    The browser must already be logged into a Google account with GSC access.
    Returns dict with screenshot paths and status info."""
    if log_fn is None:
        log_fn = print
    if pages is None:
        pages = [
            {"key": "sitemap",  "page": "sitemaps",                     "wait": 8},
            {"key": "manual",   "page": "manual-actions",               "wait": 12},
            {"key": "perf",     "page": "performance/search-analytics", "wait": 15},
            {"key": "security", "page": "security-issues",              "wait": 8},
            {"key": "removals", "page": "removals",                     "wait": 8},
        ]

    os.makedirs(out_dir, exist_ok=True)
    screenshots = {}
    statuses = {}

    def _wait_network_idle(max_wait):
        """Poll for network activity to go idle instead of always sleeping the full
        max_wait. GSC's pages are a React SPA whose charts/tables load via
        fetch/XHR after the base HTML is 'complete', so resource-count stability
        (no new network requests for ~1s) is a much safer 'actually done loading'
        signal than DOM text length here - and it's capped at the same max_wait as
        before, so a shot can only finish SOONER than the old fixed sleep, never
        later, which is what keeps this from risking a blank/partial capture."""
        last_count, stable, elapsed, step = -1, 0, 0.0, 0.5
        while elapsed < max_wait:
            try:
                count = driver.execute_script(
                    "return performance.getEntriesByType('resource').length")
            except Exception:
                break
            if count == last_count:
                stable += 1
                if stable >= 2:                        # idle for ~1s -> settled
                    return
            else:
                stable = 0
            last_count = count
            time.sleep(step)
            elapsed += step

    for p in pages:
        url = build_gsc_url(p["page"], property_url, email)
        log_fn(f"  Capturing {p['key']}...")
        try:
            driver.get(url)
            _wait_network_idle(p["wait"])
            if _looks_like_signin(driver):
                log_fn(f"  [warn] {p['key']} bounced to sign-in mid-capture - "
                       f"skipping (would be a misleading screenshot).")
                continue
            # Normalise the health-audit key names ("manual_action"/"security_issues")
            # to the status detector's base keys ("manual"/"security").
            status_key = {"manual_action": "manual",
                          "security_issues": "security"}.get(p["key"], p["key"])
            if status_key in ("manual", "security", "removals"):
                status = _detect_page_status(driver, status_key)
                if status:
                    statuses[p["key"]] = status
                    log_fn(f"  {p['key']} status: {status}")
                    if status_key == "removals" and status.lower() not in ("no removals found", ""):
                        # Rich per-URL breakdown for the alert email only - the
                        # short `status` above stays as-is for the PPTX slide's
                        # single-line "Status: ..." text.
                        detail = _extract_removals_detail(driver)
                        if detail:
                            statuses["removals_detail"] = detail
            ss_path = os.path.join(out_dir, f"gsc_{p['key']}.png")
            driver.save_screenshot(ss_path)
            screenshots[p["key"]] = ss_path
            log_fn(f"  {p['key']} captured.")
        except Exception as e:
            log_fn(f"  {p['key']} capture failed: {e}")

    screenshots["_statuses"] = statuses
    return screenshots


# ---------------------------------------------------------------------------
# PPTX Report Builders
# ---------------------------------------------------------------------------

def _init_pptx():
    """Import python-pptx and return the module."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    return Presentation, Inches, Pt, Emu, RGBColor, PP_ALIGN, MSO_ANCHOR


def _add_image_slide(prs, img_path, Inches, Pt, RGBColor, PP_ALIGN,
                     slide_w=13.33, slide_h=7.5,
                     img_x=0.5, img_y=1.8, img_w=12.33, img_h=5.1):
    """Add an image to a slide, fitting within the given bounds (contain sizing)."""
    from pptx.util import Emu
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    if img_path and os.path.exists(img_path):
        from PIL import Image
        with Image.open(img_path) as im:
            iw, ih = im.size
        aspect = iw / ih
        target_w = Inches(img_w)
        target_h = Inches(img_h)
        if aspect > (img_w / img_h):
            final_w = target_w
            final_h = int(target_w / aspect)
        else:
            final_h = target_h
            final_w = int(target_h * aspect)
        cx = Inches(img_x) + (target_w - final_w) // 2
        cy = Inches(img_y) + (target_h - final_h) // 2
        slide.shapes.add_picture(img_path, cx, cy, final_w, final_h)

    return slide


def _text(slide, text_str, x, y, w, h, font_size, font_name, color, bold=False,
          align=None, char_spacing=None, italic=False):
    """Add a text box to a slide."""
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(text_str)
    run = p.runs[0] if p.runs else p.add_run()
    run.text = str(text_str)
    run.font.size = Pt(font_size)
    run.font.name = font_name
    if isinstance(color, str):
        color = RGBColor(int(color[1:3], 16), int(color[3:5], 16), int(color[5:7], 16))
    run.font.color.rgb = color
    run.font.bold = bold
    if italic:
        run.font.italic = True
    if align:
        p.alignment = align
    if char_spacing is not None:
        run.font.spacing = Pt(char_spacing)
    return txBox


def _rect(slide, x, y, w, h, color):
    """Add a rectangle shape."""
    from pptx.util import Inches
    from pptx.dml.color import RGBColor

    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    if isinstance(color, str):
        color = RGBColor(int(color[1:3], 16), int(color[3:5], 16), int(color[5:7], 16))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _fmt_num(n):
    if n is None:
        return "0"
    if isinstance(n, float) and n < 1:
        return f"{n:.1%}"
    if isinstance(n, float):
        return f"{n:.1f}"
    return f"{n:,}" if isinstance(n, int) else str(n)


def _totals(rows):
    clicks = sum(r.get("clicks", 0) for r in rows)
    impressions = sum(r.get("impressions", 0) for r in rows)
    ctr = clicks / impressions if impressions else 0
    positions = [r.get("position", 0) for r in rows if r.get("position")]
    avg_pos = sum(positions) / len(positions) if positions else 0
    return clicks, impressions, ctr, avg_pos


# --- James Full (19 slides) ---

def build_james_full(data, out_path, log_fn=None):
    if log_fn is None:
        log_fn = print
    Presentation, Inches, Pt, Emu, RGBColor, PP_ALIGN, MSO_ANCHOR = _init_pptx()

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    C = lambda h: RGBColor(int(h[1:3],16), int(h[3:5],16), int(h[5:7],16))
    TEAL_BG = C("#165060"); TEAL_HDR = C("#0D3D4A"); GOLD = C("#C9A84C")
    WHITE = RGBColor(0xFF,0xFF,0xFF); OFF_WHITE = C("#E8F4F6")
    CARD_BG = C("#EEF6FA"); CARD_BG2 = C("#F5F8FA")
    TEXT_DARK = C("#1A1A2E"); TEXT_SOFT = C("#5A6B7A")
    KPI_ORANGE = C("#E8A020"); KPI_BLUE = C("#2980B9")
    KPI_GREEN = C("#27AE60"); KPI_PURPLE = C("#8E44AD")
    STATUS_OK = C("#27AE60"); STATUS_WARN = C("#E8A020"); STATUS_ERR = C("#C0392B")

    domain = data.get("domain", "")
    property_url = data.get("propertyUrl", data.get("property_url", ""))
    start_date = data.get("startDate", data.get("start_date", ""))
    end_date = data.get("endDate", data.get("end_date", ""))
    sitemaps = data.get("sitemaps", [])
    top_queries = data.get("topQueries", data.get("top_queries", []))
    top_pages = data.get("topPages", data.get("top_pages", []))
    perf_daily = data.get("perfDaily", data.get("perf_daily", []))
    image_perf = data.get("imagePerf", data.get("image_perf", []))
    inspections = data.get("inspections", [])
    screenshots = data.get("screenshots", {})

    def header_slide(title, desc="", status_text="", status_ok=True):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _rect(slide, 0, 0, 13.33, 1.0, "#0D3D4A")
        _rect(slide, 0, 1.0, 13.33, 0.08, "#C9A84C")
        _text(slide, title, 0.5, 0.15, 12.3, 0.55, 26, "Trebuchet MS", WHITE, bold=True)
        if desc:
            _text(slide, desc, 0.5, 0.62, 12.3, 0.3, 14, "Calibri", OFF_WHITE)
        if status_text:
            # Green when the check passes, red when it flags an issue - turns each data
            # slide from a raw table into an audit verdict the reader can scan.
            _text(slide, status_text, 0.5, 1.18, 12.3, 0.3, 14, "Calibri",
                  C("#2E7D32") if status_ok else C("#C62828"), bold=True)
        return slide

    def _verdict_ok(text):
        return (text, True)

    def _verdict_bad(text):
        return (text, False)

    def add_table(slide, headers, rows, col_widths, start_y=1.8):
        from pptx.util import Inches as I, Pt as P
        total_w = sum(col_widths)
        tbl_shape = slide.shapes.add_table(
            len(rows) + 1, len(headers), I(0.5), I(start_y), I(total_w), I(0.4 * (len(rows) + 1))
        )
        table = tbl_shape.table
        for i, w in enumerate(col_widths):
            table.columns[i].width = I(w)
        for i, h in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = h
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = P(11)
                    r.font.bold = True
                    r.font.color.rgb = WHITE
                    r.font.name = "Calibri"
            cell.fill.solid()
            cell.fill.fore_color.rgb = TEAL_HDR
        for ri, row in enumerate(rows):
            for ci, val in enumerate(row):
                cell = table.cell(ri + 1, ci)
                cell.text = str(val)
                cell.text_frame.word_wrap = True   # wrap long URLs instead of truncating
                for p in cell.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.size = P(10)
                        r.font.name = "Calibri"
                        r.font.color.rgb = TEXT_DARK

    def kpi_cards(slide, values, y=2.4):
        """values = list of (label, value, color_hex)"""
        from pptx.util import Inches as I, Pt as P
        start_x = 0.5
        card_w = 2.95
        card_h = 2.2
        gap = 0.25
        for i, (label, val, clr) in enumerate(values[:4]):
            cx = start_x + i * (card_w + gap)
            _rect(slide, cx, y, card_w, card_h, "#EEF6FA")
            _rect(slide, cx, y, card_w, 0.12, clr)
            _text(slide, _fmt_num(val), cx, y + 0.5, card_w, 0.6, 30, "Calibri",
                  C(clr), bold=True, align=PP_ALIGN.CENTER)
            _text(slide, label, cx, y + 1.3, card_w, 0.4, 12, "Calibri",
                  TEXT_SOFT, align=PP_ALIGN.CENTER)

    # --- Slide 1: Title ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s, 0, 0, 13.33, 7.5, "#165060")
    _rect(s, 0, 0, 0.4, 7.5, "#C9A84C")
    _text(s, "GSC AUDIT REPORT", 1.2, 2.3, 11, 0.8, 36, "Calibri", WHITE, bold=True)
    _rect(s, 1.2, 3.4, 5, 0.06, "#C9A84C")
    _text(s, domain, 1.2, 3.6, 11, 0.5, 22, "Calibri", GOLD)
    if data.get("generatedDate"):
        _text(s, data.get("generatedDate", ""), 1.2, 4.25, 11, 0.4, 16, "Calibri", OFF_WHITE)

    # --- Slide 2: Introduction ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s, 0, 0, 13.33, 7.5, "#F5F8FA")
    _rect(s, 0, 0, 13.33, 1.0, "#0D3D4A")
    _rect(s, 0, 1.0, 13.33, 0.08, "#C9A84C")
    _text(s, "INTRODUCTION", 0.5, 0.15, 12.3, 0.55, 26, "Trebuchet MS", WHITE, bold=True)
    _text(s, f"This report presents a full Google Search Console audit for {domain} across key areas: "
             "Sitemap Health, Indexing & Canonicalisation, Rich Results, Image Search, and Search Performance.",
          0.5, 1.2, 12.3, 0.8, 14, "Calibri", C("#5A6B7A"))
    boxes = [
        ("Property", property_url),
        ("Pages Inspected", f"{len(inspections)} URLs (homepage + top by clicks)"),
        ("Sections Covered", "Sitemap, Indexing, Rich Results, Image, Performance, Manual Action, Security"),
    ]
    for i, (label, val) in enumerate(boxes):
        bx = 0.47 + i * 4.2
        _rect(s, bx, 2.5, 4.0, 2.3, "#FFFFFF")
        _text(s, label, bx + 0.3, 2.7, 3.4, 0.4, 14, "Calibri", TEXT_SOFT, bold=True)
        _text(s, val, bx + 0.3, 3.2, 3.4, 1.0, 12, "Calibri", TEXT_DARK)

    # --- Slide 3: Sitemap Status ---
    _sm_err = sum(int(sm.get("errors", 0) or 0) for sm in sitemaps)
    _sm_warn = sum(int(sm.get("warnings", 0) or 0) for sm in sitemaps)
    if not sitemaps:
        _sv, _sok = "Issue found - no sitemaps submitted", False
    elif _sm_err:
        _sv, _sok = f"Issue found - {_sm_err} error(s) across sitemaps", False
    elif _sm_warn:
        _sv, _sok = f"{_sm_warn} warning(s) found", False
    else:
        _sv, _sok = "No issues found", True
    s = header_slide("SITEMAP SUBMITTED & STATUS",
        "Checks whether sitemaps are submitted to GSC, whether they are active (not pending), "
        "and reports any errors or warnings on each sitemap.", status_text=_sv, status_ok=_sok)
    sitemap_rows = []
    for sm in sitemaps:
        path = sm.get("path", "")
        sm_type = _detect_sitemap_type(sm)
        errors = str(sm.get("errors", 0))
        warnings = str(sm.get("warnings", 0))
        pending = "Yes" if sm.get("isPending") else "No"
        processed = "Yes" if not sm.get("isPending") else "No"
        sitemap_rows.append([path, sm_type, processed, pending, errors, warnings])
    add_table(s, ["Sitemap", "Type", "Processed", "Pending", "Errors", "Warnings"],
              sitemap_rows or [["No sitemaps found", "", "", "", "", ""]],
              [5.3, 1.7, 1.5, 1.4, 1.2, 1.2])

    # --- Slide 4: Sitemap Last Read ---
    s = header_slide("SITEMAP LAST READ & URLS DISCOVERED",
        "Reports the last date Google downloaded each sitemap and how many URLs were discovered. "
        "A fresh read date indicates healthy crawl activity.")
    lr_rows = []
    for sm in sitemaps:
        contents = sm.get("contents", [])
        submitted = contents[0].get("submitted", "N/A") if contents else "N/A"
        lr_rows.append([sm.get("path", ""), _fmt_date(sm.get("lastSubmitted", "N/A")),
                        _fmt_date(sm.get("lastDownloaded", "N/A")), submitted])
    add_table(s, ["Sitemap URL", "Last Submitted", "Last Downloaded", "URLs Submitted"],
              lr_rows or [["N/A", "N/A", "N/A", "N/A"]],
              [5.5, 2.3, 2.3, 2.2])

    # --- Slides 5-10: Inspection-based slides ---
    insp_valid = [i for i in inspections if not i.get("error")]
    # URLs whose inspection FAILED - never hidden. They show as ERROR rows on each slide
    # and are folded into the verdicts (a URL we couldn't check is not a pass).
    _err = len(inspections) - len(insp_valid)

    # Slide 5: URL Indexing
    _tot = len(insp_valid)
    _not_idx = sum(1 for ins in insp_valid
                   if (ins.get("indexStatusResult", {}).get("verdict", "") or "").upper() != "PASS")
    if len(inspections) == 0:
        _iv, _iok = "No inspection data available", False
    elif _not_idx or _err:
        _p = []
        if _not_idx: _p.append(f"{_not_idx} not indexed")
        if _err: _p.append(f"{_err} could not be inspected")
        _iv, _iok = f"Issue found - {', '.join(_p)} of {len(inspections)} URL(s)", False
    else:
        _iv, _iok = f"All {_tot} URL(s) indexed", True
    s = header_slide("URL INDEXING CHECK",
        "Verifies whether each inspected URL is indexed by Google and confirms its coverage state "
        "directly from the GSC URL Inspection API.", status_text=_iv, status_ok=_iok)
    idx_rows = []
    for ins in inspections:
        if ins.get("error"):
            idx_rows.append([ins.get("url", ""), "ERROR", str(ins.get("error", ""))[:40], ""])
            continue
        isr = ins.get("indexStatusResult", {})
        idx_rows.append([ins["url"], isr.get("verdict", "N/A"),
                         isr.get("coverageState", "N/A"), isr.get("indexingState", "N/A")])
    add_table(s, ["URL", "Verdict", "Coverage", "Indexing State"],
              idx_rows or [["No inspection data", "", "", ""]],
              [5.5, 1.5, 2.8, 2.5])

    # Slide 6: Canonical Check
    _mm = 0
    for ins in insp_valid:
        isr = ins.get("indexStatusResult", {})
        if isr.get("userCanonical", "N/A") != isr.get("googleCanonical", "N/A"):
            _mm += 1
    if not insp_valid:
        _cv, _cok = "No inspection data available", False
    elif _mm:
        _cv, _cok = f"Issue found - {_mm} canonical mismatch(es)", False
    else:
        _cv, _cok = "No mismatches found", True
    s = header_slide("CANONICAL CHECK",
        "Compares the user-declared canonical URL against the canonical URL Google has chosen. "
        "A mismatch means Google is indexing a different version than intended.", status_text=_cv, status_ok=_cok)
    canon_rows = []
    for ins in inspections:
        if ins.get("error"):
            canon_rows.append([ins.get("url", "")[:35], "ERROR", "", ""])
            continue
        isr = ins.get("indexStatusResult", {})
        user_c = isr.get("userCanonical", "N/A")
        google_c = isr.get("googleCanonical", "N/A")
        match = "Yes" if user_c == google_c else "No"
        canon_rows.append([ins["url"][:35], user_c[:35] if user_c else "N/A",
                           google_c[:35] if google_c else "N/A", match])
    add_table(s, ["URL", "User Canonical", "Google Canonical", "Match"],
              canon_rows or [["No data", "", "", ""]],
              [3.8, 4.0, 4.0, 0.5])

    # Slide 7: Robots.txt
    _blocked = sum(1 for ins in insp_valid
                   if "BLOCK" in (ins.get("indexStatusResult", {}).get("robotsTxtState", "") or "").upper())
    if not insp_valid:
        _rv, _rok = "No inspection data available", False
    elif _blocked:
        _rv, _rok = f"Issue found - {_blocked} URL(s) blocked by robots.txt", False
    else:
        _rv, _rok = "No URLs blocked", True
    s = header_slide("ROBOTS.TXT BLOCKING",
        "Checks whether Google's crawler is allowed to access each inspected URL. "
        "BLOCKED means the page cannot be crawled, which prevents indexing.", status_text=_rv, status_ok=_rok)
    robot_rows = []
    for ins in inspections:
        if ins.get("error"):
            robot_rows.append([ins.get("url", "")[:70], "ERROR"])
            continue
        isr = ins.get("indexStatusResult", {})
        robot_rows.append([ins["url"][:70], isr.get("robotsTxtState", "N/A")])
    add_table(s, ["URL", "Robots.txt State"],
              robot_rows or [["No data", ""]],
              [9.3, 3.0])

    # Slide 8: Fetchability
    _badfetch = 0
    for ins in insp_valid:
        st = (ins.get("indexStatusResult", {}).get("pageFetchState", "") or "").upper()
        if st and not st.startswith("SUCCESS"):
            _badfetch += 1
    if not insp_valid:
        _fv, _fok = "No inspection data available", False
    elif _badfetch:
        _fv, _fok = f"Issue found - {_badfetch} page(s) not fetchable", False
    else:
        _fv, _fok = "All pages fetchable", True
    s = header_slide("PAGE FETCHABILITY & LAST CRAWLED",
        "Confirms Google can successfully fetch each page and reports when Google last crawled it. "
        "Pages not crawled in 90+ days may signal crawl budget or discovery issues.", status_text=_fv, status_ok=_fok)
    fetch_rows = []
    for ins in inspections:
        if ins.get("error"):
            fetch_rows.append([ins.get("url", ""), "ERROR", ""])
            continue
        isr = ins.get("indexStatusResult", {})
        fetch_rows.append([ins["url"], isr.get("pageFetchState", "N/A"),
                           isr.get("lastCrawlTime", "N/A")[:19]])
    add_table(s, ["URL", "Fetch State", "Last Crawl"],
              fetch_rows or [["No data", "", ""]],
              [7.5, 2.5, 2.3])

    # Slide 9: Crawl Type
    _desktop = sum(1 for ins in insp_valid
                   if "DESKTOP" in (ins.get("indexStatusResult", {}).get("crawledAs", "") or "").upper())
    if not insp_valid:
        _dv, _dok = "No inspection data available", False
    elif _desktop:
        _dv, _dok = f"Issue found - {_desktop} page(s) crawled as desktop", False
    else:
        _dv, _dok = "Mobile-first (all pages)", True
    s = header_slide("CRAWL TYPE - MOBILE-FIRST CHECK",
        "Google uses mobile-first indexing. All inspected pages should show MOBILE as the crawl type. "
        "DESKTOP crawl is a flag that the page may not be mobile-optimised.", status_text=_dv, status_ok=_dok)
    crawl_rows = []
    for ins in inspections:
        if ins.get("error"):
            crawl_rows.append([ins.get("url", "")[:70], "ERROR"])
            continue
        isr = ins.get("indexStatusResult", {})
        crawl_rows.append([ins["url"][:70], isr.get("crawledAs", "N/A")])
    add_table(s, ["URL", "Crawled As"],
              crawl_rows or [["No data", ""]],
              [9.3, 3.0])

    # Slide 10: Rich Results
    _rich = sum(1 for ins in insp_valid if ins.get("richResultsResult", {}).get("detectedItems"))
    if not insp_valid:
        _rrv, _rrok = "No inspection data available", False
    elif _rich:
        _rrv, _rrok = f"{_rich} page(s) with rich results detected", True
    else:
        _rrv, _rrok = "No rich results detected", True
    s = header_slide("ACTIVE RICH RESULTS & SEARCH APPEARANCE",
        "Identifies rich result types detected in page markup via URL inspection. "
        "Rich results improve SERP visibility and CTR.", status_text=_rrv, status_ok=_rrok)
    rich_rows = []
    for ins in inspections:
        if ins.get("error"):
            rich_rows.append([ins.get("url", ""), "N/A", "ERROR"])
            continue
        rr = ins.get("richResultsResult", {})
        verdict = rr.get("verdict", "N/A")
        items = rr.get("detectedItems", [])
        item_str = ", ".join(i.get("richResultType", "") for i in items) if items else "None"
        rich_rows.append([ins["url"], item_str, verdict])
    add_table(s, ["URL", "Detected Items", "Verdict"],
              rich_rows or [["No data", "", ""]],
              [6.0, 4.3, 2.0])

    # --- Slide 11: Image Search Performance ---
    s = header_slide("IMAGE SEARCH PERFORMANCE")
    img_clicks, img_impr, img_ctr, img_pos = _totals(image_perf)
    kpi_cards(s, [
        ("Total Clicks", img_clicks, "#E8A020"),
        ("Total Impressions", img_impr, "#2980B9"),
        ("Avg CTR", img_ctr, "#27AE60"),
        ("Avg Position", round(img_pos, 1) if img_pos else 0, "#8E44AD"),
    ])
    img_rows = []
    for q in image_perf[:5]:
        keys = q.get("keys", [""])
        img_rows.append([keys[0], _fmt_num(q.get("clicks", 0)), _fmt_num(q.get("impressions", 0)),
                         f"{q.get('ctr', 0):.1%}", f"{q.get('position', 0):.1f}"])
    if img_rows:
        add_table(s, ["Top Image Query", "Clicks", "Impressions", "CTR", "Position"],
                  img_rows,
                  [5.8, 1.4, 1.8, 1.4, 1.9], start_y=5.0)

    # --- Slide 12: Performance Overview ---
    s = header_slide("PERFORMANCE OVERVIEW", f"Web search performance totals for the period: {start_date} to {end_date}")
    total_clicks, total_impr, total_ctr, avg_pos = _totals(perf_daily)
    kpi_cards(s, [
        ("Total Clicks", total_clicks, "#E8A020"),
        ("Total Impressions", total_impr, "#2980B9"),
        ("Avg CTR", total_ctr, "#27AE60"),
        ("Avg Position", round(avg_pos, 1) if avg_pos else 0, "#8E44AD"),
    ])

    # --- Slide 13: Top Queries ---
    s = header_slide("TOP QUERIES",
        "Top 10 queries driving impressions and clicks. Flag = below position benchmark "
        "(pos 4-10 should achieve >3% CTR). Several high-impression queries with low CTR "
        "represent optimisation opportunities.")
    q_rows = []
    for i, q in enumerate(top_queries[:10]):
        keys = q.get("keys", [""])
        q_rows.append([str(i+1), keys[0], _fmt_num(q.get("clicks", 0)),
                       _fmt_num(q.get("impressions", 0)),
                       f"{q.get('ctr', 0):.1%}", f"{q.get('position', 0):.1f}"])
    add_table(s, ["#", "Query", "Clicks", "Impressions", "CTR", "Position"],
              q_rows or [["", "No data", "", "", "", ""]],
              [0.5, 5.8, 1.4, 1.8, 1.4, 1.4])

    # --- Slide 14: Top Pages ---
    s = header_slide("TOP PAGES", "Top 10 pages by clicks during the period.")
    p_rows = []
    for i, pg in enumerate(top_pages[:10]):
        keys = pg.get("keys", [""])
        p_rows.append([str(i+1), keys[0], _fmt_num(pg.get("clicks", 0)),
                       _fmt_num(pg.get("impressions", 0)),
                       f"{pg.get('ctr', 0):.1%}", f"{pg.get('position', 0):.1f}"])
    add_table(s, ["#", "Page", "Clicks", "Impressions", "CTR", "Position"],
              p_rows or [["", "No data", "", "", "", ""]],
              [0.5, 5.8, 1.4, 1.8, 1.4, 1.4])

    # --- Slide 15: Keyword Opportunities ---
    s = header_slide("KEYWORD OPPORTUNITIES",
        "Queries with high impression volume but low CTR represent immediate optimisation opportunities. "
        "Improving title tags and meta descriptions for these terms can unlock clicks without new rankings.")
    opps = sorted([q for q in top_queries
                   if q.get("impressions", 0) >= 30 and q.get("ctr", 0) < 0.03],
                  key=lambda x: x.get("impressions", 0), reverse=True)[:10]
    opp_rows = []
    for q in opps:
        keys = q.get("keys", [""])
        pos = q.get("position", 0)
        if pos <= 5:
            opp = "High - Title/Meta tune"
        elif pos <= 15:
            opp = "Medium - Content boost"
        else:
            opp = "Long-term"
        opp_rows.append([keys[0], _fmt_num(q.get("impressions", 0)),
                         f"{q.get('ctr', 0):.1%}", f"{pos:.1f}", opp])
    add_table(s, ["Query", "Impressions", "CTR", "Position", "Opportunity"],
              opp_rows or [["No clear opportunities in the top queries - CTR looks healthy.", "", "", "", ""]],
              [4.5, 1.7, 1.4, 1.4, 3.3])

    # --- Screenshot slides with definitions (matches extension format) ---
    gsc_statuses = screenshots.get("_statuses", {})
    SHOT_SLIDES = [
        ("manual", "MANUAL ACTION",
         "Google issues a manual action against a site when a human reviewer at Google has determined "
         "that pages on the site are not compliant with Google's webmaster quality guidelines."),
        ("security", "SECURITY ISSUES",
         "If a Google evaluation determines that a site was hacked, or exhibits behaviour that could "
         "harm visitors, the Security Issues report will show Google's findings."),
        ("removals", "REMOVALS",
         "Checks for active temporary removal requests, outdated content removals, or SafeSearch "
         "filtering applied to this property."),
    ]
    for shot_key, title, desc in SHOT_SLIDES:
        img_path = screenshots.get(shot_key, "")
        status_text = gsc_statuses.get(shot_key, "")
        # Always emit the Manual Action / Security / Removals slide - never skip it. If a
        # screenshot couldn't be captured, show a "please check this manually" note instead
        # of silently dropping the section (Google's most-looked-at checks must appear).
        s = header_slide(title, desc)
        if status_text:
            color = C("#2E7D32") if "no issue" in status_text.lower() or "no removal" in status_text.lower() else C("#C62828")
            _text(s, f"Status: {status_text}", 0.5, 1.65, 12.3, 0.4, 16, "Calibri", color, bold=True)
        if img_path and os.path.exists(img_path):
            from PIL import Image
            with Image.open(img_path) as im:
                iw, ih = im.size
            aspect = iw / ih
            img_top = 2.4 if status_text else 2.2
            img_h_avail = 4.4 if status_text else 4.6
            target_w = Inches(12.33)
            target_h = Inches(img_h_avail)
            if aspect > (12.33 / img_h_avail):
                fw = target_w; fh = int(target_w / aspect)
            else:
                fh = target_h; fw = int(target_h * aspect)
            cx = Inches(0.5) + (target_w - fw) // 2
            cy = Inches(img_top) + (target_h - fh) // 2
            _EMU = 914400
            _rect(s, cx / _EMU - 0.03, cy / _EMU - 0.03, fw / _EMU + 0.06, fh / _EMU + 0.06, "#000000")
            s.shapes.add_picture(img_path, cx, cy, fw, fh)
        else:
            _text(s, "Please check this manually in Google Search Console.",
                  0.5, 4.0, 12.33, 0.6, 20, "Calibri", C("#C62828"), bold=True,
                  align=PP_ALIGN.CENTER)

    # --- Slide 19: Thank You ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s, 0, 0, 13.33, 7.5, "#165060")
    _rect(s, 0, 0, 0.4, 7.5, "#C9A84C")
    _text(s, "THANK YOU", 1.2, 2.8, 10, 1.0, 48, "Trebuchet MS", WHITE, bold=True)
    _rect(s, 1.2, 4.0, 4, 0.06, "#C9A84C")
    _text(s, f"GSC Audit Report - {domain}", 1.2, 4.3, 10, 0.5, 18, "Calibri", GOLD)

    prs.save(out_path)
    log_fn(f"  James Full report saved: {os.path.basename(out_path)}")
    return out_path


# --- James Short (8 slides) ---

def build_james_short(data, out_path, log_fn=None):
    if log_fn is None:
        log_fn = print
    Presentation, Inches, Pt, Emu, RGBColor, PP_ALIGN, MSO_ANCHOR = _init_pptx()

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    C = lambda h: RGBColor(int(h[1:3],16), int(h[3:5],16), int(h[5:7],16))
    BG = C("#00292E"); PANEL = C("#054854"); CARD = C("#107082")
    TEAL_MD = C("#64B2C1"); BEIGE = C("#F0CDA1")
    WHITE = RGBColor(0xFF,0xFF,0xFF); YELLOW = C("#FFFF00")

    domain = data.get("domain", "")
    screenshots = data.get("screenshots", {})

    def screenshot_slide(title, desc, shot_key, status_text=""):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        _rect(s, 0, 0, 13.33, 7.5, "#00292E")
        _rect(s, 11.5, 0, 1.83, 0.28, "#F0CDA1")
        _text(s, title, 0.689, 0.25, 12.0, 0.70, 32, "Calibri", WHITE, bold=True)
        _rect(s, 0.689, 0.87, 5.5, 0.05, "#F0CDA1")
        if desc:
            _text(s, desc, 0.622, 1.06, 12.0, 0.4, 20, "Calibri", TEAL_MD)
        if status_text:
            txb = _text(s, "Status: ", 0.622, 1.96, 2.0, 0.3, 16, "Calibri", WHITE)
            _text(s, status_text, 2.5, 1.96, 5.0, 0.3, 16, "Calibri", YELLOW, bold=True)

        # Screenshot frame
        _rect(s, 0.38, 2.55, 12.57, 4.65, "#F0CDA1")
        _rect(s, 0.42, 2.59, 12.49, 4.57, "#00292E")

        img_path = screenshots.get(shot_key, "")
        if img_path and os.path.exists(img_path):
            from PIL import Image
            with Image.open(img_path) as im:
                iw, ih = im.size
            aspect = iw / ih
            target_w = Inches(12.33); target_h = Inches(4.4)
            if aspect > (12.33/4.4):
                fw = target_w; fh = int(target_w / aspect)
            else:
                fh = target_h; fw = int(target_h * aspect)
            cx = Inches(0.5) + (target_w - fw) // 2
            cy = Inches(2.7) + (target_h - fh) // 2
            _EMU = 914400
            _rect(s, cx / _EMU - 0.03, cy / _EMU - 0.03, fw / _EMU + 0.06, fh / _EMU + 0.06, "#000000")
            s.shapes.add_picture(img_path, cx, cy, fw, fh)
        else:
            _text(s, "Please check this manually", 3, 4.5, 7, 0.5, 18, "Calibri",
                  TEAL_MD, align=PP_ALIGN.CENTER)
        return s

    # Slide 1: Title
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s, 0, 0, 13.33, 7.5, "#00292E")
    _rect(s, 0, 0, 3.2, 7.5, "#054854")
    _rect(s, 3.2, 0, 0.06, 7.5, "#F0CDA1")
    _text(s, "GOOGLE SEARCH CONSOLE", 4.0, 2.5, 9.0, 0.85, 38, "Calibri", WHITE, bold=True)
    _rect(s, 4.0, 3.43, 4.5, 0.06, "#F0CDA1")
    _text(s, "AUDIT REPORT", 4.0, 3.5, 9.0, 0.85, 38, "Calibri", WHITE, bold=True)
    _rect(s, 4.0, 4.72, 7.0, 0.65, "#107082")
    _text(s, domain, 4.1, 4.72, 6.8, 0.65, 18, "Calibri", BEIGE, bold=True, italic=True)
    _rect(s, 4.0, 5.42, 3.2, 0.05, "#F0CDA1")

    # Slide 2: Introduction
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s, 0, 0, 13.33, 7.5, "#00292E")
    _rect(s, 8.35, 0, 4.98, 7.5, "#054854")
    _rect(s, 7.0, 1.75, 1.1, 1.1, "#F0CDA1")
    _rect(s, 7.15, 1.9, 0.85, 0.85, "#107082")
    _rect(s, 1.255, 2.988, 7.549, 3.838, "#054854")
    _text(s, "INTRODUCTION", 3.089, 3.498, 5.015, 0.818, 44, "Calibri", WHITE, bold=True)
    _rect(s, 1.961, 4.738, 4.488, 0.05, "#F0CDA1")
    _text(s, "Reviewing the search console for the website is one of the major aspects of our work as we "
             "check the webmaster for all the important parameters like traffic, manual action, security issue, "
             "etc. to check if all are performing well or not so that we can take steps if required. "
             "This 3rd phase report is just for acknowledgment!",
          1.34, 4.88, 7.2, 1.82, 13, "Calibri", WHITE)
    _rect(s, 1.717, 6.457, 6.625, 0.05, "#F0CDA1")

    # Slides 3-7: Screenshots
    gsc_statuses = screenshots.get("_statuses", {})
    screenshot_slide("SITEMAP",
        "A sitemap is a file on your site that tells Google which pages we should know about.",
        "sitemap")
    screenshot_slide("PERFORMANCE",
        "The Performance report shows clicks, impressions, CTR, and average position from Google Search.",
        "perf")
    screenshot_slide("MANUAL ACTION",
        "Google issues a manual action against a site when a human reviewer at Google has determined "
        "that pages on the site are not compliant with Google's webmaster quality guidelines.",
        "manual", status_text=gsc_statuses.get("manual", ""))
    screenshot_slide("SECURITY ISSUE",
        "If a Google evaluation determines that a site was hacked, or exhibits behaviour that could "
        "harm visitors, the Security Issues report will show Google's finding.",
        "security", status_text=gsc_statuses.get("security", ""))
    screenshot_slide("CHECK ANY PAGE FOUND IN REMOVAL",
        "Checks for any pages that have been temporarily or permanently removed from Google Search results.",
        "removals", status_text=gsc_statuses.get("removals", ""))

    # Slide 8: Thank You
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s, 0, 0, 13.33, 7.5, "#00292E")
    _rect(s, 9.8, 0, 3.53, 7.5, "#054854")
    _rect(s, 9.8, 0, 0.06, 7.5, "#64B2C1")
    _text(s, "THANK YOU", 0.9, 2.6, 9.3, 1.4, 60, "Calibri", WHITE, bold=True)
    _rect(s, 1.317, 4.147, 3.796, 0.05, "#F0CDA1")

    prs.save(out_path)
    log_fn(f"  James Short report saved: {os.path.basename(out_path)}")
    return out_path


# --- Sigma (10 slides) ---

def build_sigma(data, out_path, log_fn=None):
    if log_fn is None:
        log_fn = print
    Presentation, Inches, Pt, Emu, RGBColor, PP_ALIGN, MSO_ANCHOR = _init_pptx()

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    C = lambda h: RGBColor(int(h[1:3],16), int(h[3:5],16), int(h[5:7],16))
    CREAM = C("#EFEDE3"); INK = C("#191B0E"); INK_SOFT = C("#6B6A5C")
    GOLD = C("#E6C069"); GOLD_DK = C("#B8923B"); GREEN = C("#00B050"); RED = C("#C00000")
    WHITE = RGBColor(0xFF,0xFF,0xFF); ROW_ALT = C("#F6F2E8")

    domain = data.get("domain", "")
    property_url = data.get("propertyUrl", data.get("property_url", ""))
    screenshots = data.get("screenshots", {})
    top_queries = data.get("topQueries", data.get("top_queries", []))
    top_pages = data.get("topPages", data.get("top_pages", []))
    generated = data.get("generatedDate", datetime.now().strftime("%B %d, %Y"))

    def header_slide(title, desc=""):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        _rect(s, 0, 0, 13.33, 7.5, "#EFEDE3")
        _rect(s, 11.5, 0, 1.83, 0.28, "#E6C069")
        _text(s, title, 0.689, 0.25, 12.0, 0.6, 32, "Calibri", INK, bold=True)
        _rect(s, 0.689, 0.95, 5.5, 0.05, "#E6C069")
        if desc:
            _text(s, desc, 0.689, 1.1, 12.0, 0.4, 16, "Calibri", INK_SOFT)
        return s

    def screenshot_slide(title, shot_key, desc="", status_text=""):
        s = header_slide(title, desc)
        if status_text:
            color = GREEN if "no issue" in status_text.lower() or "no removal" in status_text.lower() else RED
            _text(s, f"Status: {status_text}", 0.689, 1.65, 12.0, 0.4, 16, "Calibri", color, bold=True)
        _rect(s, 0.38, 2.55, 12.57, 4.65, "#E6C069")
        _rect(s, 0.42, 2.59, 12.49, 4.57, "#EFEDE3")
        img_path = screenshots.get(shot_key, "")
        if img_path and os.path.exists(img_path):
            from PIL import Image
            with Image.open(img_path) as im:
                iw, ih = im.size
            aspect = iw / ih
            tw = Inches(12.33); th = Inches(4.4)
            if aspect > (12.33/4.4): fw = tw; fh = int(tw / aspect)
            else: fh = th; fw = int(th * aspect)
            cx = Inches(0.5) + (tw - fw) // 2
            cy = Inches(2.7) + (th - fh) // 2
            s.shapes.add_picture(img_path, cx, cy, fw, fh)
        else:
            _text(s, "Please check this manually", 3, 4.5, 7, 0.5, 18, "Calibri",
                  INK_SOFT, align=PP_ALIGN.CENTER)
        return s

    def data_table(slide, headers, rows, col_widths, y=2.1):
        from pptx.util import Inches as I, Pt as P
        total_w = sum(col_widths)
        tbl_shape = slide.shapes.add_table(
            len(rows) + 1, len(headers), I(0.5), I(y), I(total_w), I(0.35 * (len(rows) + 1))
        )
        table = tbl_shape.table
        for i, w in enumerate(col_widths):
            table.columns[i].width = I(w)
        for i, h in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = h
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = P(11); r.font.bold = True; r.font.name = "Calibri"
                    r.font.color.rgb = INK
            cell.fill.solid()
            cell.fill.fore_color.rgb = GOLD
        for ri, row in enumerate(rows):
            for ci, val in enumerate(row):
                cell = table.cell(ri+1, ci)
                cell.text = str(val)
                for p in cell.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.size = P(10); r.font.name = "Calibri"; r.font.color.rgb = INK
                if ri % 2 == 1:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = ROW_ALT

    # Slide 1: Title
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s, 0, 0, 13.33, 7.5, "#EFEDE3")
    _rect(s, 0, 0, 13.33, 0.35, "#E6C069")
    _rect(s, 0, 7.15, 13.33, 0.35, "#E6C069")
    _text(s, "WEBMASTER AUDIT REPORT", 0, 2.8, 13.33, 0.9, 44, "Calibri", INK,
          bold=True, align=PP_ALIGN.CENTER)
    _text(s, domain, 0, 4.0, 13.33, 0.5, 22, "Calibri", GREEN, align=PP_ALIGN.CENTER)
    _text(s, f"Generated: {generated}", 0, 4.7, 13.33, 0.4, 14, "Calibri",
          INK_SOFT, align=PP_ALIGN.CENTER)

    # Slide 2: Introduction
    s = header_slide("INTRODUCTION")
    _rect(s, 0.689, 1.5, 0.12, 4.5, "#E6C069")
    _text(s, "Reviewing the Search Console for the website is one of the major aspects of our work. "
             "We check the webmaster for all the important parameters - traffic & performance, "
             "sitemaps, manual actions, security issues and removals - to confirm everything is "
             "performing well, so corrective steps can be taken if required. "
             "This report is shared for acknowledgement of the website's current Search Console health.",
          1.0, 1.6, 11.5, 3.0, 18, "Calibri", INK)

    # Slide 3: Performance screenshot
    screenshot_slide("PERFORMANCE", "perf",
        "The Performance report shows clicks, impressions, CTR and average position from Google Search over the reporting period.")

    # Slide 4: Top Queries table
    s = header_slide("TOP QUERIES", f"The top search queries driving impressions and clicks during the reporting period.")
    q_rows = []
    for q in top_queries:
        keys = q.get("keys", [""])
        q_rows.append([keys[0], _fmt_num(q.get("clicks", 0)),
                       _fmt_num(q.get("impressions", 0)),
                       f"{q.get('ctr', 0):.1%}", f"{q.get('position', 0):.1f}"])
    data_table(s, ["Query", "Clicks", "Impressions", "CTR", "Position"],
               q_rows or [["No query data available", "-", "-", "-", "-"]],
               [6.05, 1.4, 1.8, 1.3, 1.4])

    # Slide 5: Top Pages table
    s = header_slide("TOP PAGES", "The top pages by clicks during the reporting period.")
    pg_rows = []
    for pg in top_pages:
        keys = pg.get("keys", [""])
        pg_rows.append([keys[0][:50], _fmt_num(pg.get("clicks", 0)),
                        _fmt_num(pg.get("impressions", 0)),
                        f"{pg.get('ctr', 0):.1%}", f"{pg.get('position', 0):.1f}"])
    data_table(s, ["Page", "Clicks", "Impressions", "CTR", "Position"],
               pg_rows or [["No page data available", "-", "-", "-", "-"]],
               [6.05, 1.4, 1.8, 1.3, 1.4])

    # Slides 6-9: Screenshots
    gsc_statuses = screenshots.get("_statuses", {})
    screenshot_slide("SITEMAP", "sitemap",
        "A sitemap is a file on your site that tells Google which pages we should know about.")
    screenshot_slide("MANUAL ACTION", "manual",
        "Google issues a manual action against a site when a human reviewer at Google has determined "
        "that pages on the site are not compliant with Google's webmaster quality guidelines.",
        status_text=gsc_statuses.get("manual", ""))
    screenshot_slide("SECURITY ISSUE", "security",
        "If a Google evaluation determines that a site was hacked, or exhibits behaviour that could "
        "harm visitors or their computer, the Security Issues report will show Google's finding.",
        status_text=gsc_statuses.get("security", ""))
    screenshot_slide("CHECK ANY PAGE FOUND IN REMOVAL", "removals",
        "Checks for any pages that have been temporarily or permanently removed from Google Search results.",
        status_text=gsc_statuses.get("removals", ""))

    # Slide 10: Thank You
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s, 0, 0, 13.33, 7.5, "#EFEDE3")
    _rect(s, 0, 0, 13.33, 0.35, "#E6C069")
    _rect(s, 0, 7.15, 13.33, 0.35, "#E6C069")
    _text(s, "THANK YOU", 0, 3.0, 13.33, 1.0, 60, "Calibri", INK,
          bold=True, align=PP_ALIGN.CENTER)
    _rect(s, 4.66, 4.3, 4.0, 0.06, "#E6C069")

    prs.save(out_path)
    log_fn(f"  Sigma report saved: {os.path.basename(out_path)}")
    return out_path


# --- Omega (8 slides) ---

def build_omega(data, out_path, log_fn=None):
    if log_fn is None:
        log_fn = print
    Presentation, Inches, Pt, Emu, RGBColor, PP_ALIGN, MSO_ANCHOR = _init_pptx()

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    C = lambda h: RGBColor(int(h[1:3],16), int(h[3:5],16), int(h[5:7],16))
    NAVY = C("#1B2A41"); NAVY_DARK = C("#0F1A2C"); BRONZE = C("#B08D57")
    BRONZE_SOFT = C("#D4B891"); OFF_WHITE = C("#F8F7F4")
    TEXT_DARK = C("#2C2C2C"); TEXT_SOFT = C("#6B7280"); BORDER = C("#E5E0D5")
    STATUS_OK = C("#4A7C59")
    WHITE = RGBColor(0xFF,0xFF,0xFF)

    domain = data.get("domain", "")
    screenshots = data.get("screenshots", {})

    def header_slide(title, desc=""):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        _rect(s, 0, 0, 13.33, 7.5, "#F8F7F4")
        _rect(s, 0, 0, 13.33, 1.2, "#1B2A41")
        _rect(s, 0, 1.2, 13.33, 0.06, "#B08D57")
        _rect(s, 0, 7.1, 13.33, 0.4, "#1B2A41")
        _text(s, title.upper(), 0.5, 0.3, 12.3, 0.6, 28, "Georgia", OFF_WHITE,
              bold=True, char_spacing=2)
        if desc:
            _text(s, desc, 0.5, 1.5, 12.3, 0.4, 14, "Calibri", TEXT_SOFT, italic=True)
        return s

    def screenshot_slide(title, shot_key, desc="", status_text=""):
        s = header_slide(title, desc)
        if status_text:
            color = STATUS_OK if "no issue" in status_text.lower() or "no removal" in status_text.lower() else C("#C62828")
            _text(s, f"Status: {status_text}", 0.5, 1.85, 12.3, 0.35, 16, "Calibri", color, bold=True)
        # White card frame
        _rect(s, 0.7, 2.15, 11.93, 4.85, "#FFFFFF")
        img_path = screenshots.get(shot_key, "")
        if img_path and os.path.exists(img_path):
            from PIL import Image
            with Image.open(img_path) as im:
                iw, ih = im.size
            aspect = iw / ih
            tw = Inches(11.73); th = Inches(4.65)
            if aspect > (11.73/4.65): fw = tw; fh = int(tw / aspect)
            else: fh = th; fw = int(th * aspect)
            cx = Inches(0.8) + (tw - fw) // 2
            cy = Inches(2.25) + (th - fh) // 2
            s.shapes.add_picture(img_path, cx, cy, fw, fh)
        else:
            _text(s, "Please check this manually", 3, 4.5, 7, 0.5, 18, "Calibri",
                  TEXT_SOFT, align=PP_ALIGN.CENTER)
        return s

    # Slide 1: Title
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s, 0, 0, 13.33, 7.5, "#F8F7F4")
    _rect(s, 0, 0, 13.33, 0.5, "#B08D57")
    _rect(s, 0, 7.0, 13.33, 0.5, "#1B2A41")
    _text(s, "WEBMASTER AUDIT REPORT", 0, 2.8, 13.33, 0.8, 40, "Georgia", TEXT_DARK,
          bold=True, align=PP_ALIGN.CENTER)
    _rect(s, 5.66, 3.8, 2.0, 0.06, "#B08D57")
    _text(s, domain, 0, 4.2, 13.33, 0.5, 24, "Calibri", TEXT_DARK,
          italic=True, align=PP_ALIGN.CENTER)

    # Slide 2: Introduction
    s = header_slide("Introduction")
    _text(s, "Reviewing the search console for the website is one of the major aspects of our work. "
             "We check the webmaster for all the important parameters, including sitemap health, "
             "manual actions, performance, and security issues, to confirm that the website is "
             "functioning correctly and free of errors.",
          0.5, 2.0, 12.3, 2.0, 16, "Calibri", TEXT_DARK)
    _text(s, f"Domain: {domain}", 0.5, 4.5, 12.3, 0.5, 18, "Calibri", NAVY, bold=True)

    # Slides 3-6: Screenshots
    gsc_statuses = screenshots.get("_statuses", {})
    screenshot_slide("Check Error in Sitemap", "sitemap",
        "A sitemap is a record on your site that uncovers to Google which pages we should consider.")
    screenshot_slide("Check Manual Action", "manual",
        "Manual action is a penalty imposed by Google for not following the webmaster guidelines.",
        status_text=gsc_statuses.get("manual", ""))
    screenshot_slide("Check Performance Issue", "perf",
        "The Performance Report shows important metrics about how your site performs in Google Search Results.")
    screenshot_slide("Check Security Issue", "security",
        "The Security Issues report in Search Console alerts webmasters about malicious behaviour on their websites.",
        status_text=gsc_statuses.get("security", ""))

    # Slide 7: Other Issues
    s = header_slide("Check Other Issue if in Webmaster",
        "Reviewed the Removals report and any additional Search Console alerts for the property.")
    removals_status = gsc_statuses.get("removals", "")
    if removals_status:
        color = STATUS_OK if "no removal" in removals_status.lower() else C("#C62828")
        _text(s, f"Status: {removals_status}", 0.5, 1.85, 12.3, 0.35, 16, "Calibri", color, bold=True)
    img_path = screenshots.get("removals", "")
    if img_path and os.path.exists(img_path):
        # Image ends at ~6.35 so the caption below (y=6.5) never overlaps it.
        _rect(s, 0.7, 2.15, 11.93, 4.2, "#FFFFFF")
        from PIL import Image
        with Image.open(img_path) as im:
            iw, ih = im.size
        aspect = iw / ih
        tw = Inches(11.73); th = Inches(4.0)
        if aspect > (11.73/4.0): fw = tw; fh = int(tw / aspect)
        else: fh = th; fw = int(th * aspect)
        cx = Inches(0.8) + (tw - fw) // 2
        cy = Inches(2.25) + (th - fh) // 2
        s.shapes.add_picture(img_path, cx, cy, fw, fh)
    else:
        _rect(s, 3.5, 3.3, 6.3, 1.2, "#FFFFFF")
        _text(s, "Not Found", 3.5, 3.3, 6.3, 1.2, 32, "Georgia", STATUS_OK,
              bold=True, align=PP_ALIGN.CENTER)
    _text(s, "No additional issues detected in the Search Console webmaster reports for this property.",
          0.5, 6.5, 12.3, 0.45, 12, "Calibri", TEXT_DARK, align=PP_ALIGN.CENTER)

    # Slide 8: Thank You
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s, 0, 0, 13.33, 7.5, "#1B2A41")
    _rect(s, 0, 3.4, 13.33, 0.04, "#B08D57")
    _text(s, "THANK YOU", 1.5, 2.5, 10.3, 1.0, 56, "Georgia", BRONZE_SOFT,
          bold=True, align=PP_ALIGN.CENTER, char_spacing=8)
    _text(s, domain, 1.5, 3.7, 10.3, 0.5, 18, "Calibri", WHITE, align=PP_ALIGN.CENTER)

    prs.save(out_path)
    log_fn(f"  Omega report saved: {os.path.basename(out_path)}")
    return out_path


# --- Neon (8 slides) ---

def build_neon(data, out_path, log_fn=None):
    if log_fn is None:
        log_fn = print
    Presentation, Inches, Pt, Emu, RGBColor, PP_ALIGN, MSO_ANCHOR = _init_pptx()

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    C = lambda h: RGBColor(int(h[1:3],16), int(h[3:5],16), int(h[5:7],16))
    DARK = C("#0F1F2E"); DARK_MID = C("#162D40"); OFF_WHITE = C("#F4F7FA")
    TEAL = C("#B0BEC5"); TEAL_SOFT = C("#CFD8DC"); TEAL_DIM = C("#78909C")
    GOLD_N = C("#FFC100"); GOLD_SOFT = C("#FFD44D")
    TEXT_DARK = C("#1A2E3D"); TEXT_MID = C("#2D4A62")
    TEXT_SOFT = C("#5A7A8F"); TEXT_LIGHT = C("#C8D8E4"); BORDER_N = C("#D8E6F0")
    WHITE = RGBColor(0xFF,0xFF,0xFF)

    domain = data.get("domain", "")
    screenshots = data.get("screenshots", {})

    accent_colors = {"sitemap": TEAL, "manual": GOLD_N, "perf": TEAL_SOFT, "security": GOLD_SOFT}

    def header_slide(title, desc=""):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        _rect(s, 0, 0, 13.33, 7.5, "#F4F7FA")
        _rect(s, 0, 0, 13.33, 1.3, "#0F1F2E")
        _rect(s, 0, 0, 0.18, 1.3, "#B0BEC5")
        _rect(s, 0, 1.25, 13.33, 0.05, "#FFC100")
        _rect(s, 0, 7.2, 13.33, 0.3, "#0F1F2E")
        _rect(s, 0, 7.2, 0.18, 0.3, "#FFC100")
        _text(s, title.upper(), 0.5, 0.35, 12.3, 0.6, 28, "Trebuchet MS", WHITE,
              bold=True, char_spacing=3)
        if desc:
            _text(s, desc, 0.5, 1.45, 12.3, 0.4, 13, "Calibri", TEXT_SOFT, italic=True)
        return s

    def screenshot_slide(title, shot_key, desc="", accent_hex="#B0BEC5", status_text=""):
        s = header_slide(title, desc)
        if status_text:
            color = C("#2E7D32") if "no issue" in status_text.lower() or "no removal" in status_text.lower() else C("#C62828")
            _text(s, f"Status: {status_text}", 0.5, 1.85, 12.3, 0.35, 16, "Calibri", color, bold=True)
        # Frame
        _rect(s, 0.5, 2.1, 12.33, 4.85, "#FFFFFF")
        _rect(s, 0.5, 2.1, 12.33, 0.1, accent_hex)
        img_path = screenshots.get(shot_key, "")
        if img_path and os.path.exists(img_path):
            from PIL import Image
            with Image.open(img_path) as im:
                iw, ih = im.size
            aspect = iw / ih
            tw = Inches(12.21); th = Inches(4.65)
            if aspect > (12.21/4.65): fw = tw; fh = int(tw / aspect)
            else: fh = th; fw = int(th * aspect)
            cx = Inches(0.56) + (tw - fw) // 2
            cy = Inches(2.25) + (th - fh) // 2
            s.shapes.add_picture(img_path, cx, cy, fw, fh)
        else:
            _text(s, "Please check this manually", 3, 4.5, 7, 0.5, 18, "Calibri",
                  TEXT_SOFT, align=PP_ALIGN.CENTER)
        return s

    # Slide 1: Title
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s, 0, 0, 13.33, 7.5, "#0F1F2E")
    _rect(s, 0, 0, 0.18, 7.5, "#B0BEC5")
    # Gold top-right L
    _rect(s, 12.0, 0, 1.33, 0.15, "#FFC100")
    _rect(s, 13.18, 0, 0.15, 1.5, "#FFC100")
    # Teal bottom-left
    _rect(s, 0, 7.0, 1.5, 0.15, "#B0BEC5")
    _text(s, "WEBMASTER", 1.0, 2.3, 11.0, 1.0, 54, "Trebuchet MS", TEXT_LIGHT, bold=True)
    _text(s, "AUDIT REPORT", 1.0, 3.3, 11.0, 1.0, 54, "Trebuchet MS", TEAL, bold=True)
    _rect(s, 1.0, 4.6, 4.5, 0.06, "#FFC100")
    _text(s, domain, 1.0, 5.0, 11.0, 0.5, 26, "Trebuchet MS", GOLD_SOFT)

    # Slide 2: Introduction
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s, 0, 0, 13.33, 7.5, "#F4F7FA")
    _rect(s, 0, 0, 13.33, 1.3, "#0F1F2E")
    _rect(s, 0, 0, 0.18, 1.3, "#B0BEC5")
    _rect(s, 0, 1.25, 13.33, 0.05, "#FFC100")
    _text(s, "INTRODUCTION", 0.5, 0.35, 12.3, 0.6, 28, "Trebuchet MS", WHITE,
          bold=True, char_spacing=3)
    _text(s, "Reviewing the Search Console for the website is one of the major aspects of our work. "
             "We check the webmaster for all important parameters, including sitemap health, manual actions, "
             "performance issues, and security issues, to confirm the website is functioning correctly "
             "and is free of errors.",
          0.7, 1.6, 11.9, 1.8, 16, "Calibri", TEXT_MID)
    # Domain card
    _rect(s, 0.5, 3.5, 12.33, 0.8, "#0F1F2E")
    _rect(s, 0.5, 3.5, 0.12, 0.8, "#B0BEC5")
    _text(s, f"Domain: {domain}", 0.8, 3.55, 11.8, 0.7, 16, "Calibri", WHITE)
    # Section preview boxes
    sections = [("Sitemap", "#B0BEC5"), ("Manual Action", "#FFC100"),
                ("Performance", "#CFD8DC"), ("Security", "#FFD44D")]
    for i, (label, clr) in enumerate(sections):
        bx = 0.5 + i * 3.1
        _rect(s, bx, 4.8, 2.9, 1.5, "#FFFFFF")
        _rect(s, bx, 4.8, 2.9, 0.08, clr)
        _text(s, label, bx + 0.2, 5.1, 2.5, 0.4, 13, "Calibri", TEXT_MID, bold=True)

    # Slides 3-6: Screenshots
    gsc_statuses = screenshots.get("_statuses", {})
    screenshot_slide("Sitemap Check", "sitemap",
        "A sitemap guides Google on which pages to crawl and index. We review it for errors or missing URLs.",
        "#B0BEC5")
    screenshot_slide("Manual Action Check", "manual",
        "A manual action is a Google penalty for violating webmaster guidelines. This should always be clean.",
        "#FFC100", status_text=gsc_statuses.get("manual", ""))
    screenshot_slide("Performance Check", "perf",
        "The Performance report shows clicks, impressions, CTR, and average position from Google Search.",
        "#CFD8DC")
    screenshot_slide("Security Issues Check", "security",
        "Security issues alert about hacked content, malware, or deceptive pages that may harm users.",
        "#FFD44D", status_text=gsc_statuses.get("security", ""))

    # Slide 7: Other Webmaster Checks
    s = header_slide("Other Webmaster Checks",
        "Reviewed the Removals report and any additional Search Console alerts for this property.")
    removals_status = gsc_statuses.get("removals", "")
    if removals_status:
        color = C("#2E7D32") if "no removal" in removals_status.lower() else C("#C62828")
        _text(s, f"Status: {removals_status}", 0.5, 1.85, 12.3, 0.35, 16, "Calibri", color, bold=True)
    img_path = screenshots.get("removals", "")
    if img_path and os.path.exists(img_path):
        # Image ends at ~6.3 so the caption below (y=6.5) never overlaps it.
        _rect(s, 0.5, 2.1, 12.33, 4.2, "#FFFFFF")
        _rect(s, 0.5, 2.1, 12.33, 0.1, "#B0BEC5")
        from PIL import Image
        with Image.open(img_path) as im:
            iw, ih = im.size
        aspect = iw / ih
        tw = Inches(12.21); th = Inches(4.0)
        if aspect > (12.21/4.0): fw = tw; fh = int(tw / aspect)
        else: fh = th; fw = int(th * aspect)
        cx = Inches(0.56) + (tw - fw) // 2
        cy = Inches(2.2) + (th - fh) // 2
        s.shapes.add_picture(img_path, cx, cy, fw, fh)
    else:
        _rect(s, 4.16, 3.0, 5.0, 1.7, "#0F1F2E")
        _rect(s, 4.16, 3.0, 5.0, 0.08, "#B0BEC5")
        _rect(s, 4.16, 4.62, 5.0, 0.08, "#FFC100")
        _text(s, "NOT FOUND", 4.16, 3.0, 5.0, 1.7, 40, "Trebuchet MS", TEAL,
              align=PP_ALIGN.CENTER, char_spacing=8)
    _text(s, "No additional issues detected in Search Console for this property.",
          0.7, 6.5, 11.9, 0.45, 12, "Calibri", TEXT_MID, align=PP_ALIGN.CENTER)

    # Slide 8: Thank You
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s, 0, 0, 13.33, 7.5, "#0F1F2E")
    _rect(s, 13.15, 0, 0.18, 7.5, "#B0BEC5")
    # Gold top-left L
    _rect(s, 0, 0, 1.5, 0.15, "#FFC100")
    _rect(s, 0, 0, 0.15, 1.5, "#FFC100")
    _text(s, "THANK", 0.55, 1.8, 12.5, 1.2, 72, "Trebuchet MS", TEXT_LIGHT, bold=True,
          align=PP_ALIGN.CENTER)
    _text(s, "YOU", 0.55, 2.9, 12.5, 1.2, 72, "Trebuchet MS", TEAL, bold=True,
          align=PP_ALIGN.CENTER)
    _rect(s, 5.16, 4.3, 3.0, 0.05, "#FFC100")
    _text(s, domain, 0.55, 4.5, 12.5, 0.55, 20, "Calibri", GOLD_SOFT,
          align=PP_ALIGN.CENTER)

    prs.save(out_path)
    log_fn(f"  Neon report saved: {os.path.basename(out_path)}")
    return out_path


# ---------------------------------------------------------------------------
# Format registry & main entry point
# ---------------------------------------------------------------------------

GSC_FORMATS = {
    "james":       {"label": "James New (Full, 19 slides)", "builder": build_james_full},
    "james_short": {"label": "James Short (8 slides)",      "builder": build_james_short},
    "sigma":       {"label": "Sigma (10 slides)",           "builder": build_sigma},
    "omega":       {"label": "Omega (8 slides)",            "builder": build_omega},
    "neon":        {"label": "Neon (8 slides)",             "builder": build_neon},
}

# Status strings that mean "no problem" - everything else from
# _detect_page_status counts as a real issue worth emailing about.
_GSC_STATUS_OK = {"no issues detected", "no removals found", ""}
_GSC_STATUS_LABELS = {
    "manual": "Manual Action",
    "security": "Security Issue",
    "removals": "Removals",
}


def _send_gsc_alert(webapp_url, domain, statuses, email, log_fn=None,
                    extra_issues=None, access_level=None):
    """POST the Apps Script's send_alert action for any real GSC issue found
    (manual action, security issue, removals, robots.txt) - handleSendAlert
    already exists and works server-side, it just had no caller."""
    if log_fn is None:
        log_fn = print
    issues = []
    for key, label in _GSC_STATUS_LABELS.items():
        if key == "removals":
            # Prefer the rich per-URL breakdown (statuses["removals_detail"])
            # over the short "Removals Found" status used for the PPTX slide.
            status = (statuses.get("removals_detail") or statuses.get(key) or "").strip()
        else:
            status = (statuses.get(key) or "").strip()
        if status and status.lower() not in _GSC_STATUS_OK:
            issues.append({"type": label, "detail": status})
    issues.extend(extra_issues or [])
    if not issues:
        return
    payload = {
        "action": "send_alert",
        "domain": domain,
        "issues": issues,
        "timestamp": datetime.now().isoformat(),
        "accountEmail": email or "",
        "accessLevel": access_level or "",
    }
    try:
        req = urllib.request.Request(
            webapp_url, data=json.dumps(payload).encode("utf-8"),
            headers={"Content-Type": "text/plain;charset=utf-8"})
        with urllib.request.urlopen(req, timeout=30) as r:
            result = json.loads(r.read().decode())
        if result.get("sent"):
            log_fn(f"  Alert email sent for {domain} ({len(issues)} issue(s)).")
        else:
            log_fn(f"  Alert email not sent: {result.get('error') or result.get('reason')}")
    except Exception as e:
        log_fn(f"  [warn] Alert email failed: {e}")


def run_gsc_audit(domain, email, fmt="james", out_dir=None, driver=None,
                  period_days=28, end_offset=3, log_fn=None, webapp_url=None,
                  access_level=None):
    """Run a complete GSC audit: fetch data, capture screenshots, build PPTX."""
    if log_fn is None:
        log_fn = print
    # Normalize to a bare host so it's safe in the report filename (a full URL
    # like https://x.com/ produces an invalid filename).
    domain = re.sub(r'^\s*https?://', '', str(domain or '')).strip().strip('/').split('/')[0] or str(domain)

    # Build EXACTLY the selected format - fail loudly rather than silently building
    # a different one after doing all the API/screenshot work.
    fmt = str(fmt or "").strip().lower()
    if fmt not in GSC_FORMATS:
        raise ValueError(f"Unknown GSC audit format '{fmt}'. "
                         f"Available: {', '.join(sorted(GSC_FORMATS))}")

    if out_dir is None:
        out_dir = tempfile.mkdtemp(prefix="gsc_audit_")
    os.makedirs(out_dir, exist_ok=True)

    token = get_access_token(email)
    log_fn(f"  Account: {email}")

    # Resolve property
    log_fn("  Resolving GSC property...")
    property_url = resolve_property(token, domain)
    log_fn(f"  Property: {property_url}")

    # Date range
    end_date = (datetime.now() - timedelta(days=end_offset)).strftime("%Y-%m-%d")
    start_date = (datetime.now() - timedelta(days=end_offset + period_days - 1)).strftime("%Y-%m-%d")
    log_fn(f"  Period: {start_date} to {end_date}")

    report_data = {
        "domain": domain,
        "propertyUrl": property_url,
        "startDate": start_date,
        "endDate": end_date,
        "periodDays": period_days,
        "generatedDate": datetime.now().strftime("%d-%B-%Y"),
    }

    # Fetch API data for James Full and Sigma
    if fmt in ("james", "sigma"):
        log_fn("  Fetching GSC API data...")
        api_data = fetch_all_api_data(token, property_url, start_date, end_date, log_fn)
        report_data.update(api_data)

        if fmt == "james":
            log_fn("  Running URL inspections...")
            inspections = run_inspections(token, property_url,
                                          api_data.get("topPages", []), log_fn)
            report_data["inspections"] = inspections

    # Capture screenshots from the per-account LOGGED-IN browser session - NOT a shared
    # profile. The shared profile isn't signed into the account, so it would screenshot
    # the Google sign-in page instead of the real GSC data (Manual Actions / Security /
    # Removals). capture_gsc_with_session verifies the login and retries in a visible
    # window if Google bounces the headless relaunch. If no session works we leave the
    # screenshots empty so the report shows a clear "please check manually" note rather
    # than a misleading sign-in page.
    screenshots = {}
    ss_dir = os.path.join(out_dir, "gsc_screenshots")
    sessions = list_sessions()

    def _match(s):
        return bool(email) and email.lower() in [a.lower() for a in s.get("accounts", [])]
    ordered = [s for s in sessions if _match(s)] + [s for s in sessions if not _match(s)]

    for sess in ordered:
        tag = "matched" if _match(sess) else "untagged"
        log_fn(f"  Capturing GSC screenshots via session {sess.get('label', sess['id'])} ({tag})...")
        ss = capture_gsc_with_session(sess["id"], property_url, email, ss_dir, log_fn=log_fn)
        if isinstance(ss, dict) and "error" not in ss:
            screenshots = ss
            break
        if isinstance(ss, dict) and ss.get("error") == "session_expired":
            log_fn(f"  Session {sess.get('label', sess['id'])} needs re-login (GSC Audit > Browser Sessions).")
    if not screenshots:
        if not sessions:
            log_fn("  No GSC browser session found - connect the account in GSC Audit > Browser Sessions.")
        else:
            log_fn("  Could not capture from any session - slides will show a 'check manually' note.")
    log_fn(f"  {len(screenshots)} screenshot(s) captured.")
    report_data["screenshots"] = screenshots

    # Email the admin when a real issue is found (manual action, security
    # issue, removals, or robots.txt) - the Apps Script's handleSendAlert
    # already builds and sends this email, but nothing was ever calling it,
    # so these alerts never went out.
    if webapp_url:
        robots_issues = check_robots_txt(property_url, log_fn)
        _send_gsc_alert(webapp_url, domain, screenshots.get("_statuses", {}), email, log_fn,
                        extra_issues=robots_issues, access_level=access_level)

    # Build PPTX
    format_info = GSC_FORMATS[fmt]   # validated above - build the selected format
    timestamp = datetime.now().strftime("%d-%B-%Y")
    out_file = os.path.join(out_dir, f"GSC_Audit_{domain}_{timestamp}.pptx")

    log_fn(f"  Building {format_info['label']} report...")
    format_info["builder"](report_data, out_file, log_fn)

    return out_file
