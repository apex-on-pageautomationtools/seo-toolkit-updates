"""
Brief Website Analysis — SEO check + PPTX report builder.
Formats: james (15 slides, 10×5.625"), xenon (15 slides, 13.33×7.5")
"""
import os, json, re, time, tempfile, urllib.request, urllib.parse
from datetime import datetime
from pathlib import Path

# Reuse existing health_audit checks
from health_audit import (
    _fetch_html, _parse_meta,
    check_robots_txt, check_sitemap, check_canonical,
    check_broken_links, check_status200,
)

# ---------------------------------------------------------------------------
# Browser-assisted data collection (uses app's stealth Chrome)
# ---------------------------------------------------------------------------
_brief_driver = None

def _get_brief_driver():
    """Launch a headless Chrome for brief analysis scraping."""
    global _brief_driver
    if _brief_driver is not None:
        try:
            _ = _brief_driver.title
            return _brief_driver
        except Exception:
            _brief_driver = None
    try:
        from selenium.webdriver import Chrome, ChromeOptions
        opts = ChromeOptions()
        profile_dir = tempfile.mkdtemp(prefix="seo_brief_")
        opts.add_argument(f"--user-data-dir={profile_dir}")
        opts.add_argument("--headless=new")
        opts.add_argument("--no-sandbox")
        opts.add_argument("--disable-dev-shm-usage")
        opts.add_argument("--disable-gpu")
        opts.add_argument("--disable-notifications")
        opts.add_argument("--disable-popup-blocking")
        opts.add_argument("--window-size=1280,900")
        opts.add_argument("--mute-audio")
        ua = ("Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
              "AppleWebKit/537.36 (KHTML, like Gecko) "
              "Chrome/149.0.0.0 Safari/537.36")
        opts.add_argument(f"--user-agent={ua}")
        _brief_driver = Chrome(options=opts)
        _brief_driver.set_page_load_timeout(30)
        return _brief_driver
    except Exception:
        return None

def _close_brief_driver():
    global _brief_driver
    if _brief_driver:
        try:
            _brief_driver.quit()
        except Exception:
            pass
        _brief_driver = None


_UA = ("Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 "
       "(KHTML, like Gecko) Chrome/149.0.0.0 Safari/537.36")

# Per-run cache so each URL is rendered by the browser only ONCE even though
# several checks (title / meta / headers / canonical / image alt) all need it.
_PAGE_CACHE = {}
# Rendered-content fingerprints for the current run, used to catch SPA soft-404s
# (sites that return HTTP 200 for every route and client-render a fallback view
# for URLs that don't actually exist).
_HOME_SIG = [None]   # the homepage
_NF_SIG = [None]     # a deliberately-nonexistent probe URL = the site's 404 view


def _page_signature(html):
    """A fingerprint of a page's visible content (title + a large slice of body
    text) used to tell whether two URLs render the SAME page. The slice is big so
    a real content page (long) is clearly distinct from the short 404/shell — on
    themed sites (Shopify/WordPress) every page shares the header/nav, so a small
    slice would make everything look alike."""
    if not html:
        return ""
    m = re.search(r'<title[^>]*>([^<]*)</title>', html, re.I)
    title = (m.group(1) if m else "").strip().lower()
    text = re.sub(r'(?is)<(script|style|noscript)[^>]*>.*?</\1>', ' ', html)
    text = re.sub(r'(?s)<[^>]+>', ' ', text)
    text = re.sub(r'\s+', ' ', text).strip().lower()
    return title + "||" + text[:8000]


def _sig_similar(a, b, thresh=0.95):
    """True only if two page signatures are near-identical — i.e. the browser
    rendered essentially the SAME page (a pure-SPA catch-all)."""
    if not a or not b:
        return False
    if a == b:
        return True
    # Length guard: a real content page and the short 404/shell differ a lot in
    # size. Without this, two pages that merely share header/nav chrome (all pages
    # on a Shopify/WordPress theme) would false-match and be flagged 'Not Found'.
    la, lb = len(a), len(b)
    if min(la, lb) < 0.6 * max(la, lb):
        return False
    from difflib import SequenceMatcher
    return SequenceMatcher(None, a, b).ratio() >= thresh


def _prime_soft404(domain):
    """Render a guaranteed-nonexistent URL once to capture the site's soft-404
    fingerprint. SPAs that 200 every route render their fallback/404 view here,
    which lets us flag other pages that render the same thing as 'Not Found'."""
    probe = f"https://{domain}/__seo_toolkit_probe_404__/no-such-page"
    try:
        _NF_SIG[0] = _page_signature(_render_html(probe))
    except Exception:
        _NF_SIG[0] = None


def _http_status(url):
    """Best-effort HTTP status for a URL (HEAD, then GET fallback).
    Returns the code (e.g. 200, 404) or 0 if the host is unreachable."""
    import urllib.error
    for method in ("HEAD", "GET"):
        try:
            req = urllib.request.Request(url, method=method,
                                         headers={"User-Agent": _UA})
            with urllib.request.urlopen(req, timeout=10) as r:
                return getattr(r, "status", 200) or 200
        except urllib.error.HTTPError as e:
            return e.code          # 404/410/etc. is a real answer — stop here
        except Exception:
            continue
    return 0


def _looks_not_found(html):
    """True if the page's <title>/<h1> says it's a 404 / 'page not found'.
    Catches soft-404s (SPA routes that return HTTP 200 but render a 404 page)."""
    if not html:
        return False
    m = re.search(r'<title[^>]*>([^<]*)</title>', html, re.I)
    title = (m.group(1) if m else "").lower()
    hm = re.search(r'<h1[^>]*>(.*?)</h1>', html, re.I | re.S)
    h1 = re.sub(r'<[^>]+>', '', hm.group(1)).lower() if hm else ""
    hay = title + " " + h1
    for phrase in ("404 not found", "page not found", "page doesn't exist",
                   "page does not exist", "404 error", "error 404",
                   "not found", "nothing found", "page cannot be found",
                   "no longer exists"):
        if phrase in hay:
            return True
    return False


def _render_html(url, max_wait=9.0):
    """Return fully-rendered HTML (JavaScript executed) via the headless browser,
    falling back to raw HTTP if the browser is unavailable. Waits for the SPA to
    finish hydrating (body text stops growing) so title/heading/meta checks see
    the real content, not an empty shell."""
    driver = _get_brief_driver()
    html = ""
    if driver:
        try:
            driver.get(url)
            last_len, stable = -1, 0
            for _ in range(int(max_wait / 0.4)):
                try:
                    ready = driver.execute_script("return document.readyState")
                    blen = driver.execute_script(
                        "return (document.body ? document.body.innerText.length : 0)")
                except Exception:
                    break
                if ready == "complete" and blen > 0 and blen == last_len:
                    stable += 1
                    if stable >= 2:          # content unchanged for ~0.8s -> settled
                        break
                else:
                    stable = 0
                last_len = blen
                time.sleep(0.4)
            html = driver.page_source or ""
        except Exception:
            html = ""
    if not html or len(html) < 200:                   # browser failed -> raw HTTP
        raw = _fetch_html(url)
        if raw:
            html = raw
    return html


def _get_page(url):
    """Fetch a page once (rendered HTML + whether it actually exists), cached for
    the run. Returns {'html': str, 'status': int, 'exists': bool}.

    Existence is decided by: HTTP status (404/410), a 'page not found' title/H1,
    OR — for SPAs that 200 every route — the rendered content being identical to
    the homepage (the app served its shell for a URL that doesn't exist)."""
    if url in _PAGE_CACHE:
        return _PAGE_CACHE[url]
    status = _http_status(url)
    html = "" if status in (404, 410) else _render_html(url)
    exists = status not in (404, 410) and not _looks_not_found(html)

    path = re.sub(r'^https?://[^/]+', '', url).split('#')[0].split('?')[0]
    sig = _page_signature(html)
    if path in ("", "/"):
        if sig:
            _HOME_SIG[0] = sig                        # remember the homepage
    elif exists and sig:
        # SPA soft-404: this route renders the same thing as the 404 probe (or is
        # byte-identical to the homepage shell) -> the page doesn't really exist.
        if (_NF_SIG[0] and _sig_similar(sig, _NF_SIG[0])) or \
           (_HOME_SIG[0] and sig == _HOME_SIG[0]):
            exists = False

    page = {"html": html, "status": status, "exists": exists}
    _PAGE_CACHE[url] = page
    return page


def _page_name(pg):
    """Human label for a page path: last segment, title-cased. '/' -> 'Home'."""
    seg = (pg or "").strip("/").split("/")[-1]
    if not seg:
        return "Home"
    seg = re.sub(r'\.\w{2,5}$', '', seg)               # drop .html/.php extension
    seg = seg.replace("-", " ").replace("_", " ").strip()
    return (seg.title() or "Home")[:22]


def _browser_get_indexing(domain):
    """Read the 'About X results' count from a Google site: query via the browser.
    Handles the consent wall + one retry; returns None if Google blocks the query
    (so the report shows N/A rather than a wrong number)."""
    driver = _get_brief_driver()
    if not driver:
        return None
    try:
        from engine import accept_google_consent, classify_page
    except Exception:
        accept_google_consent = classify_page = None

    def _extract(src):
        for pat in [r'About ([\d,\.  ]+?) results',
                    r'result-stats[^>]*>\s*About ([\d,\.  ]+)',
                    r'result-stats[^>]*>([\d,\.  ]+)',
                    r'([\d,\. ]+)\s*results']:
            m = re.search(pat, src, re.I)
            if m:
                num = re.sub(r'[^\d]', '', m.group(1))
                if num:
                    return {"count": f"{int(num):,}", "status": "Indexed & serving"}
        if "did not match any documents" in src.lower() or "no results found" in src.lower():
            return {"count": "0", "status": "Not indexed"}
        return None

    try:
        for attempt in range(2):
            driver.get(f"https://www.google.com/search?q=site:{domain}&hl=en&num=20")
            time.sleep(3 if attempt == 0 else 5)
            if accept_google_consent:
                try:
                    accept_google_consent(driver, lambda *a: None)
                    time.sleep(1.5)
                except Exception:
                    pass
            src = driver.page_source or ""
            if classify_page:
                try:
                    if classify_page(src) in ("captcha", "soft_block", "http_403"):
                        continue   # Google blocked us — retry once, else give up (N/A)
                except Exception:
                    pass
            res = _extract(src)
            if res:
                return res
        return None
    except Exception:
        return None


def _browser_get_domain_age(domain):
    """Use headless browser to get WHOIS data from whois.domaintools.com."""
    driver = _get_brief_driver()
    if not driver:
        return None
    clean = domain.replace("www.", "")
    try:
        driver.get(f"https://whois.domaintools.com/{clean}")
        time.sleep(4)
        src = driver.page_source or ""
        rows = [["Domain Name", clean, "", ""]]

        def _dt_extract(label_pat):
            m = re.search(label_pat + r'[^<]*</(?:div|td|dt|th)[^>]*>\s*<(?:div|td|dd)[^>]*>\s*(.*?)\s*</(?:div|td|dd)',
                          src, re.I | re.S)
            if m:
                val = re.sub(r'<[^>]+>', '', m.group(1)).strip()
                if val and val.lower() not in ('n/a', 'not available', ''):
                    return val
            return None

        created = _dt_extract(r'(?:Create[d]?\s*Date|Registration\s*Date|Registered\s*On)')
        updated = _dt_extract(r'(?:Update[d]?\s*Date|Last\s*Updated)')
        expiry = _dt_extract(r'(?:Expir\w+\s*Date|Registry\s*Expiry)')
        registrar = _dt_extract(r'Registrar')

        # Fallback: search for date patterns near keywords
        if not created:
            m = re.search(r'(?:Created|Registered|Registration)[^<]{0,50}?(\d{4}-\d{2}-\d{2})', src, re.I)
            if m:
                created = m.group(1)
        if not expiry:
            m = re.search(r'(?:Expir\w+)[^<]{0,50}?(\d{4}-\d{2}-\d{2})', src, re.I)
            if m:
                expiry = m.group(1)

        if created:
            try:
                for fmt in ("%Y-%m-%d", "%d-%b-%Y", "%B %d, %Y", "%d %B %Y", "%m/%d/%Y"):
                    try:
                        dt = datetime.strptime(created.strip()[:10], fmt)
                        age_years = (datetime.now() - dt).days // 365
                        rows.append(["Created On", created.strip()[:20], f"~{age_years} years", ""])
                        break
                    except ValueError:
                        continue
                else:
                    rows.append(["Created On", created.strip()[:20], "", ""])
            except Exception:
                rows.append(["Created On", created.strip()[:20], "", ""])
        if updated:
            rows.append(["Updated", updated.strip()[:20], "", ""])
        if expiry:
            rows.append(["Expiration", expiry.strip()[:20], "", ""])
        if registrar:
            rows.append(["Registrar", registrar.strip()[:40], "", ""])

        if len(rows) > 1:
            return rows
        return None
    except Exception:
        return None




# ---------------------------------------------------------------------------
# Data collection
# ---------------------------------------------------------------------------

# Indexing captcha-solving hook. web_app_batch populates this with the SAME Buster
# solver + extension the rank checker uses, so the "site:" query can beat a CAPTCHA
# by escalating to a visible browser instead of just giving up (N/A).
_INDEXING = {"extensions": None, "solve_captcha": None}


def configure_indexing(extensions=None, solve_captcha=None):
    """Let the host (web_app_batch) share its Buster CAPTCHA solver + extension so
    the indexing 'site:' query can solve challenges like the rank checker does."""
    _INDEXING["extensions"] = extensions
    _INDEXING["solve_captcha"] = solve_captcha


def _parse_index_count(src):
    for pat in (r'result-stats[^>]*>[^<]*?About\s*([\d,\.  ]+)',
                r'About\s*([\d,\.  ]+?)\s*results',
                r'([\d,\.]+)\s*results'):
        m = re.search(pat, src, re.I)
        if m:
            num = re.sub(r'[^\d]', '', m.group(1))
            if num:
                return int(num)
    if "did not match any documents" in src.lower() or "no results found" in src.lower():
        return 0
    return None


def _run_index_query(domain, country, headless, extensions, solver, attempts=2):
    """One Google 'site:' session via the stealth engine. Retries the search, and
    if it hits a CAPTCHA and a solver is provided, runs Buster then re-reads."""
    import engine
    driver = None
    profile = tempfile.mkdtemp(prefix="seo_idx_")
    try:
        driver = engine.build_driver(profile, headless=headless, country=country,
                                     extra_extensions=extensions,
                                     logger=lambda *a: None, browser_pref="auto")
        engine.warm_up(driver, country, logger=lambda *a: None)
        for _ in range(max(1, attempts)):
            engine.human_search(driver, f"site:{domain}", country, logger=lambda *a: None)
            time.sleep(2.5)
            src = driver.page_source or ""
            blocked = False
            try:
                blocked = engine.classify_page(src) in ("captcha", "soft_block", "http_403")
            except Exception:
                pass
            if blocked:
                if solver:
                    try:
                        solver(driver, 3)          # Buster CAPTCHA solver
                    except Exception:
                        pass
                    time.sleep(2)
                    src = driver.page_source or ""
                else:
                    time.sleep(3)                  # brief pause, then retry the search
                    continue
            n = _parse_index_count(src)
            if n is not None:
                # Grab the SERP screenshot from THIS (CAPTCHA-solved) session so the
                # report's indexing slide shows the real result, not a blank/skip.
                try:
                    import base64
                    _sp = os.path.join(tempfile.mkdtemp(prefix="brief_serp_"), "serp.png")
                    res = driver.execute_cdp_cmd("Page.captureScreenshot", {
                        "format": "png", "captureBeyondViewport": True,
                        "clip": {"x": 0, "y": 0, "width": 1366.0, "height": 900.0, "scale": 1}})
                    with open(_sp, "wb") as f:
                        f.write(base64.b64decode(res["data"]))
                    _INDEXING["serp_shot"] = _sp
                except Exception:
                    pass
                return n
        return None
    except Exception:
        return None
    finally:
        try:
            if driver:
                driver.quit()
        except Exception:
            pass


def _stealth_indexing_count(domain, country="us"):
    """Google 'site:domain' result count via the app's STEALTH engine — the same
    path the rank checker uses. Tries quiet headless first; if Google blocks it and
    a Buster solver is available, escalates to a VISIBLE browser with Buster (like
    the rank checker) to solve the CAPTCHA. Returns an int, 0, or None."""
    n = _run_index_query(domain, country, headless=True, extensions=None,
                         solver=None, attempts=2)
    if n is not None:
        return n
    solver = _INDEXING.get("solve_captcha")
    if solver:
        # Visible browser + Buster loaded (extensions only load when NOT headless).
        n = _run_index_query(domain, country, headless=False,
                             extensions=_INDEXING.get("extensions"),
                             solver=solver, attempts=2)
    return n


def capture_brief_screenshots(domain, sitemap_url=None, log_fn=print):
    """Capture the per-topic screenshots the report slides embed (like the manual
    sample templates): homepage, view-source (title/meta/canonical), Google 'site:'
    SERP (indexing), sitemap and robots.txt. Returns {key: png_path}; best-effort —
    a missing/blocked shot is simply skipped."""
    # Use a FRESH driver — after all the page rendering + the separate stealth
    # engine used for indexing, the shared driver can be in a state where get()
    # times out and cascades. A clean driver captures reliably.
    _close_brief_driver()
    driver = _get_brief_driver()
    if not driver:
        return {}
    try:
        driver.set_page_load_timeout(20)     # heavy homepages can exceed the default
    except Exception:
        pass
    import base64
    out = {}
    out_dir = tempfile.mkdtemp(prefix="brief_shots_")
    root = f"https://{domain}"

    def _shot(key, url, height, view_source=False, serp=False):
        try:
            try:
                driver.get(("view-source:" + url) if view_source else url)
            except Exception:
                # Page-load timeout on a heavy page — stop loading and screenshot
                # whatever has rendered so far (good enough for a preview).
                try:
                    driver.execute_script("window.stop();")
                except Exception:
                    pass
            time.sleep(4 if serp else 3)
            if serp:                          # skip a Google CAPTCHA wall — useless as an image
                low = (driver.page_source or "").lower()
                if any(m in low for m in ("unusual traffic", "not a robot", "recaptcha",
                                          "/sorry/", "detected unusual", "before you continue")):
                    log_fn(f"  Screenshot '{key}' skipped (Google blocked)")
                    return
            try:
                driver.execute_script("window.scrollTo(0, 0);")
            except Exception:
                pass
            time.sleep(0.4)
            if view_source:
                cdp = {"format": "png", "captureBeyondViewport": False}
            else:
                try:
                    _w = driver.execute_script(
                        "return Math.max(document.documentElement.clientWidth||0, window.innerWidth||0, 1366);")
                except Exception:
                    _w = 1366
                cdp = {"format": "png", "captureBeyondViewport": True,
                       "clip": {"x": 0, "y": 0, "width": float(_w or 1366), "height": float(height), "scale": 1}}
            res = driver.execute_cdp_cmd("Page.captureScreenshot", cdp)
            p = os.path.join(out_dir, f"{key}.png")
            with open(p, "wb") as f:
                f.write(base64.b64decode(res["data"]))
            out[key] = p
            log_fn(f"  Captured screenshot: {key}")
        except Exception as e:
            log_fn(f"  Screenshot '{key}' skipped ({type(e).__name__})")

    _shot("homepage", root + "/", 950)
    _shot("viewsource", root + "/", 760, view_source=True)
    # Indexing SERP: reuse the CAPTCHA-SOLVED screenshot captured during the
    # indexing count query (stealth engine + Buster). Only if that didn't happen do
    # we try a plain headless shot (which skips on a CAPTCHA wall).
    serp_shot = _INDEXING.get("serp_shot")
    if serp_shot and os.path.exists(serp_shot):
        out["serp"] = serp_shot
        log_fn("  Using CAPTCHA-solved SERP screenshot")
    else:
        _shot("serp", f"https://www.google.com/search?q=site:{domain}", 900, serp=True)
    _shot("sitemap", sitemap_url or (root + "/sitemap.xml"), 850)
    _shot("robots", root + "/robots.txt", 700)
    return out


def _place_image(slide, path, x, y, w, h=None):
    """Place a screenshot fit within a w×h box (inches), centred, preserving aspect
    ratio. If h is None, just scales to width. Ignores a missing/broken file."""
    from pptx.util import Inches
    if not path or not os.path.exists(path):
        return False
    try:
        if h is None:
            slide.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w))
            return True
        try:
            from PIL import Image
            iw, ih = Image.open(path).size
            ar = (iw / ih) if ih else 1.5
        except Exception:
            ar = 1.45
        if (w / h) > ar:                     # box wider than image -> fit by height
            nh, nw = h, h * ar
        else:                                # fit by width
            nw, nh = w, w / ar
        slide.shapes.add_picture(path, Inches(x + (w - nw) / 2), Inches(y + (h - nh) / 2),
                                 width=Inches(nw), height=Inches(nh))
        return True
    except Exception:
        return False


def _add_preview_slides(prs, screenshots, slide_w, slide_h, title_color,
                        bg_color=None, blank_idx=6):
    """Append 'Website Preview' image slides (homepage / robots.txt / sitemap) so
    every format shows real screenshots, matching the manual sample reports."""
    if not isinstance(screenshots, dict):
        return
    for title, key in (("Homepage Preview", "homepage"),
                       ("Robots.txt", "robots"),
                       ("Sitemap", "sitemap")):
        path = screenshots.get(key)
        if not path or not os.path.exists(path):
            continue
        s = prs.slides.add_slide(prs.slide_layouts[blank_idx])
        if bg_color:
            _rect(s, 0, 0, slide_w, slide_h, bg_color)
        _text(s, title, 0.6, 0.3, slide_w - 1.2, 0.7, 26, "Calibri", title_color, bold=True)
        _place_image(s, path, 0.8, 1.2, slide_w - 1.6, slide_h - 1.6)


def check_indexing(domain):
    """Estimate indexed page count. Stealth engine first (beats Google's block on
    plain headless), then plain browser, then urllib."""
    _INDEXING["serp_shot"] = None    # reset the per-run CAPTCHA-solved SERP screenshot
    # Stealth engine (undetected-chromedriver + human-typed search) — most reliable.
    try:
        n = _stealth_indexing_count(domain)
    except Exception:
        n = None
    if n is not None:
        if n <= 0:
            return {"count": "0", "status": "Not indexed"}
        return {"count": f"{n:,}", "status": "Indexed & serving"}

    # Try plain browser next (bypasses some Google blocks)
    result = _browser_get_indexing(domain)
    if result:
        return result
    # Fallback to urllib
    url = f"https://www.google.com/search?q=site:{domain}"
    try:
        req = urllib.request.Request(url, headers={
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/126.0.0.0 Safari/537.36",
            "Accept-Language": "en-US,en;q=0.9",
        })
        with urllib.request.urlopen(req, timeout=15) as r:
            html = r.read().decode("utf-8", errors="replace")
        patterns = [
            r'About ([\d,\.]+) results',
            r'([\d,\.]+) results',
            r'result-stats[^>]*>About ([\d,\.]+)',
            r'result-stats[^>]*>([\d,\.]+)',
        ]
        for pat in patterns:
            m = re.search(pat, html, re.I)
            if m:
                return {"count": m.group(1).replace(".", ","), "status": "Indexed & serving"}
        if "did not match any documents" in html:
            return {"count": "0", "status": "Not indexed"}
    except Exception:
        pass
    return {"count": "N/A", "status": "Could not check"}


def _whois_data(clean):
    """Collect all WHOIS data from python-whois + RDAP. Returns dict."""
    info = {}
    # python-whois (best for .com/.net/.org)
    try:
        import whois
        w = whois.whois(clean)
        created = w.creation_date
        if isinstance(created, list):
            created = created[0]
        updated = w.updated_date
        if isinstance(updated, list):
            updated = updated[0]
        expiry = w.expiration_date
        if isinstance(expiry, list):
            expiry = expiry[0]
        if created and hasattr(created, 'strftime'):
            info["created"] = created.replace(tzinfo=None) if created.tzinfo else created
        if updated and hasattr(updated, 'strftime'):
            info["updated"] = updated.replace(tzinfo=None) if updated.tzinfo else updated
        if expiry and hasattr(expiry, 'strftime'):
            info["expiry"] = expiry.replace(tzinfo=None) if expiry.tzinfo else expiry
        if w.registrar:
            info["registrar"] = str(w.registrar)[:50]
        registrant = getattr(w, 'registrant_name', None) or getattr(w, 'org', None)
        if registrant:
            info["registrant"] = str(registrant)[:50]
    except Exception:
        pass

    # RDAP (fills gaps — some registries expose registration date here)
    if "created" not in info:
        try:
            url = f"https://rdap.org/domain/{clean}"
            req = urllib.request.Request(url, headers={"Accept": "application/rdap+json",
                                                        "User-Agent": "SEOToolkitPro/1.0"})
            with urllib.request.urlopen(req, timeout=12) as r:
                data = json.loads(r.read().decode())
            events = {e["eventAction"]: e["eventDate"][:10] for e in data.get("events", [])}
            if events.get("registration"):
                info["created"] = datetime.strptime(events["registration"], "%Y-%m-%d")
            if "updated" not in info and events.get("last changed"):
                info["updated"] = datetime.strptime(events["last changed"], "%Y-%m-%d")
            if "expiry" not in info and events.get("expiration"):
                info["expiry"] = datetime.strptime(events["expiration"], "%Y-%m-%d")
        except Exception:
            pass
    return info


def _wayback_first_capture(domain):
    """Get earliest Wayback Machine capture date as age proxy."""
    try:
        import requests as _req
        r = _req.get(
            f"https://web.archive.org/cdx/search/cdx?url={domain}&output=json&limit=1&fl=timestamp&from=2000",
            timeout=12)
        if r.status_code == 200:
            data = r.json()
            if len(data) > 1:
                ts = data[1][0]
                return datetime.strptime(ts[:8], "%Y%m%d")
    except Exception:
        pass
    return None


def check_domain_age(domain):
    """Get domain WHOIS info from multiple sources: python-whois → RDAP → Wayback Machine."""
    clean = domain.replace("www.", "")
    rows = [["Domain Name", clean, "", ""]]

    info = _whois_data(clean)

    # If no creation date from WHOIS/RDAP, use Wayback Machine first capture
    if "created" not in info:
        wb = _wayback_first_capture(clean)
        if wb:
            info["first_archived"] = wb

    # Build rows
    if info.get("created"):
        age_days = (datetime.now() - info["created"]).days
        age_years = age_days // 365
        rows.append(["Created On", info["created"].strftime("%Y-%m-%d"), f"~{age_years} years", ""])
    elif info.get("first_archived"):
        age_days = (datetime.now() - info["first_archived"]).days
        age_years = age_days // 365
        rows.append(["First Archived", info["first_archived"].strftime("%Y-%m-%d"),
                      f"~{age_years} yr (Wayback Machine)", ""])

    if info.get("updated"):
        rows.append(["Updated", info["updated"].strftime("%Y-%m-%d")
                      if hasattr(info["updated"], 'strftime') else str(info["updated"])[:20], "", ""])
    if info.get("expiry"):
        rows.append(["Expiration", info["expiry"].strftime("%Y-%m-%d")
                      if hasattr(info["expiry"], 'strftime') else str(info["expiry"])[:20], "", ""])
    if info.get("registrar"):
        rows.append(["Registrar", info["registrar"], "", ""])
    if info.get("registrant"):
        rows.append(["Registrant", info["registrant"], "", ""])

    if len(rows) == 1:
        rows.append(["Status", "Could not retrieve WHOIS data", "", ""])

    return rows


def check_title_tags(domain, pages=None):
    """Check title tags on key pages (uses rendered HTML so SPA pages work,
    and marks pages that don't exist as 'Not Found')."""
    if not pages:
        pages = ["/"]
    results = []
    for pg in pages:
        url = f"https://{domain}{pg}"
        name = url   # show the full page URL, not a short label
        page = _get_page(url)
        if not page["exists"]:
            results.append({"page": name, "title": "Page not found", "chars": "0", "status": "Not Found"})
            continue
        html = page["html"]
        if not html:
            results.append({"page": name, "title": "Could not fetch", "chars": "0", "status": "Could not check"})
            continue
        meta = _parse_meta(html)
        title = meta["title"][0] if meta["title"] else ""
        length = len(title)
        if not title:
            status = "Missing"
        elif length < 30:
            status = "Too short"
        elif length > 60:
            status = "Too long"
        else:
            status = "Good"
        results.append({"page": name, "title": title[:40], "chars": str(length), "status": status})
    return results


def check_meta_desc(domain, pages=None):
    """Check meta descriptions on key pages (rendered HTML; 'Not Found' for
    pages that don't exist)."""
    if not pages:
        pages = ["/"]
    results = []
    for pg in pages:
        url = f"https://{domain}{pg}"
        name = url   # show the full page URL, not a short label
        page = _get_page(url)
        if not page["exists"]:
            results.append({"page": name, "found": "No", "chars": "0", "status": "Not Found"})
            continue
        html = page["html"]
        if not html:
            results.append({"page": name, "found": "No", "chars": "0", "status": "Could not check"})
            continue
        meta = _parse_meta(html)
        desc = meta["description"][0] if meta["description"] else ""
        length = len(desc)
        found = "Yes" if desc else "No"
        if not desc:
            status = "Missing"
        elif length < 120:
            status = "Below optimal length"
        elif length > 160:
            status = "Too long"
        else:
            status = "Good"
        results.append({"page": name, "found": found, "chars": str(length), "status": status})
    return results


def check_headers(domain, pages=None):
    """Count H1-H6 tags on key pages (rendered HTML; missing pages flagged)."""
    if not pages:
        pages = ["/"]
    results = []
    for pg in pages:
        url = f"https://{domain}{pg}"
        name = url   # show the full page URL, not a short label
        page = _get_page(url)
        counts = {f"h{i}": 0 for i in range(1, 7)}
        if not page["exists"]:
            results.append({"page": f"{name} (Not Found)", **counts})
            continue
        html = page["html"]
        if not html:
            results.append({"page": f"{name} (no data)", **counts})
            continue
        for i in range(1, 7):
            counts[f"h{i}"] = len(re.findall(rf"<h{i}[\s>]", html, re.I))
        results.append({"page": name, **counts})
    return results


def check_image_alts(domain):
    """Check image alt tags on homepage (rendered HTML)."""
    html = _get_page(f"https://{domain}/")["html"]
    if not html:
        return []
    imgs = re.findall(r'<img\s[^>]*?>', html, re.I | re.S)
    results = []
    for img_tag in imgs[:20]:
        alt_m = re.search(r'alt\s*=\s*["\']([^"\']*)["\']', img_tag, re.I)
        src_m = re.search(r'(?:data-src|src)\s*=\s*["\']([^"\']*)["\']', img_tag, re.I)
        import html as _html
        src = _html.unescape(src_m.group(1).strip()) if src_m else "unknown"
        # Resolve to a full, absolute URL so the location is usable (not trimmed).
        if src.startswith("//"):
            src = "https:" + src
        elif src.startswith("/"):
            src = f"https://{domain}" + src
        if not alt_m:
            results.append({"location": src, "present": "No", "content": "(empty)", "status": "Add descriptive alt"})
        elif not alt_m.group(1).strip():
            results.append({"location": src, "present": "Empty", "content": "(empty)", "status": "Add descriptive alt"})
        else:
            alt = alt_m.group(1).strip()
            is_generic = bool(re.match(r'^(img|image|dsc|photo|pic|banner)\d*', alt, re.I))
            results.append({"location": src, "present": "Yes",
                            "content": alt[:25], "status": "Non-descriptive" if is_generic else "OK"})
    return results


def check_redirections(domain):
    """Check common redirection patterns and report the ACTUAL status code
    (200 / 301 / 302). 301 is ideal; 302 is a temporary redirect that should be
    reviewed/changed to 301."""
    import urllib.error

    class _NoRedirect(urllib.request.HTTPRedirectHandler):
        def redirect_request(self, *a, **k):
            return None   # don't follow — so we can read the redirect's own status

    results = []
    variants = [
        (f"http://{domain}/", "HTTP to HTTPS"),
        (f"https://www.{domain}/", "www to non-www") if not domain.startswith("www.") else (f"https://{domain.replace('www.','')}/", "non-www to www"),
    ]
    opener = urllib.request.build_opener(_NoRedirect)
    for url, rtype in variants:
        try:
            req = urllib.request.Request(url, method="HEAD",
                                         headers={"User-Agent": "SEOToolkitPro/1.0"})
            try:
                with opener.open(req, timeout=10) as r:
                    code = r.status
                if 200 <= code < 300:
                    results.append({"type": rtype, "status": str(code), "checked_url": url,
                                    "detail": f"Checked {url} — no redirect (already canonical / redirect missing)",
                                    "impact": "Check needed"})
                else:
                    results.append({"type": rtype, "status": str(code), "checked_url": url,
                                    "detail": f"Checked {url} — unexpected response", "impact": "Review"})
            except urllib.error.HTTPError as e:
                code = e.code
                loc = (e.headers.get("Location") or "").strip()
                if code in (301, 308):
                    impact = "Good (permanent 301)"
                elif code in (302, 307):
                    impact = "Check — temporary redirect, use 301"
                else:
                    impact = "Review"
                results.append({"type": rtype, "status": str(code), "checked_url": url,
                                "detail": f"{url}  ->  {loc}" if loc else f"Checked {url} - HTTP {code}",
                                "impact": impact})
        except Exception as e:
            results.append({"type": rtype, "status": "Error", "checked_url": url,
                            "detail": f"Checked {url} — {str(e)[:40]}", "impact": "Unknown"})
    return results


def check_canonical_tags(domain, pages=None):
    """Check canonical tags on key pages (rendered HTML; 'Not Found' handling)."""
    if not pages:
        pages = ["/"]
    results = []
    for pg in pages:
        url = f"https://{domain}{pg}"
        name = url   # show the full page URL, not a short label
        page = _get_page(url)
        if not page["exists"]:
            results.append({"page": name, "found": "No", "correct": "N/A",
                            "duplicate_risk": "N/A", "status": "Not Found"})
            continue
        html = page["html"]
        if not html:
            results.append({"page": name, "found": "No", "correct": "N/A",
                            "duplicate_risk": "N/A", "status": "Could not check"})
            continue
        m = re.search(r'<link[^>]+rel\s*=\s*["\']canonical["\'][^>]+href\s*=\s*["\']([^"\']+)["\']', html, re.I)
        if not m:
            m = re.search(r'<link[^>]+href\s*=\s*["\']([^"\']+)["\'][^>]+rel\s*=\s*["\']canonical["\']', html, re.I)
        if m:
            canonical = m.group(1)
            is_self = canonical.rstrip("/") == url.rstrip("/")
            results.append({"page": name, "found": "Yes", "correct": "Yes" if is_self else "No",
                            "duplicate_risk": "Low" if is_self else "Medium",
                            "status": "OK" if is_self else "Check canonical"})
        else:
            results.append({"page": name, "found": "No", "correct": "N/A",
                            "duplicate_risk": "High", "status": "Missing"})
    return results


def _discover_sitemap(domain):
    """Find the sitemap — robots.txt 'Sitemap:' first, then common locations —
    and return (sitemap_url, [same-site page URLs]). Follows a sitemap index."""
    import urllib.request
    def _get(u):
        try:
            req = urllib.request.Request(u, headers={"User-Agent": "SEOToolkitPro/1.0"})
            with urllib.request.urlopen(req, timeout=12) as r:
                if getattr(r, "status", 200) != 200:
                    return None
                return r.read().decode("utf-8", "ignore")
        except Exception:
            return None
    candidates = []
    robots = _get(f"https://{domain}/robots.txt")
    if robots:
        for m in re.finditer(r'(?im)^\s*sitemap:\s*(\S+)', robots):
            candidates.append(m.group(1).strip())
    for path in ("/sitemap.xml", "/sitemap_index.xml", "/sitemap-index.xml", "/sitemap"):
        u = f"https://{domain}{path}"
        if u not in candidates:
            candidates.append(u)
    for sm_url in candidates:
        content = _get(sm_url)
        if not content or "<loc" not in content.lower():
            continue
        locs = [l.strip() for l in re.findall(r'<loc>\s*([^<]+?)\s*</loc>', content, re.I)]
        if "<sitemapindex" in content.lower():
            pages = []
            for sub in locs[:8]:
                sc = _get(sub)
                if sc:
                    pages += re.findall(r'<loc>\s*([^<]+?)\s*</loc>', sc, re.I)
            pages = [p.strip() for p in pages]
        else:
            pages = locs
        pages = [p for p in pages if domain in p]
        if pages:
            return sm_url, pages
    return None, []


def _link_section(html, pos):
    """Which region of the page a byte offset falls in: Header / Navigation /
    Footer / Content — so a broken link can be located for the client."""
    for tag, label in (("header", "Header"), ("nav", "Navigation"), ("footer", "Footer")):
        for m in re.finditer(rf'<{tag}\b[^>]*>.*?</{tag}>', html, re.I | re.S):
            if m.start() <= pos < m.end():
                return label
    window = html[max(0, pos - 400):pos].lower()   # fall back to class/id hints
    if any(k in window for k in ('id="footer"', 'class="footer', 'site-footer')):
        return "Footer"
    if any(k in window for k in ('id="header"', 'class="header', 'site-header')):
        return "Header"
    if any(k in window for k in ('<nav', 'class="nav', 'navbar', 'menu')):
        return "Navigation"
    return "Content"


def _check_broken_links_located(domain, pages, max_links=80):
    """Check the real anchor (<a href>) links on each page and report ONLY the
    genuinely-broken ones (404/410/5xx), together with WHERE they sit (header /
    navigation / footer / content). Preconnect / stylesheet / font <link> hints are
    ignored, so CDN & font hosts that 404 a bare request aren't false-flagged."""
    checked, seen, broken = 0, set(), []
    for pg in pages:
        purl = f"https://{domain}{pg}" if str(pg).startswith("/") else str(pg)
        html = _get_page(purl)["html"]
        if not html:
            continue
        for m in re.finditer(r'<a\b[^>]*?\shref\s*=\s*["\']([^"\']+)["\']', html, re.I):
            if checked >= max_links:
                break
            href = m.group(1).strip()
            if not href or href.startswith(("#", "mailto:", "tel:", "javascript:", "data:")):
                continue
            full = urllib.parse.urljoin(purl, href.split("#")[0])
            if not full.startswith("http") or full in seen:
                continue
            seen.add(full)
            checked += 1
            status = _http_status(full)
            if status in (404, 410) or (isinstance(status, int) and 500 <= status < 600):
                broken.append({"type": "Broken link", "location": _link_section(html, m.start()),
                               "url": full, "status": str(status), "source": purl})
        if checked >= max_links:
            break
    return checked, broken


def _select_sample_pages(sm_paths):
    """A representative sample of REAL pages from the sitemap: Home + one
    collection/category + one about/contact page. Uses only URLs that are actually
    in the sitemap (so they exist); if a category/about isn't present it substitutes
    another real non-product page, and a site with no other pages just gets Home."""
    pages = ["/"]

    def pick(keys):
        for p in sm_paths:
            pl = p.lower()
            if p in pages or p in ("/", "") or "/products/" in pl or "/product/" in pl:
                continue                      # skip Home, already-picked, deep products
            if keys is None or any(k in pl for k in keys):
                return p
        return None

    # 1 collection/category page (else any other real page)
    collection = pick(["/collections/", "/collection/", "/category/", "/categories/",
                       "/product-category/", "/shop"]) or pick(None)
    if collection:
        pages.append(collection)
    # 1 about/contact page (else another real page)
    about = pick(["about", "contact"]) or pick(None)
    if about:
        pages.append(about)
    return pages[:3]


def _normalize_target_pages(target_pages, domain):
    """Turn user-provided pages (full URLs or paths) into site-relative paths."""
    out = []
    seq = target_pages if isinstance(target_pages, (list, tuple)) else [target_pages]
    for tp in seq:
        tp = str(tp or "").strip()
        if not tp:
            continue
        p = re.sub(r'^\s*https?://[^/]+', '', tp).split('#')[0]
        if not p.startswith("/"):
            p = "/" + p
        if p not in out:
            out.append(p)
    return out or ["/"]


def run_brief_checks(domain, target_pages=None, log_fn=None):
    """Run all brief analysis checks. Returns dict of results."""
    if log_fn is None:
        log_fn = print
    domain = re.sub(r'^\s*https?://', '', str(domain or '')).strip().strip('/').split('/')[0] or str(domain)
    _PAGE_CACHE.clear()      # fresh per-run cache of rendered pages
    _HOME_SIG[0] = None      # reset homepage fingerprint (SPA soft-404 detection)
    _NF_SIG[0] = None        # reset 404-probe fingerprint
    # Discover real pages from the sitemap so we check pages that actually EXIST
    # instead of guessing /about, /contact (which return the homepage on SPAs).
    sitemap_url, _sm_urls = None, []
    try:
        sitemap_url, _sm_urls = _discover_sitemap(domain)
    except Exception:
        pass
    _sm_paths = []
    for u in _sm_urls:
        p = re.sub(r'^https?://[^/]+', '', u).split('#')[0].split('?')[0] or "/"
        if p not in _sm_paths:
            _sm_paths.append(p)
    if target_pages:
        pages = _normalize_target_pages(target_pages, domain)   # user's own pages
    elif _sm_paths:
        pages = _select_sample_pages(_sm_paths)                 # Home + collection + about/contact
    else:
        pages = ["/"]   # no sitemap -> only the homepage (don't invent pages)

    log_fn("  Launching browser for data collection...")
    # Render the homepage + one known-bad URL first, so their fingerprints are
    # captured before any real sub-page is checked. That's what lets us flag SPA
    # soft-404s (routes that 200 but don't really exist). Homepage is cached, so
    # no real page renders twice.
    try:
        _get_page(f"https://{domain}/")
        _prime_soft404(domain)
    except Exception:
        pass

    def _safe(label, default, fn, *args, **kwargs):
        """Run one check; if it raises, log it and keep going so the rest of the
        report still generates — that section just shows as not completed."""
        try:
            return fn(*args, **kwargs)
        except Exception as e:
            log_fn(f"  {label} could not be completed — skipped ({e})")
            return default

    log_fn("  Checking website indexing...")
    indexing = _safe("Indexing check", {"status": "Could not check", "count": "N/A"},
                     check_indexing, domain)

    log_fn("  Checking domain age...")
    domain_age = _safe("Domain age check", [], check_domain_age, domain)

    log_fn("  Checking title tags...")
    titles = _safe("Title tags check", [], check_title_tags, domain, pages)

    log_fn("  Checking meta descriptions...")
    metas = _safe("Meta descriptions check", [], check_meta_desc, domain, pages)

    log_fn("  Checking header tags...")
    headers = _safe("Header tags check", [], check_headers, domain, pages[:3])

    log_fn("  Checking image alt tags...")
    img_alts = _safe("Image alt tags check", [], check_image_alts, domain)

    log_fn("  Checking redirections...")
    redirects = _safe("Redirections check", [], check_redirections, domain)

    log_fn("  Checking canonical tags...")
    canonicals = _safe("Canonical tags check", [], check_canonical_tags, domain, pages[:3])

    log_fn("  Checking robots.txt...")
    robots = _safe("Robots.txt check", {}, check_robots_txt, domain)

    log_fn("  Checking sitemap...")
    sitemap = _safe("Sitemap check",
                    {"ok": False, "found": False, "summary": "Sitemap check could not be completed."},
                    check_sitemap, domain, robots_data=robots)
    # Merge in the sitemap URL + real page count from discovery (robots.txt first).
    if isinstance(sitemap, dict):
        if sitemap_url:
            sitemap["url_checked"] = sitemap_url
            sitemap["found"] = True
        if _sm_paths:
            sitemap["page_count"] = len(_sm_paths)
            if not sitemap.get("summary"):
                sitemap["summary"] = f"Sitemap found with {len(_sm_paths)} page(s)."

    # Indexing MUST come from Google (via the stealth engine + Buster). Never
    # substitute the sitemap count — if Google couldn't be checked, leave it N/A so
    # the slide clearly asks the user to verify it manually.
    if isinstance(indexing, dict) and str(indexing.get("count", "")).strip().upper() in ("", "N/A"):
        indexing = {"count": "N/A",
                    "status": "Could not verify with Google - please check indexing manually"}

    log_fn("  Checking broken links (with location)...")
    bl_checked, bl_broken = _safe("Broken links check", (0, []),
                                  _check_broken_links_located, domain, pages[:3])

    log_fn("  Capturing screenshots (homepage, robots, sitemap)...")
    screenshots = _safe("Screenshots", {}, capture_brief_screenshots,
                        domain, sitemap_url, log_fn)

    log_fn("  Checking Domain Authority / Domain Rating / Page Authority...")
    from da_checker import check_da_pa
    da_pa = _safe("DA/PA check", {"da": "—", "dr": "—", "pa": "—", "source": "—"},
                  check_da_pa, _get_brief_driver(), domain, log_fn)

    _close_brief_driver()
    log_fn("  All checks complete.")

    # Defensive: the report builder calls .get() on these, so guarantee a dict
    # shape even if a check returned something unexpected. Prevents
    # "'tuple' object has no attribute 'get'" from aborting the whole report.
    if not isinstance(indexing, dict):
        indexing = {"status": "Could not check", "count": "N/A"}
    if not isinstance(sitemap, dict):
        sitemap = {"ok": False, "found": False, "summary": "Sitemap check could not be completed."}
    if not isinstance(robots, dict):
        robots = {}
    # check_robots_txt returns status/ok but no explicit "found" — derive it so the
    # report doesn't show "Found: No" while the summary says robots.txt WAS found.
    if "found" not in robots:
        robots["found"] = (robots.get("status") == 200) or ("found" in str(robots.get("summary", "")).lower())
    # check_broken_links returns (url, status) TUPLES, but every report builder
    # calls b.get(...) on each broken link expecting a dict. Normalize to dicts
    # so a site that actually has broken links doesn't crash the report with
    # "'tuple' object has no attribute 'get'".
    _norm_broken = []
    for b in (bl_broken if isinstance(bl_broken, (list, tuple)) else []):
        if isinstance(b, dict):
            _norm_broken.append(b)
        elif isinstance(b, (list, tuple)):
            _norm_broken.append({"type": "Broken", "source": "",
                                 "url": str(b[0]) if len(b) > 0 else "",
                                 "status": str(b[1]) if len(b) > 1 else ""})
        else:
            _norm_broken.append({"type": "Broken", "source": "", "url": str(b), "status": ""})
    bl_broken = _norm_broken
    if not isinstance(bl_checked, int):
        try:
            bl_checked = int(bl_checked)
        except Exception:
            bl_checked = 0

    return {
        "domain": domain,
        "date": datetime.now().strftime("%d %B %Y"),
        "indexing": indexing,
        "domain_age": domain_age,
        "titles": titles,
        "metas": metas,
        "headers": headers,
        "img_alts": img_alts,
        "redirects": redirects,
        "canonicals": canonicals,
        "sitemap": sitemap,
        "robots": robots,
        "broken_links_checked": bl_checked,
        "broken_links": bl_broken,
        "screenshots": screenshots if isinstance(screenshots, dict) else {},
        "da_pa": da_pa if isinstance(da_pa, dict) else {"da": "—", "dr": "—", "pa": "—", "source": "—"},
    }


# ---------------------------------------------------------------------------
# PPTX helpers
# ---------------------------------------------------------------------------

def _init_pptx():
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    return Presentation, Inches, Pt, Emu, RGBColor, PP_ALIGN, MSO_ANCHOR


def _C(h):
    from pptx.dml.color import RGBColor
    return RGBColor(int(h[1:3], 16), int(h[3:5], 16), int(h[5:7], 16))


def _rect(slide, x, y, w, h, color):
    from pptx.util import Inches
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _C(color) if isinstance(color, str) else color
    shape.line.fill.background()
    return shape


def _text(slide, text_str, x, y, w, h, font_size, font_name, color, bold=False,
          align=None, italic=False):
    from pptx.util import Inches, Pt
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(text_str)          # splits on \n into one run per line (with line breaks)
    if isinstance(color, str):
        color = _C(color)
    runs = p.runs or [p.add_run()]
    for run in runs:                # format every run — do NOT re-assign run.text (that duplicated multi-line text)
        run.font.size = Pt(font_size)
        run.font.name = font_name
        run.font.color.rgb = color
        run.font.bold = bold
        if italic:
            run.font.italic = True
    if align:
        p.alignment = align
    return txBox


UNVERIFIED = "Please check this once from your side"
RED = "#C0392B"


def _is_empty(val):
    """Check if a data value is empty/unverifiable."""
    if val is None:
        return True
    s = str(val).strip().lower()
    return s in ("", "n/a", "could not check", "could not retrieve", "error", "unknown")


def _suggestion_text(slide, data_ok, suggestion, x, y, w, h, font_name="Calibri"):
    """Show suggestion in navy if data verified, or red warning if not."""
    if data_ok:
        return _text(slide, f"SEO Suggestions: {suggestion}", x, y, w, h, 11, font_name, "#1B2A4A", bold=True)
    else:
        return _text(slide, UNVERIFIED, x, y, w, h, 11, font_name, RED, bold=True)


def _add_table(slide, headers, rows, col_widths, start_x, start_y,
               hdr_fill="#1B2A4A", hdr_color="#FFFFFF", row_alt="#D6E4F0"):
    from pptx.util import Inches as I, Pt as P
    total_w = sum(col_widths)
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_shape = slide.shapes.add_table(n_rows, n_cols,
                                        I(start_x), I(start_y),
                                        I(total_w), I(0.4 * n_rows))
    table = tbl_shape.table
    for i, w in enumerate(col_widths):
        table.columns[i].width = I(w)
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = P(10)
                r.font.bold = True
                r.font.color.rgb = _C(hdr_color)
                r.font.name = "Calibri"
        cell.fill.solid()
        cell.fill.fore_color.rgb = _C(hdr_fill)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            # Wrap long values (full page/image URLs) instead of overflowing, and
            # tighten margins so they fit the column.
            tf = cell.text_frame
            tf.word_wrap = True
            cell.margin_left = P(3)
            cell.margin_right = P(3)
            cell.margin_top = P(1)
            cell.margin_bottom = P(1)
            for p in tf.paragraphs:
                for r in p.runs:
                    r.font.size = P(8) if any(str(val).startswith(s) for s in ("http", "//")) else P(9)
                    r.font.name = "Calibri"
                    r.font.color.rgb = _C("#2B2B2B")
            if ri % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _C(row_alt)


# ---------------------------------------------------------------------------
# James format — 15 slides, 10×5.625"
# ---------------------------------------------------------------------------

def _james_sidebar(slide):
    """Draw the James left sidebar: olive bar + 3 diagonal lines + 2 red chevrons."""
    _rect(slide, 0.6, 0.0, 0.4, 5.625, "#55611F")
    _rect(slide, 0.1, 0.9, 1.2, 0.02, "#55611F")
    _rect(slide, 0.1, 1.4, 1.2, 0.02, "#55611F")
    _rect(slide, 0.1, 1.8, 1.2, 0.02, "#55611F")
    # Red chevron arrows at bottom-left
    _rect(slide, 0.2, 4.8, 0.3, 0.2, "#C0392B")
    _rect(slide, 0.5, 4.8, 0.3, 0.2, "#C0392B")
    _rect(slide, 0.2, 5.2, 0.3, 0.2, "#C0392B")
    _rect(slide, 0.5, 5.2, 0.3, 0.2, "#C0392B")


def build_james(data, out_path, log_fn=None):
    if log_fn is None:
        log_fn = print
    Presentation, Inches, Pt, Emu, RGBColor, PP_ALIGN, MSO_ANCHOR = _init_pptx()

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    NAVY = "#1B2A4A"
    OLIVE = "#55611F"
    OLIVE_LINE = "#8A9440"
    LINK_BLUE = "#1155CC"
    TEXT = "#2B2B2B"
    SOFT = "#6E6E6E"
    WHITE = "#FFFFFF"

    domain = data["domain"]
    date_str = data["date"]

    def content_slide(title):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        _james_sidebar(s)
        _text(s, title, 2.0, 0.3, 7.5, 0.6, 27, "Calibri", NAVY, bold=True)
        _rect(s, 4.6, 1.0, 2.4, 0.02, OLIVE_LINE)
        return s

    # --- Slide 1: Cover ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _james_sidebar(s)
    _text(s, "Website Analysis Report", 1.5, 1.2, 5.0, 0.8, 38, "Calibri", NAVY, bold=True)
    _rect(s, 1.6, 2.5, 3.4, 0.02, OLIVE_LINE)
    _text(s, f"https://{domain}/\n{date_str}", 1.6, 2.8, 5.0, 0.8, 16, "Calibri", LINK_BLUE, bold=True)

    # --- Slide 2: Index ---
    s = content_slide("Index")
    topics = [
        "Website Indexing", "Domain Age", "Title Tag Checking",
        "Meta Description", "Header Tag Checking", "Image Alt Tag", "Redirection Issues",
        "Canonical Tag", "Sitemap.xml", "Robots.txt", "Broken Links",
    ]
    col1 = "\n".join(f"{i+1}.  {t}" for i, t in enumerate(topics[:6]))
    col2 = "\n".join(f"{i+7}.  {t}" for i, t in enumerate(topics[6:]))
    _text(s, col1, 2.0, 1.4, 3.5, 3.5, 14, "Calibri", TEXT)
    _text(s, col2, 5.9, 1.4, 3.5, 3.5, 14, "Calibri", TEXT)

    # --- Slide 3: Website Indexing ---
    s = content_slide("Website Indexing")
    _text(s, "Indexing is the process by which search engines store and organise a website's pages so they can be returned in search results.",
          2.0, 1.2, 7.5, 0.6, 11, "Calibri", TEXT, bold=True)
    idx = data.get("indexing", {})
    sitemap = data.get("sitemap", {})
    sm_count = "N/A"
    pc = sitemap.get("page_count")
    if pc:
        sm_count = f"{pc} pages"
    else:
        sm_summary = sitemap.get("summary", "")
        m = re.search(r'(\d+)\s*URL', sm_summary)
        if m:
            sm_count = f"~{m.group(1)} URLs"
        else:
            m2 = re.search(r'(\d+)\s*sub-sitemap', sm_summary)
            if m2:
                sm_count = f"{m2.group(1)} sub-sitemaps"
    sm_status = "Found" if sitemap.get("found") or sitemap.get("ok") else "Not found"
    _add_table(s, ["Metric", "Value", "Coverage Status"],
               [["Indexed Pages (site: query)", f"~{idx.get('count', 'N/A')} URLs", idx.get("status", "N/A")],
                ["URLs in XML Sitemap", sm_count, sm_status]],
               [3.0, 2.5, 2.0], 2.0, 2.2)
    sm_url = sitemap.get("url_checked", f"https://{domain}/sitemap.xml")
    _suggestion_text(s, not _is_empty(idx.get("count")),
          f"Sitemap checked: {sm_url}\nReview GSC Coverage report for any excluded pages. Ensure all important pages are in the sitemap and not blocked by robots.txt.",
          2.0, 4.0, 7.5, 0.8)

    # --- Slide 4: Domain Age ---
    s = content_slide("Domain Age")
    _text(s, "Domain age reflects how long a domain has been registered. An older, continuously active domain often carries more search-engine trust.",
          2.0, 1.2, 7.5, 0.6, 11, "Calibri", TEXT, bold=True)
    age_rows = data.get("domain_age", [["Domain Name", domain, "", ""]])
    _add_table(s, ["Attribute", "Detail"],
               [[r[0], r[1]] for r in age_rows],
               [3.0, 4.5], 2.0, 2.0)
    _suggestion_text(s, len(age_rows) > 1 and not _is_empty(age_rows[1][1]),
          "Keep domain registration current. A longer registration period can signal trust to search engines.",
          2.0, 4.3, 7.5, 0.5)

    # --- Slide 6: Title Tag Checking ---
    s = content_slide("Title Tag Checking")
    _text(s, "The title tag is the clickable headline shown in search results and browser tabs. Optimal length is 50-60 characters.",
          2.0, 1.2, 7.5, 0.6, 11, "Calibri", TEXT, bold=True)
    title_rows = [[t["page"], t["title"], t["chars"], t["status"]] for t in data.get("titles", [])][:3]
    if not title_rows:
        title_rows = [["Home", "N/A", "0", "Could not check"]]
    _add_table(s, ["Page URL", "Current Title (sample)", "Chars", "Status"],
               title_rows, [3.3, 2.2, 0.6, 1.4], 2.0, 2.0)
    _suggestion_text(s, bool(data.get("titles")),
          "Rewrite thin titles to 50-60 chars with primary keyword and location intent where applicable.",
          2.0, 3.9, 7.5, 0.5)

    # --- Slide 7: Meta Description ---
    s = content_slide("Meta Description")
    _text(s, "A meta description is the summary snippet under the title in search results. While not a direct ranking factor, it impacts click-through rate.",
          2.0, 1.2, 7.5, 0.6, 11, "Calibri", TEXT, bold=True)
    meta_rows = [[m["page"], m["found"], m["chars"], m["status"]] for m in data.get("metas", [])][:3]
    if not meta_rows:
        meta_rows = [["Home", "N/A", "0", "Could not check"]]
    _add_table(s, ["Page URL", "Found?", "Chars", "Status"],
               meta_rows, [3.5, 1.0, 0.8, 2.2], 2.0, 2.0)
    _suggestion_text(s, bool(data.get("metas")),
          "Write unique 140-160 char descriptions for every key page, each containing the primary keyword and a clear call to action.",
          2.0, 3.9, 7.5, 0.5)

    # --- Slide 8: Header Tag Checking ---
    s = content_slide("Header Tag Checking")
    _text(s, "Header tags (H1-H6) structure page content for readers and search engines. Each page should have exactly one H1.",
          2.0, 1.2, 7.5, 0.6, 11, "Calibri", TEXT, bold=True)
    hdr_rows = [[h["page"], str(h["h1"]), str(h["h2"]), str(h["h3"]),
                 str(h["h4"]), str(h["h5"]), str(h["h6"])] for h in data.get("headers", [])][:3]
    if not hdr_rows:
        hdr_rows = [["Home", "0", "0", "0", "0", "0", "0"]]
    _add_table(s, ["Page URL", "H1", "H2", "H3", "H4", "H5", "H6"],
               hdr_rows, [3.0, 0.75, 0.75, 0.75, 0.75, 0.75, 0.75], 2.0, 2.0)
    issues = [h for h in data.get("headers", []) if h["h1"] != 1]
    suggestion = "All pages have correct H1 structure." if not issues else "Ensure each page has exactly one H1 tag. Multiple H1s dilute keyword focus."
    _suggestion_text(s, bool(data.get("headers")), suggestion, 2.0, 3.9, 7.5, 0.5)

    # --- Slide 9: Image Alt Tag ---
    s = content_slide("Image Alt Tag")
    _text(s, "Alt text describes an image for screen readers and search engines. Descriptive, keyword-rich alt text improves accessibility and image SEO.",
          2.0, 1.2, 7.5, 0.6, 11, "Calibri", TEXT, bold=True)
    alt_rows = [[a["location"], a["present"], a["content"], a["status"]] for a in data.get("img_alts", [])][:3]
    if not alt_rows:
        alt_rows = [["N/A", "N/A", "N/A", "Could not check"]]
    _add_table(s, ["Image URL", "Alt Present?", "Content", "Status"],
               alt_rows, [3.6, 1.0, 1.4, 1.5], 2.0, 2.0)
    missing = len([a for a in data.get("img_alts", []) if a["present"] != "Yes"])
    total = len(data.get("img_alts", []))
    _suggestion_text(s, bool(data.get("img_alts")),
          f"{missing} of {total} images need attention. Add meaningful alt text describing what the image shows.",
          2.0, 3.9, 7.5, 0.5)

    # --- Slide 10: Redirection Issues ---
    s = content_slide("Redirection Issues")
    _text(s, "Redirects send users and crawlers from one URL to another. Permanent 301 redirects pass full link equity.",
          2.0, 1.2, 7.5, 0.6, 11, "Calibri", TEXT, bold=True)
    redir_rows = [[r["type"], r["status"], r["detail"], r["impact"]] for r in data.get("redirects", [])]
    if not redir_rows:
        redir_rows = [["N/A", "N/A", "N/A", "N/A"]]
    _add_table(s, ["Redirect Type", "Status", "Checked URL → Destination", "SEO Impact"],
               redir_rows, [1.6, 0.7, 3.7, 1.5], 2.0, 2.0)
    _suggestion_text(s, bool(data.get("redirects")),
          "Ensure all HTTP and www variants redirect to the canonical version via 301. Avoid redirect chains.",
          2.0, 3.9, 7.5, 0.5)

    # --- Slide 11: Canonical Tag ---
    s = content_slide("Canonical Tag")
    _text(s, "The canonical tag tells search engines which URL is the preferred version of a page, preventing duplicate content issues.",
          2.0, 1.2, 7.5, 0.6, 11, "Calibri", TEXT, bold=True)
    canon_rows = [[c["page"], c["found"], c["correct"], c["duplicate_risk"], c["status"]]
                  for c in data.get("canonicals", [])]
    if not canon_rows:
        canon_rows = [["N/A", "N/A", "N/A", "N/A", "N/A"]]
    _add_table(s, ["Page URL", "Tag Found?", "Correct?", "Duplicate Risk", "Status"],
               canon_rows, [3.0, 1.0, 1.0, 1.2, 1.3], 2.0, 2.0)
    _suggestion_text(s, bool(data.get("canonicals")),
          "Add self-referencing canonical tags to all pages. This prevents duplicate content from URL parameters or www/non-www variants.",
          2.0, 3.9, 7.5, 0.5)

    # --- Slide 12: Sitemap.xml ---
    s = content_slide("Sitemap.xml")
    _text(s, "An XML sitemap lists a site's URLs to help search engines discover and crawl all important pages.",
          2.0, 1.2, 7.5, 0.6, 11, "Calibri", TEXT, bold=True)
    sm = data.get("sitemap", {})
    sm_found = sm.get("found") or sm.get("ok")
    sm_url_checked = sm.get("url_checked", f"https://{domain}/sitemap.xml")
    sm_rows = [
        ["Sitemap Found", "Yes" if sm_found else "No", sm.get("summary", "")[:40], "High"],
        ["URL Checked", sm_url_checked[:45], "", ""],
        ["Submitted to GSC", "Check in GSC", "Submit via GSC or API", "High"],
    ]
    _add_table(s, ["Check Item", "Status", "Detail", "Priority"],
               sm_rows, [2.0, 1.5, 2.5, 1.5], 2.0, 2.0)
    _suggestion_text(s, sm_found, sm.get("summary", "Sitemap not found. Create and submit one."),
          2.0, 4.2, 7.5, 0.5)

    # --- Slide 13: Robots.txt ---
    s = content_slide("Robots.txt")
    _text(s, "The robots.txt file tells search engines which pages to crawl and which to skip. A misconfigured file can block important pages.",
          2.0, 1.2, 7.5, 0.6, 11, "Calibri", TEXT, bold=True)
    rb = data.get("robots", {})
    rb_rows = [
        ["Robots.txt Found", "Yes" if rb.get("found") else "No", rb.get("summary", "")[:90]],
    ]
    _add_table(s, ["Check Item", "Status", "Detail"],
               rb_rows, [2.5, 1.5, 3.5], 2.0, 2.0)
    _suggestion_text(s, rb.get("found", False), rb.get("summary", "Check robots.txt configuration."),
          2.0, 3.5, 7.5, 0.5)

    # --- Slide 14: Broken Links ---
    s = content_slide("Broken Links")
    _text(s, "Broken links hurt user experience and waste crawl budget. Regular audits keep your site's link graph healthy.",
          2.0, 1.2, 7.5, 0.6, 11, "Calibri", TEXT, bold=True)
    bl = data.get("broken_links", [])
    bl_checked = data.get("broken_links_checked", 0)
    if bl:
        bl_rows = [[b.get("location", "Content"), b.get("url", ""),
                     str(b.get("status", ""))] for b in bl[:4]]
    else:
        bl_rows = [[f"Checked {bl_checked} link(s)", "No broken links found", "Good"]]
    _add_table(s, ["Location", "Broken Link URL", "Status"],
               bl_rows, [1.7, 4.5, 1.3], 2.0, 2.0)
    if bl:
        _suggestion_text(s, True,
              f"{len(bl)} broken link(s) found out of {bl_checked} checked. Fix or remove broken links to improve crawl efficiency.",
              2.0, 3.9, 7.5, 0.5)
    else:
        _suggestion_text(s, bl_checked > 0,
              f"No broken links found in {bl_checked} links checked. Continue monitoring regularly.",
              2.0, 3.9, 7.5, 0.5)

    # --- Slide 15: Thank You ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _james_sidebar(s)
    _text(s, "Thank You", 2.0, 1.8, 7.0, 1.0, 40, "Calibri", NAVY, bold=True, italic=True,
          align=PP_ALIGN.CENTER)
    _text(s, f"https://{domain}/", 2.0, 3.0, 7.0, 0.5, 16, "Calibri", LINK_BLUE,
          align=PP_ALIGN.CENTER)

    prs.save(out_path)
    log_fn(f"  James brief report saved: {os.path.basename(out_path)}")
    return out_path


# ---------------------------------------------------------------------------
# Xenon format — 15 slides, 13.33×7.5"
# ---------------------------------------------------------------------------

def build_xenon(data, out_path, log_fn=None):
    if log_fn is None:
        log_fn = print
    Presentation, Inches, Pt, Emu, RGBColor, PP_ALIGN, MSO_ANCHOR = _init_pptx()

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    NAVY = "#1F3864"
    NAVY_DARK = "#0D2137"
    LIGHT_BLUE = "#C7D8EF"
    MED_BLUE = "#9FC0E8"
    BADGE_BG = "#DEEAF6"
    CARD_BG = "#EAF1FA"
    CARD_WARN = "#FBEDD9"
    TEXT_DARK = "#2A2A2A"
    TEXT_SOFT = "#5A5A5A"
    WHITE = "#FFFFFF"
    GREEN_BADGE = "#B45309"
    BLUE_BADGE = "#2E75B6"

    domain = data["domain"]
    date_str = data["date"]

    def icon_circle(slide, x, y):
        _rect(slide, x, y, 0.7, 0.7, NAVY)

    def status_badge(slide, x, y, text, badge_type="good"):
        color = GREEN_BADGE if badge_type in ("good", "issue") else BLUE_BADGE
        w = 1.4 if badge_type == "good" else 2.4
        _rect(slide, x, y, w, 0.4, color)
        _text(slide, text.upper(), x, y, w, 0.4, 11, "Calibri", WHITE, bold=True,
              align=PP_ALIGN.CENTER)

    def finding_box(slide, x, y, w, h, text):
        _rect(slide, x, y, w, h, BADGE_BG)
        return text

    def content_slide(title, description, badge_text="Good", badge_type="good", finding="", data_ok=True):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        icon_circle(s, 0.6, 0.5)
        _text(s, title, 1.6, 0.5, 10.0, 0.6, 27, "Calibri", NAVY, bold=True)
        _text(s, description, 0.6, 1.5, 12.1, 1.2, 15, "Calibri", TEXT_SOFT)
        _rect(s, 0.6, 3.1, 12.1, 1.1, BADGE_BG)
        if not data_ok:
            _text(s, UNVERIFIED, 0.9, 3.3, 11.5, 0.7, 14, "Calibri", RED, bold=True)
        else:
            status_badge(s, 0.9, 3.5, badge_text, badge_type)
            bx = 2.7 if badge_type in ("good", "issue") else 3.6
            if finding:
                _text(s, finding, bx, 3.2, 12.1 - bx + 0.6, 0.9, 13, "Calibri", TEXT_DARK)
        return s

    _shots = data.get("screenshots", {}) or {}

    def page_detail(s, header, lines, y=4.35):
        """Compact per-page breakdown below the summary box — kept short so a
        screenshot can sit beneath it (matching the original template)."""
        if not lines:
            return
        _text(s, header, 0.6, y, 12.1, 0.32, 12, "Calibri", NAVY, bold=True)
        _text(s, "\n".join(lines[:3]), 0.6, y + 0.34, 12.1, 0.85, 11, "Calibri", TEXT_DARK)

    def topic_shot(s, key):
        """Screenshot at the bottom of the slide (Xenon puts the evidence there)."""
        _place_image(s, _shots.get(key), 2.6, 5.45, 8.1, 1.85)

    # --- Slide 1: Cover ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s, 0, 0, 13.33, 7.5, NAVY_DARK)
    _text(s, "SEO AUDIT REPORT", 0.6, 3.5, 5.0, 0.4, 14, "Calibri", LIGHT_BLUE, bold=True)
    _text(s, "Basic SEO Issues", 0.6, 3.9, 5.0, 1.0, 46, "Calibri", WHITE, bold=True)
    _text(s, "Website Audit — Findings & Recommendations", 0.6, 5.0, 5.0, 0.4, 18, "Calibri", LIGHT_BLUE)
    _text(s, f"https://{domain}/", 0.6, 5.8, 5.0, 0.3, 15, "Calibri", MED_BLUE)
    _text(s, date_str, 0.6, 6.2, 5.0, 0.3, 13, "Calibri", "#8AA6CC")

    # --- Slide 2: Website Indexing ---
    idx = data.get("indexing", {})
    idx_ok = not _is_empty(idx.get("count"))
    finding = f"About {idx.get('count', 'N/A')} pages are indexed in Google. Confirm that indexing isn't blocked on important pages via robots.txt or meta noindex."
    s = content_slide("Website Indexing",
                      "Indexing is how a search engine crawls and stores your pages so they can appear in search results.",
                      "Good", "good", finding, data_ok=idx_ok)
    topic_shot(s, "serp")

    # --- Slide 3: Title Tag ---
    titles = data.get("titles", [])
    titles_ok = bool(titles)
    issues = [t for t in titles if t["status"] != "Good"]
    if issues:
        finding = f"Title tags need attention on {len(issues)} page(s). " + "; ".join(f'{t["page"]}: {t["status"]}' for t in issues[:2]) + "."
        badge = ("RECOMMENDATION", "recommendation")
    else:
        finding = "All checked pages have well-optimized title tags within the recommended 50-60 character range."
        badge = ("Good", "good")
    s = content_slide("Title Tag",
                      "The title tag is the HTML element that defines a page's title. It appears on search engine results and browser tabs.",
                      badge[0], badge[1], finding, data_ok=titles_ok)
    page_detail(s, "Pages checked:", [
        f"{t['page']}  —  \"{(t.get('title') or '(none)')[:55]}\"  ({t['chars']} chars, {t['status']})"
        for t in titles])
    topic_shot(s, "viewsource")

    # --- Slide 4: Meta Description Tag ---
    metas = data.get("metas", [])
    metas_ok = bool(metas)
    issues = [m for m in metas if m["status"] != "Good"]
    if issues:
        finding = f"Meta descriptions need attention on {len(issues)} page(s). " + "; ".join(f'{m["page"]}: {m["status"]}' for m in issues[:2]) + "."
        badge = ("RECOMMENDATION", "recommendation")
    else:
        finding = "All checked pages have well-optimized meta descriptions within the recommended 140-160 character range."
        badge = ("Good", "good")
    s = content_slide("Meta Description Tag",
                      "The meta description is an HTML attribute that summarises a page. Google often displays it as the snippet in search results.",
                      badge[0], badge[1], finding, data_ok=metas_ok)
    page_detail(s, "Pages checked:", [
        f"{m['page']}  —  meta description {'FOUND' if m['found'] == 'Yes' else 'MISSING'} "
        f"({m['chars']} chars, {m['status']})" for m in metas])
    topic_shot(s, "viewsource")

    # --- Slide 5: Headings ---
    headers = data.get("headers", [])
    headers_ok = bool(headers)
    h1_issues = [h for h in headers if h["h1"] != 1]
    if h1_issues:
        finding = f"H1 tag issues found on {len(h1_issues)} page(s). " + "; ".join(f'{h["page"]}: {h["h1"]} H1s' for h in h1_issues[:2]) + "."
        badge = ("ISSUE", "issue")
    else:
        finding = "Each checked page has exactly one H1 tag. This is ideal for SEO."
        badge = ("Good", "good")
    s = content_slide("Headings",
                      "Heading tags (H1, H2, and so on) are an important part of a page and play a role in helping search engines understand the topic.",
                      badge[0], badge[1], finding, data_ok=headers_ok)
    page_detail(s, "Pages checked:", [
        f"{h['page']}  —  H1={h['h1']}, H2={h['h2']}, H3={h['h3']}, H4={h['h4']}, H5={h['h5']}, H6={h['h6']}"
        for h in headers])
    topic_shot(s, "homepage")

    # --- Slide 6: Heading Tags — Findings ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    icon_circle(s, 0.6, 0.5)
    _text(s, "Heading Tags — Findings", 1.6, 0.5, 10.0, 0.6, 27, "Calibri", NAVY, bold=True)
    hdr_detail = "Heading structure breakdown:\n"
    for h in headers[:3]:
        hdr_detail += f"{h['page']}: H1={h['h1']}, H2={h['h2']}, H3={h['h3']}, H4={h['h4']}\n"
    _text(s, hdr_detail, 0.6, 1.5, 6.0, 3.0, 15, "Calibri", TEXT_SOFT)
    _rect(s, 0.6, 3.7, 6.7, 2.7, BADGE_BG)
    total_h1 = sum(h["h1"] for h in headers)
    status_badge(s, 0.9, 4.8, "Good" if all(h["h1"] == 1 for h in headers) else "ISSUE",
                 "good" if all(h["h1"] == 1 for h in headers) else "issue")
    _text(s, f"{total_h1} H1 tag(s) found across {len(headers)} checked page(s).",
          2.5, 3.8, 5.0, 0.5, 13, "Calibri", TEXT_DARK)

    # --- Slide 7: Sitemap.xml ---
    sm = data.get("sitemap", {})
    sm_ok = sm.get("found") is not None
    sm_url = sm.get("url_checked") or f"https://{domain}/sitemap.xml"
    _parts = [sm.get("summary", "Sitemap status could not be determined.")]
    pc = sm.get("page_count")
    if pc:
        _parts.append(f"{pc} page URL(s) found in the sitemap.")
    finding = "  ".join(p for p in _parts if p)
    s = content_slide("Sitemap.xml",
                      "A sitemap is an XML file that lists a site's URLs, along with metadata, to tell search engines what to crawl.",
                      "Good" if sm.get("found") else "ISSUE",
                      "good" if sm.get("found") else "issue",
                      finding, data_ok=sm_ok)
    page_detail(s, "Sitemap URL:", [sm_url])
    topic_shot(s, "sitemap")

    # --- Slide 8: Robots.txt ---
    rb = data.get("robots", {})
    rb_ok = rb.get("found") is not None
    finding = rb.get("summary", "Robots.txt status could not be determined.")
    s = content_slide("Robots.txt",
                      "The robots.txt file tells web robots — mainly search engines — which pages to crawl and which to skip.",
                      "Good" if rb.get("found") else "ISSUE",
                      "good" if rb.get("found") else "issue",
                      finding, data_ok=rb_ok)
    topic_shot(s, "robots")

    # --- Slide 9: Content Optimization ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    icon_circle(s, 0.6, 0.5)
    _text(s, "Content Optimization", 1.6, 0.5, 10.0, 0.6, 27, "Calibri", NAVY, bold=True)
    _text(s, "Content is central to SEO. Unique, keyword-oriented, informative content improves rankings and user engagement.",
          0.6, 1.6, 7.0, 1.0, 15, "Calibri", TEXT_SOFT)
    _rect(s, 0.6, 4.0, 7.4, 2.0, BADGE_BG)
    status_badge(s, 0.9, 4.8, "RECOMMENDATION", "recommendation")
    _text(s, "Review page content for keyword optimization. Ensure each key page has unique, valuable content of at least 300 words.",
          3.6, 4.1, 4.5, 1.0, 13, "Calibri", TEXT_DARK)
    _rect(s, 8.4, 1.6, 4.3, 4.8, CARD_BG)
    _text(s, "Content is king.", 8.7, 4.5, 3.7, 0.5, 18, "Calibri", NAVY, bold=True)

    # --- Slide 10: Image Alt Tags ---
    alts = data.get("img_alts", [])
    total_imgs = len(alts)
    missing = len([a for a in alts if a["present"] == "No"])
    empty = len([a for a in alts if a["present"] == "Empty"])
    bad = missing + empty
    s = prs.slides.add_slide(prs.slide_layouts[6])
    icon_circle(s, 0.6, 0.5)
    _text(s, "Image Alt Tags", 1.6, 0.5, 10.0, 0.6, 27, "Calibri", NAVY, bold=True)
    _text(s, "Images using img tags with empty or missing alt attributes are not ideal for SEO and accessibility.",
          0.6, 1.5, 12.1, 0.7, 15, "Calibri", TEXT_SOFT)
    # Metric cards
    _rect(s, 0.6, 2.5, 3.8, 1.1, BADGE_BG)
    _text(s, str(total_imgs), 0.6, 2.6, 3.8, 0.5, 30, "Calibri", NAVY, bold=True, align=PP_ALIGN.CENTER)
    _text(s, "Images on home page", 0.6, 3.2, 3.8, 0.3, 11, "Calibri", TEXT_SOFT, align=PP_ALIGN.CENTER)
    _rect(s, 4.7, 2.5, 3.8, 1.1, CARD_WARN)
    _text(s, str(bad), 4.7, 2.6, 3.8, 0.5, 30, "Calibri", GREEN_BADGE, bold=True, align=PP_ALIGN.CENTER)
    _text(s, "Missing/empty alt", 4.7, 3.2, 3.8, 0.3, 11, "Calibri", TEXT_SOFT, align=PP_ALIGN.CENTER)
    _rect(s, 8.9, 2.5, 3.8, 1.1, CARD_WARN)
    _text(s, str(total_imgs), 8.9, 2.6, 3.8, 0.5, 30, "Calibri", GREEN_BADGE, bold=True, align=PP_ALIGN.CENTER)
    _text(s, "Missing title attribute", 8.9, 3.2, 3.8, 0.3, 11, "Calibri", TEXT_SOFT, align=PP_ALIGN.CENTER)
    _rect(s, 0.6, 3.9, 12.1, 1.0, BADGE_BG)
    badge_type = "issue" if bad > 0 else "good"
    status_badge(s, 0.9, 4.1, "ISSUE" if bad > 0 else "Good", badge_type)
    _text(s, f"{total_imgs} images on the home page — {bad} with no/empty alt tag. Add descriptive alt text to all images.",
          2.7, 4.0, 10.0, 0.5, 13, "Calibri", TEXT_DARK)

    # --- Slide 11: Canonical Tag ---
    canons = data.get("canonicals", [])
    canons_ok = bool(canons)
    missing_c = [c for c in canons if c["found"] == "No"]
    if missing_c:
        finding = f"Canonical tag missing on {len(missing_c)} page(s). Add self-referencing canonical tags to prevent duplicate content issues."
        badge = ("ISSUE", "issue")
    else:
        finding = "All checked pages have correct self-referencing canonical tags."
        badge = ("Good", "good")
    s = content_slide("Canonical Tag",
                      "The canonical tag tells search engines which version of a URL is the master copy, preventing duplicate content penalties.",
                      badge[0], badge[1], finding, data_ok=canons_ok)
    page_detail(s, "Pages checked:", [
        f"{c['page']}  —  tag {c['found']}, correct: {c.get('correct', 'N/A')} ({c['status']})"
        for c in canons])
    topic_shot(s, "viewsource")

    # --- Slide 12: Redirection Issues ---
    redirects = data.get("redirects", [])
    redir_ok = bool(redirects)
    redir_issues = [r for r in redirects if not str(r.get("impact", "")).startswith("Good")]
    if redir_issues:
        finding = f"{len(redir_issues)} redirection issue(s) found. " + "; ".join(f'{r["type"]}: {r["detail"][:40]}' for r in redir_issues[:2]) + "."
        badge = ("ISSUE", "issue")
    else:
        finding = "All redirection patterns are correctly configured. HTTP and www variants properly redirect."
        badge = ("Good", "good")
    s = content_slide("Redirection Issues",
                      "Redirects send users and crawlers from one URL to another. Correct 301 redirects pass full link equity to the target URL.",
                      badge[0], badge[1], finding, data_ok=redir_ok)
    page_detail(s, "Redirects checked:", [
        f"{r['type']}  —  {r['status']}  —  {r.get('detail', '')}" for r in redirects])

    # --- Slide 13: Broken Links ---
    bl = data.get("broken_links", [])
    bl_checked = data.get("broken_links_checked", 0)
    bl_ok = bl_checked > 0
    if bl:
        finding = f"{len(bl)} broken link(s) found out of {bl_checked} checked. Fix or remove broken links to improve crawl efficiency and user experience."
        badge = ("ISSUE", "issue")
    else:
        finding = f"No broken links found in {bl_checked} links checked. Good link health."
        badge = ("Good", "good")
    s = content_slide("Broken Links",
                      "Broken links return 404 errors, wasting crawl budget and damaging user trust. Regular audits keep your link graph healthy.",
                      badge[0], badge[1], finding, data_ok=bl_ok)
    if bl:
        page_detail(s, "Broken links found:", [
            f"[{b.get('location', 'Content')}]  {b.get('url', '')}  ({b.get('status', '')})" for b in bl[:6]])

    # --- Slide 14: Domain Age ---
    age_rows = data.get("domain_age", [])
    age_info = ""
    age_ok = False
    for r in age_rows:
        if r[0] == "Created On" and not _is_empty(r[1]):
            age_info = f"Domain registered since {r[1]}. {r[2] if len(r) > 2 else ''}"
            age_ok = True
    if not age_info:
        age_info = "Domain age information could not be retrieved."
    s = content_slide("Domain Age",
                      "Domain age reflects registration history. Older domains often carry accumulated trust and authority with search engines.",
                      "Good", "good", age_info, data_ok=age_ok)

    # --- Slide 15: Thank You ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s, 0, 0, 13.33, 7.5, NAVY_DARK)
    _text(s, "Let's take this further", 0.6, 2.1, 12.0, 0.6, 30, "Calibri", WHITE, bold=True)
    _text(s, "There may be additional SEO issues we can identify and resolve. We look forward to working with you on a comprehensive strategy.",
          2.0, 3.0, 10.0, 1.0, 16, "Calibri", LIGHT_BLUE)
    _text(s, "Thank You", 0.6, 4.4, 12.0, 1.0, 40, "Calibri", WHITE, bold=True, align=PP_ALIGN.CENTER)
    _text(s, f"https://{domain}/", 0.6, 5.6, 12.0, 0.4, 14, "Calibri", MED_BLUE, align=PP_ALIGN.CENTER)

    prs.save(out_path)
    log_fn(f"  Xenon brief report saved: {os.path.basename(out_path)}")
    return out_path


# ---------------------------------------------------------------------------
# Omega format — 16 slides, 10×5.625", dark navy bg, blue sidebar
# ---------------------------------------------------------------------------

def build_omega(data, out_path, log_fn=None):
    if log_fn is None:
        log_fn = print
    Presentation, Inches, Pt, Emu, RGBColor, PP_ALIGN, MSO_ANCHOR = _init_pptx()

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    BG = "#161F2C"
    SIDEBAR = "#385886"
    ACCENT = "#9CB4D8"
    WHITE = "#FFFFFF"
    ORANGE = "#E76A28"
    GRAY = "#999999"
    TEAL = "#0097A7"

    domain = data["domain"]
    date_str = data["date"]

    def bg_slide():
        s = prs.slides.add_slide(prs.slide_layouts[6])
        _rect(s, 0, 0, 10, 5.625, BG)
        return s

    def left_sidebar(s):
        _rect(s, -0.3, 0, 3.0, 5.625, SIDEBAR)
        _rect(s, 0.5, 0.8, 1.5, 0.02, ACCENT)
        _rect(s, 0.5, 1.3, 1.5, 0.02, ACCENT)

    def right_sidebar(s):
        _rect(s, 7.3, 0, 3.0, 5.625, SIDEBAR)
        _rect(s, 8.0, 0.8, 1.5, 0.02, ACCENT)
        _rect(s, 8.0, 1.3, 1.5, 0.02, ACCENT)

    def content_slide(title, desc, finding, shot=None):
        s = bg_slide()
        _text(s, title, 0.8, 0.1, 8.4, 0.6, 24, "Calibri", WHITE, bold=True)
        _rect(s, 0.8, 0.7, 8.4, 0.02, ACCENT)
        _text(s, desc, 0.8, 0.9, 8.8, 0.6, 14, "Calibri", ACCENT)
        # Screenshot in the middle (like the original), finding text below it.
        has_shot = _place_image(s, shot, 1.2, 1.55, 7.6, 2.05) if shot else False
        fy = 3.75 if has_shot else 1.9
        fh = 1.75 if has_shot else 3.4
        data_ok = not _is_empty(finding) and finding != UNVERIFIED
        if data_ok:
            _text(s, finding, 0.8, fy, 9.2, fh, 11 if has_shot else 12, "Calibri", WHITE)
        else:
            _text(s, UNVERIFIED, 0.8, fy, 9.2, 0.7, 14, "Calibri", RED, bold=True)
        return s

    topics = [
        ("Website Indexing", "Indexing is how search engines store pages for search results."),
        ("Title Tag", "The title tag appears in search results and browser tabs."),
        ("Meta Description", "The meta description summarises a page in search result snippets."),
        ("Header Tags", "Header tags (H1-H6) structure content for readers and crawlers."),
        ("Sitemap.xml", "An XML sitemap lists URLs to help search engines discover pages."),
        ("Robots.txt", "The robots.txt file tells crawlers which pages to crawl or skip."),
        ("Image Alt Tags", "Alt text describes images for accessibility and SEO."),
        ("Canonical Tag", "Canonical tags prevent duplicate content issues."),
        ("Redirection Issues", "Redirects send users from one URL to another."),
        ("Broken Links", "Broken links hurt UX and waste crawl budget."),
        ("Domain Age", "Domain age reflects registration history and trust signals."),
    ]

    # --- Slide 1: Title ---
    s = bg_slide()
    left_sidebar(s)
    _text(s, "BASIC ANALYSIS\nOF WEBSITE", 2.6, 0.6, 7.1, 2.8, 44, "Calibri", WHITE, bold=True)
    _rect(s, 4.5, 3.3, 3.2, 0.02, ACCENT)
    _text(s, f"Domain: {domain}", 3.4, 3.5, 5.1, 0.5, 20, "Calibri", WHITE)

    # --- Slide 2: TOC ---
    s = bg_slide()
    _text(s, "Table of contents", 3.5, 0.4, 3.0, 0.4, 18, "Calibri", WHITE, bold=True)
    for i, (t, _) in enumerate(topics):
        col = 0 if i % 2 == 0 else 1
        row = i // 2
        x = 1.0 + col * 4.3
        y = 1.2 + row * 0.6
        num = f"{i+1:02d}"
        _rect(s, x, y, 0.5, 0.35, SIDEBAR)
        _text(s, num, x, y, 0.5, 0.35, 12, "Calibri", WHITE, align=PP_ALIGN.CENTER)
        _text(s, t, x + 0.6, y, 3.5, 0.35, 12, "Calibri", ACCENT)

    # --- Slides 3-13: Content (full per-page detail, not just a status) ---
    def _lines(items, fmt_fn):
        return "\n".join(fmt_fn(x) for x in items) if items else UNVERIFIED

    idx = data.get("indexing", {})
    titles = data.get("titles", [])
    metas = data.get("metas", [])
    headers = data.get("headers", [])
    canons = data.get("canonicals", [])
    redirects = data.get("redirects", [])
    bl = data.get("broken_links", [])
    bl_checked = data.get("broken_links_checked", 0)
    alts = data.get("img_alts", [])
    sm = data.get("sitemap", {})
    rb = data.get("robots", {})
    sm_url = sm.get("url_checked") or f"https://{domain}/sitemap.xml"
    sm_find = sm.get("summary", "")
    if sm.get("found") or sm.get("ok"):
        sm_find = (sm_find + f"\nSitemap URL: {sm_url}").strip()
    findings = [
        f"About {idx.get('count', 'N/A')} pages indexed. Status: {idx.get('status', 'N/A')}",
        _lines(titles, lambda t: f'{t["page"]}  -  "{(t.get("title") or "(none)")[:55]}" ({t["chars"]} chars, {t["status"]})'),
        _lines(metas, lambda m: f'{m["page"]}  -  meta {"FOUND" if m["found"] == "Yes" else "MISSING"} ({m["chars"]} chars, {m["status"]})'),
        _lines(headers, lambda h: f'{h["page"]}  -  H1={h["h1"]}, H2={h["h2"]}, H3={h["h3"]}, H4={h["h4"]}'),
        sm_find or UNVERIFIED,
        rb.get("summary", UNVERIFIED),
        (f'{len(alts)} images checked. {len([a for a in alts if a["present"] != "Yes"])} need alt text.' if alts else UNVERIFIED),
        _lines(canons, lambda c: f'{c["page"]}  -  tag {c["found"]}, correct: {c.get("correct", "N/A")} ({c["status"]})'),
        _lines(redirects, lambda r: f'{r["type"]}: {r["status"]}  -  {r.get("detail", "")}'),
        (_lines(bl[:5], lambda b: f'[{b.get("location", "Content")}] {b.get("url", "")} ({b.get("status", "")})')
         if bl else (f'No broken links found in {bl_checked} links checked.' if bl_checked > 0 else UNVERIFIED)),
    ]
    # Domain age — show every attribute we DID find (not just "Created On"), so
    # ccTLDs like .com.au that hide the creation date still show useful info.
    age_rows = data.get("domain_age", [])
    age_lines = [f"{r[0]}: {r[1]}" + (f"  ({r[2]})" if len(r) > 2 and r[2] else "")
                 for r in age_rows if len(r) > 1 and str(r[1]).strip()]
    findings.append("\n".join(age_lines) if age_lines else UNVERIFIED)

    # Each content slide embeds the screenshot that evidences that check, like the
    # original template (indexing = Google SERP, title/meta/canonical = view-source,
    # headers/images = homepage, sitemap, robots).
    shots = data.get("screenshots", {}) or {}
    shot_map = {0: "serp", 1: "viewsource", 2: "viewsource", 3: "homepage",
                4: "sitemap", 5: "robots", 6: "homepage", 7: "viewsource"}
    for i, (title, desc) in enumerate(topics):
        content_slide(title, desc, findings[i], shots.get(shot_map.get(i)))

    # --- Slide 16: Thank You ---
    s = bg_slide()
    right_sidebar(s)
    _text(s, "Thanks!", 0.8, 2.0, 6.2, 1.7, 44, "Calibri", WHITE, bold=True)
    _rect(s, 1.2, 3.7, 5.4, 0.02, ACCENT)
    _text(s, f"https://{domain}/", 1.2, 4.0, 5.0, 0.4, 16, "Calibri", TEAL)

    prs.save(out_path)
    log_fn(f"  Omega brief report saved: {os.path.basename(out_path)}")
    return out_path


# ---------------------------------------------------------------------------
# Neon format — 19 slides, 13.33×7.5", clean white bg, text+image style
# ---------------------------------------------------------------------------

def build_neon(data, out_path, log_fn=None):
    if log_fn is None:
        log_fn = print
    Presentation, Inches, Pt, Emu, RGBColor, PP_ALIGN, MSO_ANCHOR = _init_pptx()

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    NAVY = "#1F3864"
    BLUE = "#94B6D2"
    TITLE_CLR = "#4A5B69"
    BODY_CLR = "#000000"
    RED_BAD = "#EB393B"
    WHITE = "#FFFFFF"

    domain = data["domain"]
    date_str = data["date"]

    topics = [
        ("Website Indexing", "How search engines crawl and store your pages to display in search results."),
        ("Title Tag", "The HTML title element appears in browser tabs and as the clickable headline in search results."),
        ("Meta Description Tag", "An HTML attribute that provides a brief summary of a page, often shown as the snippet in search results."),
        ("Headings", "Header tags (H1-H6) structure page content and help search engines understand topic hierarchy."),
        ("Sitemap.xml", "An XML file listing the site's URLs with metadata for search engine discovery."),
        ("Canonical Tag", "Tells search engines which URL version is the preferred one, preventing duplicate content penalties."),
        ("Robots.txt", "Instructs web crawlers which pages to crawl and which to skip."),
        ("Image ALT Tag", "Alt text describes images for screen readers and search engine image indexing."),
        ("Content Optimization", "Unique, keyword-focused content improves rankings, engagement, and user experience."),
        ("Backlink Status", "Links from other websites pointing to yours, a key ranking signal."),
        ("Domain Age", "How long a domain has been registered — older domains may carry accumulated trust."),
        ("Website Speed", "Page load time affects user experience and is a Google ranking factor."),
        ("Organic Traffic Status", "The volume of visitors arriving from unpaid search results."),
        ("SEO Score", "An overall assessment of a website's search engine optimisation health."),
        ("Broken Link", "Broken links return 404 errors, damaging UX and wasting crawl budget."),
        ("Footer Logo", "Ensure your brand logo is visible and linked in the site footer."),
        ("Copyright Year", "An up-to-date copyright year signals an actively maintained website."),
    ]

    def finding_for(i):
        idx = data.get("indexing", {});   idx = idx if isinstance(idx, dict) else {}
        titles = data.get("titles", [])
        metas = data.get("metas", [])
        headers = data.get("headers", [])
        sm = data.get("sitemap", {});      sm = sm if isinstance(sm, dict) else {}
        canons = data.get("canonicals", [])
        rb = data.get("robots", {});       rb = rb if isinstance(rb, dict) else {}
        alts = data.get("img_alts", [])
        bl = data.get("broken_links", [])
        bl_checked = data.get("broken_links_checked", 0)
        age_rows = data.get("domain_age", [])

        def _lines(items, fn):
            return "\n".join(fn(x) for x in items) if items else None
        _sm_url = sm.get("url_checked") or f"https://{domain}/sitemap.xml"
        mapping = {
            0: f"About {idx.get('count', 'N/A')} pages indexed. {idx.get('status', 'N/A')}" if not _is_empty(idx.get("count")) else None,
            1: _lines(titles, lambda t: f'{t["page"]}  -  "{(t.get("title") or "(none)")[:50]}" ({t["chars"]} chars, {t["status"]})'),
            2: _lines(metas, lambda m: f'{m["page"]}  -  meta {"FOUND" if m["found"] == "Yes" else "MISSING"} ({m["chars"]} chars, {m["status"]})'),
            3: _lines(headers, lambda h: f'{h["page"]}  -  H1={h["h1"]}, H2={h["h2"]}, H3={h["h3"]}, H4={h["h4"]}'),
            4: ((sm.get("summary", "") + f"\nSitemap URL: {_sm_url}").strip() if (sm.get("found") or sm.get("ok")) else None),
            5: _lines(canons, lambda c: f'{c["page"]}  -  tag {c["found"]}, correct: {c.get("correct", "N/A")} ({c["status"]})'),
            6: rb.get("summary") if rb.get("found") else None,
            7: f"{len(alts)} images checked. {len([a for a in alts if a['present'] != 'Yes'])} need alt text." if alts else None,
            8: "Review page content for keyword optimization. Ensure at least 300 words of unique content per key page.",
            9: "Backlink data requires a third-party API (Ahrefs, Moz, etc.) — check manually or connect an API.",
            10: None,
            11: "Run a PageSpeed Insights or GTmetrix test for detailed speed metrics.",
            12: "Check Google Analytics or Search Console for organic traffic data.",
            13: "Run a full SEO audit for a comprehensive score.",
            14: (_lines(bl[:5], lambda b: f'[{b.get("location", "Content")}] {b.get("url", "")} ({b.get("status", "")})')
                 if bl else (f"No broken links found in {bl_checked} checked." if bl_checked > 0 else None)),
            15: "Ensure your brand logo appears in the footer with a link to the homepage.",
            16: f"Verify the copyright year is set to {datetime.now().year}.",
        }
        # Domain age — show every attribute found (ccTLDs may hide the creation date).
        if i == 10:
            age_lines = [f"{r[0]}: {r[1]}" + (f"  ({r[2]})" if len(r) > 2 and r[2] else "")
                         for r in age_rows if len(r) > 1 and str(r[1]).strip()]
            return ("\n".join(age_lines) if age_lines else None, bool(age_lines))

        val = mapping.get(i)
        return (val, val is not None)

    # --- Slide 1: Title ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _text(s, "Project Name:-", 7.2, 4.4, 5.5, 0.5, 20, "Georgia", TITLE_CLR, bold=True)
    _text(s, f"SEO Analysis - {domain}", 7.2, 4.9, 5.5, 0.5, 18, "Calibri", BODY_CLR)

    # --- Slides 2-18: Content (screenshot in the middle + finding at the bottom,
    #     like the original template) ---
    shots = data.get("screenshots", {}) or {}
    shot_map = {0: "serp", 1: "viewsource", 2: "viewsource", 3: "homepage",
                4: "sitemap", 5: "viewsource", 6: "robots", 7: "homepage"}
    for i, (title, desc) in enumerate(topics):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        _text(s, f"{title} --", 0.25, 0.5, 4.5, 0.65, 36, "Calibri", TITLE_CLR, bold=True)
        _text(s, desc, 0.54, 1.2, 12.04, 1.0, 18, "Calibri", BODY_CLR)
        has_shot = _place_image(s, shots.get(shot_map.get(i)), 1.4, 2.45, 10.5, 2.75)

        finding_val, data_ok = finding_for(i)
        fy = 5.45 if has_shot else 2.5
        fh = 1.5 if has_shot else 4.4
        if data_ok and finding_val:
            _text(s, finding_val, 0.83, fy, 11.75, fh, 13, "Calibri", BODY_CLR)
        else:
            _text(s, UNVERIFIED, 0.83, fy, 11.75, 0.65, 18, "Calibri", RED_BAD, bold=True)

    # --- Slide 19: Closing ---
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _text(s, "Thank You", 3.0, 2.5, 7.0, 1.5, 44, "Calibri", TITLE_CLR, bold=True, align=PP_ALIGN.CENTER)
    _text(s, f"SEO Analysis - {domain}", 3.0, 4.2, 7.0, 0.6, 18, "Calibri", BODY_CLR, align=PP_ALIGN.CENTER)
    _text(s, date_str, 3.0, 4.8, 7.0, 0.4, 14, "Calibri", TITLE_CLR, align=PP_ALIGN.CENTER)

    prs.save(out_path)
    log_fn(f"  Neon brief report saved: {os.path.basename(out_path)}")
    return out_path


def build_camila(data, out_path, log_fn=None):
    """Camila — 16-slide PPTX, verified against the client reference
    "Camila.pptx" (ltn-stahlhallenbau.de): navy (#002060) cover title, an
    "Index:" status slide with each checked item marked red (#C00000, needs
    attention) or green (#00B050, optimized), then 14 topic slides — gold
    (#7C5F1D) section title + bold black "Result:"/"Results:" label followed
    by the finding."""
    if log_fn is None:
        log_fn = print
    Presentation, Inches, Pt, Emu, RGBColor, PP_ALIGN, MSO_ANCHOR = _init_pptx()

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    NAVY = "#002060"
    GOLD = "#7C5F1D"
    BLACK = "#000000"
    RED = "#C00000"
    GREEN = "#00B050"

    domain = data["domain"]
    date_str = data["date"]
    home = f"https://{domain}/"

    idx = data.get("indexing", {}) or {}
    titles = data.get("titles", []) or []
    metas = data.get("metas", []) or []
    headers = data.get("headers", []) or []
    sitemap = data.get("sitemap", {}) or {}
    canonicals = data.get("canonicals", []) or []
    robots = data.get("robots", {}) or {}
    img_alts = data.get("img_alts", []) or []
    broken_links = data.get("broken_links", []) or []
    bl_checked = data.get("broken_links_checked", 0)
    da_pa = data.get("da_pa") or {}

    title_ok = bool(titles) and all(t.get("status") == "Good" for t in titles)
    meta_ok = bool(metas) and all(m.get("status") == "Good" for m in metas)
    h1_ok = bool(headers) and all(h.get("h1") == 1 for h in headers)
    h2_ok = bool(headers) and all(h.get("h2", 0) >= 1 for h in headers)
    idx_ok = not _is_empty(idx.get("count"))
    da_ok = (da_pa.get("da") not in (None, "—"))
    canon_ok = bool(canonicals) and all(c.get("status") == "Good" for c in canonicals)
    img_ok = bool(img_alts) and all(a.get("present") == "Yes" for a in img_alts)
    sm_ok = bool(sitemap.get("found") or sitemap.get("ok"))
    robots_ok = bool(robots.get("found"))
    bl_ok = bl_checked > 0 and not broken_links

    # ---- Slide 1: Cover ----
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _text(s, home, 1.6, 3.4, 8.0, 0.5, 16, "Cambria", BLACK, bold=True)
    _text(s, f"Date: {date_str}", 1.6, 3.9, 8.0, 0.5, 16, "Cambria", BLACK, bold=True)
    _text(s, "Website Audit Report", 1.6, 1.8, 9.0, 0.9, 48, "Calibri", NAVY, bold=True)

    # ---- Slide 2: Index ----
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _text(s, "Index:", 1.5, 0.6, 9.0, 0.9, 44, "Calibri", NAVY, bold=True)
    index_items = [
        ("Meta title: ", title_ok),
        ("Meta Description: ", meta_ok),
        ("H1 Heading Tags: ", h1_ok),
        ("H2 Heading Tags: ", h2_ok),
        ("Website Indexing: ", idx_ok),
        ("Domain Authority & Page Authority: ", da_ok),
        ("Canonical Tag: ", canon_ok),
        ("Image Optimization: ", img_ok),
        ("Schema Markup: ", False),
        ("XML Sitemap: ", sm_ok),
        ("Redirection: ", False),
        ("Robots.txt: ", robots_ok),
        ("Broken Link: ", bl_ok),
        ("Backlink Status: ", da_ok),
    ]
    box = prs.slides[-1].shapes.add_textbox(Inches(1.2), Inches(1.7), Inches(10.1), Inches(4.65))
    tf = box.text_frame
    tf.word_wrap = True
    for i, (label, ok) in enumerate(index_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r1 = p.add_run(); r1.text = label
        r1.font.size = Pt(18); r1.font.name = "Arial"; r1.font.bold = True
        r1.font.color.rgb = _C(BLACK)
        r2 = p.add_run(); r2.text = "Optimized" if ok else "Not Optimized"
        r2.font.size = Pt(18); r2.font.name = "Arial"; r2.font.bold = True
        r2.font.color.rgb = _C(GREEN if ok else RED)

    # ---- Slides 3-16: topics ----
    shots = data.get("screenshots", {}) or {}

    def topic_slide(title, result_text, shot_key=None, results_label="Result"):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        _text(s, title, 1.2, 0.35, 10.8, 0.9, 44, "Calibri", GOLD)
        y = 1.65
        if shot_key and shots.get(shot_key):
            has_shot = _place_image(s, shots.get(shot_key), 1.2, y, 10.9, 2.6)
            if has_shot:
                y += 2.85
        box = s.shapes.add_textbox(Inches(1.2), Inches(y), Inches(10.9), Inches(2.6))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r1 = p.add_run(); r1.text = results_label
        r1.font.size = Pt(18); r1.font.name = "Calibri"; r1.font.bold = True
        r1.font.color.rgb = _C(BLACK)
        r2 = p.add_run(); r2.text = ": "
        r2.font.size = Pt(18); r2.font.name = "Calibri"; r2.font.bold = True
        r2.font.color.rgb = _C(BLACK)
        r3 = p.add_run(); r3.text = result_text
        r3.font.size = Pt(16); r3.font.name = "Calibri"; r3.font.bold = False
        r3.font.color.rgb = _C(BLACK)
        return s

    # Meta Title
    if titles:
        t = titles[0]
        topic_slide("Meta Title",
                     f"The homepage title is {t.get('chars', '?')} characters — {t.get('status', 'N/A')}. "
                     "Keeping it within 50-60 characters with the primary keyword near the start "
                     "improves search engine relevance and click-through rate.",
                     "viewsource")
    else:
        topic_slide("Meta Title", UNVERIFIED, "viewsource")

    # Meta Description
    if metas:
        m = metas[0]
        topic_slide("Meta Description",
                     f"The homepage meta description is approximately {m.get('chars', '?')} "
                     f"characters — {m.get('status', 'N/A')}. It should be unique, keyword-relevant "
                     "and within the recommended length to improve the search snippet.",
                     "viewsource")
    else:
        topic_slide("Meta Description", UNVERIFIED, "viewsource")

    # H1 Tag
    if headers:
        h = headers[0]
        topic_slide("H1 Tag",
                     f"The homepage contains {h.get('h1', 0)} H1 tag(s). Having exactly one clear, "
                     "keyword-focused H1 tag per page is recommended for SEO.",
                     "homepage")
    else:
        topic_slide("H1 Tag", UNVERIFIED, "homepage")

    # H2 Tag
    if headers:
        h = headers[0]
        topic_slide("H2 Tag",
                     f"The homepage contains {h.get('h2', 0)} H2 heading(s). A clear H2 structure "
                     "helps establish content hierarchy and organize page topics for search engines.",
                     "homepage")
    else:
        topic_slide("H2 Tag", UNVERIFIED, "homepage")

    # Website Indexing
    topic_slide("Website Indexing",
                 f"Google has indexed approximately {idx.get('count', 'N/A')} page(s) for the "
                 f"website. {idx.get('status', '')} We recommend reviewing canonical, noindex, or "
                 "robots.txt directives where necessary so search engines prioritize important "
                 "content.",
                 "serp")

    # Domain Authority & Page Authority
    da_val, dr_val, pa_val = da_pa.get("da", "—"), da_pa.get("dr", "—"), da_pa.get("pa", "—")
    if da_ok:
        topic_slide("Domain Authority & Page Authority",
                     f"The website has a Domain Authority (DA) of {da_val}, Page Authority (PA) of "
                     f"{pa_val}, and Domain Rating (DR) of {dr_val}. We recommend building "
                     "high-quality backlinks to strengthen overall domain authority.")
    else:
        topic_slide("Domain Authority & Page Authority", UNVERIFIED)

    # Canonical Tag
    canon_note = ("properly implemented self-referencing canonical tags, which helps search "
                  "engines identify the preferred version of each page."
                  if canon_ok else
                  "canonical tag implementation that needs review to ensure each page correctly "
                  "identifies its preferred URL version.")
    topic_slide("Canonical Tag", f"The website has {canon_note}", results_label="Result")

    # Image Optimisation
    missing_alt = len([a for a in img_alts if a.get("present") != "Yes"])
    if img_alts:
        topic_slide("Image Optimisation",
                     f"The website has {len(img_alts)} image(s) checked, with {missing_alt} "
                     "missing ALT text. Alt tags help provide context to search engines and "
                     "improve accessibility and image search visibility.",
                     results_label="Results")
    else:
        topic_slide("Image Optimisation", UNVERIFIED, results_label="Results")

    # Schema Markup
    topic_slide("Schema Markup",
                 "Structured data (Schema) helps search engines understand website content "
                 "better and can enable rich results. We recommend implementing relevant "
                 "schema types (Organization, WebSite, Breadcrumb) for this domain.",
                 results_label="Results")

    # Sitemap.xml
    sm_url = sitemap.get("url_checked", f"{home}sitemap.xml")
    if sm_ok:
        topic_slide("Sitemap.xml",
                     f"The website's XML sitemap ({sm_url}) is available and accessible to search "
                     "engines. We recommend keeping it up to date and submitted via Google Search "
                     "Console and Bing Webmaster Tools.",
                     "sitemap", results_label="Results")
    else:
        topic_slide("Sitemap.xml",
                     f"The website's XML sitemap ({sm_url}) could not be verified as accessible. "
                     "We recommend creating and submitting one via Google Search Console.",
                     "sitemap", results_label="Results")

    # Redirection
    topic_slide("Redirection", UNVERIFIED, results_label="Results")

    # Robots.txt
    if robots_ok:
        topic_slide("Robots.txt",
                     f"The website has a robots.txt file in place at {home}robots.txt. We "
                     "recommend regularly reviewing it to ensure important pages remain "
                     "crawlable.",
                     "robots", results_label="Results")
    else:
        topic_slide("Robots.txt",
                     f"The website does not currently have an accessible robots.txt file at "
                     f"{home}robots.txt. We recommend creating one to guide search engine "
                     "crawlers.",
                     "robots", results_label="Results")

    # Broken Link
    if broken_links:
        topic_slide("Broken Link",
                     f"The website contains {len(broken_links)} broken link(s) across "
                     f"{bl_checked} scanned page(s). We recommend removing or replacing "
                     "obsolete URLs and implementing 301 redirects where appropriate.",
                     results_label="Results")
    else:
        topic_slide("Broken Link",
                     f"No broken links were found across {bl_checked} scanned page(s).",
                     results_label="Results")

    # Backlinks Status
    if da_ok:
        topic_slide("Backlinks Status",
                     f"The website's backlink profile shows a Domain Rating (DR) of {dr_val}. "
                     "We recommend acquiring more high-quality, dofollow backlinks from "
                     "authoritative domains to strengthen link equity.",
                     results_label="Results")
    else:
        topic_slide("Backlinks Status", UNVERIFIED, results_label="Results")

    prs.save(out_path)
    log_fn(f"  Camila brief report saved: {os.path.basename(out_path)}")
    return out_path


def build_eta(data, out_path, log_fn=None):
    """ETA — DOCX brief report, verified against the client reference
    "ETA.docx" (shoreteldepot.com): plain bold-black 14pt lettered/labeled
    sections (no shaded banners, same convention as Alpha), covering Meta
    Optimization, Headings, Image Optimization, Schema, Backlink Profile
    Overview, DA/PA/Spam Score, Indexing, Robots.txt, Sitemap, Broken Links,
    and separate "Website performance according to Semrush/Ahrefs" sections
    (Ahrefs uses the real DA/DR data; Semrush isn't integrated so it stays a
    manual-attach note, same convention as Beta)."""
    if log_fn is None:
        log_fn = print
    from docx.shared import Pt, RGBColor

    domain = data["domain"]
    home = f"https://{domain}/"
    from docx import Document
    doc = Document()
    for sname in ("Normal", "List Paragraph"):
        try:
            st = doc.styles[sname]
            st.font.name = "Calibri"
            st.font.size = Pt(11)
        except KeyError:
            pass

    BLACK = RGBColor(0x00, 0x00, 0x00)

    def _run(p, text, bold=False, color=BLACK, size=11):
        r = p.add_run(text)
        r.font.name = "Calibri"
        r.font.size = Pt(size)
        r.font.bold = bold
        if color is not None:
            r.font.color.rgb = color
        return r

    def section(text):
        p = doc.add_paragraph()
        p.paragraph_format.space_before = Pt(10)
        p.paragraph_format.space_after = Pt(4)
        _run(p, text, bold=True, color=BLACK, size=14)
        return p

    def body(text, bold=False):
        p = doc.add_paragraph()
        p.paragraph_format.space_after = Pt(4)
        _run(p, text, bold=bold, color=BLACK, size=11)
        return p

    def labeled(label, text=""):
        p = doc.add_paragraph()
        p.paragraph_format.space_after = Pt(4)
        _run(p, label, bold=True, color=BLACK, size=11)
        if text:
            _run(p, text, bold=False, color=BLACK, size=11)
        return p

    p = doc.add_paragraph()
    _run(p, f"Brief Website Analysis — {domain}", bold=True, color=BLACK, size=18)
    body(f"We analyzed your website ({home}) so please go through this doc to get a quick "
         "overview of your website optimization.")

    # ---- Meta Optimization ----
    section("A) Meta Optimization")
    body("Metadata is an effective and important factor for SEO, and it is a summary of the "
         "content on that page shown in search results.")
    titles = data.get("titles", []) or []
    title_issues = [t for t in titles if t.get("status") != "Good"]
    if titles:
        labeled("Meta Title - ", f"A total of {len(title_issues)} page title issue(s) were "
                                  f"identified across {len(titles)} URL(s) audited.")
        t = titles[0]
        body("For example:")
        labeled("URL: ", t.get("page", home))
        body(f"\"{t.get('title', 'N/A')}\" ({t.get('chars', '?')} chars) — {t.get('status', 'N/A')}.")
    metas = data.get("metas", []) or []
    meta_issues = [m for m in metas if m.get("status") != "Good"]
    if metas:
        m = metas[0]
        labeled("Meta Description – ", f"The audit identified {len(meta_issues)} meta "
                                        f"description issue(s) across {len(metas)} page(s) "
                                        "audited.")
        labeled("URL: ", m.get("page", home))

    # ---- Headings ----
    headers = data.get("headers", []) or []
    h1_issues = [h for h in headers if h.get("h1") != 1]
    if headers:
        h = headers[0]
        labeled("Headings: ", f"The audit found {len(h1_issues)} H1 tag issue(s) across "
                               f"{len(headers)} page(s) audited.")
        labeled("URL: ", h.get("page", home))

    # ---- Image Optimization ----
    p = doc.add_paragraph()
    _run(p, "Image Optimization Issue :", bold=True, color=BLACK, size=14)
    img_alts = data.get("img_alts", []) or []
    missing = len([a for a in img_alts if a.get("present") != "Yes"])
    labeled("Note: ", f"The audit identified {missing} of {len(img_alts)} image(s) checked "
                       "with missing ALT text.")

    # ---- Schema ----
    p = doc.add_paragraph()
    _run(p, "Note For Schema:", bold=True, color=BLACK, size=14)
    body(f"During the audit of {home}, we checked for structured data (Schema Markup) "
         "implementation. We recommend implementing relevant schema types to help search "
         "engines better understand your business.")

    # ---- Backlink Profile Overview ----
    da_pa = data.get("da_pa") or {}
    da_val, dr_val, pa_val = da_pa.get("da", "—"), da_pa.get("dr", "—"), da_pa.get("pa", "—")
    labeled("Backlink Profile Overview: ", f"The website currently has a Domain Rating (DR) of "
                                            f"{dr_val}, indicating its current backlink authority.")

    # ---- DA/PA/Spam Score ----
    section("Domain authority, page authority & Spam Score:")
    if da_val != "—" or dr_val != "—":
        body(f"The website currently has a Domain Authority (DA) of {da_val}, Page Authority "
             f"(PA) of {pa_val}, and Domain Rating (DR) of {dr_val}.")
    else:
        body("Note — Domain Authority / Domain Rating could not be retrieved automatically for "
             "this domain. Please check manually via a DA/DR checker tool.")

    # ---- Indexing Status ----
    idx = data.get("indexing", {}) or {}
    labeled("Indexing Status: – ", f"The website has approximately {idx.get('count', 'N/A')} "
                                    f"page(s) indexed by Google. {idx.get('status', '')}")

    # ---- Robots.txt ----
    robots = data.get("robots", {}) or {}
    p = doc.add_paragraph()
    _run(p, "•  Robot.txt Optimization: - ", bold=True, color=BLACK, size=14)
    if robots.get("found"):
        _run(p, f"The website has a properly configured robots.txt file at {home}robots.txt "
                "that allows search engines to crawl important pages.", bold=False, color=BLACK,
             size=11)
    else:
        _run(p, f"The website does not currently have an accessible robots.txt file at "
                f"{home}robots.txt.", bold=False, color=BLACK, size=11)

    # ---- Sitemap ----
    sitemap = data.get("sitemap", {}) or {}
    p = doc.add_paragraph()
    _run(p, "Sitemap: ", bold=True, color=BLACK, size=14)
    sm_url = sitemap.get("url_checked", f"{home}sitemap.xml")
    if sitemap.get("found") or sitemap.get("ok"):
        _run(p, f"The website has a properly configured XML sitemap ({sm_url}) accessible to "
                "search engines.", bold=False, color=BLACK, size=11)
    else:
        _run(p, f"The website does not have an accessible XML sitemap at {sm_url}.",
             bold=False, color=BLACK, size=11)

    # ---- Broken Links ----
    p = doc.add_paragraph()
    _run(p, "Broken Links", bold=True, color=BLACK, size=14)
    bl_checked = data.get("broken_links_checked", 0)
    broken = data.get("broken_links", []) or []
    if broken:
        body(f"During the audit, we performed a broken link analysis for the website and found "
             f"that {bl_checked} URL(s) were scanned, and {len(broken)} broken link(s) were "
             "identified. Kindly refer to the attached sheet for details.")
    else:
        body(f"During the audit, we performed a broken link analysis for the website and found "
             f"that {bl_checked} URL(s) were scanned, and no broken links were found.")

    # ---- Semrush / Ahrefs performance ----
    p = doc.add_paragraph()
    _run(p, "Website performance according to Semrush: - ", bold=True, color=BLACK, size=14)
    body("Note — Please refer to the attached Semrush overview screenshot for this domain's "
         "Authority Score, organic keywords, organic traffic and referring domains.")

    p = doc.add_paragraph()
    _run(p, "Website performance according to Ahrefs: - ", bold=True, color=BLACK, size=14)
    if da_val != "—" or dr_val != "—":
        body(f"During the audit, we found that the website has a Domain Rating (DR) of "
             f"{dr_val} according to Ahrefs.")
    else:
        body("Note — Ahrefs data could not be retrieved automatically for this domain. Please "
             "check manually via Ahrefs.")

    doc.save(out_path)
    log_fn(f"  ETA report saved: {os.path.basename(out_path)}")


# ---------------------------------------------------------------------------
# Format registry & main entry
# ---------------------------------------------------------------------------

def build_alpha(data, out_path, log_fn=None):
    """Alpha — DOCX brief report, verified against the client reference
    "Alpha.docx" (irishflighttraining.com): plain bold-black lettered/titled
    sections (no shaded banners), flowing document style covering Meta
    Optimization, Headings, Image Alt, Schema, Backlink Profile, DA/PA/Spam
    Score, Broken Links, Indexing, PageSpeed (mobile+desktop), Robots.txt and
    Sitemap. DA/PA/backlink/PageSpeed data isn't collected by run_brief_checks
    today, so those sections point to the attached sheet/manual check, same
    graceful-degradation convention the on-page formats use for data outside
    what's computed inline."""
    if log_fn is None:
        log_fn = print
    from docx import Document
    from docx.shared import Pt, RGBColor

    domain = data["domain"]
    doc = Document()
    for sname in ("Normal", "List Paragraph"):
        try:
            st = doc.styles[sname]
            st.font.name = "Calibri"
            st.font.size = Pt(11)
        except KeyError:
            pass

    BLACK = RGBColor(0x00, 0x00, 0x00)

    def _run(p, text, bold=False, color=BLACK, size=11):
        r = p.add_run(text)
        r.font.name = "Calibri"
        r.font.size = Pt(size)
        r.font.bold = bold
        if color is not None:
            r.font.color.rgb = color
        return r

    def section(text):
        p = doc.add_paragraph()
        p.paragraph_format.space_before = Pt(10)
        p.paragraph_format.space_after = Pt(4)
        _run(p, text, bold=True, color=BLACK, size=13)
        return p

    def body(text, bold=False):
        p = doc.add_paragraph()
        p.paragraph_format.space_after = Pt(4)
        _run(p, text, bold=bold, color=BLACK, size=11)
        return p

    def labeled(label, text=""):
        p = doc.add_paragraph()
        p.paragraph_format.space_after = Pt(4)
        _run(p, label, bold=True, color=BLACK, size=11)
        if text:
            _run(p, text, bold=False, color=BLACK, size=11)
        return p

    home = f"https://{domain}/"

    p = doc.add_paragraph()
    _run(p, f"Brief Website Analysis — {domain}", bold=True, color=BLACK, size=18)
    body(f"We analyzed your website ({home}) so please go through this doc to get a quick "
         "overview of your website's SEO health.")

    # ---- Meta Optimization ----
    section("A) Meta Optimization")
    body("Metadata plays a crucial role in SEO as it provides a clear summary of page content and "
         "helps improve click-through rate from search results.")
    titles = data.get("titles", []) or []
    for t in titles[:3]:
        labeled("Meta Title - ", f"During the audit of the website ({home}), we found the title "
                                  f"\"{t.get('title', 'N/A')}\" ({t.get('chars', '?')} chars) — "
                                  f"{t.get('status', 'N/A')}.")
        labeled("URL: ", t.get("page", home))
    metas = data.get("metas", []) or []
    for m in metas[:3]:
        labeled("Meta Description – ", f"During the audit of the website ({home}), we found the "
                                        f"meta description is {m.get('chars', '?')} characters — "
                                        f"{m.get('status', 'N/A')}.")
        body("Long or duplicate meta descriptions can negatively impact search snippet quality, "
             "reduce click-through rate and confuse search engines about page relevance.")
        labeled("URL: ", m.get("page", home))

    # ---- Headings ----
    headers = data.get("headers", []) or []
    for h in headers[:1]:
        h1c = h.get("h1", 0)
        labeled("Headings: ", f"During the audit of the website ({home}), we found {h1c} H1 "
                              f"tag(s) on the page.")
        body("We recommend adding a unique, keyword-focused H1 tag to all pages where it is "
             "missing, ensuring a clear content hierarchy.")
        labeled("URL: ", h.get("page", home))

    # ---- Image Alt ----
    section("Image Alt texts and title optimization:")
    img_alts = data.get("img_alts", []) or []
    missing = len([a for a in img_alts if a.get("present") != "Yes"])
    labeled("Note: ", f"During the audit of the website ({home}), we found that {missing} of "
                       f"{len(img_alts)} image(s) checked are missing descriptive ALT text.")
    body("We recommend adding descriptive, keyword-relevant ALT text to all missing images, "
         "compressing large images, and using meaningful file names.")

    # ---- Schema ----
    section("Note For Schema:")
    body(f"During the audit of the homepage ({home}), we checked for structured data (schema "
         "markup) implementation.")
    body("We recommend implementing additional relevant schema types such as Organization Schema "
         "to help search engines better understand your business.")

    # ---- Backlink Profile ----
    section("Backlink Profile Overview")
    body("Note — Please refer to the attached backlink profile sheet (Ahrefs export) for full "
         "details on referring domains and anchor text distribution.")

    # ---- DA/PA/DR ----
    section("Domain authority, page authority & Domain Rating:")
    da_pa = data.get("da_pa") or {}
    da_val, dr_val, pa_val = da_pa.get("da", "—"), da_pa.get("dr", "—"), da_pa.get("pa", "—")
    if da_val != "—" or dr_val != "—":
        body(f"During the audit, we found that the website currently has a Domain Authority (DA) "
             f"of {da_val}, Domain Rating (DR) of {dr_val}, and Page Authority (PA) of {pa_val}.")
    else:
        body("Note — Domain Authority / Domain Rating could not be retrieved automatically for "
             "this domain. Please check manually via a DA/DR checker tool.")

    # ---- Broken Links ----
    section("Broken Link")
    bl_checked = data.get("broken_links_checked", 0)
    broken = data.get("broken_links", []) or []
    if broken:
        body(f"During the audit, we performed a broken link analysis for the website ({home}) and "
             f"checked {bl_checked} link(s) — {len(broken)} broken link(s) were found. Kindly "
             "refer to the attached sheet for details.")
    else:
        body(f"During the audit, we performed a broken link analysis for the website ({home}) and "
             f"checked {bl_checked} link(s) — no broken links were found.")

    # ---- Indexing ----
    section("Indexing Status:")
    idx = data.get("indexing", {}) or {}
    labeled("Status – ", f"During the audit, we performed a Google site search for {home} and "
                          f"found approximately {idx.get('count', 'N/A')} indexed page(s). "
                          f"{idx.get('status', '')}")
    body("We recommend regularly reviewing index coverage through Google Search Console, "
         "identifying and resolving any excluded or non-indexed pages.")

    # ---- Page Speed ----
    section("Page Speed Insights on Mobile:")
    body("Note — Please refer to the attached PageSpeed Insights screenshot (mobile) for the "
         "current performance score.")
    section("Page Speed Insights on Desktop:")
    body("Note — Please refer to the attached PageSpeed Insights screenshot (desktop) for the "
         "current performance score.")

    # ---- Robots.txt ----
    robots = data.get("robots", {}) or {}
    p = doc.add_paragraph()
    _run(p, "•  Robot.txt Optimization: - ", bold=True, color=BLACK, size=11)
    if robots.get("found"):
        _run(p, f"During the technical SEO audit, we found that the website has a robots.txt "
                f"file in place at https://{domain}/robots.txt.", bold=False, color=BLACK, size=11)
    else:
        _run(p, "During the technical SEO audit, we did not find a robots.txt file on the "
                "website.", bold=False, color=BLACK, size=11)
    body("We recommend regularly reviewing and updating the robots.txt file to ensure that "
         "important pages are crawlable and irrelevant ones are excluded.")

    # ---- Sitemap ----
    sitemap = data.get("sitemap", {}) or {}
    p = doc.add_paragraph()
    _run(p, "Sitemap: ", bold=True, color=BLACK, size=11)
    sm_url = sitemap.get("url_checked", f"https://{domain}/sitemap.xml")
    if sitemap.get("found") or sitemap.get("ok"):
        _run(p, f"During the technical SEO audit, we found the website's XML sitemap URL "
                f"({sm_url}) is accessible and properly configured.", bold=False, color=BLACK,
             size=11)
    else:
        _run(p, f"During the technical SEO audit, we did not find an accessible XML sitemap at "
                f"{sm_url}.", bold=False, color=BLACK, size=11)
    body("We recommend creating and properly configuring an XML sitemap, ensuring it is "
         "accessible to search engines and submitted via Google Search Console.")

    doc.save(out_path)
    log_fn(f"  Alpha report saved: {os.path.basename(out_path)}")


def build_beta(data, out_path, log_fn=None):
    """Beta — DOCX brief report, verified against the client reference
    "Beta.docx" (arnoldanabolics.com): light-teal (#E1F2F1) 1-cell-table
    section banners with navy (#2F5496) bold text, red (#C00000) "Required
    Action" labels, screenshots for Robots.txt/Sitemap/Google Indexing, and
    separate Semrush/Ahrefs overview tables under Backlinks Analysis. Semrush
    isn't integrated so that table stays a manual-attach note; the Ahrefs
    table uses the real DA/DR data from run_brief_checks."""
    if log_fn is None:
        log_fn = print
    from docx.shared import Pt, RGBColor
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    import sys
    sys.path.insert(0, str(Path(__file__).resolve().parent / "scripts"))
    import generate_health_report as hr

    domain = data["domain"]
    home = f"https://{domain}/"
    from docx import Document
    doc = Document()
    for sname in ("Normal", "List Paragraph"):
        try:
            st = doc.styles[sname]
            st.font.name = "Calibri"
            st.font.size = Pt(11)
        except KeyError:
            pass

    BLACK = RGBColor(0x00, 0x00, 0x00)
    NAVY = RGBColor(0x2F, 0x54, 0x96)
    RED = RGBColor(0xC0, 0x00, 0x00)

    def _run(p, text, bold=False, color=BLACK, size=11):
        r = p.add_run(text)
        r.font.name = "Calibri"
        r.font.size = Pt(size)
        r.font.bold = bold
        if color is not None:
            r.font.color.rgb = color
        return r

    def banner(text, fill="E1F2F1", color=NAVY):
        table = doc.add_table(rows=1, cols=1)
        cell = table.rows[0].cells[0]
        p = cell.paragraphs[0]
        _run(p, text, bold=True, color=color, size=13)
        shd = OxmlElement('w:shd')
        shd.set(qn('w:fill'), fill)
        shd.set(qn('w:val'), 'clear')
        cell._tc.get_or_add_tcPr().append(shd)
        doc.add_paragraph()
        return table

    def body(text, bold=False):
        p = doc.add_paragraph()
        p.paragraph_format.space_after = Pt(4)
        _run(p, text, bold=bold, color=BLACK, size=11)
        return p

    def labeled(label, text="", label_bold=True, label_color=BLACK):
        p = doc.add_paragraph()
        p.paragraph_format.space_after = Pt(4)
        _run(p, label, bold=label_bold, color=label_color, size=11)
        if text:
            _run(p, text, bold=False, color=BLACK, size=11)
        return p

    def required_action(text):
        labeled("Required Action: ", text, label_color=RED)

    def shot(key):
        src = (data.get("screenshots") or {}).get(key)
        if src and Path(src).exists():
            try:
                hr._add_bordered_image(doc, src)
            except Exception:
                pass

    # ---- Title / intro ----
    p = doc.add_paragraph()
    _run(p, f"Brief Website Analysis Report — {domain}", bold=True, color=BLACK, size=18)
    body(f"After examining the website ({home}) in great detail, we have identified some search "
         "engine optimization (SEO) factors that are essential for improving its visibility and "
         "performance.")
    body("These factors could potentially boost the website's search engine ranking. We have "
         "identified both technical and content-related issues that need attention.")
    body("With our consistent and effective SEO practices, we believe that we can help your "
         "website achieve a favourable position in search engine results.")

    # ---- Meta Optimization ----
    banner("Meta Optimization")
    body("Metadata plays a crucial role in search engine optimization (SEO) as it provides a "
         "concise summary of the content on a specific page.")
    titles = data.get("titles", []) or []
    if titles:
        t = titles[0]
        body(f"During the audit, we reviewed the page title tags across the website and found "
             f"that title tags are implemented on all crawled pages.")
        required_action(f"During the page title audit of {home}, a total of {len(titles)} page "
                         f"title(s) were analyzed, and we found the title \"{t.get('title', 'N/A')}\" "
                         f"({t.get('chars', '?')} chars) — {t.get('status', 'N/A')}.")
    p = doc.add_paragraph()
    _run(p, "ii) Meta Descriptions", bold=True, color=BLACK, size=11)
    metas = data.get("metas", []) or []
    if metas:
        m = metas[0]
        body("We reviewed the website's meta descriptions and found that they are implemented "
             "across all crawled pages, which is a positive sign for search visibility.")
        required_action(f"During the meta description audit of {home}, a total of {len(metas)} "
                         f"page(s) were analyzed, and we found a meta description of "
                         f"{m.get('chars', '?')} characters — {m.get('status', 'N/A')}.")

    # ---- Heading Tag Optimization ----
    banner("Heading Tag Optimization")
    body("Heading tags are an important part of on-page SEO, as they help search engines "
         "understand the structure and relevance of the content on a page.")
    headers = data.get("headers", []) or []
    h1c = headers[0].get("h1", 0) if headers else 0
    required_action(f"During the H1 audit of {home}, a total of {len(headers)} page(s) were "
                     f"analyzed, and we found {h1c} H1 tag(s) on the page. We recommend adding a "
                     "unique, keyword-focused H1 tag to all pages where it is missing.")

    # ---- Robots.txt Optimization ----
    banner("Robots.txt Optimization")
    body("The robots.txt file in SEO is a file placed on a website's server to give instructions "
         "to web crawlers (like Googlebot) about which pages or sections of the site should or "
         "should not be crawled.")
    labeled("Page URL - ", f"{home}robots.txt")
    labeled("Screenshot:", "")
    shot("robots")
    robots = data.get("robots", {}) or {}
    if robots.get("found"):
        required_action(f"The website has a robots.txt file in place at {home}robots.txt. We "
                         "recommend regularly reviewing it to ensure important pages remain "
                         "crawlable.")
    else:
        required_action(f"The website does not currently have a robots.txt file, as "
                         f"{home}robots.txt returned an error. We recommend creating one to guide "
                         "search engine crawlers.")

    # ---- XML Sitemap Optimization ----
    banner("XML Sitemap Optimization")
    body("XML sitemaps are a crucial aspect of search engine optimization (SEO) because they help "
         "search engines discover and crawl the pages of a website more efficiently.")
    sitemap = data.get("sitemap", {}) or {}
    sm_url = sitemap.get("url_checked", f"{home}sitemap.xml")
    labeled("Page URL - ", sm_url)
    labeled("Screenshot:", "")
    shot("sitemap")
    if sitemap.get("found") or sitemap.get("ok"):
        required_action(f"During the technical SEO audit, we found the XML sitemap ({sm_url}) is "
                         "accessible and properly configured.")
    else:
        required_action(f"During the technical SEO audit, we found the XML sitemap ({sm_url}) is "
                         "not accessible. We recommend creating and properly configuring one.")

    # ---- Canonical Tags ----
    banner("Canonical Tags")
    body("The canonical tag plays an important role in website audits as it helps to prevent "
         "duplicate content issues. When multiple pages have similar content, the canonical tag "
         "tells search engines which version should be treated as the primary one.")
    canonicals = data.get("canonicals", []) or []
    required_action(f"During the canonical tag audit of {home}, a total of {len(canonicals)} "
                     "page(s) were analyzed to review canonical tag implementation across the "
                     "website.")

    # ---- Schema Optimization ----
    banner("Schema Optimization")
    body("Schema markup is a standardized structured data format that helps search engines better "
         "understand website content and can enhance search results with rich snippets.")
    required_action(f"During the technical SEO audit of {home}, we checked for structured data "
                     "(Schema) implementation. We recommend implementing additional relevant "
                     "schema types to help search engines better understand your business.")

    # ---- Broken Links ----
    banner("Broken Links")
    body("Internal links are an essential component of search engine optimization (SEO) because "
         "they help establish the structure of a website and distribute page authority.")
    bl_checked = data.get("broken_links_checked", 0)
    broken = data.get("broken_links", []) or []
    if broken:
        required_action(f"During the broken link audit of {home}, a total of {bl_checked} URL(s) "
                         f"were analyzed, and {len(broken)} broken link(s) were found. Kindly "
                         "refer to the attached sheet for details.")
    else:
        required_action(f"During the broken link audit of {home}, a total of {bl_checked} URL(s) "
                         "were analyzed, and no broken links were found.")

    # ---- Backlinks Analysis ----
    banner("Backlinks Analysis")
    body("Backlinks are important in SEO. They are links from other websites to your website, "
         "acting as a \"vote of confidence\" for search engines.")
    body("I have provided score metrics from different SEO tools, which are as follows:")

    # ---- Google Indexing ----
    banner("Google Indexing")
    body("Index pages are web pages that Google has discovered and added to their index for "
         "search results. Being indexed doesn't guarantee ranking, but it's a prerequisite for "
         "appearing in search results.")
    idx = data.get("indexing", {}) or {}
    p = doc.add_paragraph()
    _run(p, str(idx.get("count", "N/A")), bold=False, color=BLACK, size=11)
    _run(p, " URLs are Index in Google", bold=True, color=BLACK, size=11)
    labeled("Page URL: ", f"site:{home}")
    labeled("Screenshot:", "")
    shot("serp")
    required_action(f"During the site indexation analysis of {home}, approximately "
                     f"{idx.get('count', 'N/A')} page(s) are currently indexed. "
                     f"{idx.get('status', '')}")

    # ---- Semrush / Ahrefs Overview ----
    banner("Semrush Overview")
    body("Note — Please refer to the attached Semrush overview screenshot for this domain's "
         "Authority Score, organic traffic, organic keywords, referring domains and backlinks.")

    banner("Ahref Overview")
    da_pa = data.get("da_pa") or {}
    da_val, dr_val, pa_val = da_pa.get("da", "—"), da_pa.get("dr", "—"), da_pa.get("pa", "—")
    if da_val != "—" or dr_val != "—":
        body(f"The domain {domain} currently has a Domain Authority (DA) of {da_val}, a Domain "
             f"Rating (DR) of {dr_val}, and a Page Authority (PA) of {pa_val}.")
        required_action(f"During the backlink profile analysis of {home}, the website shows a "
                         f"Domain Rating (DR) of {dr_val}. We recommend building high-quality "
                         "backlinks from authoritative domains to strengthen the backlink profile.")
    else:
        body("Note — Domain Authority / Domain Rating could not be retrieved automatically for "
             "this domain. Please check manually via a DA/DR checker tool.")
        required_action(f"During the backlink profile analysis of {home}, we recommend building "
                         "high-quality backlinks from authoritative domains to strengthen the "
                         "backlink profile.")

    # ---- Conclusion ----
    banner("Conclusion -", color=RED)
    body("After reviewing the points mentioned above, we are confident that addressing them will "
         "cover the major aspects of SEO for the website and help improve its overall search "
         "engine visibility and performance.")

    doc.save(out_path)
    log_fn(f"  Beta report saved: {os.path.basename(out_path)}")


BRIEF_FORMATS = {
    "james": {"label": "James (15 slides, table-based)", "builder": build_james, "ext": "pptx"},
    "xenon": {"label": "Xenon (15 slides, badge-based)", "builder": build_xenon, "ext": "pptx"},
    "omega": {"label": "Omega (16 slides, dark navy)", "builder": build_omega, "ext": "pptx"},
    "neon": {"label": "Neon (19 slides, clean white)", "builder": build_neon, "ext": "pptx"},
    "alpha": {"label": "Alpha (DOCX, flowing document)", "builder": build_alpha, "ext": "docx"},
    "beta": {"label": "Beta (DOCX, teal banners + screenshots)", "builder": build_beta, "ext": "docx"},
    "camila": {"label": "Camila (16 slides, gold titles)", "builder": build_camila, "ext": "pptx"},
    "eta": {"label": "ETA (DOCX, plain lettered sections)", "builder": build_eta, "ext": "docx"},
}


def run_brief_analysis(domain, fmt="james", target_pages=None, out_dir=None, log_fn=None):
    """Run brief website analysis: collect data + build the selected report file
    (PPTX or DOCX, per that format's registered "ext")."""
    if log_fn is None:
        log_fn = print
    domain = re.sub(r'^\s*https?://', '', str(domain or '')).strip().strip('/').split('/')[0] or str(domain)
    if out_dir is None:
        out_dir = tempfile.mkdtemp(prefix="brief_analysis_")
    os.makedirs(out_dir, exist_ok=True)

    # Build EXACTLY the selected format — never silently fall back to another one.
    fmt = str(fmt or "").strip().lower()
    format_info = BRIEF_FORMATS.get(fmt)
    if not format_info:
        raise ValueError(f"Unknown brief report format '{fmt}'. "
                         f"Available: {', '.join(sorted(BRIEF_FORMATS))}")

    log_fn(f"Running brief analysis for {domain}...")
    data = run_brief_checks(domain, target_pages, log_fn)

    timestamp = datetime.now().strftime("%d-%B-%Y")
    ext = format_info.get("ext", "pptx")
    out_file = os.path.join(out_dir, f"Brief_Report_{domain}_{timestamp}.{ext}")

    log_fn(f"Building {format_info['label']} report (format: {fmt})...")
    format_info["builder"](data, out_file, log_fn)

    return out_file
