"""
GSC AUDIT REPORT GENERATOR - v3
- Smart property format detection
- Always (re)submits the latest sitemap so the screenshot shows a fresh date
- Clipped screenshots (no wasted white space)
- Aspect-ratio-preserving placement in slides (no stretching/tilt)
- Gray border on every placed picture
- Lock-safe save

Usage: python 02_generate_report.py <domain>
"""
import sys
import re
import csv
import struct
from pathlib import Path
from urllib.parse import quote
from datetime import datetime
# patchright/playwright isn't in the bundled python. Import it lazily so this
# module stays importable (its helpers — png_dimensions etc. — are reused by the
# health/on-page reports); only the legacy live-capture functions need it.
try:
    from patchright.sync_api import sync_playwright
except Exception:
    sync_playwright = None
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

# ---------- CONFIG ----------
TEMPLATE_PATH = "template.pptx"
AUTH_DIR = Path("auth_states")             # Folder with per-account JSON files
LEGACY_AUTH_FILE = Path("auth_state.json")  # Backward compatibility
ACCOUNTS_CSV = Path("accounts.csv")         # Sheet export: Domain, Account Key, Access Level
OUTPUT_DIR = Path("output")
SCREENSHOTS_DIR = Path("screenshots")

VIEWPORT_W = 1920
VIEWPORT_H = 1080
EXTRA_WAIT_MS = 5000

# Sitemap that gets (re)submitted on every run so the screenshot shows a fresh date
DEFAULT_SITEMAP = "sitemap.xml"

PROPERTY_FORMATS = [
    "sc-domain:{domain}",
    "https://www.{domain}/",
    "https://{domain}/",
    "http://www.{domain}/",
    "http://{domain}/",
]

ACCESS_DENIED_MARKERS = [
    "don't have access to this property",
    "verify your ownership",
    "property not found",
]

# Markers that indicate the session is dead (redirected to Google sign-in)
SIGNIN_MARKERS = [
    "choose an account",
    "use another account",
    "sign in to continue",
    "to continue to search console",
]

# Clip regions per page - capture only relevant area (in pixels)
# Tuned so that the relevant content card / table is included, empty space cut off
CLIP_REGIONS = {
    "sitemap":  {"x": 0,   "y": 0,  "width": VIEWPORT_W, "height": 800},
    "manual":   {"x": 0,   "y": 0,  "width": VIEWPORT_W, "height": 450},
    "security": {"x": 0,   "y": 0,  "width": VIEWPORT_W, "height": 450},
    "removal":  {"x": 0,   "y": 0,  "width": VIEWPORT_W, "height": 800},
    # Sucuri: skip top navbar (y=80) and tighten left/right margins
    "sucuri":   {"x": 100, "y": 80, "width": 1720,       "height": 720},
}

# Available white area below the COLORED HEADER BAND on each content slide (inches)
# (left, top, width, height) - script fits image inside this maintaining aspect ratio
# Top values match the original template's picture placement (designer placed them below the colored band)
SLIDE_TARGET_BOX = {
    3: (0.40, 3.46, 12.53, 3.14),   # Sucuri      (header band ends ~3.46)
    4: (0.40, 3.16, 12.53, 3.44),   # Sitemap     (header band ends ~3.16)
    5: (0.40, 3.49, 12.53, 3.11),   # Manual      (header band ends ~3.49)
    6: (0.40, 3.54, 12.53, 3.06),   # Security    (header band ends ~3.54)
    7: (0.40, 3.14, 12.53, 3.46),   # Removal     (header band ends ~3.14)
}

# Border styling
BORDER_COLOR_RGB = (0x00, 0x00, 0x00)   # black
BORDER_WIDTH_PT  = 3

SLIDE_PICTURE_MAP = {3: "sucuri", 4: "sitemap", 5: "manual", 6: "security", 7: "removal"}


# ---------- HELPERS ----------
def clean_domain(raw):
    d = raw.strip().lower()
    d = re.sub(r'^https?://', '', d)
    d = re.sub(r'^www\.', '', d)
    return d.rstrip('/')


def looks_like_access_denied(html: str) -> bool:
    h = html.lower()
    return any(marker in h for marker in ACCESS_DENIED_MARKERS)


def looks_like_signin_page(url: str, html: str) -> bool:
    """Detect if we've been redirected to Google's sign-in / account-chooser page."""
    if "accounts.google.com" in url.lower() or "servicelogin" in url.lower():
        return True
    h = html.lower()
    return any(marker in h for marker in SIGNIN_MARKERS)


def png_dimensions(path):
    """Read PNG width/height from the file header. No PIL needed."""
    with open(path, "rb") as f:
        signature = f.read(8)
        if signature[:4] != b'\x89PNG':
            raise ValueError(f"Not a PNG file: {path}")
        f.read(4)             # chunk length
        f.read(4)             # chunk type "IHDR"
        w = struct.unpack(">I", f.read(4))[0]
        h = struct.unpack(">I", f.read(4))[0]
    return w, h


def get_auth_files():
    """Return list of all auth-state JSON files. Falls back to legacy single file."""
    files = []
    if AUTH_DIR.exists():
        files = sorted(AUTH_DIR.glob("*.json"))
    if not files and LEGACY_AUTH_FILE.exists():
        files = [LEGACY_AUTH_FILE]
    return files


def load_accounts_csv():
    """
    Load accounts.csv if present.
    Returns dict: bare_domain -> {"account": str, "access_level": str}
    Tolerant of column order and capitalization - matches headers by substring.
    """
    if not ACCOUNTS_CSV.exists():
        return {}

    mapping = {}
    try:
        with open(ACCOUNTS_CSV, encoding='utf-8-sig', newline='') as f:
            reader = csv.DictReader(f)
            headers = reader.fieldnames or []
            field_map = {}
            for h in headers:
                key = h.strip().lower()
                if 'domain' in key and 'sub' not in key:
                    field_map['domain'] = h
                elif 'account' in key and 'key' in key:
                    field_map['account'] = h
                elif 'access' in key and 'level' in key:
                    field_map['access'] = h

            if 'domain' not in field_map or 'account' not in field_map:
                print(f"  [warn] accounts.csv me 'Domain' aur 'Account Key' columns nahi mile")
                print(f"     Headers found: {headers}")
                return {}

            for row in reader:
                raw_domain = row.get(field_map['domain'], '') or ''
                d = clean_domain(raw_domain)
                if not d:
                    continue
                account = (row.get(field_map['account'], '') or '').strip().lower()
                if not account:
                    continue
                access = (row.get(field_map.get('access', ''), '') or '').strip()
                mapping[d] = {"account": account, "access_level": access}
    except Exception as e:
        print(f"  [warn] accounts.csv parse error: {type(e).__name__}: {e}")
        return {}

    return mapping


def format_hints_from_access_level(access_level: str):
    """
    Map Access Level value (e.g. 'Full (domain)' / 'Full (URL prefix)') to
    the property URL formats to try. Returns None to mean 'try all'.
    """
    if not access_level:
        return None
    level = access_level.lower()
    if 'domain' in level:
        return ["sc-domain:{domain}"]
    if 'url prefix' in level or 'prefix' in level:
        return [
            "https://www.{domain}/",
            "https://{domain}/",
            "http://www.{domain}/",
            "http://{domain}/",
        ]
    return None


def fit_in_box(img_path, box_left, box_top, box_width, box_height):
    """Return (left, top, width, height) in inches such that the image fits
    inside the box, centered, maintaining aspect ratio. No stretching."""
    img_w_px, img_h_px = png_dimensions(img_path)
    img_aspect = img_w_px / img_h_px
    box_aspect = box_width / box_height

    if img_aspect >= box_aspect:
        # Image is wider relative to box -> fit by width
        final_w = box_width
        final_h = box_width / img_aspect
    else:
        # Image is taller relative to box -> fit by height
        final_h = box_height
        final_w = box_height * img_aspect

    final_left = box_left + (box_width - final_w) / 2
    final_top  = box_top  + (box_height - final_h) / 2
    return final_left, final_top, final_w, final_h


# ---------- STEP 1: Property format detection ----------
class SessionExpired(Exception):
    pass


def find_working_property(page, domain, format_hint=None):
    print(f"  [*] Detecting property format for '{domain}'...")

    formats = PROPERTY_FORMATS[:]
    # Reorder based on hint (so the right format is tried first)
    if format_hint:
        fh = format_hint.lower()
        if "domain" in fh and "url" not in fh:
            # Hinted as Domain property - try sc-domain: first
            formats = ["sc-domain:{domain}"] + [f for f in formats if "sc-domain" not in f]
        elif "url" in fh or "prefix" in fh or "https" in fh:
            # Hinted as URL prefix - try https/http first
            formats = [f for f in formats if "sc-domain" not in f] + ["sc-domain:{domain}"]

    for fmt in formats:
        prop = fmt.format(domain=domain)
        encoded = quote(prop, safe='')
        test_url = f"https://search.google.com/search-console/sitemaps?resource_id={encoded}"
        try:
            page.goto(test_url, wait_until="domcontentloaded", timeout=30000)
            page.wait_for_timeout(3500)
            html = page.content()
            current_url = page.url

            if looks_like_signin_page(current_url, html):
                raise SessionExpired("Redirected to Google sign-in / account-chooser page")

            if looks_like_access_denied(html):
                print(f"     [no ] {prop}")
                continue
            print(f"     [YES] {prop}  <-- using this")
            return prop
        except SessionExpired:
            raise
        except Exception as e:
            print(f"     [err] {prop}: {type(e).__name__}")
    return None


def build_gsc_urls(prop):
    encoded = quote(prop, safe='')
    return {
        "sitemap":  f"https://search.google.com/search-console/sitemaps?resource_id={encoded}",
        "manual":   f"https://search.google.com/search-console/manual-actions?resource_id={encoded}",
        "security": f"https://search.google.com/search-console/security-issues?resource_id={encoded}",
        "removal":  f"https://search.google.com/search-console/removals?resource_id={encoded}",
    }


# ---------- STEP 2: Sitemap auto-submit ----------
def has_submitted_sitemap(page):
    """
    Detect if any sitemap row exists in the 'Submitted sitemaps' table.
    Looks for anchor tags whose text matches a sitemap path like '/sitemap.xml'.
    More robust than scanning innerText (which can false-match placeholder text).
    """
    try:
        count = page.evaluate("""
            () => {
                const links = document.querySelectorAll('a');
                let found = 0;
                for (const a of links) {
                    const text = (a.textContent || '').trim();
                    if (/^\\/[\\w\\-\\/\\.]*\\.xml(\\.gz)?$/i.test(text)) found++;
                }
                return found;
            }
        """)
        return count > 0
    except Exception:
        return False


def _wait_for_submit_success(page, timeout_ms=20000):
    """
    After clicking Submit, wait for GSC to confirm the submission succeeded.
    GSC shows a "Sitemap submitted successfully" snackbar/dialog and then the
    row's status becomes 'Success'. We poll for either signal instead of a
    blind fixed sleep, so we don't screenshot before the date refreshes.
    Returns True if a success signal was seen.
    """
    try:
        page.wait_for_function(
            """() => {
                const txt = (document.body.innerText || '').toLowerCase();
                return txt.includes('sitemap submitted successfully')
                    || txt.includes('submitted successfully')
                    || txt.includes('couldn')         // "Couldn't fetch" still = processed
                    || txt.includes('success');
            }""",
            timeout=timeout_ms,
        )
        return True
    except Exception:
        return False


def submit_latest_sitemap(page, sitemap=DEFAULT_SITEMAP):
    """
    ALWAYS (re)submit the latest sitemap, then wait until GSC confirms the
    submission and the table reflects the fresh 'Last read' date.

    This is intentionally unconditional: even when a sitemap row already
    exists, re-submitting makes GSC re-read the sitemap and refresh its date,
    so the screenshot always shows the most recent submission. Shared across
    all webmaster-audit formats (James / James new / Omega / Neon).
    """
    page.wait_for_timeout(3000)

    already = has_submitted_sitemap(page)
    print(f"     submitting latest sitemap '{sitemap}'"
          f"{' (re-submit to refresh date)' if already else ''}...")
    try:
        # CRITICAL: GSC page has multiple inputs (top search bar, sitemap input).
        # Target the specific sitemap input by its placeholder text.
        input_field = page.get_by_placeholder(
            re.compile(r"Enter sitemap URL", re.I)
        ).first
        input_field.wait_for(state="visible", timeout=10000)

        # Click + clear + fill - some Material inputs need this sequence
        input_field.click()
        input_field.fill("")
        input_field.fill(sitemap)
        page.wait_for_timeout(1500)

        # Find Submit button. It should now be enabled (was disabled when input empty).
        # Wait briefly for the button to become enabled before clicking.
        submit_btn = page.get_by_role(
            "button", name=re.compile(r"^\s*submit\s*$", re.I)
        ).first

        # Wait up to 5s for button to be enabled (not just visible)
        try:
            submit_btn.wait_for(state="visible", timeout=5000)
        except Exception:
            pass

        # Try click; if disabled, fall back to pressing Enter in the input
        try:
            submit_btn.click(timeout=10000)
        except Exception as e:
            print(f"     [info] submit button click failed ({type(e).__name__}), trying Enter key")
            input_field.press("Enter")

        # Wait for GSC to confirm the submission (success snackbar/dialog),
        # not just a blind sleep - this is the "wait for submit to complete" step.
        if _wait_for_submit_success(page):
            print("     submission confirmed by GSC")
        else:
            print("     [info] no explicit success signal, continuing")
        page.wait_for_timeout(2500)

        # Dismiss any confirmation dialog: "GOT IT" / "OK" / "Done"
        for label_pat in [r"got\s*it", r"^ok$", r"^done$", r"close"]:
            try:
                page.get_by_role("button", name=re.compile(label_pat, re.I)).first.click(timeout=1500)
                page.wait_for_timeout(500)
                break
            except Exception:
                continue

        # Reload so the (re)submitted row + refreshed 'Last read' date appear,
        # then verify before the caller screenshots. Retry once if slow.
        page.wait_for_timeout(2000)
        page.reload(wait_until="networkidle", timeout=45000)
        page.wait_for_timeout(5000)

        if has_submitted_sitemap(page):
            print("     sitemap submitted - row visible with refreshed date")
        else:
            print("     [info] table not yet updated, waiting longer...")
            page.wait_for_timeout(8000)
            page.reload(wait_until="networkidle", timeout=45000)
            page.wait_for_timeout(4000)
            if has_submitted_sitemap(page):
                print("     sitemap submitted - row visible (after retry)")
            else:
                print("     [warn] sitemap submitted but table still empty in screenshot")
    except Exception as e:
        print(f"     [warn] sitemap submit issue: {type(e).__name__}: {e}")


# Backward-compatible alias: older call sites used ensure_sitemap_submitted.
# The behavior is now "always submit the latest sitemap" (see above).
def ensure_sitemap_submitted(page):
    submit_latest_sitemap(page)


# ---------- STEP 3: Screenshots (with clipping) ----------
def capture_gsc_page(page, key, url, path):
    print(f"  -> [{key}]")
    try:
        page.goto(url, wait_until="networkidle", timeout=60000)
        page.wait_for_timeout(EXTRA_WAIT_MS)
    except Exception as e:
        print(f"     [warn] {type(e).__name__}, capturing anyway")
        page.wait_for_timeout(3000)

    if key == "sitemap":
        ensure_sitemap_submitted(page)

    clip = CLIP_REGIONS.get(key)
    page.screenshot(path=str(path), clip=clip, full_page=False)


def capture_sucuri(page, domain, path):
    url = f"https://sitecheck.sucuri.net/results/https/www.{domain}"
    print(f"  -> [sucuri]")
    try:
        page.goto(url, wait_until="domcontentloaded", timeout=60000)
    except Exception as e:
        print(f"     [warn] navigation: {type(e).__name__}")

    try:
        page.wait_for_function(
            """() => {
                const txt = (document.body.innerText || '').toLowerCase();
                if (txt.includes('please wait')) return false;
                return txt.includes('site is clean')
                    || txt.includes('no malware')
                    || txt.includes('malware found')
                    || txt.includes('not clean')
                    || txt.includes('blacklist')
                    || txt.includes('out-of-date')
                    || txt.includes('outdated software')
                    || txt.includes('minor security')
                    || txt.includes('low security')
                    || txt.includes('scan results');
            }""",
            timeout=45000,
        )
        print(f"     scan complete")
    except Exception:
        print(f"     [warn] scan still pending, capturing snapshot")

    page.wait_for_timeout(2500)

    # Dismiss any cookie banner so it doesn't appear in the screenshot
    try:
        for label in ["Accept", "ACCEPT", "Got it", "I Accept", "Agree"]:
            try:
                page.locator("button").filter(
                    has_text=re.compile(rf"^\s*{re.escape(label)}\s*$", re.I)
                ).first.click(timeout=1200)
                page.wait_for_timeout(500)
                break
            except Exception:
                continue
    except Exception:
        pass

    page.screenshot(path=str(path), clip=CLIP_REGIONS["sucuri"], full_page=False)


def capture_all_screenshots(domain, account_hint=None, format_hint=None):
    SCREENSHOTS_DIR.mkdir(exist_ok=True)
    out = {}

    # ---- AUTO-LOOKUP from accounts.csv (sheet) ----
    csv_map = load_accounts_csv()
    csv_entry = csv_map.get(domain)
    if csv_map:
        print(f"  [*] accounts.csv loaded: {len(csv_map)} domains mapped")
    if csv_entry:
        print(f"  [*] Sheet match: '{domain}' -> "
              f"account='{csv_entry['account']}', level='{csv_entry['access_level']}'")
        # CLI hints always override CSV
        if not account_hint:
            account_hint = csv_entry['account']
        if not format_hint:
            format_hint = csv_entry['access_level']
    elif csv_map:
        print(f"  [!] '{domain}' not found in accounts.csv - trying all accounts")

    auth_files = get_auth_files()
    if not auth_files:
        print(f"\n[ERROR] Koi auth file nahi mili.")
        print(f"  Pehle setup karo:")
        print(f"     python 01_auth_setup.py <account_name>")
        print(f"  Example:")
        print(f"     python 01_auth_setup.py default")
        sys.exit(1)

    # If account hint is provided, try that first (others as fallback)
    if account_hint:
        hint_lower = account_hint.lower()
        hinted   = [f for f in auth_files if f.stem.lower() == hint_lower]
        others   = [f for f in auth_files if f.stem.lower() != hint_lower]
        if hinted:
            auth_files = hinted + others
            print(f"  [*] Account hint: '{account_hint}' -> trying first")
        else:
            print(f"  [!] Account hint '{account_hint}' not found, trying all accounts")

    print(f"  [*] {len(auth_files)} account(s) found: {[f.stem for f in auth_files]}")

    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)

        working_context = None
        working_page = None
        working_prop = None
        working_account = None

        # Try each saved account until one has access to this domain
        for auth_file in auth_files:
            print(f"\n  [*] Trying account: '{auth_file.stem}'")
            context = browser.new_context(
                storage_state=str(auth_file),
                viewport={"width": VIEWPORT_W, "height": VIEWPORT_H},
            )
            page = context.new_page()

            try:
                prop = find_working_property(page, domain, format_hint=format_hint)
                if prop:
                    working_context = context
                    working_page = page
                    working_prop = prop
                    working_account = auth_file.stem
                    break
                else:
                    context.close()
            except SessionExpired:
                print(f"     [session expired] '{auth_file.stem}' - refresh with: python 01_auth_setup.py {auth_file.stem}")
                context.close()
                continue
            except Exception as e:
                print(f"     [error] {type(e).__name__}: {e}")
                context.close()
                continue

        if not working_prop:
            browser.close()
            print(f"\n{'='*55}")
            print(f"[ERROR] '{domain}' did not match any account")
            print(f"{'='*55}")
            print(f"  Tried accounts: {[f.stem for f in auth_files]}")
            print(f"")
            print(f"  Possible reasons:")
            print(f"   - Domain is not added to any of the accounts")
            print(f"   - All account sessions have expired")
            print(f"     -> Refresh: python 01_auth_setup.py <account_name>")
            print(f"   - You need to add a new account")
            print(f"     -> New auth: python 01_auth_setup.py <new_account_name>")
            print(f"{'='*55}\n")
            sys.exit(2)

        # Continue with the matched account's page
        page = working_page

        for key, url in build_gsc_urls(working_prop).items():
            path = SCREENSHOTS_DIR / f"{domain}_{key}.png"
            capture_gsc_page(page, key, url, path)
            out[key] = str(path)

        sucuri_path = SCREENSHOTS_DIR / f"{domain}_sucuri.png"
        capture_sucuri(page, domain, sucuri_path)
        out["sucuri"] = str(sucuri_path)

        browser.close()

    return out, working_prop, working_account


# ---------- STEP 4: PPTX build with proper sizing + borders ----------
def remove_existing_picture(slide):
    for shape in list(slide.shapes):
        if shape.shape_type == 13:  # PICTURE
            sp = shape._element
            sp.getparent().remove(sp)
            return True
    return False


def add_picture_with_border(slide, slide_idx, image_path):
    """Place image in the slide's target box, aspect-preserving, with gray border."""
    box_left, box_top, box_w, box_h = SLIDE_TARGET_BOX[slide_idx]
    fl, ft, fw, fh = fit_in_box(image_path, box_left, box_top, box_w, box_h)

    pic = slide.shapes.add_picture(
        image_path,
        Inches(fl), Inches(ft),
        Inches(fw), Inches(fh),
    )

    # Border (gray, 1pt)
    line = pic.line
    line.color.rgb = RGBColor(*BORDER_COLOR_RGB)
    line.width = Pt(BORDER_WIDTH_PT)

    return fl, ft, fw, fh


def replace_text_in_slide(slide, replacements: dict):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                for old, new in replacements.items():
                    if old in run.text:
                        run.text = run.text.replace(old, new)


def is_file_locked(path: Path) -> bool:
    if not path.exists():
        return False
    try:
        with open(path, "a"):
            return False
    except (PermissionError, OSError):
        return True


def safe_save_path(base_path: Path) -> Path:
    if not is_file_locked(base_path):
        return base_path
    stamp = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
    new_path = base_path.with_name(f"{base_path.stem}_{stamp}{base_path.suffix}")
    print(f"  [!] {base_path.name} is open in PowerPoint (locked).")
    print(f"  [!] Using new name: {new_path.name}")
    return new_path


def build_report(domain, screenshots):
    OUTPUT_DIR.mkdir(exist_ok=True)
    prs = Presentation(TEMPLATE_PATH)

    OLD_DOMAIN = "richardrosalaw.com"
    replacements = {
        f"https://www.{OLD_DOMAIN}/": f"https://www.{domain}/",
        OLD_DOMAIN: domain,
    }

    for slide in prs.slides:
        replace_text_in_slide(slide, replacements)

    for slide_idx, key in SLIDE_PICTURE_MAP.items():
        slide = prs.slides[slide_idx - 1]
        remove_existing_picture(slide)
        fl, ft, fw, fh = add_picture_with_border(slide, slide_idx, screenshots[key])
        print(f"  -> Slide {slide_idx} ({key:8s}): "
              f"placed at L={fl:.2f} T={ft:.2f} W={fw:.2f} H={fh:.2f} + border")

    safe_domain = domain.replace('.', '_').replace('/', '_')
    base_path = OUTPUT_DIR / f"GSC_Audit_{safe_domain}.pptx"
    out_path = safe_save_path(base_path)

    try:
        prs.save(out_path)
    except PermissionError:
        stamp = datetime.now().strftime("%Y%m%d_%H%M%S_%f")
        fallback = OUTPUT_DIR / f"GSC_Audit_{safe_domain}_{stamp}.pptx"
        print(f"  [!] Fallback name: {fallback.name}")
        prs.save(fallback)
        out_path = fallback

    return out_path


# ---------- MAIN ----------
def parse_cli_args(args):
    """Parse: <domain> [--account name] [--format hint]"""
    if not args or args[0].startswith("--"):
        return None, None, None
    domain = args[0]
    account_hint = None
    format_hint = None
    i = 1
    while i < len(args):
        if args[i] == "--account" and i + 1 < len(args):
            account_hint = args[i+1]
            i += 2
        elif args[i] == "--format" and i + 1 < len(args):
            format_hint = args[i+1]
            i += 2
        else:
            i += 1
    return domain, account_hint, format_hint


def main():
    domain, account_hint, format_hint = parse_cli_args(sys.argv[1:])
    if not domain:
        print("Usage:   python 02_generate_report.py <domain> [--account NAME] [--format HINT]")
        print("Example: python 02_generate_report.py removalistauction.com.au")
        print("         python 02_generate_report.py example.com --account default")
        print("         python 02_generate_report.py example.com --account default --format \"URL prefix\"")
        sys.exit(1)

    domain = clean_domain(domain)
    print(f"\n{'='*55}")
    print(f"  GSC Audit Report Generator v3.2")
    print(f"  Domain : {domain}")
    if account_hint:
        print(f"  Hint   : account='{account_hint}'")
    if format_hint:
        print(f"  Hint   : format='{format_hint}'")
    print(f"{'='*55}")

    print(f"\n[1/2] Capturing screenshots...")
    shots, working_prop, working_account = capture_all_screenshots(domain, account_hint, format_hint)

    print(f"\n[2/2] Building PPTX...")
    out = build_report(domain, shots)

    print(f"\n{'='*55}")
    print(f"  [DONE] Report ready!")
    print(f"  File          : {out.resolve()}")
    print(f"  Account used  : {working_account}")
    print(f"  Property used : {working_prop}")
    print(f"{'='*55}\n")


if __name__ == "__main__":
    main()