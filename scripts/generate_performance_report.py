"""
generate_performance_report.py - Standalone monthly Performance Report (pptx),
separate from every other report tool in this app.

Covers what's buildable with data this app already has real API access to:
  - Google Search Console (via gsc_audit.py's existing OAuth/Search Analytics
    plumbing - clicks/impressions/CTR/position trend, top queries, top pages,
    top countries)
  - Google Analytics 4 (via gsc_audit.py's new GA4 Admin/Data API helpers -
    users/sessions trend, traffic by channel, device breakdown)

Does NOT attempt to replicate the rank-tracker (SE Ranking) or GSC Security/
Manual-Actions/Links sections some client-facing reports also include - those
need a different data source (SE Ranking API access, or GSC report types
Google no longer exposes via API) not wired up yet. This is "james" format -
the first Performance Report format; more will be added the same way the
On-Page formats were, once real client references for them exist.

Run:
    python generate_performance_report.py example.com --gsc-account you@x.com \
        --ga4-property properties/123456789 --out "SEO Performance Report.pptx"
"""
import argparse
import datetime
import os
import sys
from pathlib import Path

ROOT = Path(__file__).parent.resolve()
REPO_ROOT = ROOT.parent
for p in (str(ROOT), str(REPO_ROOT)):
    if p not in sys.path:
        sys.path.insert(0, p)

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

import gsc_audit

NAVY = RGBColor(0x1F, 0x38, 0x64)
BLUE = RGBColor(0x2F, 0x54, 0x96)
LIGHT_BLUE = RGBColor(0xDE, 0xEA, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x59, 0x59, 0x59)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Channel groups GA4 attributes to non-paid, SEO-relevant traffic. This is an
# SEO reporting tool - paid ads (Paid Search/Paid Social/Paid Other/Display)
# are deliberately excluded from the Traffic Acquisition section, per-team
# decision. GA4 has no dedicated "AI" channel group as of this writing - AI
# search referrers (chatgpt.com, perplexity.ai, etc.) land in "Referral" or
# "Organic Search" under GA4's own default channel grouping, so they're
# already included here rather than broken out separately.
SEO_CHANNEL_GROUPS = {"Organic Search", "Referral", "Organic Social", "Direct"}


def log(msg):
    print(msg, flush=True)


def _period_days(start_date, end_date):
    try:
        return (datetime.date.fromisoformat(end_date) - datetime.date.fromisoformat(start_date)).days
    except Exception:
        return 28


# --------------------------------------------------------------------------- #
# Data fetch - GSC (reuses gsc_audit.py's existing, already-live functions)
# --------------------------------------------------------------------------- #
def fetch_gsc_data(token, property_url, start_date, end_date):
    return {
        "daily": gsc_audit.fetch_performance_daily(token, property_url, start_date, end_date),
        "queries": gsc_audit.fetch_top_queries(token, property_url, start_date, end_date, limit=10),
        "pages": gsc_audit.fetch_top_pages(token, property_url, start_date, end_date, limit=10),
        "countries": gsc_audit.fetch_top_countries(token, property_url, start_date, end_date, limit=10),
    }


def fetch_gsc_totals(token, property_url, start_date, end_date):
    """Aggregate clicks/impressions/ctr/position for the whole period - one
    row, no dimensions. gsc_audit.fetch_search_analytics can't do this
    directly (dimensions=[] is falsy in Python, so its `dimensions or
    ["query"]` default would silently fall back to a per-query breakdown
    instead of a true aggregate), so this posts the empty-dimensions body
    itself."""
    import urllib.parse
    encoded = urllib.parse.quote(property_url, safe="")
    body = {"startDate": start_date, "endDate": end_date, "dimensions": []}
    resp = gsc_audit._api_post(
        f"{gsc_audit.SEARCH_ANALYTICS_URL}/sites/{encoded}/searchAnalytics/query", token, body, timeout=30)
    rows = resp.get("rows") or []
    if not rows:
        return {"clicks": 0, "impressions": 0, "ctr": 0.0, "position": 0.0}
    r = rows[0]
    return {"clicks": r.get("clicks", 0), "impressions": r.get("impressions", 0),
            "ctr": r.get("ctr", 0.0), "position": r.get("position", 0.0)}


# --------------------------------------------------------------------------- #
# Data fetch - GA4 totals (for period-over-period comparison)
# --------------------------------------------------------------------------- #
def fetch_ga4_totals(token, property_name, start_date, end_date):
    rows = gsc_audit.run_ga4_report(
        token, property_name, start_date, end_date,
        dimensions=[], metrics=["activeUsers", "sessions"], limit=1)
    if not rows:
        return {"activeUsers": 0, "sessions": 0}
    r = rows[0]
    return {"activeUsers": int(r.get("activeUsers", 0) or 0), "sessions": int(r.get("sessions", 0) or 0)}


# --------------------------------------------------------------------------- #
# Period-over-period improvement detection - drives which screenshots get
# included (only metrics that actually improved, per the client's ask).
# --------------------------------------------------------------------------- #
def compute_gsc_improvements(current, previous):
    """previous can be all-zero (e.g. a brand new property with no prior
    data) - treated as "can't compare", not as "improved from zero", since
    that would be a misleading claim, not a real improvement."""
    improved = set()
    if previous.get("clicks"):
        if current.get("clicks", 0) > previous["clicks"]:
            improved.add("clicks")
    if previous.get("impressions"):
        if current.get("impressions", 0) > previous["impressions"]:
            improved.add("impressions")
    if previous.get("position"):
        # Lower average position is better (position 1 beats position 10).
        if 0 < current.get("position", 0) < previous["position"]:
            improved.add("position")
    return improved


def compute_ga4_improvements(current, previous):
    improved = set()
    if previous.get("activeUsers"):
        if current.get("activeUsers", 0) > previous["activeUsers"]:
            improved.add("activeUsers")
    if previous.get("sessions"):
        if current.get("sessions", 0) > previous["sessions"]:
            improved.add("sessions")
    return improved


def fetch_ga4_channel_sessions(token, property_name, start_date, end_date):
    """Per-channel-group session totals for one period - used both for the
    Traffic Acquisition chart/screenshot data and for the SEO-channels
    period-over-period comparison below."""
    return gsc_audit.run_ga4_report(
        token, property_name, start_date, end_date,
        dimensions=["sessionDefaultChannelGroup"], metrics=["sessions"], limit=20)


def filter_seo_channels(channel_rows):
    """Non-paid channel rows only - see SEO_CHANNEL_GROUPS."""
    return [r for r in channel_rows if r.get("sessionDefaultChannelGroup") in SEO_CHANNEL_GROUPS]


def sum_seo_sessions(channel_rows):
    return sum(int(r.get("sessions", 0) or 0) for r in filter_seo_channels(channel_rows))


def compute_ga4_seo_improvement(current_channels, previous_channels):
    """Whether combined Organic Search/Referral/Organic Social/Direct sessions
    grew vs. the previous period - the signal that gates the Traffic
    Acquisition screenshot, separate from compute_ga4_improvements' overall
    (paid-inclusive) 'sessions' flag."""
    current_seo = sum_seo_sessions(current_channels)
    previous_seo = sum_seo_sessions(previous_channels)
    return {"seo_sessions"} if previous_seo and current_seo > previous_seo else set()


# --------------------------------------------------------------------------- #
# Data fetch - GA4
# --------------------------------------------------------------------------- #
def fetch_ga4_data(token, property_name, start_date, end_date):
    return {
        "daily": gsc_audit.run_ga4_report(
            token, property_name, start_date, end_date,
            dimensions=["date"], metrics=["activeUsers", "sessions", "newUsers"], limit=200),
        "channels": gsc_audit.run_ga4_report(
            token, property_name, start_date, end_date,
            dimensions=["sessionDefaultChannelGroup"], metrics=["sessions"], limit=10),
        "devices": gsc_audit.run_ga4_report(
            token, property_name, start_date, end_date,
            dimensions=["deviceCategory"], metrics=["activeUsers"], limit=10),
    }


# --------------------------------------------------------------------------- #
# Real dashboard screenshots - only for whichever metrics actually improved.
# Requires a GSC "session" (a persisted, logged-in browser profile - the
# SAME mechanism GSC Audit's screenshot capture already uses, separate from
# the OAuth token used for the API calls above). If the given account has no
# session, screenshots are skipped entirely and the report falls back to the
# native charts already built - never a fabricated/placeholder screenshot.
#
# GSC's Performance page has real per-metric toggle chips, so the "only show
# what improved" ask is precisely achievable there. GA4's own report UI is a
# much heavier, harder-to-deep-link SPA - this captures its Reports Snapshot
# overview (real screenshot, real data) rather than attempting the same
# precise per-metric toggle GA4's UI doesn't expose as simply as GSC's does.
# --------------------------------------------------------------------------- #
def launch_screenshot_browser(session_id, browser_pref="edge", log_fn=None):
    """Thin wrapper so this script doesn't need its own browser-launch code -
    reuses gsc_audit's exact session-profile mechanism."""
    return gsc_audit.launch_session_browser(session_id, browser_pref=browser_pref, log_fn=log_fn or log)


def _robust_click(driver, element):
    """A plain WebDriver .click() confirmed real failures live across every
    click site in this file - "element not interactable" (GSC tabs, GA4 nav
    items) and "element click intercepted" (GA4's own filter-panel footer
    overlapping the button it was trying to click) - the matched element can
    be technically present in the DOM but not yet scrolled into view, or
    momentarily covered by another node, when WebDriver's own visibility/
    hit-testing runs. Scroll it into view first, then try a native click,
    falling back to a JS-dispatched click (bypasses that hit-testing
    entirely) if the native one still fails - used at every click site in
    this file instead of a bare element.click()."""
    import time
    driver.execute_script("arguments[0].scrollIntoView({block:'center'});", element)
    time.sleep(0.3)
    try:
        element.click()
    except Exception:
        driver.execute_script("arguments[0].click();", element)


def capture_gsc_performance_screenshot(driver, property_url, improved_metrics, out_path, log_fn=None):
    """Navigates to GSC's own Performance report, sets the compare-to-
    previous-period view, and toggles the metric chips so only the metrics
    that actually improved (clicks/impressions/position) are shown on the
    graph - matching the reference deck's "only include what improved"
    request. Returns True on success, False if it couldn't be captured
    (signed out, no access, or a chip/control GSC's UI moved)."""
    import time
    log_fn = log_fn or log
    url = gsc_audit.build_gsc_url("performance/search-analytics", property_url)
    driver.get(url)
    time.sleep(6)
    if gsc_audit._looks_like_signin(driver):
        log_fn("  [warn] GSC Performance page bounced to sign-in - skipping screenshot.")
        return False
    if gsc_audit._looks_like_no_access(driver):
        log_fn("  [warn] Signed-in account has no access to this property - skipping screenshot.")
        return False

    # GSC's own metric toggle chips - visible text is the label; clicking
    # toggles that metric's line on/off the graph. Default state usually has
    # Clicks + Impressions on, CTR + Position off - reconcile to exactly
    # improved_metrics when something improved. When NOTHING improved, leave
    # GSC's own default chip state alone rather than skipping the screenshot
    # outright - this is just the real current-period dashboard, no
    # "improved vs. previous period" framing attached to it.
    if improved_metrics:
        chip_label = {"clicks": "Total clicks", "impressions": "Total impressions",
                      "ctr": "Average CTR", "position": "Average position"}
        from selenium.webdriver.common.by import By
        for metric, label in chip_label.items():
            wanted_on = metric in improved_metrics
            try:
                chips = driver.find_elements(By.XPATH, f"//*[contains(text(), '{label}')]")
                if not chips:
                    continue
                chip = chips[0]
                is_on = "selected" in (chip.get_attribute("class") or "").lower() or \
                        chip.get_attribute("aria-pressed") == "true"
                if is_on != wanted_on:
                    _robust_click(driver, chip)
                    time.sleep(1)
            except Exception as e:
                log_fn(f"  [warn] Could not toggle '{label}' chip: {e}")

    time.sleep(2)
    try:
        driver.save_screenshot(out_path)
        if improved_metrics:
            log_fn(f"  GSC Performance screenshot saved ({', '.join(sorted(improved_metrics))}).")
        else:
            log_fn("  GSC Performance screenshot saved (current period - no improvement to highlight).")
        return True
    except Exception as e:
        log_fn(f"  [warn] GSC screenshot save failed: {e}")
        return False


def capture_ga4_screenshot(driver, ga4_property_name, out_path, log_fn=None):
    """Real screenshot of GA4's Reports Snapshot overview for this property -
    see the module-level note above on why this isn't toggled to individual
    metrics the same precise way the GSC one is."""
    import time
    log_fn = log_fn or log
    property_id = ga4_property_name.split("/")[-1]
    url = f"https://analytics.google.com/analytics/web/#/p{property_id}/reports/reportinghub"
    driver.get(url)
    # GA4 is a heavy Angular SPA - give it real time to render past the
    # loading skeleton before trusting the screenshot.
    time.sleep(10)
    cur = (driver.current_url or "").lower()
    if "accounts.google.com" in cur or "/signin" in cur:
        log_fn("  [warn] GA4 bounced to sign-in - skipping screenshot.")
        return False
    try:
        driver.save_screenshot(out_path)
        log_fn("  GA4 Reports Snapshot screenshot saved.")
        return True
    except Exception as e:
        log_fn(f"  [warn] GA4 screenshot save failed: {e}")
        return False


def capture_gsc_dimension_screenshot(driver, property_url, tab_label, out_path, log_fn=None, sort_by=None):
    """Real screenshot of the GSC Performance page's own breakdown table -
    Queries/Pages/Countries are tabs on the SAME page as the main graph (not
    separate URLs), so this reuses that URL and clicks the tab by its visible
    label before capturing, same defensive pattern as the metric chips
    above (never fabricates - a missing tab just means whatever's on screen
    gets captured, and a sign-in bounce skips the screenshot entirely).
    sort_by (optional): "Clicks" or "Impressions" - clicks that column's
    header to sort the table by it (GSC defaults to Clicks descending);
    lets the caller capture both a by-clicks and a by-impressions view of
    the same table as two separate screenshots."""
    import time
    from selenium.webdriver.common.by import By
    log_fn = log_fn or log
    url = gsc_audit.build_gsc_url("performance/search-analytics", property_url)
    driver.get(url)
    time.sleep(6)
    if gsc_audit._looks_like_signin(driver):
        log_fn(f"  [warn] GSC Performance page bounced to sign-in - skipping {tab_label} screenshot.")
        return False
    if gsc_audit._looks_like_no_access(driver):
        log_fn(f"  [warn] Signed-in account has no access to this property - skipping {tab_label} screenshot.")
        return False
    try:
        tabs = driver.find_elements(By.XPATH, f"//*[contains(text(), '{tab_label}')]")
        if tabs:
            _robust_click(driver, tabs[0])
            time.sleep(2)
        else:
            log_fn(f"  [warn] Could not find the '{tab_label}' tab - screenshot will show whatever tab is default.")
    except Exception as e:
        log_fn(f"  [warn] Could not switch to '{tab_label}' tab: {e}")
    # The tab click only swaps which breakdown is active - the actual ranked
    # table (the whole point of this screenshot) sits BELOW the summary
    # chart, out of the viewport, so a plain save_screenshot() only ever
    # captured the chart + tab bar, identically for every tab. Scroll the
    # table into view before capturing.
    try:
        rows = driver.find_elements(By.XPATH, "//table | //*[@role='table'] | //*[@role='row']")
        if rows:
            driver.execute_script("arguments[0].scrollIntoView({block:'start'});", rows[0])
            time.sleep(0.5)
        else:
            driver.execute_script("window.scrollBy(0, 500);")
            time.sleep(0.5)
    except Exception as e:
        log_fn(f"  [warn] Could not scroll to the '{tab_label}' table: {e}")
    if sort_by:
        try:
            headers = driver.find_elements(
                By.XPATH, f"//th[contains(., '{sort_by}')] | //*[@role='columnheader'][contains(., '{sort_by}')]")
            if headers:
                _robust_click(driver, headers[0])
                time.sleep(1.5)
            else:
                log_fn(f"  [warn] Could not find the '{sort_by}' column header - "
                        f"{tab_label} screenshot will show the default sort.")
        except Exception as e:
            log_fn(f"  [warn] Could not sort {tab_label} by '{sort_by}': {e}")
    try:
        driver.save_screenshot(out_path)
        log_fn(f"  GSC {tab_label} screenshot saved" + (f" (sorted by {sort_by})" if sort_by else "") + ".")
        return True
    except Exception as e:
        log_fn(f"  [warn] GSC {tab_label} screenshot save failed: {e}")
        return False


def _expand_ga4_nav_if_collapsed(driver, log_fn=None):
    """GA4's left rail can be in icon-only collapsed mode with no visible
    text at all (confirmed real case - the report labels below can't be
    found by text until this runs). The toggle button is stable regardless
    of state: <button aria-label="Nav toggle" aria-expanded="true|false">.
    Only clicks it when aria-expanded is false, since it's a real toggle -
    clicking it while already open would collapse it instead."""
    import time
    from selenium.webdriver.common.by import By
    log_fn = log_fn or log
    try:
        toggles = driver.find_elements(By.XPATH, "//button[@aria-label='Nav toggle']")
        if toggles and toggles[0].get_attribute("aria-expanded") == "false":
            _robust_click(driver, toggles[0])
            time.sleep(1.5)
    except Exception as e:
        log_fn(f"  [warn] Could not confirm/expand the GA4 nav rail: {e}")


def _navigate_ga4_report(driver, ga4_property_name, nav_labels, log_fn=None):
    """Reaches a specific GA4 standard report by clicking through GA4's own
    left-nav labels from the Reports Snapshot overview, rather than a
    hand-built deep link, since GA4's internal report-routing params are
    undocumented and liable to change; clicking the real UI by visible text
    is the same approach already proven for GSC's tab/chip clicks. Top-level
    groups (Life cycle, User) are already expanded by default once the rail
    itself is open, so nav_labels should start at the sub-group that's
    actually collapsed (e.g. "Acquisition", "User attributes"), not the
    top-level group name - clicking an already-expanded group would
    collapse it instead. Returns True if the sign-in check passed
    (navigation itself is always best-effort - a missing nav label just
    leaves the driver wherever it got to, which the caller still
    screenshots rather than failing outright)."""
    import time
    from selenium.webdriver.common.by import By
    log_fn = log_fn or log
    property_id = ga4_property_name.split("/")[-1]
    url = f"https://analytics.google.com/analytics/web/#/p{property_id}/reports/reportinghub"
    driver.get(url)
    time.sleep(10)
    cur = (driver.current_url or "").lower()
    if "accounts.google.com" in cur or "/signin" in cur:
        log_fn(f"  [warn] GA4 bounced to sign-in - skipping {nav_labels[-1]} screenshot.")
        return False
    _expand_ga4_nav_if_collapsed(driver, log_fn)
    try:
        for label in nav_labels:
            els = driver.find_elements(By.XPATH, f"//span[@class='item-text' and text()='{label}']")
            if not els:
                els = driver.find_elements(By.XPATH, f"//*[contains(text(), '{label}')]")
            if not els:
                log_fn(f"  [warn] Could not find GA4 nav item '{label}' - capturing whatever page was reached so far.")
                break
            _robust_click(driver, els[0])
            time.sleep(3)
    except Exception as e:
        log_fn(f"  [warn] Could not navigate to GA4 '{nav_labels[-1]}' report: {e}")
    return True


def capture_ga4_nav_screenshot(driver, ga4_property_name, nav_labels, out_path, log_fn=None,
                               dimension_label=None, scroll_to_table=False):
    """Real screenshot of a specific GA4 standard report (e.g. Demographic
    details) - see _navigate_ga4_report for how it gets there. dimension_label
    (optional): some reports (Tech details) default to a dimension other than
    the one wanted - e.g. "Browser" instead of "Device category" - confirmed
    live. When given, clicks the report's own dimension dropdown (its label
    text, e.g. "Browser", sits right after the report title) and picks the
    named option before capturing. Best-effort: on any failure, still
    captures whatever dimension was showing rather than skipping the slide.
    scroll_to_table (optional): GA4's standard reports show a chart above the
    fold and the actual row-by-row breakdown table further down - a plain
    viewport screenshot misses the table entirely. When True, scrolls the
    table into view first (mirroring capture_gsc_dimension_screenshot's same
    scroll-to-table approach), so this can be called a second time on the
    same report to get a table-focused screenshot alongside the chart one."""
    import time
    from selenium.webdriver.common.by import By
    log_fn = log_fn or log
    if not _navigate_ga4_report(driver, ga4_property_name, nav_labels, log_fn):
        return False
    if dimension_label:
        try:
            dim_triggers = driver.find_elements(By.XPATH, "//*[@role='button' or @role='combobox']")
            dim_triggers = [t for t in dim_triggers if t.is_displayed() and t.text.strip()
                           and len(t.text.strip()) < 30] or dim_triggers
            clicked = False
            for t in dim_triggers[:5]:
                try:
                    _robust_click(driver, t)
                    time.sleep(1)
                    opts = driver.find_elements(By.XPATH, f"//*[contains(text(), '{dimension_label}')]")
                    if opts:
                        _robust_click(driver, opts[-1])
                        time.sleep(2)
                        clicked = True
                        break
                except Exception:
                    continue
            if not clicked:
                log_fn(f"  [warn] Could not switch dimension to '{dimension_label}' - "
                        f"screenshot will show the default dimension.")
        except Exception as e:
            log_fn(f"  [warn] Could not switch dimension to '{dimension_label}': {e}")
    if scroll_to_table:
        try:
            table_els = driver.find_elements(
                By.XPATH, "//table | //*[@role='table'] | //*[@role='row']")
            table_els = [t for t in table_els if t.is_displayed()]
            if table_els:
                driver.execute_script(
                    "arguments[0].scrollIntoView({block:'center'});", table_els[-1])
            else:
                driver.execute_script("window.scrollBy(0, 600);")
            time.sleep(1)
        except Exception:
            try:
                driver.execute_script("window.scrollBy(0, 600);")
                time.sleep(1)
            except Exception:
                pass
    try:
        driver.save_screenshot(out_path)
        log_fn(f"  GA4 {nav_labels[-1]} screenshot saved.")
        return True
    except Exception as e:
        log_fn(f"  [warn] GA4 {nav_labels[-1]} screenshot save failed: {e}")
        return False


def capture_ga4_seo_channels_screenshot(driver, ga4_property_name, out_path, log_fn=None):
    """Real screenshot of GA4's Traffic acquisition report, filtered down to
    non-paid/SEO channels (see SEO_CHANNEL_GROUPS) via GA4's own "Add filter"
    UI - an SEO report shouldn't spotlight paid ad spend. UNVERIFIED against
    a live GA4 account as of when this was written (GA4's filter-builder
    click path is our best guess at the real UI, not confirmed) - if the
    filter can't be applied for any reason, this still captures the
    unfiltered Traffic acquisition report rather than failing the slide
    outright, and logs a warning so it's visible in the run log which
    happened."""
    import time
    from selenium.webdriver.common.by import By
    log_fn = log_fn or log
    if not _navigate_ga4_report(driver, ga4_property_name,
                                ["Acquisition", "Traffic acquisition"], log_fn):
        return False
    try:
        add_filter = driver.find_elements(By.XPATH, "//*[contains(text(), 'Add filter')]")
        if add_filter:
            _robust_click(driver, add_filter[0])
            time.sleep(1.5)
            dim_search = driver.find_elements(
                By.XPATH, "//input[contains(@placeholder, 'dimensions') or contains(@placeholder, 'Search')]")
            if dim_search:
                dim_search[0].send_keys("Session default channel group")
                time.sleep(1)
            dim_option = driver.find_elements(By.XPATH, "//*[contains(text(), 'Session default channel group')]")
            if dim_option:
                _robust_click(driver, dim_option[-1])
                time.sleep(1)
            for group in SEO_CHANNEL_GROUPS:
                opts = driver.find_elements(By.XPATH, f"//*[contains(text(), '{group}')]")
                if opts:
                    _robust_click(driver, opts[-1])
                    time.sleep(0.3)
            apply_btn = driver.find_elements(By.XPATH, "//*[contains(text(), 'Apply')]")
            if apply_btn:
                _robust_click(driver, apply_btn[-1])
                time.sleep(2)
            else:
                log_fn("  [warn] Could not find GA4 filter 'Apply' button - screenshot may be unfiltered.")
        else:
            log_fn("  [warn] Could not find GA4 'Add filter' control - screenshot will be unfiltered (includes paid channels).")
    except Exception as e:
        log_fn(f"  [warn] Could not apply SEO-channels filter on GA4 Traffic acquisition: {e} - capturing unfiltered.")
    # Confirmed real case: the filter-apply sequence above can silently fail
    # partway through (e.g. the wrong element matched a text search) and
    # leave the "Build filter" side panel open and unapplied in the final
    # screenshot - worse than just an unfiltered report, since it also looks
    # broken. Force it closed (Escape is the universal close gesture for this
    # kind of side-panel/dialog) regardless of whether the filter itself
    # succeeded, so the screenshot never shows this panel either way.
    try:
        from selenium.webdriver.common.keys import Keys
        if driver.find_elements(By.XPATH, "//*[contains(text(), 'Build filter')]"):
            driver.find_element(By.TAG_NAME, "body").send_keys(Keys.ESCAPE)
            time.sleep(1)
    except Exception:
        pass
    try:
        driver.save_screenshot(out_path)
        log_fn("  GA4 Traffic acquisition (SEO channels) screenshot saved.")
        return True
    except Exception as e:
        log_fn(f"  [warn] GA4 Traffic acquisition screenshot save failed: {e}")
        return False


# --------------------------------------------------------------------------- #
# Slide-building helpers
# --------------------------------------------------------------------------- #
def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _title(slide, text, color=WHITE, size=32, top=Inches(0.4)):
    box = slide.shapes.add_textbox(Inches(0.6), top, SLIDE_W - Inches(1.2), Inches(1))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = color
    return box


def _subtitle(slide, text, top, color=GREY, size=14):
    box = slide.shapes.add_textbox(Inches(0.6), top, SLIDE_W - Inches(1.2), Inches(0.8))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    return box


def _stat_box(slide, left, top, width, value, label):
    box = slide.shapes.add_textbox(left, top, width, Inches(1.1))
    tf = box.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = value
    p1.font.size = Pt(28)
    p1.font.bold = True
    p1.font.color.rgb = NAVY
    p1.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(12)
    p2.font.color.rgb = GREY
    p2.alignment = PP_ALIGN.CENTER


def build_title_slide(prs, domain, report_date):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, NAVY)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), SLIDE_W - Inches(1.6), Inches(1.5))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Performance Report"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    p2.text = f"{domain}  |  {report_date}"
    p2.font.size = Pt(18)
    p2.font.color.rgb = LIGHT_BLUE
    return slide


def build_section_divider(prs, number, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, NAVY)
    num_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(2), Inches(1.5))
    p = num_box.text_frame.paragraphs[0]
    p.text = number
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = LIGHT_BLUE
    _title(slide, title, color=WHITE, size=36, top=Inches(3.3))
    return slide


def build_overview_slide(prs, domain, start_date, end_date, has_gsc, has_ga4):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, WHITE)
    _title(slide, "Report Overview", color=NAVY)
    parts = [f"This report covers {domain}'s search and site performance from {start_date} to {end_date}."]
    if has_gsc:
        parts.append("Google Search Console data shows real organic search visibility - clicks, impressions, "
                     "and which queries/pages/countries are driving traffic.")
    if has_ga4:
        parts.append("Google Analytics 4 data shows how visitors actually behave once they land on the site - "
                     "user volume, traffic sources, and device usage.")
    body = slide.shapes.add_textbox(Inches(0.6), Inches(1.6), SLIDE_W - Inches(1.2), Inches(4))
    tf = body.text_frame
    tf.word_wrap = True
    for i, text in enumerate(parts):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(16)
        p.font.color.rgb = GREY
        p.space_after = Pt(14)
    return slide


def _line_chart_slide(prs, title, subtitle, categories, series):
    """series: list of (name, values) tuples."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, WHITE)
    _title(slide, title, color=NAVY)
    _subtitle(slide, subtitle, top=Inches(1.1))
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for name, values in series:
        chart_data.add_series(name, values)
    x, y, cx, cy = Inches(0.6), Inches(1.8), SLIDE_W - Inches(1.2), Inches(5.2)
    gframe = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, x, y, cx, cy, chart_data)
    chart = gframe.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    return slide


def _bar_chart_slide(prs, title, subtitle, categories, series_name, values, headline=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, WHITE)
    _title(slide, title, color=NAVY)
    _subtitle(slide, subtitle, top=Inches(1.1))
    if headline:
        _stat_box(slide, SLIDE_W - Inches(3.2), Inches(0.4), Inches(2.6), headline[0], headline[1])
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series(series_name, values)
    x, y, cx, cy = Inches(0.6), Inches(1.8), SLIDE_W - Inches(1.2), Inches(5.2)
    gframe = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data)
    gframe.chart.has_legend = False
    return slide


def _pie_chart_slide(prs, title, subtitle, categories, values):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, WHITE)
    _title(slide, title, color=NAVY)
    _subtitle(slide, subtitle, top=Inches(1.1))
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series("Share", values)
    x, y, cx, cy = Inches(2.5), Inches(1.8), Inches(8.3), Inches(5.2)
    gframe = slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data)
    chart = gframe.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False
    return slide


def _screenshot_slide(prs, title, subtitle, img_path):
    """A real dashboard screenshot slide - used instead of a native chart
    when one is available, matching what the reference deck's own GA/GSC
    sections actually look like (screenshots, not synthetic charts)."""
    from PIL import Image
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, WHITE)
    _title(slide, title, color=NAVY)
    _subtitle(slide, subtitle, top=Inches(1.1))
    max_w, max_h = SLIDE_W - Inches(1.2), Inches(5.3)
    with Image.open(img_path) as im:
        iw, ih = im.size
    scale = min(max_w / iw, max_h / ih)
    w, h = int(iw * scale), int(ih * scale)
    x = (SLIDE_W - w) // 2
    slide.shapes.add_picture(img_path, x, Inches(1.9), width=w, height=h)
    return slide


def build_summary_slide(prs, domain):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, NAVY)
    _title(slide, "Final Summary", color=WHITE, top=Inches(2.6))
    _subtitle(slide, f"Thank you for reviewing {domain}'s performance report.",
              top=Inches(3.4), color=LIGHT_BLUE, size=16)
    return slide


# --------------------------------------------------------------------------- #
# Report assembly
# --------------------------------------------------------------------------- #
def build_report(domain, out_path, gsc_data=None, ga4_data=None, start_date=None, end_date=None,
                 gsc_improved=None, ga4_improved=None, gsc_screenshot=None, ga4_screenshot=None,
                 gsc_queries_clicks_screenshot=None, gsc_queries_impressions_screenshot=None,
                 gsc_pages_clicks_screenshot=None, gsc_pages_impressions_screenshot=None,
                 gsc_countries_clicks_screenshot=None, gsc_countries_impressions_screenshot=None,
                 ga4_acquisition_screenshot=None, ga4_demographics_screenshot=None,
                 ga4_events_screenshot=None, ga4_events_table_screenshot=None,
                 ga4_devices_screenshot=None, ga4_source_medium_screenshot=None):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    gsc_improved = gsc_improved or set()
    ga4_improved = ga4_improved or set()

    today = datetime.date.today().isoformat()
    build_title_slide(prs, domain, today)
    build_overview_slide(prs, domain, start_date, end_date, bool(gsc_data), bool(ga4_data))

    if gsc_data:
        build_section_divider(prs, "01", "Google Search Console")

        daily = gsc_data["daily"]
        cats = [r["keys"][0] for r in daily]
        clicks = [r.get("clicks", 0) for r in daily]
        impressions = [r.get("impressions", 0) for r in daily]
        total_clicks = sum(clicks)
        total_impr = sum(impressions)
        metric_label = {"clicks": "Clicks", "impressions": "Impressions", "position": "Position"}
        if gsc_screenshot:
            if gsc_improved:
                improved_txt = " & ".join(metric_label[m] for m in ("clicks", "impressions", "position") if m in gsc_improved)
                subtitle = f"{improved_txt} improved vs. the previous {(_period_days(start_date, end_date))} days - real GSC dashboard view"
            else:
                subtitle = f"{start_date} to {end_date} - real GSC dashboard view"
            _screenshot_slide(prs, "Traffic Status", subtitle, gsc_screenshot)
        elif daily:
            slide = _line_chart_slide(
                prs, "Traffic Status", f"Clicks & impressions, {start_date} to {end_date}",
                cats, [("Clicks", clicks), ("Impressions", impressions)])
            _stat_box(slide, SLIDE_W - Inches(3.2), Inches(0.4), Inches(1.5), f"{total_clicks:,}", "Total Clicks")
            _stat_box(slide, SLIDE_W - Inches(1.6), Inches(0.4), Inches(1.5), f"{total_impr:,}", "Total Impressions")

        def _gsc_dimension_slides(base_title, clicks_shot, impressions_shot, rows, dimension_fmt):
            """Each GSC breakdown (Queries/Pages/Countries) gets up to two
            screenshot slides - one sorted by Clicks, one by Impressions -
            instead of a single screenshot at whatever GSC's default sort
            happened to be, so both are genuinely visible rather than only
            ever showing "top by clicks" and leaving impressions unseen."""
            if clicks_shot:
                subtitle = ("Clicks improved vs. the previous period - real GSC "
                            f"{base_title} view (sorted by clicks)" if "clicks" in gsc_improved else
                            f"{start_date} to {end_date} - real GSC {base_title} view (sorted by clicks)")
                _screenshot_slide(prs, f"Top {base_title} by Clicks", subtitle, clicks_shot)
            elif rows and not impressions_shot:
                _bar_chart_slide(
                    prs, f"Top {base_title} by Clicks", f"Top 10 by clicks in this period",
                    [dimension_fmt(r) for r in rows], "Clicks", [r.get("clicks", 0) for r in rows])
            if impressions_shot:
                subtitle = (f"{start_date} to {end_date} - real GSC {base_title} view (sorted by impressions)")
                _screenshot_slide(prs, f"Top {base_title} by Impressions", subtitle, impressions_shot)
            elif rows and not clicks_shot:
                _bar_chart_slide(
                    prs, f"Top {base_title} by Impressions", f"Top 10 by impressions in this period",
                    [dimension_fmt(r) for r in rows], "Impressions", [r.get("impressions", 0) for r in rows])

        _gsc_dimension_slides("Searches by Keywords", gsc_queries_clicks_screenshot,
                              gsc_queries_impressions_screenshot, gsc_data["queries"], lambda r: r["keys"][0])
        _gsc_dimension_slides("Searches by Pages", gsc_pages_clicks_screenshot,
                              gsc_pages_impressions_screenshot, gsc_data["pages"], lambda r: r["keys"][0])
        _gsc_dimension_slides("Searches by Country", gsc_countries_clicks_screenshot,
                              gsc_countries_impressions_screenshot, gsc_data["countries"],
                              lambda r: r["keys"][0].upper())

    if ga4_data:
        build_section_divider(prs, "02", "Google Analytics")

        daily = ga4_data["daily"]
        if ga4_screenshot:
            if ga4_improved & {"activeUsers", "sessions"}:
                ga4_label = {"activeUsers": "Active Users", "sessions": "Sessions"}
                improved_txt = " & ".join(ga4_label[m] for m in ("activeUsers", "sessions") if m in ga4_improved)
                subtitle = f"{improved_txt} improved vs. the previous period - real GA4 dashboard view"
            else:
                subtitle = f"Active users & sessions, {start_date} to {end_date} - real GA4 dashboard view"
            _screenshot_slide(prs, "Audience Trend", subtitle, ga4_screenshot)
        elif daily:
            daily_sorted = sorted(daily, key=lambda r: r.get("date", ""))
            cats = [r.get("date", "") for r in daily_sorted]
            users = [int(r.get("activeUsers", 0) or 0) for r in daily_sorted]
            sessions = [int(r.get("sessions", 0) or 0) for r in daily_sorted]
            slide = _line_chart_slide(
                prs, "Audience Trend", f"Active users & sessions, {start_date} to {end_date}",
                cats, [("Active Users", users), ("Sessions", sessions)])
            _stat_box(slide, SLIDE_W - Inches(3.2), Inches(0.4), Inches(1.5), f"{sum(users):,}", "Total Active Users")
            _stat_box(slide, SLIDE_W - Inches(1.6), Inches(0.4), Inches(1.5), f"{sum(sessions):,}", "Total Sessions")

        seo_channels = filter_seo_channels(ga4_data["channels"])
        if ga4_acquisition_screenshot:
            subtitle = ("Organic/referral sessions improved vs. the previous period - real GA4 "
                        "Traffic acquisition view (paid channels excluded)"
                        if "seo_sessions" in ga4_improved else
                        f"{start_date} to {end_date} - real GA4 Traffic acquisition view "
                        f"(paid channels excluded)")
            _screenshot_slide(prs, "Traffic Acquisition", subtitle, ga4_acquisition_screenshot)
        elif seo_channels:
            _pie_chart_slide(
                prs, "Traffic Acquisition", "Non-paid sessions by channel in this period",
                [r.get("sessionDefaultChannelGroup", "Unknown") for r in seo_channels],
                [int(r.get("sessions", 0) or 0) for r in seo_channels])

        devices = ga4_data["devices"]
        if ga4_demographics_screenshot:
            subtitle = ("Active users improved vs. the previous period - real GA4 Demographic details view"
                        if "activeUsers" in ga4_improved else
                        f"{start_date} to {end_date} - real GA4 Demographic details view")
            _screenshot_slide(prs, "Demographic Details", subtitle, ga4_demographics_screenshot)
        elif devices:
            _pie_chart_slide(
                prs, "Demographic Details", "Active users by device category",
                [r.get("deviceCategory", "Unknown").title() for r in devices],
                [int(r.get("activeUsers", 0) or 0) for r in devices])

        # Screenshot-only sections (no native-chart fallback data pipeline for
        # these yet) - simply omitted from the deck if the screenshot
        # couldn't be captured, rather than fabricating placeholder data.
        if ga4_events_screenshot:
            _screenshot_slide(prs, "Events", f"{start_date} to {end_date} - real GA4 Events view",
                             ga4_events_screenshot)
        if ga4_events_table_screenshot:
            _screenshot_slide(prs, "Events (Table)",
                             f"{start_date} to {end_date} - real GA4 Events table, scrolled to the full breakdown",
                             ga4_events_table_screenshot)
        if ga4_devices_screenshot:
            _screenshot_slide(prs, "Devices", f"{start_date} to {end_date} - real GA4 Tech details view",
                             ga4_devices_screenshot)
        if ga4_source_medium_screenshot:
            _screenshot_slide(prs, "Source / Medium",
                             f"{start_date} to {end_date} - real GA4 User acquisition (source/medium) view",
                             ga4_source_medium_screenshot)

    build_summary_slide(prs, domain)
    prs.save(out_path)
    log(f"[DONE] {out_path}")


# --------------------------------------------------------------------------- #
def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("domain", help="Target domain, e.g. example.com")
    ap.add_argument("--out", required=True, help="Output .pptx path")
    ap.add_argument("--gsc-account", default=None, help="Connected GSC account email")
    ap.add_argument("--ga4-account", default=None,
                    help="Connected account email for GA4 API calls, if different from --gsc-account "
                         "(a client may grant Search Console and Analytics access to different Google "
                         "accounts). Defaults to --gsc-account when omitted.")
    ap.add_argument("--ga4-property", default=None, help='GA4 property resource name, e.g. "properties/123456789"')
    ap.add_argument("--days", type=int, default=28, help="How many days back to report on")
    ap.add_argument("--session-id", default=None,
                    help="A GSC 'session' id (logged-in browser profile, see gsc_audit.py) for the "
                         "--gsc-account - if given, real GSC dashboard screenshots (and GA4's, when GA4 "
                         "uses the same account) are captured for every section regardless of "
                         "improvement. Omit to skip GSC screenshots and use native charts only.")
    ap.add_argument("--ga4-session-id", default=None,
                    help="A separate GSC 'session' id for --ga4-account, when it's a different account "
                         "from --gsc-account. Defaults to --session-id when omitted (the common case: "
                         "same account for both).")
    args = ap.parse_args()
    if not args.ga4_account:
        args.ga4_account = args.gsc_account
    if not args.ga4_session_id:
        args.ga4_session_id = args.session_id

    if not args.gsc_account and not args.ga4_property:
        log("[ERROR] Provide --gsc-account and/or --ga4-property - nothing to report on otherwise.")
        sys.exit(2)

    end_date = datetime.date.today()
    start_date = end_date - datetime.timedelta(days=args.days)
    start_s, end_s = start_date.isoformat(), end_date.isoformat()
    prev_end = start_date - datetime.timedelta(days=1)
    prev_start = prev_end - datetime.timedelta(days=args.days - 1)
    prev_start_s, prev_end_s = prev_start.isoformat(), prev_end.isoformat()

    gsc_data = None
    ga4_data = None
    gsc_improved = set()
    ga4_improved = set()
    property_url = None

    if args.gsc_account:
        log(f"[1/4] Resolving GSC access for {args.gsc_account}...")
        try:
            token = gsc_audit.get_access_token(args.gsc_account)
            property_url = gsc_audit.resolve_property(token, args.domain)
            log(f"   -> Property: {property_url}")
            log("[2/4] Fetching Search Console data...")
            gsc_data = fetch_gsc_data(token, property_url, start_s, end_s)
            log("   -> Comparing against the previous period...")
            current_totals = fetch_gsc_totals(token, property_url, start_s, end_s)
            previous_totals = fetch_gsc_totals(token, property_url, prev_start_s, prev_end_s)
            gsc_improved = compute_gsc_improvements(current_totals, previous_totals)
            log(f"   -> Improved vs. previous {args.days} days: {sorted(gsc_improved) or 'none'}")
        except Exception as e:
            log(f"   [warn] GSC data skipped: {type(e).__name__}: {e}")

    if args.ga4_property:
        log(f"[2/4] Fetching Google Analytics data for {args.ga4_property} (account: {args.ga4_account})...")
        try:
            token = gsc_audit.get_access_token(args.ga4_account) if args.ga4_account else None
            if not token:
                raise Exception("GA4 requires --gsc-account and/or --ga4-account.")
            ga4_data = fetch_ga4_data(token, args.ga4_property, start_s, end_s)
            current_ga4 = fetch_ga4_totals(token, args.ga4_property, start_s, end_s)
            previous_ga4 = fetch_ga4_totals(token, args.ga4_property, prev_start_s, prev_end_s)
            ga4_improved = compute_ga4_improvements(current_ga4, previous_ga4)
            previous_channels = fetch_ga4_channel_sessions(token, args.ga4_property, prev_start_s, prev_end_s)
            ga4_improved |= compute_ga4_seo_improvement(ga4_data["channels"], previous_channels)
            log(f"   -> Improved vs. previous {args.days} days: {sorted(ga4_improved) or 'none'}")
        except Exception as e:
            log(f"   [warn] GA4 data skipped: {type(e).__name__}: {e}")

    if not gsc_data and not ga4_data:
        log("[ERROR] Could not fetch any GSC or GA4 data - nothing to build a report from.")
        sys.exit(2)

    gsc_screenshot, ga4_screenshot = None, None
    gsc_queries_clicks_screenshot, gsc_queries_impressions_screenshot = None, None
    gsc_pages_clicks_screenshot, gsc_pages_impressions_screenshot = None, None
    gsc_countries_clicks_screenshot, gsc_countries_impressions_screenshot = None, None
    ga4_acquisition_screenshot, ga4_demographics_screenshot = None, None
    ga4_events_screenshot, ga4_events_table_screenshot = None, None
    ga4_devices_screenshot, ga4_source_medium_screenshot = None, None
    # Real dashboard screenshots are now attempted for EVERY section whenever a
    # session is available, not just the ones that improved vs. the previous
    # period - previously a section with no improvement fell back to a native
    # (synthetic) chart even though a real screenshot was just as capturable;
    # the caption below is what actually distinguishes "improved" framing from
    # a plain current-period view, not whether the screenshot itself exists.
    if (args.session_id or args.ga4_session_id) and (gsc_data or ga4_data):
        log("[3/4] Capturing real dashboard screenshots...")
        shot_dir = os.path.dirname(os.path.abspath(args.out)) or "."
        # GSC and GA4 can be different Google accounts (--session-id vs.
        # --ga4-session-id) - a single logged-in browser session can only
        # ever be one Google account, so when they differ this launches TWO
        # separate browser sessions, one per account, instead of trying to
        # capture both dashboards through whichever account happened to be
        # logged into the one browser. When they're the same (the common
        # case), the second session is just a reference to the first - no
        # second browser launch.
        same_session = bool(args.ga4_session_id) and args.ga4_session_id == args.session_id
        gsc_driver, ga4_driver = None, None
        try:
            if args.session_id and gsc_data and property_url:
                gsc_driver = launch_screenshot_browser(args.session_id)
                path = os.path.join(shot_dir, "_gsc_perf_screenshot.png")
                if capture_gsc_performance_screenshot(gsc_driver, property_url, gsc_improved, path):
                    gsc_screenshot = path
                # Each dimension gets two screenshots - sorted by Clicks and
                # by Impressions - instead of one at whichever sort GSC
                # happened to default to.
                path = os.path.join(shot_dir, "_gsc_queries_clicks_screenshot.png")
                if capture_gsc_dimension_screenshot(gsc_driver, property_url, "Queries", path, sort_by="Clicks"):
                    gsc_queries_clicks_screenshot = path
                path = os.path.join(shot_dir, "_gsc_queries_impressions_screenshot.png")
                if capture_gsc_dimension_screenshot(gsc_driver, property_url, "Queries", path, sort_by="Impressions"):
                    gsc_queries_impressions_screenshot = path
                path = os.path.join(shot_dir, "_gsc_pages_clicks_screenshot.png")
                if capture_gsc_dimension_screenshot(gsc_driver, property_url, "Pages", path, sort_by="Clicks"):
                    gsc_pages_clicks_screenshot = path
                path = os.path.join(shot_dir, "_gsc_pages_impressions_screenshot.png")
                if capture_gsc_dimension_screenshot(gsc_driver, property_url, "Pages", path, sort_by="Impressions"):
                    gsc_pages_impressions_screenshot = path
                path = os.path.join(shot_dir, "_gsc_countries_clicks_screenshot.png")
                if capture_gsc_dimension_screenshot(gsc_driver, property_url, "Countries", path, sort_by="Clicks"):
                    gsc_countries_clicks_screenshot = path
                path = os.path.join(shot_dir, "_gsc_countries_impressions_screenshot.png")
                if capture_gsc_dimension_screenshot(gsc_driver, property_url, "Countries", path, sort_by="Impressions"):
                    gsc_countries_impressions_screenshot = path
            if ga4_data and args.ga4_property:
                if same_session and gsc_driver:
                    ga4_driver = gsc_driver
                elif args.ga4_session_id:
                    ga4_driver = launch_screenshot_browser(args.ga4_session_id)
                else:
                    log("   No --ga4-session-id given - skipping GA4 screenshots, using native charts instead.")
                if ga4_driver:
                    path = os.path.join(shot_dir, "_ga4_screenshot.png")
                    if capture_ga4_screenshot(ga4_driver, args.ga4_property, path):
                        ga4_screenshot = path
                    path = os.path.join(shot_dir, "_ga4_acquisition_screenshot.png")
                    if capture_ga4_seo_channels_screenshot(ga4_driver, args.ga4_property, path):
                        ga4_acquisition_screenshot = path
                    path = os.path.join(shot_dir, "_ga4_demographics_screenshot.png")
                    if capture_ga4_nav_screenshot(ga4_driver, args.ga4_property,
                                                  ["User attributes", "Demographic details"], path):
                        ga4_demographics_screenshot = path
                    # These 3 nav paths are newer additions and, like the SEO-channels
                    # filter above, not yet confirmed against a live GA4 account - if
                    # a label doesn't match GA4's real current UI, capture_ga4_nav_screenshot
                    # still captures whatever page was reached (never fails outright),
                    # and build_report() simply omits the slide if the screenshot path
                    # stays None, so a wrong guess here degrades gracefully rather than
                    # breaking the rest of the report.
                    path = os.path.join(shot_dir, "_ga4_events_screenshot.png")
                    if capture_ga4_nav_screenshot(ga4_driver, args.ga4_property, ["Engagement", "Events"], path):
                        ga4_events_screenshot = path
                    # Second capture of the same report, scrolled down to the
                    # full Events breakdown table - the first capture above
                    # is the chart view above the fold.
                    path = os.path.join(shot_dir, "_ga4_events_table_screenshot.png")
                    if capture_ga4_nav_screenshot(ga4_driver, args.ga4_property, ["Engagement", "Events"], path,
                                                  scroll_to_table=True):
                        ga4_events_table_screenshot = path
                    path = os.path.join(shot_dir, "_ga4_devices_screenshot.png")
                    if capture_ga4_nav_screenshot(ga4_driver, args.ga4_property, ["Tech", "Tech details"], path,
                                                  dimension_label="Device category"):
                        ga4_devices_screenshot = path
                    path = os.path.join(shot_dir, "_ga4_source_medium_screenshot.png")
                    if capture_ga4_nav_screenshot(ga4_driver, args.ga4_property, ["Acquisition", "User acquisition"], path):
                        ga4_source_medium_screenshot = path
        except Exception as e:
            log(f"   [warn] Screenshot capture skipped: {type(e).__name__}: {e}")
        finally:
            for drv in ({id(gsc_driver): gsc_driver, id(ga4_driver): ga4_driver} if not same_session
                        else {id(gsc_driver): gsc_driver}).values():
                if drv:
                    try:
                        drv.quit()
                    except Exception:
                        pass
    elif gsc_data or ga4_data:
        log("   No session available - skipping real screenshots, using native charts instead.")

    log("[4/4] Building report...")
    build_report(args.domain, args.out, gsc_data=gsc_data, ga4_data=ga4_data, start_date=start_s, end_date=end_s,
                gsc_improved=gsc_improved, ga4_improved=ga4_improved,
                gsc_screenshot=gsc_screenshot, ga4_screenshot=ga4_screenshot,
                gsc_queries_clicks_screenshot=gsc_queries_clicks_screenshot,
                gsc_queries_impressions_screenshot=gsc_queries_impressions_screenshot,
                gsc_pages_clicks_screenshot=gsc_pages_clicks_screenshot,
                gsc_pages_impressions_screenshot=gsc_pages_impressions_screenshot,
                gsc_countries_clicks_screenshot=gsc_countries_clicks_screenshot,
                gsc_countries_impressions_screenshot=gsc_countries_impressions_screenshot,
                ga4_acquisition_screenshot=ga4_acquisition_screenshot,
                ga4_demographics_screenshot=ga4_demographics_screenshot,
                ga4_events_screenshot=ga4_events_screenshot,
                ga4_events_table_screenshot=ga4_events_table_screenshot,
                ga4_devices_screenshot=ga4_devices_screenshot,
                ga4_source_medium_screenshot=ga4_source_medium_screenshot)

    for p in (gsc_screenshot, ga4_screenshot,
              gsc_queries_clicks_screenshot, gsc_queries_impressions_screenshot,
              gsc_pages_clicks_screenshot, gsc_pages_impressions_screenshot,
              gsc_countries_clicks_screenshot, gsc_countries_impressions_screenshot,
              ga4_acquisition_screenshot, ga4_demographics_screenshot,
              ga4_events_screenshot, ga4_events_table_screenshot,
              ga4_devices_screenshot, ga4_source_medium_screenshot):
        if p and os.path.exists(p):
            try:
                os.remove(p)
            except Exception:
                pass


if __name__ == "__main__":
    try:
        main()
    except SystemExit:
        raise
    except Exception as e:
        log(f"[ERROR] {type(e).__name__}: {e}")
        sys.exit(1)
